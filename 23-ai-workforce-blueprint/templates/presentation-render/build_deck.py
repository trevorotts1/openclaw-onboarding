#!/usr/bin/env python3
"""
build_deck.py — DETERMINISTIC, NO-AI deck builder for the Presentations department.

This is the PROVEN deterministic deck pipeline, generalized for the fleet. It is a
companion to render_deck.py in this same directory: render_deck.py is the canonical
prompt-driven render module used by the full webinar pipeline; build_deck.py is the
single-command, zero-AI-at-runtime path that takes a slides.json (the agent's only
creative output) and produces a finished .pptx with no further model judgement.

WHAT IT DOES (zero AI judgement at runtime):
    1. Reads a slides.json (see slides.schema.json) and an output .pptx path.
    2. For EACH slide, composes a KIE.ai text-to-image prompt MECHANICALLY by concatenating:
         scene description  +  exact English copy (verbatim from slides.json)  +
         optional logo wordmark + layout hint  +  the MANDATORY English/Latin-only pin.
       No model decides wording — the copy is whatever the slides.json author wrote.
    3. Calls KIE.ai (gpt-image-2-text-to-image) per the ONLY-allowed recipe:
         POST https://api.kie.ai/api/v1/jobs/createTask  -> data.taskId
         GET  https://api.kie.ai/api/v1/jobs/recordInfo?taskId=<id>  -> data.state
         on success: parse data.resultJson (JSON string) -> resultUrls[0]
       Poll every ~8s up to ~3 minutes. On HTTP 429, sleep 20s and retry.
    4. Downloads resultUrls[0] UNAUTHENTICATED to <renders_dir>/slide-NN.png.
    5. VERIFIES each PNG: real PNG magic bytes AND non-zero size.
    6. Retries a failing slide up to 3 times (re-submit from scratch).
    7. Assembles ALL slide PNGs into a 16:9 .pptx via python-pptx — ONE full-bleed
       picture per slide, NO text boxes (text is baked into the KIE image, matching the
       designed architecture).
    8. Prints a JSON summary: {slidesRendered, kieTaskIds, outputPath, failures}.

HARD CONTRACTS:
    * It NEVER calls any native/built-in image tool (image_generate, openai, etc.).
      KIE createTask is the ONLY image path.
    * The dead endpoint /api/v1/image/gpt-image is NEVER used (and is refused if seen).
    * If KIE is unreachable / a slide cannot be rendered after retries, the script FAILS
      LOUDLY with a non-zero exit code. It NEVER silently substitutes a placeholder image.

USAGE:
    python3 build_deck.py <slides.json> <out.pptx> [renders_dir]

    renders_dir defaults to "<out.pptx dir>/renders".

EXIT CODES:
    0 — every slide rendered and the .pptx was written.
    1 — one or more slides failed after retries (NO .pptx written), or assembly failed.
    2 — fatal config error (no API key, bad slides.json, missing python-pptx, etc.).

ENVIRONMENT:
    KIE_API_KEY — the CLIENT's own KIE.ai key (never the operator's, never shared).
    Read from env, else from one of the client's standard env stores:
        ~/.openclaw/workspace/.env
        ~/clawd/secrets/.env
        ~/.openclaw/secrets/.env
    (paths expanded against the current user's HOME).
"""

import json
import os
import sys
import time
import urllib.request
import urllib.error
from pathlib import Path
from typing import Optional

# ---------------------------------------------------------------------------
# Constants — the ONLY allowed KIE.ai endpoints/model (verified live 2026-06-16)
# ---------------------------------------------------------------------------

CREATE_URL = "https://api.kie.ai/api/v1/jobs/createTask"
POLL_URL   = "https://api.kie.ai/api/v1/jobs/recordInfo"
MODEL_T2I  = "gpt-image-2-text-to-image"

ASPECT_RATIO = "16:9"
RESOLUTION   = "2K"

# DEAD endpoint — refuse to ever touch it.
DEAD_ENDPOINT_FRAGMENT = "/api/v1/image/gpt-image"

# Poll cadence per the recipe: ~8s interval, up to ~3 minutes total.
POLL_INTERVAL_S = 8
POLL_MAX_SECONDS = 180
POLL_MAX_PASSES = POLL_MAX_SECONDS // POLL_INTERVAL_S  # ~22 passes

# 429 backoff.
RATE_LIMIT_SLEEP_S = 20

# Per-slide retries (full re-submit) on any failure.
SLIDE_MAX_ATTEMPTS = 3

# The MANDATORY trailing pin appended to EVERY prompt.
ENGLISH_PIN = (
    "All text rendered in the image MUST be in English, Latin alphabet ONLY. "
    "NO Chinese/CJK or non-Latin characters anywhere. Render the copy spelled "
    "correctly, letter-for-letter. No garbled, misspelled, or invented text."
)

# Client's own env stores (HOME-relative — works on any client box).
SECRETS_CANDIDATES = [
    os.path.expanduser("~/.openclaw/workspace/.env"),
    os.path.expanduser("~/clawd/secrets/.env"),
    os.path.expanduser("~/.openclaw/secrets/.env"),
]


# ---------------------------------------------------------------------------
# API key
# ---------------------------------------------------------------------------

def load_api_key() -> str:
    key = os.environ.get("KIE_API_KEY", "").strip()
    if key:
        return key.strip("'\"")
    for path in SECRETS_CANDIDATES:
        p = Path(path)
        if not p.exists():
            continue
        for line in p.read_text().splitlines():
            line = line.strip()
            if line.startswith("KIE_API_KEY="):
                value = line[len("KIE_API_KEY="):].strip().strip("'\"")
                if value:
                    return value
    print("FATAL: KIE_API_KEY not found in env or any of:", file=sys.stderr)
    for path in SECRETS_CANDIDATES:
        print("   ", path, file=sys.stderr)
    sys.exit(2)


# ---------------------------------------------------------------------------
# Prompt composition — PURELY MECHANICAL (no AI)
# ---------------------------------------------------------------------------

def compose_prompt(slide: dict) -> str:
    """
    Deterministically build the KIE prompt for one slide by concatenating, in order:
      scene  ->  exact copy lines (verbatim)  ->  logo  ->  layout  ->  English pin.
    No wording is invented; every text token comes straight from slides.json.
    """
    scene = str(slide["scene"]).strip()
    copy_lines = [str(c).strip() for c in slide["copy"] if str(c).strip()]
    if not copy_lines:
        raise ValueError(f"slide {slide.get('slide')}: 'copy' has no non-empty lines.")

    headline = copy_lines[0]
    rest = copy_lines[1:]

    layout = str(slide.get("layout", "")).strip()
    if not layout:
        layout = (
            "headline in the upper area, remaining copy in the lower third, "
            "high-contrast legible placement against the background"
        )

    logo = str(slide.get("logo", "")).strip()

    parts = []
    parts.append(f"A 16:9 presentation slide background. SCENE: {scene}")
    # The copy is rendered AS TEXT into the image — quote it verbatim so the model
    # spells it exactly.
    parts.append(f'Render this exact headline text on the slide, spelled exactly: "{headline}".')
    if rest:
        joined = " | ".join(f'"{line}"' for line in rest)
        parts.append(
            f"Also render this exact supporting copy, each on its own line, "
            f"spelled exactly: {joined}."
        )
    if logo:
        parts.append(
            f'Render a small, clean brand wordmark reading exactly "{logo}" as a logo lockup.'
        )
    parts.append(f"LAYOUT: {layout}.")
    parts.append(
        "Typography must be crisp, modern, and perfectly legible; do not crop or cut off "
        "any letters; keep text inside the safe area away from the slide edges."
    )
    parts.append(ENGLISH_PIN)

    prompt = " ".join(parts)

    # Self-guard: never let the dead endpoint string sneak into a prompt payload.
    if DEAD_ENDPOINT_FRAGMENT in prompt:
        raise ValueError(
            f"slide {slide.get('slide')}: composed prompt contains the dead endpoint "
            f"fragment '{DEAD_ENDPOINT_FRAGMENT}'. Refusing."
        )
    return prompt


# ---------------------------------------------------------------------------
# HTTP helpers
# ---------------------------------------------------------------------------

class RateLimited(Exception):
    pass


def _http_json(method: str, url: str, api_key: str, body: Optional[dict] = None) -> dict:
    if DEAD_ENDPOINT_FRAGMENT in url:
        raise RuntimeError(
            f"REFUSED: attempted to call the dead endpoint {DEAD_ENDPOINT_FRAGMENT}. "
            "This pipeline only uses /api/v1/jobs/createTask and /api/v1/jobs/recordInfo."
        )
    data = json.dumps(body).encode() if body is not None else None
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    req = urllib.request.Request(url, data=data, headers=headers, method=method)
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            return json.loads(resp.read().decode())
    except urllib.error.HTTPError as exc:
        if exc.code == 429:
            raise RateLimited(f"HTTP 429 from {url}")
        body_text = exc.read().decode(errors="replace")
        raise RuntimeError(f"HTTP {exc.code} {method} {url}\nResponse: {body_text}") from exc
    except urllib.error.URLError as exc:
        # Network unreachable — fail loud (no substitution).
        raise RuntimeError(f"NETWORK ERROR reaching {url}: {exc}. KIE is unreachable.") from exc


def submit_task(prompt: str, api_key: str) -> str:
    payload = {
        "model": MODEL_T2I,
        "input": {
            "prompt": prompt,
            "aspect_ratio": ASPECT_RATIO,
            "resolution": RESOLUTION,
        },
    }
    resp = _http_json("POST", CREATE_URL, api_key, body=payload)
    if resp.get("code") != 200:
        raise RuntimeError(f"createTask non-200 code. Full response: {json.dumps(resp)}")
    task_id = (resp.get("data") or {}).get("taskId")
    if not task_id:
        raise RuntimeError(f"createTask 200 but no taskId. Full response: {json.dumps(resp)}")
    return task_id


def poll_task(task_id: str, api_key: str) -> str:
    """Poll recordInfo every ~8s up to ~3 min. Returns resultUrls[0] on success."""
    url = f"{POLL_URL}?taskId={task_id}"
    deadline = time.time() + POLL_MAX_SECONDS
    passes = 0
    while time.time() < deadline and passes < POLL_MAX_PASSES + 2:
        passes += 1
        try:
            resp = _http_json("GET", url, api_key)
        except RateLimited:
            print(f"    [poll] 429 — sleeping {RATE_LIMIT_SLEEP_S}s", flush=True)
            time.sleep(RATE_LIMIT_SLEEP_S)
            continue
        data = resp.get("data") or {}
        state = str(data.get("state", "")).lower()

        if state == "success":
            result_json_str = data.get("resultJson")
            if not result_json_str:
                raise RuntimeError(f"taskId {task_id}: success but no resultJson: {json.dumps(resp)}")
            result_obj = json.loads(result_json_str)
            urls = result_obj.get("resultUrls", [])
            if not urls:
                raise RuntimeError(f"taskId {task_id}: empty resultUrls: {json.dumps(result_obj)}")
            return urls[0]

        if state in ("fail", "failed", "error", "cancelled"):
            raise RuntimeError(
                f"taskId {task_id}: terminal state '{state}' "
                f"failCode={data.get('failCode')} failMsg={data.get('failMsg')}"
            )

        print(f"    [poll {passes}] taskId={task_id} state={state!r}; sleep {POLL_INTERVAL_S}s", flush=True)
        time.sleep(POLL_INTERVAL_S)

    raise RuntimeError(f"taskId {task_id}: not complete within {POLL_MAX_SECONDS}s.")


def download_unauthenticated(url: str, dest: Path) -> None:
    """UNAUTHENTICATED GET of the CDN result URL (Bearer token causes 403)."""
    req = urllib.request.Request(url, headers={"User-Agent": "build_deck/1.0"})
    with urllib.request.urlopen(req, timeout=180) as resp, open(dest, "wb") as f:
        f.write(resp.read())


def verify_png(path: Path) -> None:
    """Verify the downloaded file is a real, non-empty PNG."""
    if not path.exists():
        raise RuntimeError(f"{path}: file was not written.")
    size = path.stat().st_size
    if size == 0:
        raise RuntimeError(f"{path}: file is zero bytes.")
    with open(path, "rb") as f:
        magic = f.read(8)
    if magic[:4] != b"\x89PNG":
        raise RuntimeError(f"{path}: not a PNG (magic={magic[:8].hex()}, size={size}).")


# ---------------------------------------------------------------------------
# Per-slide render with retry
# ---------------------------------------------------------------------------

def render_slide(slide: dict, api_key: str, renders_dir: Path) -> dict:
    """
    Render a single slide. Returns {"slide", "file", "taskId"} on success.
    Raises RuntimeError if all SLIDE_MAX_ATTEMPTS fail.
    """
    ordinal = int(slide["slide"])
    name = f"slide-{ordinal:02d}"
    out_path = renders_dir / f"{name}.png"
    prompt = compose_prompt(slide)

    last_err = None
    for attempt in range(1, SLIDE_MAX_ATTEMPTS + 1):
        print(f"  [{name}] attempt {attempt}/{SLIDE_MAX_ATTEMPTS}", flush=True)
        try:
            # submit (429-aware)
            while True:
                try:
                    task_id = submit_task(prompt, api_key)
                    break
                except RateLimited:
                    print(f"    [submit] 429 — sleeping {RATE_LIMIT_SLEEP_S}s", flush=True)
                    time.sleep(RATE_LIMIT_SLEEP_S)
            print(f"    submitted -> taskId={task_id}", flush=True)
            result_url = poll_task(task_id, api_key)
            print(f"    success resultUrls[0]={result_url}", flush=True)
            download_unauthenticated(result_url, out_path)
            verify_png(out_path)
            size = out_path.stat().st_size
            print(f"    downloaded+verified -> {out_path} ({size:,} bytes)", flush=True)
            return {"slide": ordinal, "file": str(out_path), "taskId": task_id}
        except Exception as exc:  # noqa: BLE001 — deliberately catch to retry
            last_err = exc
            print(f"    FAIL attempt {attempt}: {exc}", file=sys.stderr, flush=True)
            if attempt < SLIDE_MAX_ATTEMPTS:
                time.sleep(3)

    raise RuntimeError(f"{name}: failed after {SLIDE_MAX_ATTEMPTS} attempts. Last error: {last_err}")


# ---------------------------------------------------------------------------
# pptx assembly — full-bleed picture per slide, NO text boxes
# ---------------------------------------------------------------------------

def assemble_pptx(rendered: list, out_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("FATAL: python-pptx is not installed (pip install python-pptx).", file=sys.stderr)
        sys.exit(2)

    prs = Presentation()
    # 16:9 at the dept-documented dimensions.
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]  # blank layout — no placeholders

    for item in sorted(rendered, key=lambda r: r["slide"]):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            item["file"], 0, 0, width=Inches(10), height=Inches(5.625)
        )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------

def main():
    if len(sys.argv) not in (3, 4):
        print("Usage: python3 build_deck.py <slides.json> <out.pptx> [renders_dir]", file=sys.stderr)
        sys.exit(2)

    slides_path = Path(sys.argv[1])
    out_path = Path(sys.argv[2])
    renders_dir = Path(sys.argv[3]) if len(sys.argv) == 4 else out_path.parent / "renders"

    if not slides_path.exists():
        print(f"FATAL: slides.json not found: {slides_path}", file=sys.stderr)
        sys.exit(2)

    try:
        slides = json.loads(slides_path.read_text())
    except json.JSONDecodeError as exc:
        print(f"FATAL: slides.json is not valid JSON: {exc}", file=sys.stderr)
        sys.exit(2)

    if not isinstance(slides, list) or not slides:
        print("FATAL: slides.json must be a non-empty JSON array.", file=sys.stderr)
        sys.exit(2)

    # Basic schema validation (deterministic, fail loud).
    seen = set()
    for s in slides:
        if not isinstance(s, dict):
            print("FATAL: every slide must be an object.", file=sys.stderr)
            sys.exit(2)
        for req in ("slide", "scene", "copy"):
            if req not in s:
                print(f"FATAL: slide missing required field '{req}': {json.dumps(s)}", file=sys.stderr)
                sys.exit(2)
        if not isinstance(s["copy"], list) or not s["copy"]:
            print(f"FATAL: slide {s.get('slide')}: 'copy' must be a non-empty array.", file=sys.stderr)
            sys.exit(2)
        ordn = s["slide"]
        if not isinstance(ordn, int) or ordn in seen:
            print(f"FATAL: slide ordinal must be a unique integer: {ordn}", file=sys.stderr)
            sys.exit(2)
        seen.add(ordn)

    api_key = load_api_key()
    renders_dir.mkdir(parents=True, exist_ok=True)

    print(f"=== build_deck — {len(slides)} slides ===", flush=True)
    print(f"slides.json: {slides_path}", flush=True)
    print(f"renders_dir: {renders_dir}", flush=True)
    print(f"out.pptx:    {out_path}", flush=True)
    print(f"endpoint:    {CREATE_URL}  model={MODEL_T2I}\n", flush=True)

    rendered = []
    task_ids = []
    failures = []

    for slide in sorted(slides, key=lambda s: s["slide"]):
        try:
            result = render_slide(slide, api_key, renders_dir)
            rendered.append(result)
            task_ids.append(result["taskId"])
        except Exception as exc:  # noqa: BLE001
            failures.append({"slide": slide.get("slide"), "error": str(exc)})
            print(f"  SLIDE FAILED: {exc}", file=sys.stderr, flush=True)

    # FAIL LOUD: never produce a partial deck silently.
    if failures:
        summary = {
            "slidesRendered": len(rendered),
            "kieTaskIds": task_ids,
            "outputPath": None,
            "failures": failures,
        }
        print("\n=== SUMMARY (FAILED) ===", file=sys.stderr)
        print(json.dumps(summary, indent=2))
        sys.exit(1)

    # Assemble.
    try:
        assemble_pptx(rendered, out_path)
    except SystemExit:
        raise
    except Exception as exc:  # noqa: BLE001
        print(f"FATAL: pptx assembly failed: {exc}", file=sys.stderr)
        summary = {
            "slidesRendered": len(rendered),
            "kieTaskIds": task_ids,
            "outputPath": None,
            "failures": [{"slide": "ASSEMBLY", "error": str(exc)}],
        }
        print(json.dumps(summary, indent=2))
        sys.exit(1)

    summary = {
        "slidesRendered": len(rendered),
        "kieTaskIds": task_ids,
        "outputPath": str(out_path),
        "failures": [],
    }
    print("\n=== SUMMARY (OK) ===", flush=True)
    print(json.dumps(summary, indent=2))
    sys.exit(0)


if __name__ == "__main__":
    main()
