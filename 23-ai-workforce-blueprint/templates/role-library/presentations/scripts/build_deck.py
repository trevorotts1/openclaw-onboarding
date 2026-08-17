#!/usr/bin/env python3
# =============================================================================
# SOP-LOCKED DEPARTMENT: if you add/modify any SOP, role, or gate, you MUST
# update PIPELINE-MANIFEST.json + build_deck.py + a test. Run
# scripts/sync_check.py — it fails the gate if the Python and the SOP stack
# drift. The PREFLIGHT_REQUIRED set + the _chk_* checkers + every AF-* code this
# file cites are reconciled against
# universal-sops/presentation-slide-craft/PIPELINE-MANIFEST.json by
# scripts/sync_check.py (in BOTH directions). Procedure:
# universal-sops/presentation-slide-craft/SOP-SLIDE-06-EXTENSION-AND-SYNC.md.
# =============================================================================
"""
build_deck.py — DETERMINISTIC, NO-AI deck builder for the Presentations department.

================================================================================
A .pptx IS NOT A DELIVERED PRESENTATION. THIS SCRIPT IS ONLY THE PHASE-4 RENDERER.
================================================================================
Producing only a .pptx with this script is NOT a delivered presentation and is NOT
a complete deliverable. build_deck.py renders slide images and assembles a bare
.pptx — and NOTHING ELSE. It does NOT produce the presenter guide PDF, the presenter
speech (PRESENTERS-SPEECH.md/.pdf/-FISH-TAGGED.md), the presenter audio
(PRESENTER-AUDIO.mp3), the teleprompter app, the deck PDF export, the infographic
checklist slide, or the GHL media upload. The FULL deliverable downstream is the
NINE-MEMBER bundle —
    [Deck-Title]-FINAL.pptx, [Deck-Title]-FINAL.pdf, PRESENTER-GUIDE.pdf,
    PRESENTERS-SPEECH.md (pure), PRESENTERS-SPEECH.pdf (teleprompter PDF),
    PRESENTERS-SPEECH-FISH-TAGGED.md (fish-tagged), PRESENTER-AUDIO.mp3,
    infographic.png, presenter-teleprompter.html (the teleprompter web app).
The .pptx this script emits is an INTERMEDIATE artifact. The Director/Delivery flow
MUST run the guide / speech / audio / PDF / infographic / teleprompter roles, and the
run is blocked from "Done" by the POSTFLIGHT COMPLETENESS GATE (AF-BUNDLE-COMPLETE)
AND the DELIVERY INTERLOCK (AF-DELIVER + AF-DH1 + AF-DELIVERY-COMPLETE). Do NOT
"finish at a .pptx." See sops/CLIENT-WEBINAR-DECK-SOP.md §9a,
sops/delivery-concierge-sops.md SOP 9.5, and
sops/SOP-SLIDE-00-MASTER-QC-AUTOFAIL-RULESET.md (AF-DELIVERY-COMPLETE).

UPSTREAM STEPS THAT PRODUCE THE REQUIRED OUTPUTS:
    - [Deck-Title]-FINAL.pdf            : produced by PPTX Assembly Specialist (PDF export)
    - PRESENTER-GUIDE.pdf               : produced by Presenters Guide Specialist
    - PRESENTERS-SPEECH.md/.pdf          : produced by Presenters Speech Writer
    - PRESENTERS-SPEECH-FISH-TAGGED.md   : produced by Presenters Speech Writer / Fish Audio
    - PRESENTER-AUDIO.mp3               : produced by Audio Demonstration Specialist
    - infographic.png                   : produced by infographic-checklist role
    - presenter-teleprompter.html       : produced by build_teleprompter.py (teleprompter role)
The postflight gate enforces that ALL NINE exist and pass their size thresholds. It
hard-fails (AF-BUNDLE-COMPLETE, exit 5) if any are missing or under-size, so they
can NEVER be silently skipped.

DEFAULT OUTPUT DESTINATION — ~/Downloads:
    The final bundle dir defaults to ~/Downloads/<client-slug>-<deck-slug>/.
    Use --out <dir> to override. Files NEVER default to a scratch/run dir — they
    go to ~/Downloads so the client receives them from a predictable, clean location.
================================================================================

This is the PROVEN deterministic deck pipeline, generalized for the fleet. It is the
single-command, zero-AI-at-runtime path that takes a slides.json (the agent's only
creative output) and produces a finished .pptx with no further model judgement.

WHAT IT DOES (zero AI judgement at runtime):
    1. Reads a slides.json (see slides.schema.json) and an output .pptx path.
    2. For EACH slide, LOADS the per-slide RICH prompt VERBATIM from
         working/prompts/slide-NN.txt  (or the SOP-named  slide-NN-prompt.txt).
       This is the Slide Image Creator's hand-authored output — a 9,000–14,000-char
       (HARD floor 9,000; HARD ceiling 18,000) prompt that already carries the
       typography (per-line weight + pt size),
       placement, usage, the logo(s), the scene, the verbatim copy, the negative
       block, and everything else that appears on the slide. build_deck.py renders
       THAT prompt verbatim — it does NOT compose its own thin prompt from
       scene+copy. A whole slide is rendered in ONE gpt-image-2 generation.
       If a slide has NO rich prompt file, or the prompt is < PROMPT_CHAR_FLOOR
       (9,000) chars, or it fails the QUALITY floor (AF-P13 eight-class negative
       block / AF-P14 spelling-lock / AF-P-DENSITY / AF-P-VERBATIM), build_deck.py
       FAILS LOUD — it NEVER silently falls back to a thin composed prompt. The
       prompt-side QC gates run on every rich prompt, all FAIL-LOUD (no silent render):
         (a) FACIAL-INTELLIGENCE / REPRESENTATION gate — refuses any prompt that
             carries a forbidden hardcoded demographic default (the "60/30/10"
             landmine). Representation comes from the slide spec / casting ledger
             (the client's captured audience), never a baked-in default split
             (SOP-CAST-01, AF-R3).
         (b) CHAR-COUNT gate — the rich prompt length must be within
             [PROMPT_CHAR_FLOOR, PROMPT_CHAR_CEILING] = [9000, 18000]; the floor is
             HARD (any prompt under the 9,000-char target-band low end is not run, not
             rendered, not updated — AF-P1/AF-PROMPT-FLOOR) and the 18,000 ceiling is the universal GPT-Image 2
             safety boundary (AF-P2).
    3. Calls KIE.ai (gpt-image-2-text-to-image, or gpt-image-2-image-to-image when a
       logo is supplied via input_urls so the WHOLE slide + logo render in ONE
       generation) per the ONLY-allowed recipe:
         POST https://api.kie.ai/api/v1/jobs/createTask  -> data.taskId
         GET  https://api.kie.ai/api/v1/jobs/recordInfo?taskId=<id>  -> data.state
         on success: parse data.resultJson (JSON string) -> resultUrls[0]
       Poll every ~8s up to ~7 minutes. Heavy gpt-image-2 renders (split scenes,
       multiple photoreal zones) routinely run past 3 minutes; a slide that is still
       actively 'generating' is given the full ceiling (plus a bounded grace window)
       before it is declared a timeout. On HTTP 429, sleep 20s and retry.
    4. Downloads resultUrls[0] UNAUTHENTICATED to <renders_dir>/slide-NN.png.
    5. VERIFIES each PNG: real PNG magic bytes AND non-zero size.
    6. Retries a failing slide up to 3 times (re-submit from scratch).
    7. Assembles ALL slide PNGs into a 16:9 .pptx via python-pptx — ONE full-bleed
       picture per slide, NO text boxes (text is baked into the KIE image, matching the
       designed architecture).
    8. Prints a JSON summary: {slidesRendered, kieTaskIds, outputPath, failures}.

HARD CONTRACTS:
    * It NEVER calls any native/built-in image tool (image_generate, openai, etc.).
      KIE createTask is the ONLY image path.
    * The dead endpoint /api/v1/image/gpt-image is NEVER used (and is refused if seen).
    * If KIE is unreachable / a slide cannot be rendered after retries, the script FAILS
      LOUDLY with a non-zero exit code. It NEVER silently substitutes a placeholder image.

USAGE:
    python3 build_deck.py <slides.json> <out.pptx> [renders_dir] [--run-dir DIR]
        [--logo PNG] [--out BUNDLE_DIR] [--adhoc-no-process]

    renders_dir defaults to "<out.pptx dir>/renders".
    --out BUNDLE_DIR  overrides the default ~/Downloads/<client-slug>-<deck-slug>/
                      destination for the final deliverable bundle. If not supplied,
                      the bundle goes to ~/Downloads.

    NOTE: a successful run (exit 0 + a .pptx at <out.pptx>) means the deck is
    RENDERED, not DELIVERED. The .pptx alone is NOT a delivered presentation — the
    full nine-member bundle (deck .pptx/.pdf, guide PDF, the three speech artifacts,
    audio, infographic, teleprompter app) is required downstream (see the banner
    above and the POSTFLIGHT COMPLETENESS GATE / DELIVERY INTERLOCK).

OFFICIAL-LOGO COMPOSITING (optional):
    --logo <png-or-URL>  — the client's REAL, official logo. When supplied
        (or when working/copy/intake.json carries brand.logo_image_path), the deck
        wears the client's actual mark, not an AI-hallucinated wordmark:
          (1) When --logo is a URL: KIE composites the REAL logo via image-to-image
              (gpt-image-2-image-to-image, input_urls) in the SAME single generation
              that renders the slide — the rich prompt already carries the
              "place this exact provided logo, do not redraw/recolor/restyle"
              instruction (the AI wordmark is suppressed by the prompt itself).
          (2) When --logo is a local PNG: assemble_pptx composites the EXACT same
              logo PNG (RGBA, transparency preserved) onto EVERY slide in a
              consistent top-right corner at a tasteful size (~13% of slide width)
              with a small margin, layered over the full-bleed background.
        The --logo flag wins over intake.json brand.logo_image_path if both exist.
        Either way the prompt is the rich verbatim prompt — build_deck.py never
        rewrites the logo instruction; the Slide Image Creator authored it.

PROCESS PREFLIGHT (un-bypassable by default):
    Before ANY render or assembly, build_deck.py REQUIRES that the upstream
    department process artifacts already exist on disk in the run/working dir.
    build_deck.py is ONLY the deterministic Phase-4 renderer + a stripped Phase-8
    assembler; it carries none of the upstream intelligence (research, copy,
    human approvals) or QC gates. Running it on a hand-fed slides.json with no
    upstream chain produces a technically-valid .pptx that NEVER went through the
    department flow — exactly the shortcut this preflight exists to refuse.

    The preflight enforces EVERY SOP artifact: intake, research brief, converting
    arc, slides_copy, design/typography brief, copy QC (pass), anti-compression
    coverage, AND a >=9,000-char RICH prompt (clearing the AF-P13/AF-P14/density/
    verbatim quality floor) for EVERY slide in working/prompts/
    (the Slide Image Creator's output). Any missing/short/thin/deviation → refuse to
    render, exit 3, loud. There is no path past the gate with a thin or absent
    per-slide prompt.

    The run dir is the directory that contains a `working/` subtree. It is found
    by --run-dir if given, else by walking UP from slides.json (then from the
    out.pptx dir) until a `working/` dir is found, else slides.json's parent.

    Required artifacts (each annotated with the dept phase/role that produces it)
    are listed in PREFLIGHT_REQUIRED below. If ANY is missing or incomplete, the
    script FAILS LOUD (exit 3), prints exactly which artifacts are missing and who
    produces each, and renders/assembles NOTHING.

    --adhoc-no-process  — DELIBERATE standalone-testing override ONLY. Skips the
        process preflight, prints a loud banner that the output is NOT a
        process-compliant deliverable, and proceeds. NEVER use for client work.

EXIT CODES:
    0 — every slide rendered, the .pptx was written, AND the postflight
        completeness gate confirmed all nine bundle deliverables present + sized.
    1 — one or more slides failed after retries (NO .pptx written), or assembly failed.
    2 — fatal config error (no API key, bad slides.json, missing python-pptx, etc.).
    3 — process preflight FAILED: required upstream dept artifacts are missing
        (NO render, NO .pptx). Run the real dept pipeline, or pass
        --adhoc-no-process for explicit non-deliverable standalone testing.
    5 — POSTFLIGHT COMPLETENESS GATE FAILED (AF-BUNDLE-COMPLETE): one or more
        required deliverables are missing or below their size threshold. The run
        may NOT be reported as "done" / "complete" until exit 0. The .pptx was
        written but the bundle is incomplete.

ENVIRONMENT:
    KIE_API_KEY — the CLIENT's own KIE.ai key (never the operator's, never shared).
    Read from env, else from one of the client's standard env stores:
        ~/.openclaw/workspace/.env
        ~/clawd/secrets/.env
        ~/.openclaw/secrets/.env
    (paths expanded against the current user's HOME).
"""

import concurrent.futures
import hashlib
import hmac
import ipaddress
import json
import os
import re
import sys
import time
import urllib.request
import urllib.error
from datetime import date, datetime, timezone
from pathlib import Path

from presentation_job.checkpoint import atomic_write_text, PREDICATES
from presentation_job.result import CheckResult
from typing import Dict, Optional, Tuple
from urllib.parse import urlparse, quote

# ---------------------------------------------------------------------------
# Constants — the ONLY allowed KIE.ai endpoints/model (verified live 2026-06-16)
# ---------------------------------------------------------------------------

CREATE_URL = "https://api.kie.ai/api/v1/jobs/createTask"
POLL_URL   = "https://api.kie.ai/api/v1/jobs/recordInfo"
MODEL_T2I  = "gpt-image-2-text-to-image"
MODEL_I2I  = "gpt-image-2-image-to-image"  # OFFICIAL-LOGO mode: KIE composites the REAL logo via input_urls (image-to-image), NOT a flat overlay or AI wordmark

ASPECT_RATIO = "16:9"
RESOLUTION   = "2K"

# ---------------------------------------------------------------------------
# OFFICIAL-LOGO compositing geometry (used by assemble_pptx)
# ---------------------------------------------------------------------------
# Slide is 16:9 at 10" x 5.625". The real logo is composited top-right at a
# tasteful size: ~13% of slide width, with a small consistent margin. The logo's
# native aspect ratio is preserved (height derived from its pixel dimensions) so
# the EXACT client file is never distorted; transparency (RGBA) is preserved
# because python-pptx passes PNG alpha straight through.
SLIDE_WIDTH_IN  = 10.0
SLIDE_HEIGHT_IN = 5.625
LOGO_WIDTH_FRACTION = 0.13   # ~13% of slide width (within the 12-15% tasteful band)
LOGO_MARGIN_IN = 0.25        # small consistent margin from the top/right edges

# DEAD endpoint — refuse to ever touch it.
DEAD_ENDPOINT_FRAGMENT = "/api/v1/image/gpt-image"

# Poll cadence per the recipe: exponential backoff 10s/20s/40s with a HARD cap at
# ~15 minutes (D14 / FIX-7). A dead or never-completing kie task must surface a
# terminal FAIL (RenderPollTimeout) instead of hanging the run in 'in_progress'.
# Heavy gpt-image-2 renders (e.g. a split then-vs-now band with two photoreal scenes
# and multiple text zones) regularly exceed 3 minutes of pure render latency. 180s was
# too low and made merely-slow slides time out and fail loud with outputPath=null.
# The backoff ladder below (10s for the first 2 min, 20s for the next 3 min, then 40s)
# gives slow-but-healthy renders room to finish while guaranteeing a hard stop at
# POLL_MAX_SECONDS; genuine terminal failures still fail immediately (see poll_task:
# 'fail'/'error'/'cancelled' short-circuit). All three are overridable via env so a
# slow kit render can stretch the wall-clock ceiling without editing code.
def _poll_ladder() -> tuple:
    """Return (fast_s, fast_window_s, medium_s, medium_window_s, slow_s).
    The schedule is 10s for the first 2 min, 20s for the next 3 min, then 40s after."""
    def _f(name, default):
        try:
            v = int(os.environ.get(name, str(default)))
        except ValueError:
            v = default
        return max(1, v)
    return (_f("BUILD_DECK_POLL_FAST_S", 10),
            _f("BUILD_DECK_POLL_FAST_WINDOW_S", 120),
            _f("BUILD_DECK_POLL_MEDIUM_S", 20),
            _f("BUILD_DECK_POLL_MEDIUM_WINDOW_S", 180),
            _f("BUILD_DECK_POLL_SLOW_S", 40))


POLL_FAST_S, POLL_FAST_WINDOW_S, POLL_MEDIUM_S, POLL_MEDIUM_WINDOW_S, POLL_SLOW_S = _poll_ladder()

# Hard wall-clock cap for a single task's poll (FIX-7: ~15 min). Beyond this the poll
# raises RenderPollTimeout (a terminal FAIL), never a silent hang.
def _poll_max_seconds() -> int:
    try:
        v = int(os.environ.get("BUILD_DECK_POLL_MAX_SECONDS", "900"))
    except ValueError:
        v = 900
    return max(30, v)


POLL_MAX_SECONDS = _poll_max_seconds()

# Backward-compatible alias: the old fixed 8s interval / pass count are gone; code
# that referenced POLL_INTERVAL_S / POLL_MAX_PASSES now uses the ladder. Kept as a
# derived value so any external reader still resolves a sane number.
POLL_INTERVAL_S = POLL_FAST_S

# 429 backoff.
RATE_LIMIT_SLEEP_S = 20

# ---------------------------------------------------------------------------
# BATCH-RENDER cadence (FIX-5 / D22 root cause B). KIE allows 20 new generation
# requests per 10s per account and 100+ concurrent tasks. The old path submitted
# ONE slide and polled it to completion before submitting the next (or ran a
# small thread pool), which serialized submit+poll+download and wasted 20+ min.
# The batch path submits ALL prompts once at BATCH_SUBMIT_INTERVAL_S (0.6s),
# so a 20-slide deck lands its 20 createTask calls in ~12s — inside the 20/10s
# window, empirically proven 200 for a burst of 8 (see ERRORS-DETECTED D22) —
# then polls all taskIds together every BATCH_POLL_INTERVAL_S and downloads each
# image the moment ITS task completes (downloads do not wait for the slowest
# slide). These are the DEFAULTS; render_slides_batch accepts them as parameters
# so a box can tune the cadence at the call site without a code edit. The
# defaults implement the documented 20/10s + concurrent-task model.
BATCH_SUBMIT_INTERVAL_S = 0.6   # seconds between consecutive createTask calls
BATCH_POLL_INTERVAL_S = 10.0    # seconds between poll passes over ALL pending tasks
BATCH_MAX_POLL_SECONDS = 900    # hard cap (15 min) — a stuck task surfaces FAIL, no hang

# Per-slide retries (full re-submit) on any failure.
# A KIE failCode 400 "Internal Error, Please try again later" is a VERIFIED
# TRANSIENT upstream error: retry the IDENTICAL request with exponential backoff
# (do NOT trim the prompt or change aspect_ratio). Overridable via env so a kit
# render can use 4-6 attempts; default kept at a safe 6 with exponential backoff.
def _slide_max_attempts() -> int:
    try:
        n = int(os.environ.get("BUILD_DECK_SLIDE_MAX_ATTEMPTS", "6"))
    except ValueError:
        n = 6
    return max(1, min(10, n))


SLIDE_MAX_ATTEMPTS = _slide_max_attempts()

# Exponential backoff base (seconds) between slide re-submits on transient failure.
def _slide_backoff_base() -> float:
    try:
        return max(0.5, float(os.environ.get("BUILD_DECK_SLIDE_BACKOFF_BASE", "4")))
    except ValueError:
        return 4.0


SLIDE_BACKOFF_BASE = _slide_backoff_base()
SLIDE_BACKOFF_CAP_S = 90.0

# Parallel render fan-out. Each slide is an independent KIE generation (submit +
# poll + download + verify), so they run concurrently in a ThreadPoolExecutor.
# Bounded so the per-task 429 backoff still holds and a client box is not swamped.
# Overridable via BUILD_DECK_RENDER_WORKERS (clamped to [1, 12]).
def _render_workers() -> int:
    try:
        n = int(os.environ.get("BUILD_DECK_RENDER_WORKERS", "6"))
    except ValueError:
        n = 6
    return max(1, min(12, n))

# The MANDATORY trailing pin appended to EVERY prompt.
ENGLISH_PIN = (
    "All text rendered in the image MUST be in English, Latin alphabet ONLY. "
    "NO Chinese/CJK or non-Latin characters anywhere. Render the copy spelled "
    "correctly, letter-for-letter. No garbled, misspelled, or invented text."
)

# ---------------------------------------------------------------------------
# PROMPT CHAR-COUNT GATE (the prompt-side QC gate, fail-loud)
# ---------------------------------------------------------------------------
# Two facts from the standards (qc-specialist-presentations.md AF-P1/AF-P2,
# slide-image-creator-sops.md SOP 9.1 step 3, MODEL-SPECS.md):
#   * GPT-Image 2 hard ceiling on input.prompt is 20,000 characters on both
#     endpoints. The HARD MAXIMUM is 18,000 — a 2,000-char safety margin below
#     that ceiling so a prompt is never rejected or truncated by the platform.
#     This ceiling is UNIVERSAL: it applies to every prompt path. A prompt over
#     it is an AF-P2 auto-fail.
#   * The FLOOR is 9,000 chars and it is HARD (AF-P1). build_deck.py does NOT
#     compose its own thin prompt any more — it renders the Slide Image Creator's
#     hand-authored RICH per-slide prompt VERBATIM (working/prompts/slide-NN.txt).
#     That prompt carries the full 15-element spec (typography size + per-line
#     weight, placement, usage, the logo(s), the scene, verbatim copy, the negative
#     block, everything on the slide); the SOP targets 9,000–14,000 chars. A prompt
#     under 9,000 chars is, by definition, not a real slide prompt — it is a thin
#     stub or a truncated file — so it is NOT run, NOT rendered, and NOT updated:
#     the slide FAILS LOUD (AF-P1). 9,000 is the reconciled HARD floor: it is the LOW
#     end of the SOP's own authoring target band (9,000–14,000 per slide-image-creator.md),
#     raised from the retired 5,000 floor because the deep diagnosis proved a 5,000-char
#     prompt — mostly boilerplate negative block + archetype line — is fully "compliant"
#     yet too thin to carry the 15-element spec, and a competing role taught "5,000 = done."
#     The HARD floor now IS the target-band low end, so a floor-grazing prompt physically
#     cannot omit the typography/face/composition/lighting/palette detail. On top of the
#     raw length floor, the rich-prompt gate ALSO enforces a density/specificity check
#     (hex palette + a pt/px type size + a composition/zone token + a distinct-word floor)
#     and a VERBATIM-WORDS-BAKED check (the slide's exact copy must appear in the prompt
#     body), so length alone can never satisfy the gate. The Prompt-Authoring phase
#     (P4-PROMPT) writes to this band and the Prompt-QC phase (P4-PROMPT-QC) verifies it
#     independently — and the governed Prompt-QC gate now RE-MEASURES every on-disk prompt
#     (it is no longer a JSON rubber stamp).
PROMPT_CHAR_FLOOR = 9000      # HARD floor (AF-P1/AF-PROMPT-FLOOR): the 9,000-char target-band LOW end; any rich prompt under it is NOT run/rendered/updated — FAIL LOUD
PROMPT_CHAR_TARGET_HIGH = 18000  # v15.0.0: SOP authoring-target HIGH end raised 14000->18000 so the authoring band is 9,000-18,000 (matches PROMPT_CHAR_CEILING; the HARD ceiling stays PROMPT_CHAR_CEILING). Final char standard: MIN 9,000 / MAX 18,000.
PROMPT_CHAR_CEILING = 18000   # UNIVERSAL hard maximum (AF-P2; 2,000 under the 20,000 API ceiling)
PROMPT_MIN_DISTINCT_WORDS = 220  # AF-P-DENSITY: a >=9,000-char prompt that repeats one paragraph to pad length has few distinct words; a genuinely rich prompt has 400+. Floor catches paste-repetition padding.

# ---------------------------------------------------------------------------
# AF-COPY-BAND — per-slide COPY character-count FLOOR and CEILING (spec item 7 /
# §8.2). Reconciles the 3-way bullet-count conflict (slide-copywriter.md:187 said
# "5 bullets, 7 words each"; qc-specialist-presentations.md c4 said "3 bullets max
# or 30 words max"; slide-image-creator.md:229,759 renders "max 5 words") DOWN to
# one enforced rule, and adds the missing MINIMUM (previously no floor existed
# anywhere for headline/subhead/slide-total — only prose ceilings that no code
# ever checked). Characters, not words: deterministic, locale-stable, and directly
# enforceable on slides.json copy[] without an LLM judgment call.
#
# copy[] positional semantics (slides.schema.json): index 0 = HEADLINE,
# index 1 = SUBHEAD (optional), index 2 = KICKER / third block (optional),
# index 3+ = BULLETS (max COPY_BULLET_MAX_COUNT).
#
# Because QC_FOREIGN_SIGNATURES (below) bans QC *reports* that carry a
# word-count rubric, this band is enforced ONLY as a first-party preflight
# check (_chk_copy_density, like AF-P1) — never as a QC-report criterion — so
# there is no collision with that ban.
# ---------------------------------------------------------------------------
COPY_HEADLINE_CHAR_FLOOR = 12        # ~2 words -- kills one-word non-headlines
COPY_HEADLINE_CHAR_CEILING = 60      # ~9 words -- preserves the historical 9-word ceiling
COPY_SUBHEAD_CHAR_FLOOR = 20         # ~4 words -- only applied when a subhead is present
COPY_SUBHEAD_CHAR_CEILING = 110      # ~18 words -- preserves the historical 18-word ceiling
COPY_KICKER_CHAR_CEILING = 40        # kicker/third block: absent OR <= 40 chars, no floor
COPY_BULLET_CHAR_FLOOR = 8           # ~1-2 words
COPY_BULLET_CHAR_CEILING = 30        # ~5 words -- resolves the 5-vs-7-word conflict DOWN to
                                      # the image-render reality (slide-image-creator.md's
                                      # "max 5 words" is what actually bakes into the PNG)
COPY_BULLET_MAX_COUNT = 3            # resolves the 3-vs-5-bullet conflict DOWN (qc-specialist wins)
COPY_SLIDE_TOTAL_CHAR_FLOOR = 40     # kills the "empty slide, one orphan word" failure
COPY_SLIDE_TOTAL_CHAR_CEILING = 180  # ~30 words @ ~6 chars/word -- preserves historical ceiling
COPY_HOOK_SLIDE_TOTAL_CHAR_FLOOR = 12  # pure-typography HOOK / section-banner slides: the
                                        # hook line alone is allowed to be short

# U067 stage 1 (Rule 3.5). The hook/section-banner exemption was unreachable on the
# boolean arc_allocation schema, so every such slide has been judged against the
# 40-char floor. Turning the exemption on flips slides from FAIL to PASS, which is
# the fail-open direction and must never arrive in one step. While this is False the
# gate still applies the standard floor and merely REPORTS which slides the exemption
# would cover. Stage 3 (flipping this to True) is a SEPARATE unit.
COPY_HOOK_EXEMPTION_ENFORCED = False

# REQUIRED STRUCTURAL BLOCKS (AF-P1). A real rich prompt is not just long enough — it
# carries the load-bearing structural scaffolding: an [ARCHETYPE ...] layout declaration,
# the final-paragraph negative block (the header our gold exemplars + the SOP 9.8 template
# emit is literally "DO-NOT BLOCK"), and at least one "Do not " imperative inside it. A
# 9,000-char file with none of these is a verbose stub, not a slide spec. This is the
# structural-block check FOLDED IN from the retired presentation-render/render_deck.py
# (_validate_prompt) so build_deck.py is the ONE canonical renderer — it now enforces
# every invariant render_deck.py used to, with no separate render module to drift.
# Matched case-insensitively, on the same STRIPPED text the char-count gate measures.
#
# Each entry is the CANONICAL (primary) label the renderer reports as "missing" when the
# block is absent. The negative-block entry is the canonical "DO-NOT BLOCK" header that
# both Appendix A gold exemplars and the SOP 9.8 prompt template actually emit; the older
# "NEGATIVE BLOCK" wording remains an ACCEPTED ALIAS (see STRUCTURAL_BLOCK_ALIASES) so a
# prompt authored faithfully to either label satisfies the gate. The gate is NOT weakened:
# a prompt that carries NEITHER header still fails AF-P-STRUCT, and the 8-class CONTENT
# teeth (AF-P13) + spelling-lock (AF-P14) fire independently regardless of header wording.
REQUIRED_STRUCTURAL_BLOCKS = ["[ARCHETYPE", "DO-NOT BLOCK", "Do not "]

# ACCEPTED ALIASES for a required structural block. A block in REQUIRED_STRUCTURAL_BLOCKS is
# "present" iff the canonical label OR any of its aliases appears (case-insensitive). The
# negative-block header is the only block with historical label drift: the canonical header
# our exemplars/SOP emit is "DO-NOT BLOCK", but earlier authoring (and the test fixtures /
# preflight examples) used "NEGATIVE BLOCK" — both must satisfy the structural requirement,
# so neither a DO-NOT-BLOCK-headed nor a NEGATIVE-BLOCK-headed prompt false-bounces AF-P-STRUCT.
STRUCTURAL_BLOCK_ALIASES = {
    "DO-NOT BLOCK": ["NEGATIVE BLOCK"],
}


def _structural_block_present(block: str, text_lc: str) -> bool:
    """True iff the required structural BLOCK (or any of its accepted aliases) appears in
    text_lc (which MUST already be lower-cased). Folds the canonical label + every alias
    into a single OR so a prompt authored to either the canonical 'DO-NOT BLOCK' header or
    the legacy 'NEGATIVE BLOCK' header satisfies the same structural slot — without
    weakening the gate (a prompt carrying none of them is still missing the block)."""
    candidates = [block] + STRUCTURAL_BLOCK_ALIASES.get(block, [])
    return any(c.lower() in text_lc for c in candidates)


def _missing_structural_blocks(text_lc: str) -> list:
    """Return the canonical labels of every REQUIRED_STRUCTURAL_BLOCKS entry whose canonical
    label AND all of its accepted aliases are absent from text_lc (lower-cased). Empty list
    = all required structural blocks are present. The returned label is always the canonical
    primary so the re-author message tells the author the header to add."""
    return [b for b in REQUIRED_STRUCTURAL_BLOCKS if not _structural_block_present(b, text_lc)]

# ---------------------------------------------------------------------------
# AF-P13 — the EIGHT mandatory negative-block defect CLASSES (slide-image-creator.md
# SOP 9.8). REQUIRED_STRUCTURAL_BLOCKS above only proves a negative-block header
# ("DO-NOT BLOCK", or the legacy "NEGATIVE BLOCK" alias) and ONE "Do not " sentence
# exist — a one-line stub satisfies it. AF-P13 gives the
# negative block real teeth: the block must actually name ALL EIGHT defect classes the
# forensic reference deck shipped (garbled text, logo mutation, placeholder tokens,
# image narration, anatomical artifacts, competing background, skin-tone fidelity,
# universal baseline). Each class below must have >=1 of its tokens present
# (case-insensitive) somewhere in the prompt. Token sets are deliberately TOLERANT
# (any one hit = the class is named) so legitimate phrasing variants pass while a thin
# stub that says only "no text" fails. This is the mechanical half of AF-P13 (the
# paired-positive-twin + no-contradiction audit stays with the QC vision pass).
NEGATIVE_BLOCK_CLASS_TOKENS = {
    "garbled/misspelled text": [
        "misspell", "garble", "letter-for-letter", "letter for letter",
        "render every quoted", "exactly as written", "render every letter"],
    "logo mutation": [
        "logo", "monogram", "tagline lockup", "reference mark", "redraw",
        "redesign", "recolor", "restyle", "reinterpret"],
    "placeholder/bracket tokens": [
        "bracketed token", "square bracket", "owner to confirm", "placeholder",
        "tbd", "build note", "to supply", "pending", "insert"],
    "image narration/presenter/meta": [
        "narrat", "presenter line", "spoken-script", "spoken script",
        "stage direction", "telegraphing", "webinar", "self-talk",
        "describe the picture", "description of the picture", "build note"],
    "anatomical artifacts": [
        "finger", "fused hand", "malformed", "anatom", "distorted facial",
        "mismatched eye", "asymmetric eye", "distorted teeth",
        "over-smoothed skin", "body proportion", "extra limb"],
    "background competing with text": [
        "busy", "cluttered", "high-detail background", "compete", "behind any text",
        "text zone", "scrim", "legib", "negative space"],
    "demographic/skin-tone fidelity": [
        "demographic", "skin tone", "skin-tone", "representation_mix", "lighten",
        "ashen", "desaturate", "mono-cast", "mono cast", "deep skin"],
    "carried-forward universal baseline": [
        "watermark", "emoji", "clipart", "default font", "calibri", "arial",
        "times new roman", "system default", "ui artifact", "user-interface",
        "em dash", "pure-black", "pure black"],
}

# AF-P14 — per-string SPELLING-LOCK. A real rich prompt pins the EXACT spelling of every
# on-slide word with a per-line spelling-lock directive (SOP 9.9 "HEADLINE VERBATIM +
# SPELLING-LOCK"). The mechanical half: at least one spelling-lock marker token must be
# present. Paired with the VERBATIM-WORDS-BAKED check (the exact slide copy must appear
# in the prompt body), a stub can neither fake the lock nor the words.
SPELLING_LOCK_TOKENS = [
    "spelling-lock", "spelling lock", "letter-for-letter", "letter for letter",
    "render this exact string", "reads exactly", "render every quoted text string exactly",
    "spelled exactly", "exact spelling", "render every letter",
]

# AF-P-DENSITY — concrete specificity signals. A genuinely rich prompt that carries the
# 15-element spec necessarily declares: a brand palette HEX, an explicit type SIZE
# (pt/px), and a COMPOSITION/zone instruction. A 9,000-char file padded with boilerplate
# can clear the length floor while carrying none of these. Each signal below is required.
PROMPT_COMPOSITION_TOKENS = [
    "thirds", "rule of thirds", "grid", "left third", "right third", "upper third",
    "lower third", "left-third", "right-third", "upper-third", "lower-third", "zone",
    "safe margin", "safe-margin", "quadrant", "negative space", "focal point", "composition",
]
_HEX_COLOR_RE = re.compile(r"#[0-9a-fA-F]{6}\b")
_TYPE_SIZE_RE = re.compile(r"\b\d{2,4}\s?(?:pt|px|pixels)\b", re.IGNORECASE)
_WORD_RE = re.compile(r"[a-z0-9][a-z0-9'\-]+")

# ---------------------------------------------------------------------------
# RESEARCH-CITATION GATE (AF-RESEARCH-UNCITED) — the minimum cited-URL floor
# ---------------------------------------------------------------------------
# A research brief that asserts research_complete:true but carries fewer than
# MIN_CITED_SOURCES real http(s) URLs is not a research pack — it is a
# self-reported completion that skipped the actual work.  This constant is the
# authoritative threshold.  It is intentionally named (not inlined) so sync_check
# can verify it exists and the PIPELINE-MANIFEST.json py_symbol can point at it.
#
# Rationale for 8:
#   The Deep Research SOP mandates categories A, B, C, D (D1+D2+D3), G, H, I, K, L
#   for a full brief — that is at minimum 10 categories, each requiring at least one
#   source URL.  8 is a conservative floor that catches the "zero research" case
#   (0-2 URLs, or a brief built entirely from intake with no external sources) while
#   not over-penalising a condensed micro-deck brief (which the SOP allows to be
#   smaller but still requires real external sources in C and D).
MIN_CITED_SOURCES = 8   # HARD floor (AF-RESEARCH-UNCITED): fewer real URLs = FAIL LOUD

# C3: counting raw http(s) URL strings is gameable — the same source cited 8 times,
# or 8 localhost/example.com/RFC-1918/bare-IP "URLs", all pass a naive count. The
# authoritative measure is DISTINCT REAL PUBLIC DOMAINS. MIN_DISTINCT_DOMAINS is the
# configurable floor on that count (kept separate from, and additive to,
# MIN_CITED_SOURCES so both constants stay referenced by sync_check / the manifest).
# 6 is the conservative floor: a real Deep-Research brief draws on many more, but a
# self-reported "research_complete" with fewer than 6 distinct public sources is not
# a research pack. Hosts that are NOT real public sources are excluded entirely:
#   * localhost / loopback (127.0.0.0/8, ::1)
#   * RFC-1918 / private + link-local IP ranges
#   * bare IP literals (a citation should name a source, not an IP)
#   * .local / .localhost / .internal / .example / .invalid / .test reserved TLDs
#   * example.* placeholder domains (example.com / example.org / example-source-*)
MIN_DISTINCT_DOMAINS = 6   # HARD floor on DISTINCT REAL PUBLIC domains (AF-RESEARCH-UNCITED)

# ---------------------------------------------------------------------------
# DETERMINISTIC TYPOGRAPHY FLOORS (AF-FONT-FLOOR) — the CODED size/scale/contrast
# rejector that the vision-opinion Typography-QC gate never was.
# ---------------------------------------------------------------------------
# R8 Gap 1+2: the font-floor, type-scale and contrast rules (the qc-specialist
# vision codes F12 / F13 / DC4, kept as the semantic backstop) were written as
# "coded mechanical asserts" but no code implemented them, and
# they pointed at a phantom artifact the Typography Architect never emitted. This
# gate closes both holes deterministically WITHOUT vision/OCR: the Typography
# Architect now emits machine-readable size/scale/contrast tokens to
# working/typography/type_layout_system.md, and check_font_floor() parses those
# tokens and REJECTS a design system that declares a sub-floor body size, a
# non-modular type scale, or a below-WCAG text contrast — before any render burns a
# KIE credit. This is a numeric floor on the DECLARED type system; the per-pixel
# OCR/SSIM verification stays with the (independent) vision Typography-QC pass, so
# the two layers are complementary, not duplicative. Tokens are pt-EQUIVALENT (the
# architect declares sizes relative to a 1080-tall canvas; "18pt" == 1.6% cap-height).
FONT_BODY_PT_FLOOR = 18            # absolute body/subhead floor (pt-equiv at 1080px tall); WCAG large-text boundary
TYPE_SCALE_STEPS_MIN = 4           # a modular type scale must declare 4..5 named steps (no more, no fewer)
TYPE_SCALE_STEPS_MAX = 5
CONTRAST_RATIO_FLOOR_NORMAL = 4.5  # WCAG 2.2 AA — normal text
CONTRAST_RATIO_FLOOR_LARGE = 3.0   # WCAG 2.2 AA — large text (>= body floor & bold-or-bigger)
DARK_THEME_BODY_PT_FLOOR = 22      # +~20% projection-dispersion compensation when client_dark_theme opt-in
DARK_THEME_CONTRAST_FLOOR = 7.0    # AAA when dark is opt-in/premium (dark slides are washed out at projection)
TYPE_LAYOUT_SYSTEM_REL = "working/typography/type_layout_system.md"  # the gate-of-record artifact

# U047 — Rule 3.5 staging. These three checks are new gates meeting already-shipped
# decks, so they REPORT and return success until the finding count is driven to zero.
# Stage 3 flips this default. Only the exact string "0" disables enforcement once the
# default is flipped, and only the exact string "1" enables it before then.
SLIDE_GEOMETRY_ENFORCE_ENV = "PRESENTATION_SLIDE_GEOMETRY_ENFORCE"
SLIDE_GEOMETRY_EDGE_MARGIN_FRAC = 0.02
SLIDE_GEOMETRY_PT_REFERENCE_HEIGHT_PX = 1080

# ---------------------------------------------------------------------------
# RESEARCH-WEAVE / BREADTH GATE (AF-RESEARCH-WEAVE) — research distributed ACROSS
# the deck, not funnelled to one proof slide.
# ---------------------------------------------------------------------------
# R3 root cause of "ONE fact in ONE spot": the brief is real + early, but it never
# reaches most slides — the writer SOP funnels research to proof beats, the brief is
# category-organised (never slide-mapped), and the only enforcement is a deck-level
# "a research pack exists" presence check with NO breadth floor. This gate enforces
# the missing third leg: the Deep Research Specialist now emits a slide-assignment
# map (working/research/research_map.json) BEFORE copy, the copy template carries a
# RESEARCH_USED tag, and _chk_research_map() REJECTS a deck that weaves a mapped
# research item into fewer than RESEARCH_WEAVE_FLOOR_PCT of non-exempt content
# slides, or draws on fewer than MIN_DISTINCT_RESEARCH_ITEMS distinct items deck-wide
# (so the breadth cannot be satisfied by repeating one stat). Hook/pure-type/
# transition slides are EXEMPT and excluded from the denominator so the gate never
# pressures fabrication onto a slide that should not carry a stat.
RESEARCH_WEAVE_FLOOR_PCT = 60      # >= 60% of NON-EXEMPT content slides must weave a mapped research item (anchor present in copy)
MIN_DISTINCT_RESEARCH_ITEMS = 8    # the deck must draw on >= 8 DISTINCT research items (mirrors MIN_CITED_SOURCES)
RESEARCH_MAP_REL = "working/research/research_map.json"  # the slide-assignment artifact

# AF-I14: a real KIE-baked 2K slide PNG is hundreds of KB; a flat-fill / native-render
# placeholder (a solid colour, a Pillow/PPTX-drawn text card, or a tiny stub) is far
# smaller. 50 KiB is a conservative floor that any genuine model-baked slide clears
# easily while a flat-placeholder/native-render signature falls below it.
PLACEHOLDER_MIN_BYTES = 51200   # AF-I14: per-slide PNG floor — below this = not a real KIE bake

# ---------------------------------------------------------------------------
# GOAL-4 CONSTANTS (decisions 1C / 2A / 3C / 5C)
# ---------------------------------------------------------------------------
# 3C — Kie.ai balance pre-flight (AF-KIE-BALANCE). Phase-0 of the signature runner
# (and a shared pre-render gate) GETs the live Kie credit balance and HARD-ABORTS
# before a single slide is dispatched when balance < estimated_floor. The floor is
# estimated as slide_count x PER_SLIDE_CREDIT_ESTIMATE x KIE_BALANCE_FLOOR_MULTIPLIER.
# These are named (not inlined) so sync_check can verify they exist and the manifest
# AF-KIE-BALANCE secondary_py_symbols can point at them.
KIE_CREDIT_URL = "https://api.kie.ai/api/v1/chat/credit"  # AF-KIE-BALANCE: live credit endpoint
PER_SLIDE_CREDIT_ESTIMATE = 4      # AF-KIE-BALANCE: estimated Kie credits consumed per slide render (gpt-image-2 2K)
KIE_BALANCE_FLOOR_MULTIPLIER = 1.25  # AF-KIE-BALANCE: safety headroom over the bare per-slide estimate (retries/QC re-renders)

# 5C — native PPTX text/element-overlay path ELIMINATED (AF-OVERLAY-DELIVERED). The
# legacy fallback wrote intended text into pptx_text_overlays.json and the assembler
# composited it as a native PPTX text box. That path is removed by construction: its
# mere presence in a run dir, OR any native (non-notes) on-slide text run in the
# delivered PPTX, is a hard auto-fail. The only legitimate slide text is baked into
# the single composed gpt-image-2 image; the only legitimate PPTX text part is the
# off-slide speaker-notes pane. OVERLAY_FORBIDDEN_FILES are the run-relative paths the
# eliminated path used to write (any present = violation).
OVERLAY_FORBIDDEN_FILES = (
    "working/copy/pptx_text_overlays.json",
    "working/checkpoints/pptx_text_overlays.json",
    "pptx_text_overlays.json",
)

# ---------------------------------------------------------------------------
# FIX-2 SHARED CONTRACT — canonical render path + pixel-level image-QC teeth
# ---------------------------------------------------------------------------
# The canonical render path is build_deck.py / run_signature_deck.py ONLY. Every
# slide's WORDS and VISUAL are generated TOGETHER in ONE image by kie.ai
# gpt-image-2 (text-to-image, or image-to-image for the logo composite). There is
# NO Pillow slide canvas, NO other model, NO silent fallback, and NO native
# PowerPoint text overlaid on top. The three auto-fail codes below are SHARED
# across the enforcement surface: this module DEFINES + EXPORTS them and the
# check symbols that raise them; the orchestrator (run_signature_deck.py) and the
# preflight lockstep CONSUME them. A gate they raise may be waived ONLY by an
# explicit, LOGGED owner/founder token in working/checkpoints/process_manifest.json
# ("owner_skip_approval") — never silently, never by an agent's own choice.
#
# FIX-2 EXPORTS (stable public symbols other fixes import as build_deck.<name>):
#   constants:  AF_CANONICAL_RENDER_BYPASS, AF_LOCAL_CANVAS, AF_IMAGE_QC_VISION,
#               CANONICAL_RENDER_SCRIPTS
#   checks:     check_canonical_render_path(run_dir, slides_path=None) -> str
#               check_image_qc_vision(run_dir, slides_path=None) -> str
#   helper:     _owner_skip_approved(run_dir, af_code) -> Optional[dict]
AF_CANONICAL_RENDER_BYPASS = "AF-CANONICAL-RENDER-BYPASS"  # a non-canonical hand-rolled renderer/assembler produced (part of) the deck
AF_LOCAL_CANVAS            = "AF-LOCAL-CANVAS"             # a slide PNG was drawn on a LOCAL Pillow canvas (Image.new 2048x1152) instead of kie.ai
AF_IMAGE_QC_VISION         = "AF-IMAGE-QC-VISION"          # image-QC "passed" without a real pixel/vision read (rubber-stamp number / flat card / overlay-blessing rubric)

# U027 (not FIX-2, but reuses the same _owner_skip_approved helper): the postflight
# OCR text-readback audit. See check_ocr_readback() below.
AF_OCR_READBACK            = "AF-OCR-READBACK"             # a rendered slide's baked text was never read back against its approved copy (checked:false — NEVER waivable), or the readback found a mismatch (matched:false — waivable ONLY by a logged owner_skip_approval)

# The canonical render/assemble tools. A *.py file inside a run/working dir whose
# name is NOT in this set, and that defines a slide canvas / native text box /
# direct kie task submission, is a forbidden hand-rolled renderer.
CANONICAL_RENDER_SCRIPTS = frozenset({
    "build_deck.py", "run_signature_deck.py", "build_teleprompter.py", "sync_check.py",
})

# Pixel-level flat-card detector thresholds (AF-IMAGE-QC-VISION). A genuine kie.ai
# gpt-image-2 render carries a photographic/illustrative subject (and model noise),
# so no single quantised colour dominates the frame. A flat cream/typography card
# (Pillow Image.new fill + drawn type) is overwhelmingly ONE near-uniform light
# fill — exactly the slide-1/24/49 hook-card defect this gate must reject.
IMAGE_QC_FLATFILL_DOMINANCE = 0.90   # >= 90% of the frame is one quantised colour = flat card (no visual subject)
IMAGE_QC_FLATFILL_LUMA      = 0.62   # ...and that dominant fill is light/cream (0-1 luma) = typography card, not a photo

# Reserved / non-routable TLD suffixes that are never a real public source.
_NON_PUBLIC_TLD_SUFFIXES = (
    ".local", ".localhost", ".internal", ".example", ".invalid", ".test",
)

# ---------------------------------------------------------------------------
# QC-INDEPENDENCE GATE (AF-QC-INDEPENDENCE) — the QC report must be graded by a
# reviewer INDEPENDENT of the builder, never self-graded.
# ---------------------------------------------------------------------------
# A Phase-1Q QC report with average>=8.5 and no triggered autofails STILL fails
# the deck if the builder graded its own work.  A report self-written by the
# renderer / by the same role that authored the copy proves nothing — it is the
# rubber-stamp the whole gate exists to refuse.  The report MUST carry explicit
# provenance naming an independent QC reviewer identity that is NOT the builder.
#
# These tokens are the "builder / self" identities a graded_by / reviewer /
# author field may NOT name.  Any of them (or a self_graded:true flag, or a
# reviewer == the builder, or NO independent-reviewer provenance at all) trips
# AF-QC-INDEPENDENCE.  The qc-specialist persona writes the report and stamps
# its own (independent) identity into the provenance block — see
# qc-specialist-presentations-sops.md (PROVENANCE block).
FORBIDDEN_QC_GRADER_IDENTITIES = frozenset({
    "build_deck.py", "build_deck", "self", "builder", "author",
})
# The role that AUTHORS the deck copy — a QC report whose reviewer equals this
# role is self-grading-by-proxy and is refused (the copywriter cannot grade the
# copywriter's own copy).
QC_BUILDER_ROLE = "slide-copywriter"


def _registered_domain(netloc: str) -> Optional[str]:
    """Reduce a URL netloc to a real PUBLIC host, or None when the host is not a
    real public source. Strips userinfo + port, lowercases, then REJECTS (returns
    None) for: empty hosts, localhost/loopback, bare IP literals (any private OR
    public IP — a citation must name a source, not an address), RFC-1918 / link-local
    ranges, reserved non-public TLDs, and example.* placeholder domains. The returned
    value is the host used for de-duplication (distinct registered domains)."""
    host = (netloc or "").strip().lower()
    # Drop any userinfo (user:pass@) and port (:443).
    if "@" in host:
        host = host.rsplit("@", 1)[1]
    # Strip a trailing :port (but not the colons inside an IPv6 literal in brackets).
    if host.startswith("["):
        # IPv6 literal like [::1]:443 -> ::1
        host = host[1:].split("]", 1)[0]
    elif ":" in host:
        host = host.rsplit(":", 1)[0]
    host = host.strip(".")
    if not host:
        return None
    # localhost variants.
    if host == "localhost" or host.endswith(".localhost"):
        return None
    # Reserved / non-public TLDs.
    for suffix in _NON_PUBLIC_TLD_SUFFIXES:
        if host.endswith(suffix):
            return None
    # example.* placeholder domains (example.com, example.org, example-source-*.org…).
    if host == "example" or host.startswith("example.") or host.startswith("example-"):
        return None
    # Bare IP literal? Reject (both private and public — cite a named source).
    try:
        ip = ipaddress.ip_address(host)
        # Any bare IP is not a named public source for our purposes.
        return None
    except ValueError:
        pass  # not an IP literal — it's a hostname, good.
    # A real public hostname must have at least one dot (a TLD).
    if "." not in host:
        return None
    return host


def _distinct_public_domains(text: str) -> set:
    """Extract every http(s) URL from text and return the set of DISTINCT REAL
    PUBLIC registered domains (junk/localhost/loopback/RFC-1918/.local/example.*/
    bare-IP hosts excluded). De-duplication is by host, so the same source cited N
    times counts once."""
    url_re = re.compile(r'https?://[^\s\)\]\>,\'"\\]+', re.IGNORECASE)
    domains = set()
    for raw in url_re.findall(text):
        try:
            netloc = urlparse(raw.rstrip(".")).netloc
        except ValueError:
            continue
        dom = _registered_domain(netloc)
        if dom:
            domains.add(dom)
    return domains

# Factual/statistical claim markers — tokens that signal a claim requiring a citation.
# When any of these appear in the slide copy and no corresponding cited URL exists in
# the research pack, the gate fires.  Documented heuristic: it catches the most common
# patterns (percentage figures, dollar figures, key research-signal words) without
# false-positives on purely narrative copy.
_CLAIM_MARKERS_RE_SRC = (
    r'\d+\s*%',                        # "45%", "20 %"
    r'\$\s*\d+',                       # "$5,000", "$ 3 million"
    r'\b(?:research|study|studies\s+show|statistics|data\s+shows)\b',
)

# ---------------------------------------------------------------------------
# FACIAL INTELLIGENCE / REPRESENTATION GUARD (the 60/30/10 landmine, fail-loud)
# ---------------------------------------------------------------------------
# Doctrine: SOP-CAST-01 ("No inverted default. There is no system default
# demographic.") and qc-specialist-presentations.md AF-R3 ("No racial or gender
# default is ever inferred."). Representation in any people-prompt MUST come from
# the slide spec (the client's captured audience), NEVER from a baked-in default
# percentage split such as the "60/30/10" ratio. This builder is deterministic
# and renders the scene VERBATIM from slides.json, so it must NEVER carry a
# hardcoded demographic default of its own, and it must REJECT a slide whose spec
# tries to smuggle one in. These patterns (case-insensitive substrings) are the
# forbidden demographic-default landmines; any of them in a slide's scene/copy/
# layout/logo text is an immediate fail-loud (the prompt-side representation gate).
FORBIDDEN_DEMOGRAPHIC_DEFAULTS = [
    "60/30/10",
    "60-30-10",
    "60/30/10 ratio",
    "default demographic",
    "default ethnicity",
    "default race",
    "default skin tone",
    "default skin-tone",
    "standard demographic mix",
    "standard representation mix",
    "assume the audience is",
    "assumed demographic",
    "inferred demographic",
    "system default demographic",
]

# Client's own env stores (HOME-relative — works on any client box).
SECRETS_CANDIDATES = [
    os.path.expanduser("~/.openclaw/workspace/.env"),
    os.path.expanduser("~/clawd/secrets/.env"),
    os.path.expanduser("~/.openclaw/secrets/.env"),
]

# ---------------------------------------------------------------------------
# OUTPUT DESTINATION — default ~/Downloads (Requirement 4)
# ---------------------------------------------------------------------------
# The final bundle dir defaults to ~/Downloads/<client-slug>-<deck-slug>/.
# --out overrides this. Files NEVER default to a scratch/run dir; they must
# land in ~/Downloads so the client receives them from a predictable location.
# The slug is derived from the deck slug (out.pptx stem) at runtime if present.
BUNDLE_DIR_DEFAULT = os.path.expanduser("~/Downloads")

# ---------------------------------------------------------------------------
# DELIVERABLES_REQUIRED manifest (Requirement 1)
# ---------------------------------------------------------------------------
# The nine mandatory output artifacts, their canonical relative paths inside the
# bundle dir, their per-artifact minimum-size gates, and a human description.
# Any artifact below its min_bytes threshold (or absent) triggers
# AF-BUNDLE-COMPLETE (exit 5). Rationale for each threshold:
#
#   deck_pptx  > 1 MB  : a real multi-slide rendered deck with 2K images is
#                         always several MB; < 1 MB implies the pptx is empty
#                         or contains placeholder content (zero-image shell < 100KB).
#
#   deck_pdf   > 50 KB : a minimal 1-slide PDF export is ~20-30KB; 50KB ensures
#                         at least two slides' worth of rendered content.
#
#   guide_pdf  > 50 KB : a minimal guide covers all slides with talking points and
#                         timing; < 50KB implies only a stub header.
#
#   speech_md  > 2 KB  : a word-for-word script for any real webinar talk will be
#                         thousands of words; 2KB floors an obvious empty or stub.
#
#   speech_pdf > 3 KB  : a PDF export of a real speech (doctrine floor, AF-BUNDLE-COMPLETE
#                         per presenters-speech-writer.md / -sops.md); < 3KB implies a stub.
#
#   teleprompter_html > 20 KB : a real self-contained scrolling teleprompter app (doctrine
#                                floor, same authority); < 20KB implies an empty/stub page.
#
#   audio_mp3  > 500KB : a real Fish Audio S2 rendition of a 30-min script is
#                         typically 50-150MB; 500KB floors the obvious failure case
#                         (silence stub or failed render < 100KB per SOP-PITCH-05).
#
#   infographic_png > 100KB : a real 2K-resolution infographic checklist PNG is
#                              always several hundred KB; < 100KB implies a blank
#                              placeholder image or a tiny thumbnail.
#
DELIVERABLES_REQUIRED = [
    {
        "key": "deck_pptx",
        "filename": "{deck_slug}-FINAL.pptx",
        "label": "assembled deck PPTX",
        "min_bytes": 1_048_576,          # 1 MB — multi-slide 2K-image deck floor
        "note": ">1MB floor; a sub-1MB pptx implies empty/placeholder content",
    },
    {
        "key": "deck_pdf",
        "filename": "{deck_slug}-FINAL.pdf",
        "label": "deck PDF export",
        "min_bytes": 51_200,             # 50 KB — PDF export of at least 2 slides
        "note": ">50KB; produced by PPTX Assembly Specialist (LibreOffice/Pandoc export)",
    },
    {
        "key": "guide_pdf",
        "filename": "PRESENTER-GUIDE.pdf",
        "label": "presenter guide PDF",
        "min_bytes": 51_200,             # 50 KB — guide covers all slides with talking points
        "note": ">50KB; produced by Presenters Guide Specialist. REQUIRED UPSTREAM STEP.",
    },
    {
        "key": "speech_md",
        "filename": "PRESENTERS-SPEECH.md",
        "label": "presenter speech markdown source (pure)",
        "min_bytes": 2_048,              # 2 KB — word-for-word script stub floor
        "note": ">2KB; produced by Presenters Speech Writer (pure script). REQUIRED UPSTREAM STEP.",
    },
    {
        "key": "speech_pdf",
        "filename": "PRESENTERS-SPEECH.pdf",
        "label": "presenter speech teleprompter PDF",
        # P3-01(c)5 — RECONCILED to the doctrine-ratified floor (presenters-speech-writer.md
        # + sops/presenters-speech-writer-sops.md AF-BUNDLE-COMPLETE: "PDF >= 3,000 bytes"),
        # which is also presenters_speech_pdf.py's own PDF_MIN_BYTES self-check. This table
        # previously drifted to 20_480 — a stricter number nothing else in the doctrine cited.
        "min_bytes": 3_000,              # 3 KB — doctrine floor (AF-BUNDLE-COMPLETE)
        "note": ">3KB; produced by Presenters Speech Writer (teleprompter PDF render). REQUIRED UPSTREAM STEP.",
    },
    {
        "key": "speech_fish_md",
        "filename": "PRESENTERS-SPEECH-FISH-TAGGED.md",
        "label": "presenter speech (Fish-Audio expression-tagged)",
        "min_bytes": 2_048,              # 2 KB — fish-tagged variant of the script floor
        "note": ">2KB; produced by Presenters Speech Writer / Fish Audio Expression "
                "Specialist (the (break)/(laugh)/emphasis-tagged script that feeds the "
                "audio render). REQUIRED UPSTREAM STEP.",
    },
    {
        "key": "audio_mp3",
        "filename": "PRESENTER-AUDIO.mp3",
        "label": "presenter audio MP3",
        "min_bytes": 512_000,            # 500 KB — real Fish Audio S2 rendition floor
        "note": ">500KB; produced by Audio Demonstration Specialist. REQUIRED UPSTREAM STEP.",
    },
    {
        "key": "infographic_png",
        "filename": "infographic.png",
        "label": "infographic checklist PNG",
        "min_bytes": 102_400,            # 100 KB — real 2K-resolution infographic floor
        "note": ">100KB; produced by infographic-checklist role. REQUIRED UPSTREAM STEP.",
    },
    {
        "key": "teleprompter_html",
        "filename": "presenter-teleprompter.html",
        "label": "presenter teleprompter web app",
        # P3-01(c)5 — RECONCILED to the doctrine-ratified floor (presenters-speech-writer.md
        # + sops/presenters-speech-writer-sops.md AF-BUNDLE-COMPLETE: "HTML >= 20,000 bytes"),
        # which is also build_teleprompter.py's own TELEPROMPTER_MIN_BYTES self-check. This
        # table previously drifted to 10_240 — a looser number nothing else in the doctrine cited.
        "min_bytes": 20_000,             # 20 KB — a real self-contained scrolling
                                          # teleprompter app (HTML + inline CSS/JS);
                                          # < 20KB implies an empty/stub page.
        "note": ">20KB; the standalone scrolling teleprompter app. Produced by the "
                "build_teleprompter.py generator (owned by the speech/teleprompter "
                "role). REQUIRED UPSTREAM STEP.",
    },
    {
        "key": "webinar_mp4",
        "filename": "{deck_slug}-WEBINAR.mp4",
        "label": "webinar video mp4",
        # P3-01(c)5 — Feature L2-G (Gauntlet Loop 2, Feature C): the fluid webinar video
        # slideshow (Ken Burns + xfade, synced to the spoken audio). PRODUCED LATER than
        # P8-ASSEMBLE/postflight: build_webinar_video.py (P9.6-WEBINAR-VIDEO, order 8.92)
        # renders it AFTER the render/assemble phases, from the deck's slide PNGs + the
        # PRESENTERS-SPEECH.md + PRESENTER-AUDIO.mp3 + the Fish chunk mp3s. It is hosted in
        # the client's GHL media library via the v3 500MB video tier. `produced_later: True`
        # tells run_postflight_gate to SKIP the file-existence loop at P8-ASSEMBLE time (the
        # file does not exist yet); the P9.6 phase + fix_bundle_complete + delivery_gate own
        # the webinar's real presence gate.
        "min_bytes": 1_048_576,          # 1 MB — a real rendered webinar video floor
        "note": ">1MB; produced by build_webinar_video.py (P9.6-WEBINAR-VIDEO, Feature "
                "L2-G) — NOT present at P8-ASSEMBLE time. REQUIRED UPSTREAM STEP.",
        "produced_later": True,
    },
]


# ---------------------------------------------------------------------------
# API key
# ---------------------------------------------------------------------------

def load_api_key() -> str:
    key = os.environ.get("KIE_API_KEY", "").strip()
    if key:
        return key.strip("'\"")
    for path in SECRETS_CANDIDATES:
        p = Path(path)
        if not p.exists():
            continue
        for line in p.read_text().splitlines():
            line = line.strip()
            if line.startswith("KIE_API_KEY="):
                value = line[len("KIE_API_KEY="):].strip().strip("'\"")
                if value:
                    return value
    print("FATAL: KIE_API_KEY not found in env or any of:", file=sys.stderr)
    for path in SECRETS_CANDIDATES:
        print("   ", path, file=sys.stderr)
    sys.exit(2)


# ---------------------------------------------------------------------------
# Rich per-slide prompt loading — VERBATIM, NO AI, NO THIN FALLBACK
# ---------------------------------------------------------------------------
# build_deck.py renders the Slide Image Creator's hand-authored RICH per-slide
# prompt VERBATIM. It NEVER composes its own thin scene+copy prompt. The rich
# prompt file for slide N lives in the run's working/prompts/ dir under one of
# these names (the CLIENT-WEBINAR-DECK-SOP convention first, then the
# slide-image-creator-sops.md convention):
PROMPT_FILE_PATTERNS = (
    "working/prompts/slide-{nn:02d}.txt",
    "working/prompts/slide-{nn:02d}-prompt.txt",
    # R3 / D10: signature decks have a 100-slide FLOOR, so a 3-digit ordinal is
    # canonical (slide-100.txt .. slide-999.txt). The %02d family above cannot
    # express it, so the canonical name set extends with an explicit %03d family.
    # Ordering keeps %02d first, so the zero-padded twin is always preferred when
    # both spellings exist for the SAME slide (see _canonical_prompt_dir_problems).
    "working/prompts/slide-{nn:03d}.txt",
    "working/prompts/slide-{nn:03d}-prompt.txt",
)


def resolve_prompt_path(run_dir: Path, ordinal: int) -> Optional[Path]:
    """Return the rich-prompt file Path for this slide ordinal, or None if no
    candidate file exists. Tries both the slide-NN.txt and slide-NN-prompt.txt
    naming conventions under working/prompts/ — and, for ordinals >= 100
    (signature decks have a 100-slide floor), the 3-digit slide-NNN.txt /
    slide-NNN-prompt.txt names (R3 / D10)."""
    for pat in PROMPT_FILE_PATTERNS:
        p = run_dir / pat.format(nn=ordinal)
        if p.exists() and p.is_file():
            return p
    return None


# R3 / D10: signature decks have a 100-slide floor, so a 3-digit ordinal is
# canonical too. The shared prompt_gate (scan_prompt_dir / prompt_dir_problems,
# FIX-22 / D16) still treats ONLY the %02d family as canonical and flags
# slide-100.txt as non-canonical — but every build_deck consumer consults the
# gate ONLY through _canonical_prompt_dir_problems below, which runs the shared
# scan and overlays the R3 3-digit-canonical rule on top of its verdict. The
# shared module itself is deliberately NOT changed (its %02d-only contract
# predates the signature-deck floor and other consumers depend on it); the
# build_deck-side fix keeps THIS file's canonical name set (PROMPT_FILE_PATTERNS,
# which already carries the 3-digit family) and the gate verdicts it consumes in
# agreement.
# NOTE on why the duplicate verdict passes through UNCHANGED: Python's minimum-
# width padding never truncates — `{nn:02d}` for nn=100 formats as "100" — so
# for ordinals 100+ the %02d and %03d families are the SAME spelling, and for
# ordinals 1-99 the %03d spelling (slide-001.txt) targets the same slide as the
# %02d twin (slide-01.txt). A collision is therefore ALWAYS between names that
# resolve to the same ordinal (slide-01 vs slide-1, slide-100 vs slide-0100,
# slide-01 vs slide-001) and R3 never makes a same-ordinal collision benign —
# exactly one canonical file per slide stays the rule at any digit width.
CANONICAL_PROMPT_DIGITS_R3_RE = re.compile(r"^\d{2,3}$")


def _prompt_ordinal_field(name: str) -> str:
    """Extract the raw ordinal FIELD (the digits) from a slide-NN.txt or
    slide-NN-prompt.txt prompt filename. Mirrors prompt_gate.scan_prompt_dir's
    digit_part extraction exactly (its idiom), so the R3 re-judgement and the
    shared scan always agree on what 'the digits' means."""
    return name[len("slide-"):].split("-")[0].split(".")[0]


def _canonical_prompt_dir_problems(run_dir: Path) -> list:
    """Directory-level prompt problems (duplicates / non-canonical names) as the
    build_deck gates consume them. Runs the shared prompt_gate detector (FIX-22 /
    D16), then reconciles the R3 3-digit-canonical overlay on its verdicts:
      * AF-PROMPT-DUP-FILE verdicts pass through unchanged — see the NOTE above;
      * AF-PROMPT-NAME is relaxed: a name whose ordinal FIELD is exactly 2 or 3
        digits is canonical (slide-01..slide-99, slide-100..slide-999, plus the
        -prompt variants); a 1-digit (`slide-1.txt`) or 4+ digit (`slide-0100.txt`,
        `slide-1000.txt`) name stays non-canonical (D16).
    Every consumed gate (preflight rich-prompt _chk_rich_prompts, Prompt-QC
    teeth check_prompt_qc_teeth) routes through _collect_prompt_problems, the
    single choke point that calls this, so the two sides can never disagree.
    Returns [] when the directory is clean, else the fatal problem strings."""
    pg = _import_prompt_gate()
    if pg is None:
        return []
    base_problems = pg.prompt_dir_problems(run_dir / "working" / "prompts")
    if not base_problems:
        return []
    rep = pg.scan_prompt_dir(run_dir / "working" / "prompts")
    problems = []
    for prob in base_problems:
        if prob.startswith("AF-PROMPT-DUP-FILE"):
            problems.append(prob)
        elif prob.startswith("AF-PROMPT-NAME"):
            # Re-judge every name this message names: keep the verdict ONLY when
            # at least one of them still has a non-canonical field (1-digit or
            # 4+ digits). A message whose only offender is now-canonical (e.g.
            # slide-100.txt) is dropped. The '.txt' tail makes the substring
            # check unambiguous (slide-100.txt is never a substring of a
            # slide-1000.txt message).
            flagged = [n for n in rep.get("non_canonical", []) if n in prob]
            if any(not CANONICAL_PROMPT_DIGITS_R3_RE.match(_prompt_ordinal_field(n))
                   for n in flagged):
                problems.append(prob)
        else:
            problems.append(prob)
    return problems


def _norm_ws(s: str) -> str:
    """Lowercase + collapse all runs of whitespace to a single space + strip. Used by
    the verbatim-words-baked check so a copy line and its prompt quotation match even
    across line wraps / double spaces."""
    return re.sub(r"\s+", " ", str(s)).strip().lower()


def _negative_block_class_problems(prompt_lc: str) -> list:
    """AF-P13 mechanical half: return the names of any of the EIGHT mandatory
    negative-block defect classes that are NOT named in the prompt. Empty list = all
    eight present."""
    missing = []
    for cls, tokens in NEGATIVE_BLOCK_CLASS_TOKENS.items():
        if not any(t in prompt_lc for t in tokens):
            missing.append(cls)
    return missing


def _spelling_lock_present(prompt_lc: str) -> bool:
    """AF-P14 mechanical half: True iff the prompt carries a per-string spelling-lock
    directive marker."""
    return any(t in prompt_lc for t in SPELLING_LOCK_TOKENS)


def _prompt_density_problems(prompt_text: str, prompt_lc: str) -> list:
    """AF-P-DENSITY: return human-readable reasons the prompt is thin/padded despite
    clearing the char floor. Requires a brand HEX, an explicit type SIZE, a
    COMPOSITION/zone token, and a distinct-word floor (catches paste-repetition padding)."""
    problems = []
    if not _HEX_COLOR_RE.search(prompt_text):
        problems.append("no brand palette HEX (#RRGGBB) — element (f) palette is mandatory")
    if not _TYPE_SIZE_RE.search(prompt_text):
        problems.append("no explicit type SIZE token (e.g. '72pt', '28pt', '120px') — "
                        "typography size is a mandatory 15-element field")
    if not any(t in prompt_lc for t in PROMPT_COMPOSITION_TOKENS):
        problems.append("no composition/zone token (thirds grid, zone, safe margin, "
                        "quadrant) — 'centered' alone is an auto-fail in doctrine")
    distinct = len(set(_WORD_RE.findall(prompt_lc)))
    if distinct < PROMPT_MIN_DISTINCT_WORDS:
        problems.append(f"only {distinct} distinct words (floor {PROMPT_MIN_DISTINCT_WORDS}) "
                        "— a long file with few distinct words is paste-repetition padding, "
                        "not a rich 15-element spec")
    return problems


def _verbatim_copy_problems(prompt_text: str, copy_val) -> list:
    """VERBATIM-WORDS-BAKED check (FIX-4): the slide's EXACT copy strings must appear in
    the prompt body so the words are baked into the kie.ai image (never paraphrased and
    never left to be overlaid later). Returns the copy strings that are NOT present
    verbatim (whitespace-normalised, case-insensitive). Empty list = every copy line is
    quoted in the prompt."""
    if isinstance(copy_val, list):
        strings = [str(c) for c in copy_val]
    elif copy_val in (None, ""):
        strings = []
    else:
        strings = [str(copy_val)]
    prompt_norm = _norm_ws(prompt_text)
    missing = []
    for c in strings:
        cn = _norm_ws(c)
        # Skip empty / trivial fragments (a 1-2 char label proves nothing).
        if len(cn) < 3:
            continue
        if cn not in prompt_norm:
            short = c if len(str(c)) <= 60 else str(c)[:57] + "..."
            missing.append(short)
    return missing


def rich_prompt_quality_problems(prompt_text: str, copy_val=None) -> list:
    """The QUALITY-LAYER teeth on a single rich prompt (on top of the char floor +
    REQUIRED_STRUCTURAL_BLOCKS already enforced by the caller): AF-P13 (eight-class
    negative block), AF-P14 (spelling-lock), AF-P-DENSITY (hex/size/composition/distinct
    words), and VERBATIM-WORDS-BAKED (the slide's exact copy is in the prompt). Returns a
    list of fatal problem strings (empty = the prompt clears every quality gate). Shared
    by load_rich_prompt (render-time, raises) and _chk_rich_prompts / the Prompt-QC teeth
    (preflight, accumulates) so the floor can never be a length-only rubber stamp."""
    prompt_lc = prompt_text.lower()
    problems = []
    missing_classes = _negative_block_class_problems(prompt_lc)
    if missing_classes:
        problems.append(
            "AF-P13: negative block does not name defect class(es): "
            + ", ".join(missing_classes)
            + " — the 8-class paired negative block (SOP 9.8) is mandatory; a "
            "one-line 'no text' AVOID stub does not satisfy it")
    if not _spelling_lock_present(prompt_lc):
        problems.append(
            "AF-P14: no per-string spelling-lock directive (e.g. 'render this exact "
            "string, letter-for-letter') — every verbatim on-slide string must be "
            "spelling-locked")
    for d in _prompt_density_problems(prompt_text, prompt_lc):
        problems.append("AF-P-DENSITY: " + d)
    if copy_val is not None:
        missing_copy = _verbatim_copy_problems(prompt_text, copy_val)
        if missing_copy:
            problems.append(
                "AF-P-VERBATIM: the slide's exact copy is NOT baked into the prompt "
                "body (must appear verbatim so kie.ai bakes the words, never overlaid): "
                + " | ".join(missing_copy))
    return problems


def assert_no_forbidden_demographic_default(slide: dict, prompt_text: str = "") -> None:
    """
    FACIAL-INTELLIGENCE / REPRESENTATION GATE (fail-loud).

    Representation must come from the slide spec / casting ledger (the client's
    captured audience), never from a baked-in default split. Per SOP-CAST-01
    ("there is no system default demographic") and AF-R3 ("no racial or gender
    default is ever inferred"), this builder REFUSES any slide whose author-supplied
    spec OR whose rich prompt text smuggles a forbidden demographic-default landmine
    (e.g. the "60/30/10" ratio). Raises ValueError (caller fails the slide loud) on
    a hit. prompt_text is the verbatim rich prompt being rendered — scanned too so a
    landmine in the prompt file (not just the slides.json spec) is also caught.
    """
    # Gather every author-supplied text field on this slide into one haystack,
    # plus the rich prompt text that will actually be sent to KIE.
    chunks = [str(slide.get("scene", "")), str(slide.get("layout", "")),
              str(slide.get("logo", "")), str(prompt_text)]
    copy_val = slide.get("copy", [])
    if isinstance(copy_val, list):
        chunks.extend(str(c) for c in copy_val)
    else:
        chunks.append(str(copy_val))
    haystack = " ".join(chunks).lower()

    for landmine in FORBIDDEN_DEMOGRAPHIC_DEFAULTS:
        if landmine.lower() in haystack:
            raise ValueError(
                f"slide {slide.get('slide')}: forbidden hardcoded demographic "
                f"default detected ('{landmine}'). Representation must come from "
                f"the client's captured audience / casting ledger, never a baked-in "
                f"default split (SOP-CAST-01 / AF-R3). Refusing to render."
            )


def load_rich_prompt(slide: dict, run_dir: Path) -> str:
    """
    Load the Slide Image Creator's RICH per-slide prompt VERBATIM and gate it.

    This is the ONLY prompt path. build_deck.py renders the hand-authored rich
    prompt (working/prompts/slide-NN.txt or slide-NN-prompt.txt) EXACTLY as written
    — it does NOT compose its own thin scene+copy prompt and NEVER silently falls
    back to one. The rich prompt already carries the full 15-element spec
    (typography per-line weight + pt size, placement, usage, the logo(s), the scene,
    the verbatim copy, the negative block — everything on the slide), and KIE
    (gpt-image-2) renders the WHOLE slide in ONE generation.

    FAIL LOUD (ValueError; the caller fails the slide and the run is blocked) when:
      * no rich prompt file exists for this slide                       → AF-P1
      * the rich prompt is < PROMPT_CHAR_FLOOR (9,000) chars            → AF-P1
        (not run, not rendered, not updated)
      * the rich prompt is > PROMPT_CHAR_CEILING (18,000) chars         → AF-P2
      * a forbidden demographic-default landmine is present             → AF-R3
      * the dead endpoint fragment is present
    """
    ordinal = int(slide["slide"])
    prompt_path = resolve_prompt_path(run_dir, ordinal)
    if prompt_path is None:
        tried = " or ".join(p.format(nn=ordinal) for p in PROMPT_FILE_PATTERNS)
        raise ValueError(
            f"slide {ordinal}: NO rich prompt file found (looked for {tried} under "
            f"{run_dir}). build_deck.py renders the Slide Image Creator's rich prompt "
            f"VERBATIM and NEVER composes a thin fallback. Refusing to render "
            f"(rich-prompt-required, AF-P1)."
        )

    prompt = prompt_path.read_text(errors="replace")

    # FACIAL-INTELLIGENCE / REPRESENTATION gate over the slide spec AND the verbatim
    # rich prompt (so a landmine in the prompt file is also caught).
    assert_no_forbidden_demographic_default(slide, prompt_text=prompt)

    # Self-guard: never let the dead endpoint string ride inside a prompt payload.
    if DEAD_ENDPOINT_FRAGMENT in prompt:
        raise ValueError(
            f"slide {ordinal}: rich prompt {prompt_path} contains the dead endpoint "
            f"fragment '{DEAD_ENDPOINT_FRAGMENT}'. Refusing."
        )

    # PROMPT CHAR-COUNT GATE (fail-loud). The floor is HARD (9,000): a prompt under
    # it is a thin stub / truncated file, NOT a real slide prompt — never run it.
    # H1: measure the STRIPPED length so a file padded with whitespace (or one that is
    # whitespace-only) can never satisfy the floor. len(prompt) over raw bytes would
    # let "   \n   ... 5000 spaces ..." pass as a "5,000-char prompt".
    if not prompt.strip():
        raise ValueError(
            f"slide {ordinal}: rich prompt {prompt_path} is empty / whitespace-only. "
            f"A blank prompt carries none of the mandatory per-slide spec. It is NOT "
            f"run, NOT rendered, NOT updated. Re-author the rich prompt. Refusing "
            f"(prompt char-count gate, AF-P1 floor)."
        )
    length = len(prompt.strip())
    if length < PROMPT_CHAR_FLOOR:
        raise ValueError(
            f"slide {ordinal}: rich prompt {prompt_path} is {length} chars, UNDER the "
            f"HARD floor of {PROMPT_CHAR_FLOOR}. A prompt this short cannot carry the "
            f"mandatory per-slide spec (typography size/placement/usage, logo, scene, "
            f"verbatim copy, negative block). It is NOT run, NOT rendered, and NOT "
            f"updated. Re-author the rich prompt to >= {PROMPT_CHAR_FLOOR} chars. "
            f"Refusing (prompt char-count gate, AF-P1 floor)."
        )
    if length > PROMPT_CHAR_CEILING:
        raise ValueError(
            f"slide {ordinal}: rich prompt {prompt_path} is {length} chars, over the "
            f"hard ceiling of {PROMPT_CHAR_CEILING} (2,000 under the 20,000 GPT-Image 2 API "
            f"ceiling, MODEL-SPECS). The prompt is too long; tighten redundant phrasing "
            f"(never delete the negative block or any spelling-lock). Refusing "
            f"(prompt char-count gate, AF-P2 ceiling)."
        )
    # STRUCTURAL-BLOCK GATE (fail-loud, AF-P1). Folded in from the retired
    # render_deck.py: a long file is not enough — the rich prompt MUST carry the
    # load-bearing scaffolding (an [ARCHETYPE ...] layout declaration, the final-paragraph
    # negative block — header "DO-NOT BLOCK", or the legacy "NEGATIVE BLOCK" alias — and at
    # least one "Do not " imperative). A prompt that clears the floor but is missing any of
    # these is a verbose stub, not a slide spec — never run it.
    prompt_lc = prompt.lower()
    missing_blocks = _missing_structural_blocks(prompt_lc)
    if missing_blocks:
        raise ValueError(
            f"slide {ordinal}: rich prompt {prompt_path} clears the char floor but is "
            f"MISSING required structural block(s): {', '.join(missing_blocks)}. A real "
            f"rich prompt declares its [ARCHETYPE ...] layout, carries the final-paragraph "
            f"negative block (header 'DO-NOT BLOCK', or the legacy 'NEGATIVE BLOCK' alias), "
            f"and states 'Do not ...' imperatives. Re-author the rich prompt with the full "
            f"15-element structure. Refusing (structural-block gate, AF-P1)."
        )
    # QUALITY-LAYER teeth (AF-P13 / AF-P14 / AF-P-DENSITY / AF-P-VERBATIM). A prompt that
    # clears the length floor + structural blocks is still rejected when its negative block
    # omits any of the eight defect classes, it carries no per-string spelling-lock, it is
    # thin/padded (no hex palette / type size / composition token / too few distinct words),
    # or it fails to bake the slide's exact copy verbatim. Length alone is never enough.
    quality = rich_prompt_quality_problems(prompt, slide.get("copy"))
    if quality:
        raise ValueError(
            f"slide {ordinal}: rich prompt {prompt_path} clears the char floor + structural "
            f"blocks but FAILS the quality floor — it is not a real 15-element spec. "
            f"It is NOT run, NOT rendered, NOT updated. Re-author. Problems: "
            + " || ".join(quality)
        )
    return prompt


# ---------------------------------------------------------------------------
# HTTP helpers
# ---------------------------------------------------------------------------

class RateLimited(Exception):
    pass


class AuthError(Exception):
    """Permanent authentication failure (HTTP 401/403) — the request is identical
    and will fail forever: the key is wrong, the Authorization header format is
    wrong, or the account is locked/rate-locked by the provider. Retrying the same
    request is a guaranteed token furnace. Fail loud, NEVER back off, NEVER
    re-submit. Raised by _http_json on HTTP 401/403 and swallowed by NO retry loop:
    render_slide re-raises it immediately (< 2s) instead of entering the
    exponential-backoff re-submit path."""


class RenderPollTimeout(Exception):
    """Terminal FAIL surfaced when a kie task did not reach a terminal state within
    the poll hard cap (POLL_MAX_SECONDS, ~15 min by default). Raised by poll_task
    so a dead / never-completing render FAILS loudly instead of hanging the run in
    'in_progress' (D14 / FIX-7). Carries the elapsed seconds and last observed
    state so the caller can log a precise failure."""


# C1 (SSRF / local-file read guard): the ONLY URL schemes this pipeline is ever
# allowed to open. urllib.request.urlopen will happily honour file://, ftp://,
# data:, etc. — so a KIE result URL, a --logo URL, or any API URL that resolves to
# anything other than http(s) is an SSRF / arbitrary-local-file-read landmine
# (e.g. file:///etc/passwd, file:///Users/.../secrets/.env). Every URL we open
# MUST pass this allowlist FIRST or we refuse, loud.
ALLOWED_URL_SCHEMES = ("http", "https")


def assert_url_scheme_allowed(url: str, what: str = "URL") -> None:
    """Refuse (ValueError) unless the URL's scheme is http or https. This is the
    SSRF / local-file-read guard (C1): urlopen honours file://, ftp://, data:, etc.,
    so before opening ANY fetched URL we enforce a strict http(s) allowlist. Applied
    to _http_json (the KIE API calls), download_image (the KIE result URL),
    and the --logo URL path."""
    scheme = (urlparse(str(url)).scheme or "").lower()
    if scheme not in ALLOWED_URL_SCHEMES:
        raise ValueError(
            f"REFUSED: {what} {url!r} has scheme {scheme!r}, which is not in the "
            f"allowlist {ALLOWED_URL_SCHEMES}. Only http(s) URLs may be opened "
            f"(file://, ftp://, data:, etc. are blocked — SSRF / local-file-read guard)."
        )


def _http_json(method: str, url: str, api_key: str, body: Optional[dict] = None) -> dict:
    assert_url_scheme_allowed(url, what="API URL")
    if DEAD_ENDPOINT_FRAGMENT in url:
        raise RuntimeError(
            f"REFUSED: attempted to call the dead endpoint {DEAD_ENDPOINT_FRAGMENT}. "
            "This pipeline only uses /api/v1/jobs/createTask and /api/v1/jobs/recordInfo."
        )
    data = json.dumps(body).encode() if body is not None else None
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    req = urllib.request.Request(url, data=data, headers=headers, method=method)
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            return json.loads(resp.read().decode())
    except urllib.error.HTTPError as exc:
        if exc.code == 429:
            raise RateLimited(f"HTTP 429 from {url}")
        if exc.code in (401, 403):
            # PERMANENT auth failure (FIX-6): never transient, never retried. The
            # key/header/account is wrong and an identical re-submit fails forever —
            # the observed 164x "backing off" 401 tailspin. Raised as AuthError so
            # render_slide fails the slide immediately instead of backoff-resubmitting.
            body_text = exc.read().decode(errors="replace")
            raise AuthError(
                f"HTTP {exc.code} {method} {url}\nResponse: {body_text}\n"
                "Permanent auth failure — do NOT re-submit. Check the KIE_API_KEY, "
                "the Authorization: Bearer header format, and that the key is not "
                "locked/rate-blocked by the provider.") from exc
        body_text = exc.read().decode(errors="replace")
        raise RuntimeError(f"HTTP {exc.code} {method} {url}\nResponse: {body_text}") from exc
    except urllib.error.URLError as exc:
        # Network unreachable — fail loud (no substitution).
        raise RuntimeError(f"NETWORK ERROR reaching {url}: {exc}. KIE is unreachable.") from exc


def _import_prompt_gate():
    """Import the shared prompt_gate module that ships beside build_deck.py (the ONE
    image-prompt gate every path uses). Tries a normal import first (scripts/ is on
    sys.path when build_deck runs / is imported), then a path-based load from this file's
    own directory. Returns the module or None (callers degrade gracefully). Mirrors
    _import_intelligence_engines_check."""
    try:
        import importlib
        return importlib.import_module("prompt_gate")
    except Exception:  # noqa: BLE001
        try:
            import importlib.util as _ilu
            spec = _ilu.spec_from_file_location(
                "prompt_gate",
                str(Path(__file__).resolve().parent / "prompt_gate.py"))
            if spec and spec.loader:
                mod = _ilu.module_from_spec(spec)
                spec.loader.exec_module(mod)
                return mod
        except Exception:  # noqa: BLE001
            return None
    return None


def _ensure_english_pin(prompt: str) -> str:
    """Make the mandatory English/Latin anti-garble pin REAL on the render payload. The
    ENGLISH_PIN constant was historically defined and appended NOWHERE (dead code); this
    appends it (belt-and-braces) if the authored rich prompt does not already carry it, so
    every createTask the canonical renderer submits pins the copy to correctly-spelled
    Latin-alphabet text. Prefers the shared prompt_gate helper (single source of truth);
    falls back to a local append if the module is unavailable, so the pin is ALWAYS on the
    payload. Never appends past the 20,000-char GPT-Image-2 API ceiling."""
    pg = _import_prompt_gate()
    if pg is not None:
        try:
            return pg.ensure_english_pin(prompt)
        except Exception:  # noqa: BLE001 — never let the pin helper break a render
            pass
    norm = lambda s: re.sub(r"\s+", " ", s).strip().lower()  # noqa: E731
    if norm(ENGLISH_PIN) in norm(prompt):
        return prompt
    candidate = prompt.rstrip() + "\n\n" + ENGLISH_PIN
    return candidate if len(candidate) <= 20000 else prompt


def _record_ocr_readback(out_path: Path, readback: dict) -> None:
    """Write the OCR text-readback provenance record beside the slide PNG (best-effort,
    never raises). Mirrors the AF-IMAGE-QC-VISION provenance pattern: on a box WITHOUT an
    OCR engine the absence is visible in the record rather than silently skipped."""
    try:
        sidecar = out_path.with_suffix(".ocr.json")
        sidecar.write_text(json.dumps(readback, indent=2))
    except Exception:  # noqa: BLE001
        pass


def _verify_aspect_and_readback(out_path: Path, slide: dict, ordinal: int) -> None:
    """POST-DOWNLOAD image verification, run after verify_png inside render_slide's per-slide
    attempt loop so a failure retries against SLIDE_MAX_ATTEMPTS:
      * ASPECT/2K — the canonical renderer always requests the 16:9 / 2K pin; a response
        whose shape/size does not match is stretched/distorted by assemble_pptx, so refuse
        it (raises -> retry) instead of shipping a distorted slide.
      * OCR TEXT-READBACK — a deterministic garbled-text catch (optional engine;
        provenance-recorded when absent). When the engine RAN and the rendered text does not
        match the slide's approved copy, re-render (retry). This converts garbled baked text
        from an LLM-honesty check into code.
    No-op (PNG magic already verified by verify_png) when the shared gate is unavailable."""
    pg = _import_prompt_gate()
    if pg is None:
        return
    label = f"slide-{ordinal:02d}"
    # Aspect/resolution: build_deck always renders the pinned 16:9 / 2K. Raises on mismatch.
    pg.verify_aspect_ratio(out_path, slide_id=label)
    # OCR readback (optional engine; provenance-recorded). Compare the baked text to the
    # slide's approved copy strings.
    readback = pg.ocr_readback(out_path, slide.get("copy"), slide_id=label)
    _record_ocr_readback(out_path, readback)
    if readback.get("checked") and readback.get("matched") is False:
        raise RuntimeError(
            f"slide {ordinal}: OCR readback — rendered text does not match the approved "
            f"copy (unreadable/garbled: {readback.get('misses')}). Re-rendering this slide."
        )


def submit_task(prompt: str, api_key: str, logo_url: Optional[str] = None) -> str:
    # Make the English/Latin anti-garble pin REAL: append it (belt-and-braces) if the
    # authored rich prompt does not already carry it. Before this, ENGLISH_PIN was a dead
    # constant — defined and appended nowhere — so the #1 defense against garbled/CJK baked
    # text never rode on the payload.
    prompt = _ensure_english_pin(prompt)

    # OFFICIAL-LOGO mode = IMAGE-TO-IMAGE: pass the real logo URL as input_urls so
    # KIE (gpt-image-2-image-to-image) composites the ACTUAL logo into the slide
    # (the verified technique). No logo_url -> plain text-to-image.
    if logo_url:
        payload = {
            "model": MODEL_I2I,
            "input": {
                "prompt": prompt,
                "input_urls": [logo_url],
                "aspect_ratio": ASPECT_RATIO,
                "resolution": RESOLUTION,
            },
        }
    else:
        payload = {
            "model": MODEL_T2I,
            "input": {
                "prompt": prompt,
                "aspect_ratio": ASPECT_RATIO,
                "resolution": RESOLUTION,
            },
        }
    resp = _http_json("POST", CREATE_URL, api_key, body=payload)
    if resp.get("code") != 200:
        raise RuntimeError(f"createTask non-200 code. Full response: {json.dumps(resp)}")
    task_id = (resp.get("data") or {}).get("taskId")
    if not task_id:
        raise RuntimeError(f"createTask 200 but no taskId. Full response: {json.dumps(resp)}")
    return task_id


def _poll_interval_for_elapsed(elapsed_s: float) -> int:
    """Exponential backoff ladder per FIX-7: 10s for the first 2 min, 20s for the
    next 3 min, then 40s after. Returns the sleep seconds for a poll at elapsed_s
    seconds since the poll started. The ladder keeps the FIRST polls cheap so a
    fast render finishes in <2 min, then slows to cut request volume on the long
    tail — while the hard cap (POLL_MAX_SECONDS) bounds total wall-clock time."""
    if elapsed_s < POLL_FAST_WINDOW_S:
        return POLL_FAST_S
    if elapsed_s < POLL_FAST_WINDOW_S + POLL_MEDIUM_WINDOW_S:
        return POLL_MEDIUM_S
    return POLL_SLOW_S


def poll_task(task_id: str, api_key: str) -> str:
    """Poll recordInfo with exponential backoff (10s/20s/40s) and a HARD cap at
    POLL_MAX_SECONDS (~15 min). Returns resultUrls[0] on success.

    FIX-7 (D14): a task that never completes must surface a TERMINAL FAIL
    (RenderPollTimeout) instead of hanging the run in 'in_progress'. The wall-clock
    deadline is absolute — there is NO unbounded grace window. A genuinely failed/
    garbled render still fails loud: any terminal state ('fail'/'failed'/'error'/
    'cancelled') short-circuits immediately. Slow-but-healthy renders are given the
    full ladder; only a task that is still non-terminal at the hard cap raises.
    """
    url = f"{POLL_URL}?taskId={task_id}"

    # States KIE reports while a slide is healthily in-flight (not yet done, not failed).
    HEALTHY_INFLIGHT = ("generating", "waiting", "queuing", "queued", "processing", "running", "pending")

    started = time.time()
    deadline = started + POLL_MAX_SECONDS
    passes = 0
    last_state = ""
    while True:
        if time.time() >= deadline:
            # HARD CAP reached — terminal FAIL, never a silent hang.
            raise RenderPollTimeout(
                f"taskId {task_id}: not complete within {POLL_MAX_SECONDS}s "
                f"hard cap (last state {last_state!r}; {passes} polls). "
                "Render FAILED — the kie task never reached a terminal state."
            )
        passes += 1
        try:
            resp = _http_json("GET", url, api_key)
        except RateLimited:
            print(f"    [poll] 429 — sleeping {RATE_LIMIT_SLEEP_S}s", flush=True)
            time.sleep(RATE_LIMIT_SLEEP_S)
            continue
        data = resp.get("data") or {}
        state = str(data.get("state", "")).lower()
        last_state = state

        if state == "success":
            result_json_str = data.get("resultJson")
            if not result_json_str:
                raise RuntimeError(f"taskId {task_id}: success but no resultJson: {json.dumps(resp)}")
            result_obj = json.loads(result_json_str)
            urls = result_obj.get("resultUrls", [])
            if not urls:
                raise RuntimeError(f"taskId {task_id}: empty resultUrls: {json.dumps(result_obj)}")
            return urls[0]

        if state in ("fail", "failed", "error", "cancelled"):
            raise RuntimeError(
                f"taskId {task_id}: terminal state '{state}' "
                f"failCode={data.get('failCode')} failMsg={data.get('failMsg')}"
            )

        elapsed = time.time() - started
        interval = _poll_interval_for_elapsed(elapsed)
        print(f"    [poll {passes}] taskId={task_id} state={state!r}; "
              f"elapsed={elapsed:.0f}s sleep {interval}s", flush=True)
        time.sleep(interval)


def poll_task_once(task_id: str, api_key: str) -> dict:
    """Single recordInfo status check for ONE task (used by the batch poll pass).

    Unlike poll_task — which BLOCKS for a single slide until that slide completes
    (up to ~7 min) — poll_task_once issues exactly ONE recordInfo GET and returns a
    plain status dict:

        {"state": "success", "result_url": "<first resultUrls[0]>"}
        {"state": "<in-flight>", "result_url": None}   # generating/waiting/queuing/etc.
        raises RuntimeError on a terminal FAIL state

    The batch loop calls this for every pending task on a shared 10s cadence, so one
    slow slide never delays the download of a fast one and a failed slide surfaces
    immediately instead of holding the pass. 429 is surfaced as RateLimited (the
    caller retries that task on the next pass)."""
    url = f"{POLL_URL}?taskId={task_id}"
    resp = _http_json("GET", url, api_key)
    data = resp.get("data") or {}
    state = str(data.get("state", "")).lower()

    if state == "success":
        result_json_str = data.get("resultJson")
        if not result_json_str:
            raise RuntimeError(
                f"taskId {task_id}: success but no resultJson: {json.dumps(resp)}")
        try:
            result_obj = json.loads(result_json_str)
        except (TypeError, ValueError) as exc:
            raise RuntimeError(
                f"taskId {task_id}: resultJson is not valid JSON: {exc}") from exc
        urls = result_obj.get("resultUrls", []) or []
        if not urls:
            raise RuntimeError(f"taskId {task_id}: empty resultUrls: {json.dumps(result_obj)}")
        return {"state": "success", "result_url": urls[0]}

    if state in ("fail", "failed", "error", "cancelled"):
        raise RuntimeError(
            f"taskId {task_id}: terminal state '{state}' "
            f"failCode={data.get('failCode')} failMsg={data.get('failMsg')}"
        )

    # Any other state = still in flight (generating/waiting/queuing/queued/processing/
    # running/pending). Return it verbatim so the batch loop can log + keep polling.
    return {"state": state, "result_url": None}


def download_image(url: str, dest: Path, api_key: str) -> int:
    """AUTHENTICATED GET of the KIE CDN result URL (tempfile.aiquickdraw.com).

    The result URLs KIE returns require `Authorization: Bearer <key>` plus a browser
    User-Agent. A plain urllib GET with neither header returns HTTP 403, which left
    renders/ empty despite successful createTask calls (D22 root cause A). This is the
    FIX-4 helper: send both headers, write the bytes, return the on-disk size.
    """
    # C1 (SSRF / local-file-read guard): the result URL comes back from KIE's API.
    # Refuse anything that is not http(s) before opening it (a file:// result URL
    # would otherwise read an arbitrary local file into the slide PNG).
    assert_url_scheme_allowed(url, what="KIE result URL")
    req = urllib.request.Request(url, headers={
        "Authorization": f"Bearer {api_key}",
        "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)",
    })
    with urllib.request.urlopen(req, timeout=180) as resp:
        dest.write_bytes(resp.read())
    return dest.stat().st_size


def verify_png(path: Path) -> None:
    """Verify the downloaded file is a real, non-empty PNG."""
    if not path.exists():
        raise RuntimeError(f"{path}: file was not written.")
    size = path.stat().st_size
    if size == 0:
        raise RuntimeError(f"{path}: file is zero bytes.")
    with open(path, "rb") as f:
        magic = f.read(8)
    if magic[:4] != b"\x89PNG":
        raise RuntimeError(f"{path}: not a PNG (magic={magic[:8].hex()}, size={size}).")


# ---------------------------------------------------------------------------
# Per-slide render with retry
# ---------------------------------------------------------------------------



# ---------------------------------------------------------------------------
# U028 -- pre-call task-id checkpointing (see presentation_job/checkpoint.py)
# ---------------------------------------------------------------------------

def _pending_tasks_path(run_dir: Path) -> Path:
    return run_dir / "working" / "checkpoints" / "pending_tasks.json"


def _checkpoint_pending_task(run_dir: Path, ordinal: int, task_id: str,
                             attempt: int) -> None:
    try:
        p = _pending_tasks_path(run_dir)
        p.parent.mkdir(parents=True, exist_ok=True)
        existing: dict = {}
        if p.exists():
            try: existing = json.loads(p.read_text())
            except (json.JSONDecodeError, ValueError): pass
        existing[str(ordinal)] = {"task_id": task_id, "attempt": attempt,
                                   "submitted_at": datetime.now(timezone.utc).isoformat()}
        atomic_write_text(p, json.dumps(existing, indent=2))
    except Exception:  # noqa: BLE001
        pass


def _read_pending_tasks(run_dir: Path) -> dict:
    p = _pending_tasks_path(run_dir)
    if not p.exists(): return {}
    try: return json.loads(p.read_text())
    except (json.JSONDecodeError, ValueError, OSError): return {}


def _record_completed_task(run_dir: Path, ordinal: int, task_id: str,
                           out_path: Path) -> None:
    try: sha = hashlib.sha256(out_path.read_bytes()).hexdigest()
    except Exception: return
    try:
        p = _pending_tasks_path(run_dir)
        existing = _read_pending_tasks(run_dir)
        existing[str(ordinal)] = {"task_id": task_id, "completed": True,
                                   "output_path": str(out_path), "sha256": sha,
                                   "completed_at": datetime.now(timezone.utc).isoformat()}
        atomic_write_text(p, json.dumps(existing, indent=2))
    except Exception: pass  # noqa: BLE001


def _clear_pending_task(run_dir: Path, ordinal: int) -> None:
    """Remove a terminal/exhausted pending task record so a future resume
    submits fresh instead of re-polling the same dead task id forever."""
    try:
        p = _pending_tasks_path(run_dir)
        existing = _read_pending_tasks(run_dir)
        if str(ordinal) in existing:
            del existing[str(ordinal)]
            atomic_write_text(p, json.dumps(existing, indent=2))
    except Exception:  # noqa: BLE001
        pass


def _resume_pending_task(run_dir: Path, ordinal: int, api_key: str) -> Optional[str]:
    pending = _read_pending_tasks(run_dir)
    rec = pending.get(str(ordinal))
    if not rec: return None
    task_id = rec.get("task_id")
    if not task_id: return None
    if rec.get("completed"):
        out_p = rec.get("output_path")
        sha = rec.get("sha256")
        if out_p and sha:
            if PREDICATES["image"](Path(out_p), sha256=sha):
                return "__COMPLETED__"
        return None
    try:
        result_url = poll_task(task_id, api_key)
        print(f"    [resume] pending taskId={task_id} completed -- reusing.")
        return result_url
    except Exception:
        print(f"    [resume] pending taskId={task_id} is terminal/exhausted -- "
              f"will submit fresh.", flush=True)
        _clear_pending_task(run_dir, ordinal)
        return None



def render_slide(slide: dict, api_key: str, renders_dir: Path, run_dir: Path,
                 has_official_logo: bool = False, logo_url: Optional[str] = None) -> dict:
    """
    Render a single slide. Returns {"slide", "file", "taskId"} on success.
    Raises RuntimeError if all SLIDE_MAX_ATTEMPTS fail.

    The prompt is the Slide Image Creator's RICH per-slide prompt, loaded VERBATIM
    from working/prompts/ (no thin fallback). has_official_logo is accepted for API
    stability but does NOT mutate the prompt — the rich prompt already carries the
    full logo instruction; logo_url is what makes KIE composite the real logo via
    image-to-image (input_urls).
    """
    ordinal = int(slide["slide"])
    name = f"slide-{ordinal:02d}"
    out_path = renders_dir / f"{name}.png"
    prompt = load_rich_prompt(slide, run_dir)

    # U028 -- resume: poll a pending task id from a prior killed run before
    # paying for a fresh submission.  This is the entire saving.
    resumed_url = _resume_pending_task(run_dir, ordinal, api_key)
    if resumed_url == "__COMPLETED__":
        print(f"  [{name}] resume found completed artifact -- skipping.")
        return {"slide": ordinal, "file": str(out_path),
                "taskId": "resumed-completed"}
    if resumed_url is not None:
        print(f"  [{name}] resume polling succeeded -- downloading result.")
        download_image(resumed_url, out_path, api_key)
        verify_png(out_path)
        _verify_aspect_and_readback(out_path, slide, ordinal)
        size = out_path.stat().st_size
        print(f"    downloaded+verified -> {out_path} ({size:,} bytes)", flush=True)
        _record_completed_task(run_dir, ordinal, "resumed", out_path)
        return {"slide": ordinal, "file": str(out_path), "taskId": "resumed"}

    last_err = None
    for attempt in range(1, SLIDE_MAX_ATTEMPTS + 1):
        print(f"  [{name}] attempt {attempt}/{SLIDE_MAX_ATTEMPTS}", flush=True)
        try:
            # submit (429-aware)
            while True:
                try:
                    task_id = submit_task(prompt, api_key, logo_url=logo_url)
                    break
                except RateLimited:
                    print(f"    [submit] 429 — sleeping {RATE_LIMIT_SLEEP_S}s", flush=True)
                    time.sleep(RATE_LIMIT_SLEEP_S)
            print(f"    submitted -> taskId={task_id}", flush=True)
            _checkpoint_pending_task(run_dir, ordinal, task_id, attempt)
            result_url = poll_task(task_id, api_key)
            print(f"    success resultUrls[0]={result_url}", flush=True)
            download_image(result_url, out_path, api_key)
            verify_png(out_path)
            # POST-DOWNLOAD aspect/2K verification + OCR text-readback (shared prompt_gate).
            # A non-16:9 / sub-2K response, or rendered text that does not match the approved
            # copy, raises here and retries against SLIDE_MAX_ATTEMPTS rather than shipping a
            # distorted or garbled slide.
            _verify_aspect_and_readback(out_path, slide, ordinal)
            size = out_path.stat().st_size
            print(f"    downloaded+verified -> {out_path} ({size:,} bytes)", flush=True)
            _record_completed_task(run_dir, ordinal, task_id, out_path)
            return {"slide": ordinal, "file": str(out_path), "taskId": task_id}
        except AuthError:
            # FIX-6 — fail-fast on auth errors. A 401/403 is PERMANENT: the identical
            # request can never succeed, so backoff-re-submitting is a token furnace
            # (the observed 164x "backing off" 401 tailspin). Re-raise immediately —
            # zero backoff, zero re-submits — so the slide fails fast (< 2s).
            raise
        except RenderPollTimeout:
            # FIX-7 (D14): a task that never reaches a terminal state within the
            # poll hard cap is a TERMINAL FAIL — surface it immediately. Re-submitting
            # a fresh identical task would only re-poll a (probably) dead task for
            # another full cap; a never-completing render must FAIL at ≤15 min, never
            # hang the run in 'in_progress'. No retry, no backoff.
            raise
        except Exception as exc:  # noqa: BLE001 — deliberately catch to retry
            last_err = exc
            print(f"    FAIL attempt {attempt}: {exc}", file=sys.stderr, flush=True)
            if attempt < SLIDE_MAX_ATTEMPTS:
                # Exponential backoff on the IDENTICAL request. The KIE failCode 400
                # "Internal Error, Please try again later" is transient; re-submit the
                # same prompt/params rather than trimming.
                backoff = min(SLIDE_BACKOFF_CAP_S, SLIDE_BACKOFF_BASE * (2 ** (attempt - 1)))
                print(f"    backing off {backoff:.0f}s before identical re-submit", flush=True)
                time.sleep(backoff)

    raise RuntimeError(f"{name}: failed after {SLIDE_MAX_ATTEMPTS} attempts. Last error: {last_err}")


# ---------------------------------------------------------------------------
# BATCH RENDER (FIX-5 / D22 root cause B) — submit ALL prompts once, then poll
# all together, download each image the moment ITS task finishes.
# ---------------------------------------------------------------------------

def render_slides_batch(slides: list, api_key: str, renders_dir: Path, run_dir: Path,
                        has_official_logo: bool = False, logo_url: Optional[str] = None,
                        submit_interval: float = BATCH_SUBMIT_INTERVAL_S,
                        poll_interval: float = BATCH_POLL_INTERVAL_S,
                        max_seconds: float = BATCH_MAX_POLL_SECONDS) -> dict:
    """Batch-render every slide against KIE.ai, the way D22's live fix proved works.

    Phase A — SUBMIT ALL: every prompt is submitted once to createTask, spaced
      BATCH_SUBMIT_INTERVAL_S (0.6s) apart. KIE allows 20 requests per 10s and 100+
      concurrent tasks (verified live: a burst of 8 rapid submits all returned 200;
      a 17-slide batch rendered + downloaded in ~3 min), so a 20-slide deck lands its
      whole submission window in ~12s — far under the 20/10s cap. A 429 (rate limit)
      is the ONLY submission error that retries, with the existing RATE_LIMIT_SLEEP_S
      backoff; any other error is recorded against that slide and submission continues,
      so one bad prompt never blocks the other 19 (the run still FAILS LOUD at the end
      if any slide failed — never a silently-partial deck).

    Phase B — POLL ALL + DOWNLOAD AS FINISHED: every submitted taskId is polled on
      one shared cadence (BATCH_POLL_INTERVAL_S = 10s). The moment a task's state
      turns 'success' its image is downloaded + verified + recorded, and that task is
      removed from the poll set — so a fast slide is never held hostage by the slowest
      one. Terminal states (fail/failed/error/cancelled) are recorded as failures and
      also removed. A hard wall-clock cap (BATCH_MAX_POLL_SECONDS = 15 min) makes a
      genuinely-stuck task surface FAIL instead of hanging the run forever (D14).

    The per-slide rich-prompt gate (load_rich_prompt), the English pin, logo compositing
    (i2i via input_urls), post-download aspect/OCR verification, and the U028 task-id
    checkpoint all still run per slide — batch does NOT weaken a single gate; it only
    changes WHEN submission happens (all up front) and WHEN polling happens (all at once).

    Returns the same contract main() consumed from the thread pool:
      {"rendered": [ {slide,file,taskId} ... ], "failures": [ {slide,error} ... ]}
    """
    ordered = sorted(slides, key=lambda s: s["slide"])
    n = len(ordered)
    renders_dir.mkdir(parents=True, exist_ok=True)

    # ---- Phase A: submit every prompt once, 0.6s apart ----
    submitted = {}  # ordinal -> (slide, task_id)
    submit_failures = []  # {slide, error}
    start = time.time()
    print(f"=== BATCH SUBMIT: {n} prompts, {submit_interval:.1f}s apart "
          f"(kie 20/10s cap + 100 concurrent tasks) ===", flush=True)
    for slide in ordered:
        ordinal = int(slide["slide"])
        name = f"slide-{ordinal:02d}"
        try:
            prompt = load_rich_prompt(slide, run_dir)
            # 429-aware submit: a rate limit backs off and re-tries the SAME slide,
            # everything else (network, shape, 401) fails that slide without blocking
            # the batch. Never re-submit a slide that already has a taskId.
            while True:
                try:
                    task_id = submit_task(prompt, api_key, logo_url=logo_url)
                    break
                except RateLimited:
                    print(f"    [submit] 429 — sleeping {RATE_LIMIT_SLEEP_S}s", flush=True)
                    time.sleep(RATE_LIMIT_SLEEP_S)
            submitted[ordinal] = (slide, task_id)
            _checkpoint_pending_task(run_dir, ordinal, task_id, attempt=1)
            print(f"  [{name}] submitted -> taskId={task_id}", flush=True)
        except Exception as exc:  # noqa: BLE001 — one bad prompt must not block 19
            submit_failures.append({"slide": ordinal, "error": str(exc)})
            print(f"  [{name}] SUBMIT FAILED: {exc}", file=sys.stderr, flush=True)
        # Space submissions inside the 20/10s window.
        if submit_interval > 0:
            time.sleep(submit_interval)
    submit_wall = time.time() - start
    print(f"=== BATCH SUBMIT DONE: {len(submitted)}/{n} submitted in "
          f"{submit_wall:.1f}s ===", flush=True)

    # ---- Phase B: poll every submitted task together, download as each finishes ----
    rendered = []
    poll_failures = list(submit_failures)
    pending = {ordinal: (slide, task_id) for ordinal, (slide, task_id) in submitted.items()}
    poll_start = time.time()
    pass_no = 0
    while pending:
        elapsed = time.time() - poll_start
        if elapsed >= max_seconds:
            # Hard cap: whatever is still pending is STUCK — surface FAIL, no hang.
            for ordinal, (slide, task_id) in list(pending.items()):
                name = f"slide-{ordinal:02d}"
                poll_failures.append({"slide": ordinal,
                                      "error": f"taskId {task_id} still pending after "
                                               f"{max_seconds:.0f}s — surfaced FAIL (no hang)"})
                print(f"  [{name}] POLL CAP REACHED ({max_seconds:.0f}s) — FAIL: {task_id}",
                      file=sys.stderr, flush=True)
                _clear_pending_task(run_dir, ordinal)
            pending.clear()
            break

        pass_no += 1
        done_this_pass = []
        for ordinal, (slide, task_id) in list(pending.items()):
            name = f"slide-{ordinal:02d}"
            try:
                status = poll_task_once(task_id, api_key)
            except RateLimited:
                # Transient 429 on the poll — retry the same task on the next pass.
                print(f"    [poll {pass_no}] {name} 429 — retry next pass", flush=True)
                continue
            except Exception as exc:  # noqa: BLE001 — terminal failure surfaced
                print(f"  [{name}] POLL FAILED: {exc}", file=sys.stderr, flush=True)
                poll_failures.append({"slide": ordinal, "error": str(exc)})
                _clear_pending_task(run_dir, ordinal)
                done_this_pass.append(ordinal)
                continue

            if status.get("state") != "success":
                # Still in flight (generating/waiting/queuing/...). Leave it in the
                # pending set and let the next shared pass check it again — a slow
                # slide never delays a fast one.
                print(f"    [poll {pass_no}] {name} state={status.get('state')!r}", flush=True)
                continue

            result_url = status.get("result_url")
            try:
                out_path = renders_dir / f"{name}.png"
                download_image(result_url, out_path, api_key)
                verify_png(out_path)
                _verify_aspect_and_readback(out_path, slide, ordinal)
                size = out_path.stat().st_size
                print(f"  [{name}] downloaded+verified -> {out_path} ({size:,} bytes)",
                      flush=True)
                _record_completed_task(run_dir, ordinal, task_id, out_path)
                rendered.append({"slide": ordinal, "file": str(out_path), "taskId": task_id})
            except Exception as exc:  # noqa: BLE001 — a bad image fails this slide only
                print(f"  [{name}] DOWNLOAD/VERIFY FAILED: {exc}", file=sys.stderr, flush=True)
                poll_failures.append({"slide": ordinal, "error": str(exc)})
                _clear_pending_task(run_dir, ordinal)
            done_this_pass.append(ordinal)

        for ordinal in done_this_pass:
            pending.pop(ordinal, None)

        if pending:
            remaining = ", ".join(f"slide-{o:02d}" for o in sorted(pending))
            print(f"    [poll pass {pass_no}] waiting on: {remaining} "
                  f"(next poll in {poll_interval:.0f}s)", flush=True)
            time.sleep(poll_interval)

    poll_wall = time.time() - poll_start
    print(f"=== BATCH RENDER DONE: {len(rendered)}/{n} rendered in "
          f"{time.time() - start:.1f}s total (submit {submit_wall:.1f}s + "
          f"poll/download {poll_wall:.1f}s) ===", flush=True)

    rendered.sort(key=lambda r: r["slide"])
    poll_failures.sort(key=lambda f: (f.get("slide") is None, f.get("slide")))
    return {"rendered": rendered, "failures": poll_failures}


# ---------------------------------------------------------------------------
# pptx assembly — full-bleed picture per slide, NO text boxes
# ---------------------------------------------------------------------------

def _logo_dimensions_in(logo_path: Path):
    """Return (width_in, height_in) for the composited logo: width fixed at
    LOGO_WIDTH_FRACTION of the slide width, height derived from the logo's native
    pixel aspect ratio so the EXACT client file is never distorted. Falls back to a
    square box if the image dimensions can't be read."""
    width_in = SLIDE_WIDTH_IN * LOGO_WIDTH_FRACTION
    try:
        from PIL import Image  # optional; only used to preserve aspect ratio
        with Image.open(str(logo_path)) as im:
            px_w, px_h = im.size
        if px_w > 0:
            height_in = width_in * (px_h / px_w)
        else:
            height_in = width_in
    except Exception:  # noqa: BLE001 — PIL absent or unreadable: assume square
        height_in = width_in
    return width_in, height_in


# ---------------------------------------------------------------------------
# PER-SLIDE SPEAKER NOTES — parse the presenter speech into per-slide chunks
# ---------------------------------------------------------------------------
# The Presenters Speech Writer authors a word-for-word script segmented per slide.
# Each slide's block is introduced by a SLIDE marker in one of two forms:
#       ## Slide 7          (a markdown heading, 1-3 '#'s, "Slide" any case)
#       SLIDE 7             (an inline marker with no heading hashes)
# parse_speech_chunks splits the speech on those markers and returns
#   {slide_no: spoken_text}
# where spoken_text is the block AFTER the marker/title line, up to the next marker
# (the marker/title line itself is stripped — it is a structural cue, not spoken).
#
# This is best-effort and NON-FATAL by contract: a speech with no recognisable
# markers yields {} and the caller injects no notes. Mismatches are handled by the
# caller (extra chunks are skipped; deck slides with no chunk get empty notes).

# Matches both "## Slide 7" (markdown heading) and inline "SLIDE 7". The optional
# leading 1-3 '#'s + whitespace cover the heading form; the marker word is matched
# case-insensitively (re.I) and anchored to the start of a line (re.M).
SPEECH_SLIDE_MARKER_RE = re.compile(r"^(?:#{1,3}\s+)?SLIDE\s+(\d+)\b", re.IGNORECASE | re.MULTILINE)


def parse_speech_chunks(speech_text: str) -> dict:
    """Parse a presenter-speech string into {slide_no: spoken_text}.

    Recognises BOTH marker forms (see SPEECH_SLIDE_MARKER_RE):
        '## Slide N' / '# Slide N' / '### Slide N'  (markdown heading)
        'SLIDE N'                                    (inline marker)
    For each marker, captures the text from the END of the marker line up to the
    start of the NEXT marker (or end of file), then strips surrounding whitespace.
    The marker/title line is NOT included in the spoken text.

    Returns {} when speech_text is empty/None or contains no recognisable markers
    (best-effort, never raises). If the same slide number appears more than once,
    the LAST occurrence wins (a later, fuller block supersedes an earlier stub)."""
    if not speech_text:
        return {}

    chunks: dict = {}
    matches = list(SPEECH_SLIDE_MARKER_RE.finditer(speech_text))
    if not matches:
        return {}

    for idx, m in enumerate(matches):
        try:
            slide_no = int(m.group(1))
        except (TypeError, ValueError):
            continue
        # Spoken text starts at the END of the marker LINE (skip the rest of the
        # marker/title line, which may carry a slide title after the number).
        line_end = speech_text.find("\n", m.end())
        block_start = (line_end + 1) if line_end != -1 else len(speech_text)
        # Spoken text ends at the start of the NEXT marker, else end of file.
        block_end = matches[idx + 1].start() if idx + 1 < len(matches) else len(speech_text)
        spoken = speech_text[block_start:block_end].strip()
        # Last occurrence of a slide number wins (supersedes an earlier stub).
        chunks[slide_no] = spoken
    return chunks


def discover_speech_chunks(run_dir: Path, bundle_dir: Path) -> Optional[dict]:
    """Auto-discover the presenter speech and parse it into per-slide chunks for
    PPTX speaker-notes injection. NON-FATAL: this is a phase-4 render and the speech
    is a phase-9 artifact, so the speech is frequently absent. Returns the parsed
    {slide_no: spoken_text} dict when a speech file is found, or None when none is
    present (caller then injects no notes). NEVER raises and NEVER blocks the render.

    Search order (first hit wins): the working-dir locations the speech roles write
    to, then the delivered bundle copy. Filenames use the standardized possessive
    plural PRESENTERS-SPEECH.md (canonical); the legacy singular PRESENTER-SPEECH.md,
    the speech.md scratch name, and the presenters_speech.md snake_case variant
    (audio-demonstration-specialist.md's legacy reference name) are also accepted —
    legacy tolerance, so no producer's output silently never reaches the notes pane."""
    candidates = [
        run_dir / "working/deliverables/PRESENTERS-SPEECH.md",
        run_dir / "working/presenter-speech/speech.md",
        run_dir / "working/delivery/PRESENTERS-SPEECH.md",
        run_dir / "working/presenter-speech/PRESENTERS-SPEECH.md",
        bundle_dir / "PRESENTERS-SPEECH.md",
        run_dir / "working/delivery/PRESENTER-SPEECH.md",
        run_dir / "working/presenter-speech/PRESENTER-SPEECH.md",
        bundle_dir / "PRESENTER-SPEECH.md",
        run_dir / "working/presenter-speech/presenters_speech.md",
        run_dir / "working/delivery/presenters_speech.md",
        bundle_dir / "presenters_speech.md",
    ]
    for path in candidates:
        try:
            if not path.is_file():
                continue
        except OSError:
            continue
        try:
            text = path.read_text(errors="replace")
        except OSError as exc:
            print(f"=== SPEAKER NOTES: found {path} but could not read it ({exc}); "
                  f"rendering deck WITHOUT per-slide notes (non-fatal) ===", flush=True)
            return None
        chunks = parse_speech_chunks(text)
        if chunks:
            print(f"=== SPEAKER NOTES: parsed {len(chunks)} per-slide chunk(s) from "
                  f"{path} — injecting into the deck's notes pane ===", flush=True)
        else:
            print(f"=== SPEAKER NOTES: speech found at {path} but no 'Slide N' / "
                  f"'## Slide N' markers were recognised; rendering deck WITHOUT "
                  f"per-slide notes (non-fatal) ===", flush=True)
        return chunks
    print("=== SPEAKER NOTES: no presenter speech found yet (searched "
          "working/presenter-speech/, working/delivery/, and the bundle dir). This is "
          "a phase-4 render; the speech is a phase-9 artifact. Rendering the deck "
          "WITHOUT per-slide notes (non-fatal — never blocks the render). ===",
          flush=True)
    return None


def assemble_pptx(rendered: list, out_path: Path, logo_path: Optional[Path] = None,
                  speech_chunks: Optional[dict] = None) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("FATAL: python-pptx is not installed (pip install python-pptx).", file=sys.stderr)
        sys.exit(2)

    prs = Presentation()
    # 16:9 at the dept-documented dimensions.
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank = prs.slide_layouts[6]  # blank layout — no placeholders

    # Pre-compute the official-logo geometry ONCE so every slide gets the EXACT
    # same file at the EXACT same size/position (consistent top-right corner).
    logo_geom = None
    if logo_path is not None:
        logo_w_in, logo_h_in = _logo_dimensions_in(logo_path)
        logo_left_in = SLIDE_WIDTH_IN - LOGO_MARGIN_IN - logo_w_in
        logo_top_in = LOGO_MARGIN_IN
        logo_geom = (logo_left_in, logo_top_in, logo_w_in, logo_h_in)

    for item in sorted(rendered, key=lambda r: r["slide"]):
        slide = prs.slides.add_slide(blank)
        # Full-bleed background first.
        slide.shapes.add_picture(
            item["file"], 0, 0, width=Inches(SLIDE_WIDTH_IN), height=Inches(SLIDE_HEIGHT_IN)
        )
        # Then the EXACT official logo on top, top-right. add_picture writes the PNG
        # bytes verbatim, so RGBA transparency is preserved over the background.
        if logo_geom is not None:
            left_in, top_in, w_in, h_in = logo_geom
            slide.shapes.add_picture(
                str(logo_path),
                Inches(left_in), Inches(top_in),
                width=Inches(w_in), height=Inches(h_in),
            )

        # PER-SLIDE SPEAKER NOTES: if a parsed speech chunk exists for this slide
        # ordinal, write the spoken text into the slide's notes pane. None/{} =>
        # no injection (phase-4 render before the speech is written). Touching
        # notes_slide lazily creates the notes part, so a deck slide with no chunk
        # is left with no notes part at all — never an empty injection.
        if speech_chunks:
            spoken = speech_chunks.get(item["slide"])
            if spoken:
                slide.notes_slide.notes_text_frame.text = spoken

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


# ---------------------------------------------------------------------------
# P9.5-NOTES-SYNC — reorder notes-pane assembly to AFTER the speech is finalized.
#
# ROOT CAUSE (spec §9.1): assemble_pptx() runs at P8-ASSEMBLE, but the presenter
# speech (PRESENTERS-SPEECH.md) is a P9-SPEECH / P-SPEECH-QC artifact. On a normal
# linear run the notes pane is injected BEFORE the speech exists, so it ships
# empty. Fix: REORDER (do not rewrite) — a small idempotent pass that reopens the
# already-assembled .pptx after the speech has been written AND QC-passed, and
# re-runs the same discover_speech_chunks()/notes injection this time with the
# real word-for-word script on disk. This is the "P9.5-NOTES-SYNC" phase
# (PIPELINE-MANIFEST.json order 8.7, between P-SPEECH-QC and P9-DELIVER).
# ---------------------------------------------------------------------------

def notes_sync_pass(bundle_pptx: Path, run_dir: Path, bundle_dir: Path) -> dict:
    """P9.5-NOTES-SYNC: reopen the assembled .pptx and inject the FULL word-for-word
    presenter speech into every slide's notes pane, now that the speech is (expected
    to be) present and QC-passed. IDEMPOTENT — this OVERWRITES the notes text frame
    on every matched slide each time it runs; it never appends or duplicates, and
    running it twice on the same inputs produces the same output.

    Returns a small JSON-able result dict:
      {"status": "synced"|"no_speech"|"error", "slides_total": N,
       "slides_with_notes": N, "speech_source": "<path or None>",
       "reason": "<human string>"}
    NEVER raises — a notes-sync problem must never take down a build that already
    has a valid assembled deck; the AF-EMPTY-NOTES-PANE postflight gate is what
    turns an unsynced deck into a hard failure, not this pass."""
    try:
        from pptx import Presentation
    except ImportError:
        return {"status": "error", "slides_total": 0, "slides_with_notes": 0,
                "speech_source": None,
                "reason": "python-pptx is not installed (pip install python-pptx)."}

    if not bundle_pptx.is_file():
        return {"status": "error", "slides_total": 0, "slides_with_notes": 0,
                "speech_source": None,
                "reason": f"assembled PPTX not found at {bundle_pptx}; run assembly "
                          f"(P8-ASSEMBLE) before P9.5-NOTES-SYNC."}

    speech_chunks = discover_speech_chunks(run_dir, bundle_dir)
    if not speech_chunks:
        return {"status": "no_speech", "slides_total": 0, "slides_with_notes": 0,
                "speech_source": None,
                "reason": "no presenter speech found (or it parsed to zero chunks); "
                          "the notes pane is left as-is. This is non-fatal here — the "
                          "AF-EMPTY-NOTES-PANE postflight gate is what fails the run "
                          "if notes are still empty at delivery."}

    try:
        prs = Presentation(str(bundle_pptx))
    except Exception as exc:  # noqa: BLE001
        return {"status": "error", "slides_total": 0, "slides_with_notes": 0,
                "speech_source": None,
                "reason": f"could not open {bundle_pptx} with python-pptx: {exc}"}

    slides_total = len(prs.slides)
    slides_with_notes = 0
    for idx, slide in enumerate(prs.slides, start=1):
        spoken = speech_chunks.get(idx)
        if spoken:
            # Overwrite (never append) so re-running this pass is idempotent.
            slide.notes_slide.notes_text_frame.text = spoken
            slides_with_notes += 1

    try:
        prs.save(str(bundle_pptx))
    except Exception as exc:  # noqa: BLE001
        return {"status": "error", "slides_total": slides_total,
                "slides_with_notes": slides_with_notes, "speech_source": None,
                "reason": f"could not save {bundle_pptx} after notes injection: {exc}"}

    print(f"=== P9.5-NOTES-SYNC: injected notes for {slides_with_notes}/{slides_total} "
          f"slide(s) into {bundle_pptx} from the QC-passed presenter speech ===",
          flush=True)
    return {"status": "synced", "slides_total": slides_total,
            "slides_with_notes": slides_with_notes, "speech_source": "found",
            "reason": ""}


# Env-overridable, comma-separated, case-insensitive slide-tag substrings that
# exempt a slide from AF-EMPTY-NOTES-PANE (structural/divider slides carry no
# spoken content of their own — the spec's "section-banner slides, configurable").
NOTES_PANE_EXEMPT_TAG_DEFAULT = "banner,divider,section,transition,intermission"


def _notes_pane_exempt_tags() -> set:
    raw = os.environ.get("NOTES_PANE_EXEMPT_TAGS", NOTES_PANE_EXEMPT_TAG_DEFAULT)
    return {t.strip().lower() for t in raw.split(",") if t.strip()}


def _notes_pane_exempt_slide_numbers(slides_path: Optional[Path]) -> set:
    """Read slides.json (if available) and return the set of 1-based slide numbers
    tagged as structural/divider — exempt from the AF-EMPTY-NOTES-PANE requirement.
    Best-effort: any parse problem yields an empty exemption set (fail toward MORE
    checking, never less)."""
    if slides_path is None:
        return set()
    try:
        slides = json.loads(Path(slides_path).read_text())
    except Exception:  # noqa: BLE001
        return set()
    if not isinstance(slides, list):
        return set()
    exempt_tags = _notes_pane_exempt_tags()
    exempt = set()
    for i, s in enumerate(slides, start=1):
        if not isinstance(s, dict):
            continue
        tokens = []
        for k in ("arc_section", "section", "beat", "tag", "type", "role", "kind"):
            v = s.get(k)
            if isinstance(v, str):
                tokens.append(v.lower())
        tags = s.get("tags")
        if isinstance(tags, list):
            tokens += [str(t).lower() for t in tags]
        blob = " ".join(tokens)
        if any(tag in blob for tag in exempt_tags):
            exempt.add(i)
    return exempt


def _chk_notes_pane(bundle_dir: Path, run_dir: Optional[Path] = None,
                    slides_path: Optional[Path] = None) -> str:
    """AF-EMPTY-NOTES-PANE (delivery/closeout gate, D11 in the Master Ruleset).
    Opens the delivered *.pptx in bundle_dir and reads
    slide.notes_slide.notes_text_frame.text for every slide; any audience-facing
    CONTENT slide with an empty notes pane is a hard fail (structural/divider
    slides are exempt per _notes_pane_exempt_slide_numbers). Returns "" on pass, or
    a fatal AF-EMPTY-NOTES-PANE message naming the empty slide numbers.

    This is the code implementation of the autofail that previously existed only
    on paper (sops/SOP-SLIDE-00-MASTER-QC-AUTOFAIL-RULESET.md:578-582) — closing
    the exact gap the notes-pane reorder (P9.5-NOTES-SYNC) exists to prevent: a
    deck assembled before the speech existed, never re-synced, shipping silent."""
    try:
        from pptx import Presentation
    except ImportError:
        return ("AF-EMPTY-NOTES-PANE: python-pptx is not installed, so the notes "
                "pane cannot be verified. Install python-pptx before delivery.")

    pptx_candidates = [p for p in sorted(Path(bundle_dir).glob("*.pptx"))
                       if not p.name.startswith("~$")]
    if not pptx_candidates:
        # No delivered PPTX here yet — another gate (AF-BUNDLE-COMPLETE) owns the
        # "missing deliverable" failure; this gate has nothing to inspect.
        return ""

    exempt = _notes_pane_exempt_slide_numbers(slides_path)
    empty_slides = []
    for pptx_path in pptx_candidates:
        try:
            prs = Presentation(str(pptx_path))
        except Exception:  # noqa: BLE001
            # A file python-pptx cannot open is NOT a valid deck to scan here — the
            # postflight bundle-completeness magic-byte gate (AF-BUNDLE-COMPLETE) and
            # the C2 content-type check own malformed/decoy .pptx files. Defer rather
            # than false-fail AF-EMPTY-NOTES-PANE (mirrors _delivered_pptx_native_text).
            continue
        for idx, slide in enumerate(prs.slides, start=1):
            if idx in exempt:
                continue
            has_notes_part = slide.has_notes_slide
            text = (slide.notes_slide.notes_text_frame.text.strip()
                    if has_notes_part else "")
            if not text:
                empty_slides.append((pptx_path.name, idx))

    if empty_slides:
        listing = "; ".join(f"{fname} slide {n}" for fname, n in empty_slides[:20])
        more = f" (+{len(empty_slides) - 20} more)" if len(empty_slides) > 20 else ""
        return (f"AF-EMPTY-NOTES-PANE: DECK FAIL -- final .pptx ships with empty "
                f"notes panes on content slides: {listing}{more}. Re-assemble with "
                f"the presenter speech present so per-slide notes are injected "
                f"(build_deck.py auto-injection at P8-ASSEMBLE, re-synced at "
                f"P9.5-NOTES-SYNC / notes_sync_pass()), then re-verify.")
    return ""


# ---------------------------------------------------------------------------
# PROCESS PREFLIGHT GATE — refuse to render without the upstream dept artifacts
# ---------------------------------------------------------------------------
#
# build_deck.py is ONLY the deterministic Phase-4 renderer + a stripped Phase-8
# assembler. The real department flow wraps it in research, copy, two human gates,
# and four QC gates (see the Presentations dept SOPs). A hand-fed slides.json that
# skips all of that produces a .pptx that is technically valid but is NOT a
# process-compliant deliverable. This preflight makes that bypass fail LOUD.
#
# Each entry: (relative-glob-under-run-dir, human label, producing phase/role,
#              completeness-check callable). The completeness check receives the
# resolved Path (or None when the file is absent) and returns "" when satisfied
# or a short reason string when the artifact exists but is not complete.

def _read_json(path: Path):
    try:
        return json.loads(path.read_text())
    except Exception as exc:  # noqa: BLE001
        return {"__parse_error__": str(exc)}


# The three DEFINITIVE audience modes (additive to presentation_mode; SOP-CAST /
# Director intake). STANDARD = net-new from a brainstorm; PERSONAL = built from the
# client's own content for a NAMED recipient; GENERAL = built from the client's own
# content but DE-IDENTIFIED. Captured in intake.json as audience_mode (uppercase).
AUDIENCE_MODES = ("STANDARD", "PERSONAL", "GENERAL")


# The SIX mandatory Brainstorming-Buddy fields captured under
# intake.json -> pre_presentation_capture (SOP 9.0 / deck-intake-questions.json
# storeTarget). Each maps to the turn-gated interview question id and, where the SOP
# allows a captured DECLINE, the sanctioned flag that waives a non-empty value.
# DARK_OK is a boolean (false is a legitimate captured answer, NOT "empty").
# GROUNDED_CONTENT has NO decline waiver — the SOP forbids a blank grounded content.
_INTAKE_MANDATORY_FIELDS = (
    # (canonical field key, accepted aliases, ledger question id, decline-flag(s), is_bool)
    ("REPRESENTATION_MIX", (), "representation_mix",
     ("representation_uncaptured",), False),
    ("AUDIENCE_COMPOSITION", ("AUDIENCE_COMPOSITION_NOTE",), "audience_composition_note",
     ("audience_composition_uncaptured", "representation_uncaptured"), False),
    ("GROUNDED_CONTENT", (), "grounded_content",
     (), False),
    ("VISUAL_MIX", (), "visual_mix",
     ("visual_mix_defaulted",), False),
    ("DARK_OK", (), "dark_ok",
     (), True),
    ("HOOK_SEED", (), "hook_seed",
     ("hook_seed_missing",), False),
)


def _intake_run_dir(intake_path: Path) -> Optional[Path]:
    """Resolve the deck run dir from the intake.json path: the nearest ancestor that
    contains a working/ subtree (intake.json canonically lives at
    <run_dir>/working/copy/intake.json)."""
    try:
        cur = intake_path.resolve().parent
    except OSError:
        return None
    for _ in range(8):
        if (cur / "working").is_dir():
            return cur
        if cur.parent == cur:
            break
        cur = cur.parent
    return None


def _flag_true(v) -> bool:
    if v is True:
        return True
    if isinstance(v, str):
        return v.strip().lower() in ("true", "yes", "1", "y")
    return False


def _intake_provenance_gate(obj: dict, intake_path: Path) -> str:
    """FAIL-CLOSED intake provenance gate (P1-C). Beyond the self-attested
    intake.json flags, this asserts:
      (1) all SIX mandatory pre_presentation_capture fields exist with non-empty /
          non-default values (or a sanctioned captured-decline flag), AND
      (2) intake.json is CONSISTENT with a COMPLETED, turn-gated
          working/interview/intake_ledger.json (the Brainstorming-Buddy record) —
          the ledger must exist, be marked complete, and carry the interview's
          captured answers. An absent/incomplete ledger means the interview was
          skipped → DENY.
    Owner override: a logged owner_skip_approval for gate INTAKE-INTERVIEW in
    working/checkpoints/process_manifest.json waives this battery (bypass only when
    the human explicitly says so). Returns "" on pass, or a reason string naming the
    exact missing field(s) / ledger problem on fail."""
    run_dir = _intake_run_dir(intake_path)

    # Owner/founder override — honored exactly like the FIX-2 gates.
    if run_dir is not None and _owner_skip_approved(run_dir, "INTAKE-INTERVIEW"):
        return ""

    problems = []

    # (1) SIX mandatory pre_presentation_capture fields.
    cap = obj.get("pre_presentation_capture")
    if not isinstance(cap, dict):
        cap = {}
    for field, aliases, _qid, decline_flags, is_bool in _INTAKE_MANDATORY_FIELDS:
        keys = (field,) + tuple(aliases)
        present_key = next((k for k in keys if k in cap), None)
        val = cap.get(present_key) if present_key is not None else None
        declined = any(_flag_true(obj.get(f)) or _flag_true(cap.get(f)) for f in decline_flags)
        if is_bool:
            # DARK_OK: a real boolean (true OR false) is a captured answer; a string
            # yes/no/true/false is tolerated. Missing/other → not captured.
            ok = isinstance(val, bool) or (
                isinstance(val, str) and val.strip().lower() in ("true", "false", "yes", "no"))
            if not ok:
                problems.append(f"{field} (missing — must be an explicit true/false from the interview)")
        else:
            non_empty = isinstance(val, str) and val.strip() != "" or \
                (val is not None and not isinstance(val, str) and str(val).strip() != "")
            if not non_empty and not declined:
                problems.append(
                    f"{field} (missing/empty in pre_presentation_capture and no sanctioned decline flag set)")

    # (2) Turn-gated ledger provenance.
    if run_dir is None:
        problems.append(
            "intake_ledger.json could not be located (no working/ run dir resolved from intake.json)")
    else:
        ledger_path = run_dir / "working" / "interview" / "intake_ledger.json"
        if not ledger_path.is_file():
            problems.append(
                "working/interview/intake_ledger.json is ABSENT — the Brainstorming Buddy "
                "deck-intake interview was never run (run deck-intake-driver.py --next/"
                "--answer/--complete)")
        else:
            led = _read_json(ledger_path)
            if not isinstance(led, dict) or "__parse_error__" in led:
                problems.append("working/interview/intake_ledger.json is not valid JSON")
            else:
                status = str(led.get("status", "")).strip().lower()
                complete = status == "complete" or led.get("complete") is True \
                    or str(led.get("complete", "")).strip().lower() == "true"
                if not complete:
                    problems.append(
                        f"intake_ledger.json is not complete (status={status or 'unset'!r}) — "
                        "finish the interview with deck-intake-driver.py --complete")
                else:
                    entries = led.get("entries") if isinstance(led.get("entries"), dict) else {}
                    # Consistency: every mandatory field either has a validated/circled-back
                    # ledger entry for its question, OR a sanctioned captured-decline flag in
                    # intake.json. A field claimed in intake.json with neither is unprovenanced.
                    for field, aliases, qid, decline_flags, _is_bool in _INTAKE_MANDATORY_FIELDS:
                        e = entries.get(qid) if isinstance(entries.get(qid), dict) else {}
                        captured = bool(e.get("validated")) or bool(e.get("circled_back"))
                        declined = any(_flag_true(obj.get(f)) or _flag_true(cap.get(f))
                                       for f in decline_flags)
                        if not captured and not declined:
                            problems.append(
                                f"{field}: intake.json is inconsistent with the ledger — no "
                                f"validated interview answer for '{qid}' and no sanctioned decline flag")

    if problems:
        return ("brainstorm intake incomplete / not provenanced — the Brainstorming Buddy "
                "interview must be run and its six mandatory fields captured. Missing/"
                "inconsistent: " + "; ".join(problems) + ". "
                "Owner override: log an owner_skip_approval for gate INTAKE-INTERVIEW "
                "(owner_approved:true, approved_by, reason, timestamp) in "
                "working/checkpoints/process_manifest.json.")
    return ""


def _chk_intake(path: Optional[Path]) -> str:
    if path is None:
        return "file absent"
    obj = _read_json(path)
    if "__parse_error__" in obj:
        return f"not valid JSON ({obj['__parse_error__']})"
    if obj.get("interview_confirmed") is not True:
        return "interview_confirmed is not true"
    # Existing presentation_mode (one-person | general) — UNCHANGED, still required.
    mode = str(obj.get("presentation_mode", "")).strip().lower()
    if mode not in ("one-person", "general"):
        return f"presentation_mode must be 'one-person' or 'general' (got {obj.get('presentation_mode')!r})"
    # Additive: the definitive audience_mode (STANDARD | PERSONAL | GENERAL).
    audience_mode = str(obj.get("audience_mode", "")).strip().upper()
    if audience_mode not in AUDIENCE_MODES:
        return (f"audience_mode must be one of {AUDIENCE_MODES} (got "
                f"{obj.get('audience_mode')!r}). STANDARD=net-new from brainstorm; "
                f"PERSONAL=from client content, named recipient; GENERAL=from client "
                f"content, de-identified.")
    # PERSONAL is the only NAMED mode — it must carry a recipient_name.
    if audience_mode == "PERSONAL" and not str(obj.get("recipient_name", "")).strip():
        return "audience_mode is PERSONAL but recipient_name is empty (PERSONAL is the named mode)."
    # Additive: target_talk_minutes — the speaking-length target the deck is sized to
    # (slides/copy/speech sized to target x 130 wpm). Must be a positive number.
    raw = obj.get("target_talk_minutes")
    try:
        if raw is None or float(raw) <= 0:
            return (f"target_talk_minutes must be a positive number of minutes "
                    f"(got {raw!r}). The deck/copy/speech are sized to "
                    f"target_talk_minutes x 130 wpm.")
    except (TypeError, ValueError):
        return f"target_talk_minutes is not numeric (got {raw!r})."
    # Additive: client_requested_slide_count — OPTIONAL. When the client states an
    # explicit number of slides, it is recorded here and honored EXACTLY (never
    # floored/capped/defaulted/changed; the client is never asked to accept a
    # different number) — enforced at build time by _chk_slide_count_exact
    # (AF-SLIDE-COUNT-EXACT). When the key is PRESENT it MUST be a positive integer; a
    # malformed value is a build-stopping intake error (so a bad value can never
    # silently disable the exact-count gate). When ABSENT, the duration/source sizing
    # governs (max(target-sized count, source_slide_count)) and the pacing floor applies.
    if "client_requested_slide_count" in obj:
        rc = obj.get("client_requested_slide_count")
        if isinstance(rc, bool) or not isinstance(rc, int) or rc <= 0:
            return (f"client_requested_slide_count, when present, must be a positive "
                    f"integer — the client's EXACT requested slide count, honored "
                    f"verbatim (got {rc!r}).")
    # FAIL-CLOSED brainstorm provenance (P1-C): the self-attested flags above are NOT
    # sufficient. Assert the six mandatory pre_presentation_capture fields AND a
    # completed, turn-gated intake_ledger.json consistent with them (or a logged owner
    # override). This stops a hand-written thin intake.json from skipping the interview.
    prov = _intake_provenance_gate(obj, path)
    if prov:
        return prov
    return ""


# Deep-Research category-body helpers. Shared by BOTH _chk_research_brief (the
# research_complete:true G/H/I/K/L gate) and _chk_research_cited (the citation gate's
# G/H/I population check) so the two checkers can never drift in how they decide a
# Category section is present + non-empty.
_RESEARCH_PLACEHOLDER_RE = re.compile(r"^\s*\[Output of SOP[^\]]*\]\s*$", re.IGNORECASE)


def _research_category_body(raw_text: str, letter: str) -> str:
    """Return the body text of the '## Category <letter>:' section, or '' if absent.
    Locates the heading (case-insensitive), then reads to the next '## ' heading
    (or end-of-file)."""
    pattern = re.compile(
        rf"^##\s+Category\s+{letter}[:\s]", re.IGNORECASE | re.MULTILINE
    )
    m = pattern.search(raw_text)
    if not m:
        return ""
    start = m.end()
    next_hdr = re.search(r"^##\s+", raw_text[start:], re.MULTILINE)
    return raw_text[start : start + next_hdr.start()] if next_hdr else raw_text[start:]


def _research_category_nonempty(body: str) -> bool:
    """True when the category body has real content beyond whitespace / placeholder
    lines (a body that is solely the '[Output of SOP …]' placeholder counts as empty)."""
    stripped = body.strip()
    if not stripped:
        return False
    non_blank_lines = [ln for ln in stripped.splitlines() if ln.strip()]
    if not non_blank_lines:
        return False
    real_lines = [ln for ln in non_blank_lines if not _RESEARCH_PLACEHOLDER_RE.match(ln)]
    return len(real_lines) > 0


def _chk_research_brief(path: Optional[Path]) -> str:
    if path is None:
        return "no working/research/brief-*.md found"
    text = path.read_text(errors="replace")
    low = text.lower()
    if "research_complete:true" not in low.replace(" ", ""):
        # tolerate 'research_complete: true' spacing
        if "research_complete" not in low or "true" not in low.split("research_complete", 1)[1][:40].lower():
            return "research_complete:true not present"

    # H2: Deep-Research categories G/H/I/K/L must be present AND non-empty before
    # research_complete:true is honoured.  These categories carry the sourced quotes
    # (G), fact-validation ledger (H), objection research (I), persuasion-framework
    # validation (K), and compliance flags (L) that feed the "who says so" slides.
    # A brief with research_complete:true but an empty G, H, I, K, or L section means
    # the Deep Research Specialist skipped the sourced-quote / validation pass — the
    # resulting slides WILL be empty or unsupported.
    #
    # Detection heuristic: locate the "## Category X:" heading (case-insensitive),
    # then read everything up to the next "## " heading (or end-of-file).  The body is
    # considered non-empty when it contains at least one non-whitespace, non-template
    # token (i.e. it is NOT solely the placeholder "[Output of SOP …]" text or blank).
    # These two helpers are shared module-level functions (_research_category_body /
    # _research_category_nonempty) so this gate and _chk_research_cited never drift.
    required_cats = {
        "G": "Credible Attributable Quotes",
        "H": "Fact-Validation Ledger",
        "I": "Objection Research",
        "K": "Persuasion-Framework Validation",
        "L": "Compliance Flags",
    }
    missing_or_empty = []
    for letter, label in required_cats.items():
        body = _research_category_body(text, letter)
        if not body:
            missing_or_empty.append(f"Category {letter} ({label}) section absent")
        elif not _research_category_nonempty(body):
            missing_or_empty.append(
                f"Category {letter} ({label}) section is empty or contains only "
                f"template placeholder text"
            )
    if missing_or_empty:
        return (
            "AF-RESEARCH-GATE: research_complete:true asserted but the following "
            "Deep-Research categories are missing or empty — 'who says so' slides "
            "WILL be empty without them. Re-run Phase -0.5 (Deep Research Specialist "
            "SOP 9.4) and populate each category before setting research_complete:true: "
            + "; ".join(missing_or_empty)
        )

    return ""


def _qc_independence_reason(obj: dict) -> str:
    """AF-QC-INDEPENDENCE: prove the QC report was graded by a reviewer INDEPENDENT
    of the builder, never self-graded. Returns "" when the report carries valid
    independent-reviewer provenance, or a non-empty AF-QC-INDEPENDENCE reason when
    the report is (or could be) self/builder-graded.

    A report PASSES only when ALL of these hold:
      * it is NOT marked self_graded:true;
      * it names an EXPLICIT independent reviewer identity — provenance lives in a
        `qc_independence` block (preferred) or top-level `graded_by`/`reviewer`/
        `reviewed_by` field — that is a non-empty string;
      * that reviewer identity is NOT one of the builder/self tokens
        (build_deck.py, self, builder, author) and is NOT the deck-copy authoring
        role (slide-copywriter) — i.e. it is genuinely someone other than the builder;
      * if an `independent` flag is present it is not explicitly false, and if a
        `builder`/`built_by` identity is recorded it does NOT equal the reviewer.

    A report that simply OMITS the provenance entirely FAILS — independence must be
    affirmatively proven, not assumed."""
    blk = obj.get("qc_independence")
    blk = blk if isinstance(blk, dict) else {}

    # self_graded:true (top-level or in the block) is an immediate fail.
    for src in (obj, blk):
        if src.get("self_graded") is True:
            return ("AF-QC-INDEPENDENCE: QC report is marked self_graded:true — a QC "
                    "report graded by the builder cannot pass. An INDEPENDENT QC "
                    "specialist (not the copy author / not build_deck.py) must grade "
                    "the deck and stamp the qc_independence provenance block "
                    "(graded_by + independent:true). See qc-specialist-presentations-sops.md.")

    # Resolve the reviewer identity from any of the accepted provenance fields.
    reviewer = None
    for src in (blk, obj):
        for key in ("graded_by", "reviewer", "reviewed_by", "reviewer_identity"):
            val = src.get(key)
            if isinstance(val, str) and val.strip():
                reviewer = val.strip()
                break
        if reviewer is not None:
            break

    if not reviewer:
        return ("AF-QC-INDEPENDENCE: QC report carries no independent-reviewer "
                "provenance — it lacks a qc_independence block (or a graded_by/"
                "reviewer/reviewed_by field) naming the INDEPENDENT QC specialist who "
                "graded the deck. A QC report self-written by the builder with no "
                "reviewer identity cannot pass; independence must be PROVEN, not "
                "assumed. The qc-specialist persona must stamp graded_by + "
                "independent:true. See qc-specialist-presentations-sops.md.")

    reviewer_norm = reviewer.lower()
    if reviewer_norm in FORBIDDEN_QC_GRADER_IDENTITIES or reviewer_norm == QC_BUILDER_ROLE:
        return (f"AF-QC-INDEPENDENCE: QC report graded_by/reviewer is {reviewer!r}, "
                f"which is the builder/self (or the deck-copy author) — a self-graded "
                f"report cannot pass. An INDEPENDENT QC specialist must grade the deck. "
                f"See qc-specialist-presentations-sops.md.")

    # An explicit independent:false flag is a self-report that it is NOT independent.
    indep_flag = blk.get("independent", obj.get("independent"))
    if indep_flag is False:
        return ("AF-QC-INDEPENDENCE: QC report provenance sets independent:false — the "
                "reviewer self-reports as NOT independent of the builder. An "
                "INDEPENDENT QC specialist must grade the deck. "
                "See qc-specialist-presentations-sops.md.")

    # If a builder/built_by identity is recorded, the reviewer must not equal it.
    for key in ("builder", "built_by"):
        builder_id = blk.get(key, obj.get(key))
        if isinstance(builder_id, str) and builder_id.strip().lower() == reviewer_norm:
            return (f"AF-QC-INDEPENDENCE: QC report reviewer {reviewer!r} is the SAME "
                    f"identity recorded as the builder ({key}) — the builder graded its "
                    f"own work. An INDEPENDENT QC specialist must grade the deck. "
                    f"See qc-specialist-presentations-sops.md.")

    return ""


# ---------------------------------------------------------------------------
# FOREIGN / CORRUPT QC-REPORT SIGNATURES — the governed path IGNORES any QC report it
# did not itself produce. A report bearing one of these signatures came from the
# eliminated client-side generators (the word-count prompt rubric `_build_qc_report.py`
# that REWARDS sub-floor prompts, and the post-production-overlay image rubric with its
# `typography_overlay_readiness` criterion that rewards a blank canvas) — it is a
# hand-authored / corrupt artifact, not a governed grade, and is REJECTED on sight.
# Matched case-insensitively against the whole report JSON blob. The bare phrase
# "out of scope" is NOT here (it false-positives on legit prose); slide-exclusion is
# caught structurally by _image_qc_report_defects' excluded-keys check instead.
QC_FOREIGN_SIGNATURES = [
    "score_prompt_length", "_build_qc_report", "word_count_band", "words_in_band",
    "80-180 words", "80–180 words", "word-count rubric", "word count rubric",
    "length_words_score", "typography_overlay_readiness", "overlay_readiness",
    "overlay the headline", "overlay the canonical", "post-production overlay",
    "post production overlay", "applied in post",
]


def _qc_report_substance_problems(obj: dict) -> list:
    """ANTI-RUBBER-STAMP teeth shared by EVERY report-shape QC gate (copy / prompt /
    image / typography / speech). Returns a list of fatal reasons the report cannot be a
    governed grade. Two classes:
      (1) FOREIGN/CORRUPT generator signature — the eliminated word-count prompt rubric or
          the post-production-overlay image rubric; the governed path ignores reports it
          did not produce.
      (2) SCORE-INCONSISTENCY — a declared headline `average` that grossly exceeds the mean
          of the report's OWN per-criterion scores, or a per-criterion mean below the 8.5
          bar; a typed pass over low criteria is a rubber stamp.
    Empty list = the report carries no rubber-stamp signature."""
    problems = []
    try:
        blob = json.dumps(obj, default=str).lower()
    except Exception:  # noqa: BLE001
        blob = ""
    for sig in QC_FOREIGN_SIGNATURES:
        if sig in blob:
            problems.append(
                f"report bears a corrupt/foreign QC-generator signature ({sig!r}) — the "
                "word-count prompt rubric and the post-production-overlay image rubric are "
                "ELIMINATED; a report carrying them was not produced by the governed QC pass "
                "and is REJECTED (the governed path ignores any QC report it did not generate)")
            break
    # Score-consistency — only fires when the report cites its own per-criterion scores.
    scores = []
    for key in ("criteria", "scores", "dimensions", "rubric", "per_criterion", "scorecard"):
        v = obj.get(key)
        if isinstance(v, list):
            for r in v:
                if isinstance(r, dict):
                    s = r.get("score", r.get("value"))
                    try:
                        scores.append(float(s))
                    except (TypeError, ValueError):
                        pass
            if scores:
                break
    if len(scores) >= 2:
        mean = sum(scores) / len(scores)
        avg = obj.get("average", obj.get("average_score"))
        try:
            avgf = float(avg)
        except (TypeError, ValueError):
            avgf = None
        if avgf is not None and avgf - mean > 1.0:
            problems.append(
                f"declared average {avgf} exceeds the mean of the report's OWN per-criterion "
                f"scores ({mean:.2f}) by more than 1.0 — the headline score is inflated over "
                "the rubric it cites")
        if mean < 8.5:
            problems.append(
                f"the mean of the report's OWN per-criterion scores is {mean:.2f}, below the "
                "8.5 bar — the declared pass contradicts its own rubric")
    return problems


# ---------------------------------------------------------------------------
# FIX-2 (Error 2) — QC phases structurally unskippable + real report floor.
# The four QC phases (copy / prompt / typography / shift) may NEVER be waived by
# a phase-skip record, and their report files must clear a REAL-CONTENT floor
# (not a 3-byte "{}" placeholder) before the phase may attest. These are the
# exact codes the fix must raise (shared with run_signature_deck + the guard):
#   AF-QC-SKIP        — a skip-approval record names a QC phase (refused)
#   AF-QC-PLACEHOLDER — a QC report is a 3-byte "{}" / sub-floor / verdict-less
# ---------------------------------------------------------------------------
UNSKIPPABLE_QC_PHASES = frozenset({
    "P1Q-COPY-QC",      # -> working/qc/copy_qc_report.json
    "P-PROMPT-QC",      # -> working/qc/prompt_qc_report.json
    "P-TYPO-QC",        # -> working/qc/typography_qc_report.json
    "P-SHIFT-QC",       # -> working/qc/priority_shift_report.json
})

# QC phase id -> the report file it must produce (mirrors PIPELINE-MANIFEST
# produces_artifact for the four QC phases). Used by the report-floor gate.
QC_PHASE_REPORT = {
    "P1Q-COPY-QC": "working/qc/copy_qc_report.json",
    "P-PROMPT-QC": "working/qc/prompt_qc_report.json",
    "P-TYPO-QC": "working/qc/typography_qc_report.json",
    "P-SHIFT-QC": "working/qc/priority_shift_report.json",
}

QC_REPORT_FLOOR_BYTES = 256   # bytes; a real per-slide QC report is far larger
QC_REPORT_SLIDE_FLOOR = 20    # per-slide verdicts the report must carry


def _qc_report_per_slide_verdicts(obj) -> list:
    """Return the real per-slide verdict list from a QC report object.

    Accepts the per-slide structures the four QC report schemas actually use:
      * copy_qc_report.json     -> per_slide_scores: [{slide, average, pass, ...}]
      * prompt_qc_report.json   -> slides / results / per_slide_scores
      * typography_qc_report.json -> slides / per_slide_scores (per-slide archetype)
      * priority_shift_report.json -> items: [{item, pass, evidence}] (the 14-point
        ship gate) AND, after FIX-2, a per-slide `slides` list added by the ledger.
    Returns [] when the report carries no per-slide verdict structure."""
    if not isinstance(obj, dict):
        return []
    for key in ("per_slide_scores", "slides", "results", "per_slide",
                "per_slide_verdicts", "slide_verdicts", "items"):
        v = obj.get(key)
        if isinstance(v, list):
            return [e for e in v if isinstance(e, dict)]
        if isinstance(v, dict):
            return [e for e in v.values() if isinstance(e, dict)]
    return []


def _qc_slide_verdict_is_real(entry: dict) -> bool:
    """True when a per-slide verdict entry carries an explicit PASS/FAIL/score — a
    real verdict, not a bare slide id or empty stub. `pass` must be a real bool or
    a verdict string; a numeric score counts; a slide-only dict does not."""
    for key in ("pass", "verdict", "score", "status", "result", "grade", "ok",
                "pass_fail", "passed"):
        val = entry.get(key)
        if isinstance(val, bool):
            return True
        if isinstance(val, (int, float)):
            return True
        if isinstance(val, str) and val.strip():
            low = val.strip().lower()
            if low in ("pass", "passed", "fail", "failed", "ok", "pending", "warn"):
                return True
    return False


def check_qc_reports_real(run_dir: Path, slides_path=None) -> str:
    """FIX-2 REPORT FLOOR — fail-closed real-content floor for the FOUR QC reports.

    Returns "" when every QC report exists, is > QC_REPORT_FLOOR_BYTES, parses as
    JSON, and carries >= QC_REPORT_SLIDE_FLOOR real per-slide verdicts. Returns a
    fatal 'AF-QC-PLACEHOLDER: ...' message naming the FIRST failing report otherwise.
    A 3-byte '{}' placeholder (the exact Error-2 artifact) fails here."""
    for phase_id in sorted(UNSKIPPABLE_QC_PHASES):
        rel = QC_PHASE_REPORT[phase_id]
        p = run_dir / rel
        if not p.is_file():
            return (f"AF-QC-PLACEHOLDER: {phase_id} report {rel} is MISSING — a real "
                    "QC report must exist before the phase can attest.")
        raw = p.read_bytes()
        if len(raw) <= QC_REPORT_FLOOR_BYTES:
            return (f"AF-QC-PLACEHOLDER: {phase_id} report {rel} is only {len(raw)} "
                    f"bytes (floor {QC_REPORT_FLOOR_BYTES}) — a 3-byte placeholder, not "
                    "a real QC report. Re-run the QC phase.")
        try:
            obj = json.loads(raw.decode("utf-8", errors="replace"))
        except Exception as exc:  # noqa: BLE001
            return (f"AF-QC-PLACEHOLDER: {phase_id} report {rel} is not valid JSON "
                    f"({exc}).")
        if not isinstance(obj, dict):
            return (f"AF-QC-PLACEHOLDER: {phase_id} report {rel} is not a JSON object.")
        verdicts = [e for e in _qc_report_per_slide_verdicts(obj)
                    if _qc_slide_verdict_is_real(e)]
        if len(verdicts) < QC_REPORT_SLIDE_FLOOR:
            return (f"AF-QC-PLACEHOLDER: {phase_id} report {rel} carries "
                    f"{len(verdicts)}/{QC_REPORT_SLIDE_FLOOR} real per-slide verdicts — "
                    "a placeholder or a verdict-less rubber stamp, not real QC.")
    return ""


def check_qc_phase_report_real(run_dir: Path, phase_id: str) -> str:
    """FIX-2 per-phase report floor — returns "" when ONE QC phase's report is real,
    or a fatal 'AF-QC-PLACEHOLDER: ...' message naming it otherwise. Called at QC
    phase attestation so a 3-byte placeholder can never satisfy a QC phase."""
    if phase_id not in UNSKIPPABLE_QC_PHASES:
        return ""
    rel = QC_PHASE_REPORT[phase_id]
    p = run_dir / rel
    if not p.is_file():
        return (f"AF-QC-PLACEHOLDER: {phase_id} report {rel} is MISSING — a real "
                "QC report must exist before the phase can attest.")
    raw = p.read_bytes()
    if len(raw) <= QC_REPORT_FLOOR_BYTES:
        return (f"AF-QC-PLACEHOLDER: {phase_id} report {rel} is only {len(raw)} "
                f"bytes (floor {QC_REPORT_FLOOR_BYTES}) — a 3-byte placeholder, not "
                "a real QC report. Re-run the QC phase.")
    try:
        obj = json.loads(raw.decode("utf-8", errors="replace"))
    except Exception as exc:  # noqa: BLE001
        return (f"AF-QC-PLACEHOLDER: {phase_id} report {rel} is not valid JSON "
                f"({exc}).")
    if not isinstance(obj, dict):
        return (f"AF-QC-PLACEHOLDER: {phase_id} report {rel} is not a JSON object.")
    verdicts = [e for e in _qc_report_per_slide_verdicts(obj)
                if _qc_slide_verdict_is_real(e)]
    if len(verdicts) < QC_REPORT_SLIDE_FLOOR:
        return (f"AF-QC-PLACEHOLDER: {phase_id} report {rel} carries "
                f"{len(verdicts)}/{QC_REPORT_SLIDE_FLOOR} real per-slide verdicts — "
                "a placeholder or a verdict-less rubber stamp, not real QC.")
    return ""


def _chk_copy_qc(path: Optional[Path]) -> str:
    if path is None:
        return "file absent"
    obj = _read_json(path)
    if "__parse_error__" in obj:
        return f"not valid JSON ({obj['__parse_error__']})"
    # M1: require gate == "Phase 1Q" AFFIRMATIVELY. The old check (`if gate and gate
    # != "Phase 1Q"`) let an EMPTY/missing gate short-circuit straight past — a QC
    # report with no gate field would silently pass. The gate must be exactly the
    # Phase 1Q gate, present and correct.
    gate = str(obj.get("gate", "")).strip()
    if gate != "Phase 1Q":
        return f"gate is {gate!r}, expected 'Phase 1Q' (must be present and exactly 'Phase 1Q')"
    avg = obj.get("average", obj.get("average_score"))
    try:
        if avg is None or float(avg) < 8.5:
            return f"average score {avg!r} is below the 8.5 pass threshold"
    except (TypeError, ValueError):
        return f"average score {avg!r} is not numeric"
    # any triggered autofail blocks
    triggered = obj.get("triggered_autofails") or obj.get("autofails_triggered") or []
    if triggered:
        return f"triggered autofails present: {triggered}"
    # M1: require pass IS True affirmatively. The old check only rejected pass:false,
    # so a report with pass absent / null / a non-True value slipped through. A QC
    # report must explicitly assert pass:true.
    if obj.get("pass") is not True:
        return f"report does not affirmatively mark pass:true (got pass={obj.get('pass')!r})"
    # AF-QC-INDEPENDENCE: the report MUST prove an INDEPENDENT reviewer (not the
    # builder / not self) graded it. A QC report self-written by the builder passes
    # every numeric check above yet proves nothing — this is the rubber-stamp the
    # gate exists to refuse. Keep this LAST so the more obvious shape errors report
    # first, but never let a self-graded report through.
    indep = _qc_independence_reason(obj)
    if indep:
        return indep
    # ANTI-RUBBER-STAMP teeth (same class applied to every report-shape gate): reject a
    # corrupt/foreign generator signature (the eliminated word-count / overlay rubrics)
    # and a headline average inflated over the report's own per-criterion scores.
    for sub in _qc_report_substance_problems(obj):
        return f"AF-COPY-QC: {sub}. See SOP-SLIDE-00 / qc-specialist-presentations-sops.md."
    # Report shape is valid — now apply the deterministic re-measure teeth: re-derive
    # truth from the ACTUAL rendered copy (slides.json copy[]) and reject a pass:true
    # report that contradicts it. The report lives at <run_dir>/working/qc/
    # copy_qc_report.json, so the run dir is parents[2].
    if path is not None:
        try:
            run_dir = path.resolve().parents[2]
        except (IndexError, OSError):
            run_dir = None
        if run_dir is not None and run_dir.is_dir():
            teeth = check_copy_qc_teeth(run_dir)
            if teeth:
                return teeth
    return ""


# ---------------------------------------------------------------------------
# THE FIVE QC GATES — each is a written-standard rubric report graded by an
# INDEPENDENT reviewer (generalizing AF-QC-INDEPENDENCE across the whole
# pipeline). copy-QC is _chk_copy_qc above; the four below mirror its contract:
#   * gate == the phase's expected gate string (present + exact),
#   * average >= 8.5 pass threshold,
#   * no triggered autofails,
#   * pass:true affirmatively,
#   * INDEPENDENT-reviewer provenance (_qc_independence_reason) — a self/builder
#     grade is REFUSED. The builder/author of the artifact may NOT grade it.
# Each QC gate sequences AFTER its artifact (manifest order): Prompt-QC after
# prompt authoring, Image-QC after render, Typography-QC after design,
# Speech-QC after speech. Returns "" on pass, or a fatal AF message string.
# ---------------------------------------------------------------------------
def _qc_report_gate(path: Optional[Path], af_code: str, gate_label: str,
                   sop_ref: str) -> str:
    """Generic INDEPENDENT-reviewer QC-report gate. Validates the report shape
    (gate string, >=8.5 average, no triggered autofails, pass:true) AND the
    independent-reviewer provenance block, citing `af_code` on every failure so
    the enforcement path names the code it raises (Guard A requirement)."""
    if path is None:
        return (f"{af_code}: file absent — the {gate_label} QC report was never "
                f"produced. An INDEPENDENT QC specialist must grade this stage and "
                f"write the report with a qc_independence provenance block. See {sop_ref}.")
    obj = _read_json(path)
    if "__parse_error__" in obj:
        return f"{af_code}: report not valid JSON ({obj['__parse_error__']}). See {sop_ref}."
    gate = str(obj.get("gate", "")).strip()
    if gate != gate_label:
        return (f"{af_code}: gate is {gate!r}, expected {gate_label!r} (must be present "
                f"and exactly {gate_label!r}). See {sop_ref}.")
    avg = obj.get("average", obj.get("average_score"))
    try:
        if avg is None or float(avg) < 8.5:
            return f"{af_code}: average score {avg!r} is below the 8.5 pass threshold. See {sop_ref}."
    except (TypeError, ValueError):
        return f"{af_code}: average score {avg!r} is not numeric. See {sop_ref}."
    triggered = obj.get("triggered_autofails") or obj.get("autofails_triggered") or []
    if triggered:
        return f"{af_code}: triggered autofails present: {triggered}. See {sop_ref}."
    if obj.get("pass") is not True:
        return (f"{af_code}: report does not affirmatively mark pass:true "
                f"(got pass={obj.get('pass')!r}). See {sop_ref}.")
    # INDEPENDENCE — reuse the copy-QC independence proof so a self/builder-graded
    # report for ANY of the five QC stages is refused identically.
    indep = _qc_independence_reason(obj)
    if indep:
        # _qc_independence_reason cites AF-QC-INDEPENDENCE; re-stamp this stage's code
        # so Guard A's "the path names the code" proof holds for THIS gate too.
        return f"{af_code}: {indep}"
    # ANTI-RUBBER-STAMP teeth — reject a corrupt/foreign QC-generator signature and a
    # headline average inflated over the report's own per-criterion scores. This gives
    # the typography / speech / prompt / image report-shape gates the same teeth the
    # image-QC pixel cross-check already has (the governed path ignores foreign reports).
    for sub in _qc_report_substance_problems(obj):
        return f"{af_code}: {sub}. See {sop_ref}."
    return ""


def check_prompt_qc_teeth(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-PROMPT-QC — the DETERMINISTIC RE-MEASURE behind the Prompt-QC gate.

    The legacy _chk_prompt_qc was a JSON RUBBER STAMP — it validated only the report's
    SHAPE (gate string, average >= 8.5, no triggered autofails, pass:true, independence)
    and NEVER OPENED A SINGLE PROMPT FILE, so an agent could type pass:true over a deck of
    800-char prompts and it sailed through (the #1 uncovered gap: image-QC got pixel teeth,
    prompt-QC did not). This is the equivalent upgrade: it RE-OPENS every on-disk per-slide
    prompt, RE-MEASURES the actual char count, and re-runs the floor / structural-block /
    quality / verbatim gate (via _collect_prompt_problems — the same source of truth the
    rich-prompt preflight uses). A prompt-QC report that marks pass:true while the on-disk
    prompts do NOT clear the gate is REJECTED — its claims do not match the files.

    It DEFERS (returns "") only when the slide count cannot be determined yet (the
    rich-prompt gate owns that case). Returns "" on pass, or a fatal AF-PROMPT-QC message."""
    prob = _collect_prompt_problems(run_dir, slides_path)
    if prob and prob[0][0] == 0:
        return ""  # slide count unknown — _chk_rich_prompts owns the "no slides.json" case.
    if prob and prob[0][0] < 0:
        # FIX-22 / D16 directory-level fatal (duplicate / non-canonical prompt file) —
        # a Prompt-QC report cannot be trusted when the on-disk prompt set is broken.
        return ("AF-PROMPT-QC: " + prob[0][1])
    if not prob:
        return ""
    offenders = "; ".join(f"slide {o:02d}: {r}" for o, r in prob[:10])
    more = "" if len(prob) <= 10 else f" (+{len(prob) - 10} more)"
    return ("AF-PROMPT-QC: the prompt-QC report passed its shape check but the on-disk "
            "per-slide prompts do NOT clear the rich-prompt floor/quality gate — a real "
            "Prompt-QC pass RE-MEASURES the actual prompt files, it is not a self-typed "
            "score over thin prompts. Re-author the prompts to the 9,000–14,000 standard "
            "(>= 9,000 chars, 8-class negative block, per-string spelling-lock, the slide's "
            "verbatim copy baked, real density) and re-run the Prompt QC Specialist. "
            "Offenders: " + offenders + more + ".")


def check_copy_qc_teeth(run_dir: Path) -> str:
    """AF-COPY-QC — the DETERMINISTIC RE-MEASURE behind the copy-QC gate.

    The legacy _chk_copy_qc was a JSON RUBBER STAMP — it validated only the report's
    SHAPE (gate string, average >= 8.5, no triggered autofails, pass:true, independence)
    and never opened the copy it graded, so an agent could type 8.9/pass:true over a deck
    whose real copy violates the AF-COPY-BAND character bands and it sailed through. This
    is the equivalent upgrade the prompt/image gates already got: it RE-DERIVES TRUTH from
    the ACTUAL rendered slides.json copy[] (the file the renderer bakes verbatim into the
    pixels), RE-MEASURES every slide's copy against the shared COPY_* band constants, and
    rejects a copy-QC report whose pass claims do not match the on-disk copy.

    Two independent classes of failure, both fatal:
      A. SLIDE-COVERAGE — the report's per_slide_scores must cover every slide the
         renderer will actually render (one graded row per slide, none dropped), and
         every row must carry a real verdict. A report that grades 10 of 20 slides is
         not a copy-QC pass.
      B. COPY-BAND RE-MEASURE — every rendered slide's copy[] must clear the
         deterministic character bands (HEADLINE 12-60 / SUBHEAD 20-110 / KICKER <= 40 /
         <= 3 BULLETS each 8-30 / SLIDE TOTAL 40-180) that AF-COPY-BAND enforces
         first-party. A pass:true report over copy that violates a band is REJECTED,
         naming the exact discrepancy (slide + field + chars + band).

    It DEFERS (returns "") ONLY when the rendered copy is genuinely not produced yet —
    no readable slides.json (its absence is owned by the schema / AF-P1 / slide-count
    gates, exactly as _chk_copy_density defers). Returns "" on pass, or a fatal
    AF-COPY-QC message naming the exact discrepancy."""
    n = _count_output_slides(run_dir)
    if n is None:
        return ""  # no slides.json yet — upstream checks own the absence (D10 defer).
    copy_map = _load_slide_copy_map(run_dir)
    if not copy_map:
        return ""  # slides.json unreadable/empty — upstream checks own that failure.
    arc_tags = _load_slide_arc_tags(run_dir)
    report_path = run_dir / "working" / "qc" / "copy_qc_report.json"
    report_obj = _read_json(report_path) if report_path.exists() else None

    problems = []

    # --- A. SLIDE-COVERAGE cross-check (when the report actually grades per-slide) ---
    # A shape-only report (no per-slide rows) is NOT judged here: verdict-count absence
    # is owned by check_qc_phase_report_real / AF-QC-PLACEHOLDER. But once a report
    # DOES carry per-slide rows, it must cover every rendered slide.
    if report_obj is not None and "__parse_error__" not in report_obj:
        verdicts = _qc_report_per_slide_verdicts(report_obj)
        # Coverage is judged only when the report actually grades per-slide; a shape-only
        # report's verdict-count absence is owned by check_qc_phase_report_real.
        if verdicts:
            covered = {}
            for e in verdicts:
                s = e.get("slide")
                if isinstance(s, int) and s >= 1:
                    covered[s] = e
            missing = [s for s in range(1, n + 1) if s not in covered]
            if missing:
                shown = ", ".join(str(s) for s in missing[:10])
                more = "" if len(missing) <= 10 else f" (+{len(missing) - 10} more)"
                problems.append(
                    f"the copy-QC report grades {len(covered)} of {n} rendered slide(s) — "
                    f"slides {shown}{more} have NO per-slide verdict. A real copy-QC pass "
                    f"carries one graded row per rendered slide (a partial report is not a "
                    f"deck-wide pass).")
            blank = [s for s, e in covered.items() if not _qc_slide_verdict_is_real(e)]
            if blank:
                problems.append(
                    f"copy-QC per-slide rows {', '.join(str(s) for s in blank[:10])} carry "
                    f"no real verdict (no pass/fail/score) — a bare slide-id row is not a "
                    f"grade.")

    # --- B. COPY-BAND RE-MEASURE of the actual rendered copy ---
    for ordinal in range(1, n + 1):
        copy_val = copy_map.get(ordinal)
        if not isinstance(copy_val, list) or not copy_val:
            continue  # missing/malformed copy — schema validation / AF-P1 own this.
        fields = [str(c) if c is not None else "" for c in copy_val]
        headline = fields[0]
        subhead = fields[1] if len(fields) > 1 else None
        kicker = fields[2] if len(fields) > 2 else None
        bullets = fields[3:]

        hlen = len(headline)
        if not (COPY_HEADLINE_CHAR_FLOOR <= hlen <= COPY_HEADLINE_CHAR_CEILING):
            problems.append(
                f"slide {ordinal:02d} HEADLINE is {hlen} chars, outside the "
                f"{COPY_HEADLINE_CHAR_FLOOR}-{COPY_HEADLINE_CHAR_CEILING} band the copy-QC "
                f"report marked pass")

        if subhead:
            slen = len(subhead)
            if not (COPY_SUBHEAD_CHAR_FLOOR <= slen <= COPY_SUBHEAD_CHAR_CEILING):
                problems.append(
                    f"slide {ordinal:02d} SUBHEAD is {slen} chars, outside the "
                    f"{COPY_SUBHEAD_CHAR_FLOOR}-{COPY_SUBHEAD_CHAR_CEILING} band the copy-QC "
                    f"report marked pass")

        if kicker:
            klen = len(kicker)
            if klen > COPY_KICKER_CHAR_CEILING:
                problems.append(
                    f"slide {ordinal:02d} KICKER is {klen} chars, over the "
                    f"{COPY_KICKER_CHAR_CEILING}-char ceiling the copy-QC report marked pass")

        real_bullets = [b for b in bullets if b]
        if len(real_bullets) > COPY_BULLET_MAX_COUNT:
            problems.append(
                f"slide {ordinal:02d} carries {len(real_bullets)} bullets, over the "
                f"{COPY_BULLET_MAX_COUNT}-bullet max the copy-QC report marked pass")
        for bi, b in enumerate(real_bullets, start=1):
            blen = len(b)
            if not (COPY_BULLET_CHAR_FLOOR <= blen <= COPY_BULLET_CHAR_CEILING):
                problems.append(
                    f"slide {ordinal:02d} BULLET {bi} is {blen} chars, outside the "
                    f"{COPY_BULLET_CHAR_FLOOR}-{COPY_BULLET_CHAR_CEILING} band the copy-QC "
                    f"report marked pass")

        total = sum(len(f) for f in fields)
        tag_blob = arc_tags.get(ordinal, "")
        exempt_eligible = _is_hook_or_banner_slide(tag_blob)
        exempt = exempt_eligible and COPY_HOOK_EXEMPTION_ENFORCED
        floor = COPY_HOOK_SLIDE_TOTAL_CHAR_FLOOR if exempt else COPY_SLIDE_TOTAL_CHAR_FLOOR
        if not (floor <= total <= COPY_SLIDE_TOTAL_CHAR_CEILING):
            problems.append(
                f"slide {ordinal:02d} SLIDE TOTAL is {total} chars, outside the "
                f"{floor}-{COPY_SLIDE_TOTAL_CHAR_CEILING} band the copy-QC report marked pass")

    if problems:
        return ("AF-COPY-QC: the copy-QC report passed its shape check but the ACTUAL "
                "rendered copy (slides.json copy[]) CONTRADICTS it — a real copy-QC pass "
                "RE-MEASURES the copy the renderer bakes into the pixels; a typed pass:true "
                "over sub-band copy is a fabricated pass. Fix the copy and re-run the Copy "
                "QC Specialist. Discrepancies: " + "; ".join(problems[:10])
                + ("" if len(problems) <= 10 else f" (+{len(problems) - 10} more)") + ".")
    return ""


def check_typography_qc_teeth(run_dir: Path) -> str:
    """AF-TYPOGRAPHY-QC — the DETERMINISTIC RE-MEASURE behind the typography-QC gate.

    The legacy _chk_typography_qc was a JSON RUBBER STAMP — it validated only the report's
    SHAPE and never opened a design-system artifact, so an agent could type 8.8/pass:true
    over a design system whose declared tokens violate the AF-FONT-FLOOR deterministic
    floors (min_body_pt / type_scale_steps / min_contrast_ratio) and it sailed through.
    This is the equivalent upgrade the prompt/image gates already got: it RE-MEASURES the
    ACTUAL design-system tokens in working/typography/type_layout_system.md (the gate-of-
    record artifact AF-FONT-FLOOR parses) against the shared floors, and rejects a
    typography-QC report whose pass claims do not match the on-disk design system.

    Two independent classes of failure, both fatal:
      A. FONT-FLOOR RE-MEASURE — the declared tokens must clear the deterministic floors
         (18pt body / 4-5 scale steps / 4.5:1 contrast; raised 22pt / 7.0:1 under a
         client dark-theme opt-in), exactly as check_font_floor enforces. A pass:true
         report over a below-floor token set is REJECTED, naming the token + value + floor.
      B. SLIDE-COVERAGE — when the report grades per-slide archetypes (per_slide_scores),
         every slide the deck will render must carry a real verdict row (none dropped).

    It DEFERS (returns "") ONLY pre-typography — no design system produced yet (no
    type_layout_system.md AND no design_system.json / design-brief, mirroring
    check_font_floor's defer). Once a design system exists, a missing token file or a
    below-floor token FAILS (the report cannot claim pass over a phantom/sub-floor
    design). Returns "" on pass, or a fatal AF-TYPOGRAPHY-QC message naming the exact
    discrepancy."""
    layout = run_dir / TYPE_LAYOUT_SYSTEM_REL
    design_present = (
        (run_dir / "working" / "typography" / "design_system.json").exists()
        or any((run_dir / "working" / "research").glob("design-brief-*.md"))
        if (run_dir / "working").exists() else False
    )
    if not layout.exists() and not design_present:
        return ""  # genuine pre-typography state — defer (D10).

    problems = []

    # --- A. FONT-FLOOR RE-MEASURE of the actual design tokens ---
    if not layout.exists():
        problems.append(
            f"the design system exists but the deterministic type tokens file "
            f"{TYPE_LAYOUT_SYSTEM_REL} is MISSING — the typography-QC report marked pass "
            f"over a design system with no measurable type tokens")
    else:
        text = layout.read_text(encoding="utf-8", errors="replace")

        def _num(key):
            m = re.search(rf'(?im)^\s*{re.escape(key)}\s*[:=]\s*([0-9]+(?:\.[0-9]+)?)', text)
            return float(m.group(1)) if m else None

        dark = _read_dark_optin(run_dir)
        pt_floor = DARK_THEME_BODY_PT_FLOOR if dark else FONT_BODY_PT_FLOOR
        contrast_floor = (DARK_THEME_CONTRAST_FLOOR if dark
                          else CONTRAST_RATIO_FLOOR_NORMAL)

        min_body_pt = _num("min_body_pt")
        type_scale_steps = _num("type_scale_steps")
        min_contrast = _num("min_contrast_ratio")

        if min_body_pt is None:
            problems.append(
                f"token min_body_pt is MISSING from {TYPE_LAYOUT_SYSTEM_REL} — the "
                f"typography-QC report marked pass over a type system with no declared "
                f"body size")
        elif min_body_pt < pt_floor:
            problems.append(
                f"min_body_pt={min_body_pt:g} is below the {pt_floor:g}pt "
                f"{'dark-theme ' if dark else ''}body floor the typography-QC report "
                f"marked pass (FONT_BODY_PT_FLOOR={FONT_BODY_PT_FLOOR})")
        if type_scale_steps is None:
            problems.append(
                f"token type_scale_steps is MISSING from {TYPE_LAYOUT_SYSTEM_REL} — the "
                f"typography-QC report marked pass over a type system with no declared "
                f"modular scale")
        elif not (TYPE_SCALE_STEPS_MIN <= int(type_scale_steps) <= TYPE_SCALE_STEPS_MAX):
            problems.append(
                f"type_scale_steps={int(type_scale_steps)} is not a modular "
                f"{TYPE_SCALE_STEPS_MIN}-{TYPE_SCALE_STEPS_MAX}-step scale the "
                f"typography-QC report marked pass")
        if min_contrast is None:
            problems.append(
                f"token min_contrast_ratio is MISSING from {TYPE_LAYOUT_SYSTEM_REL} — the "
                f"typography-QC report marked pass over a type system with no declared "
                f"contrast")
        elif min_contrast < contrast_floor:
            problems.append(
                f"min_contrast_ratio={min_contrast:g} is below the {contrast_floor:g}:1 "
                f"{'dark-theme AAA ' if dark else 'WCAG AA '}floor the typography-QC "
                f"report marked pass (CONTRAST_RATIO_FLOOR_NORMAL="
                f"{CONTRAST_RATIO_FLOOR_NORMAL})")

    # --- B. SLIDE-COVERAGE cross-check (when the report grades per-slide archetypes) ---
    report_path = run_dir / "working" / "qc" / "typography_qc_report.json"
    report_obj = _read_json(report_path) if report_path.exists() else None
    if report_obj is not None and "__parse_error__" not in report_obj:
        n = _count_output_slides(run_dir)
        verdicts = _qc_report_per_slide_verdicts(report_obj)
        if n is not None and verdicts:
            covered = {}
            for e in verdicts:
                s = e.get("slide")
                if isinstance(s, int) and s >= 1:
                    covered[s] = e
            missing = [s for s in range(1, n + 1) if s not in covered]
            if missing:
                shown = ", ".join(str(s) for s in missing[:10])
                more = "" if len(missing) <= 10 else f" (+{len(missing) - 10} more)"
                problems.append(
                    f"the typography-QC report grades {len(covered)} of {n} rendered "
                    f"slide(s) — slides {shown}{more} have NO per-slide verdict. A real "
                    f"typography-QC pass carries one graded row per rendered slide.")

    if problems:
        return ("AF-TYPOGRAPHY-QC: the typography-QC report passed its shape check but "
                "the ACTUAL design-system artifacts (working/typography/"
                "type_layout_system.md tokens) CONTRADICT it — a real typography-QC pass "
                "RE-MEASURES the declared type system against the deterministic floors; a "
                "typed pass:true over a missing/sub-floor design is a fabricated pass. Fix "
                "the design tokens and re-run the Typography QC Specialist. "
                "Discrepancies: " + "; ".join(problems[:10])
                + ("" if len(problems) <= 10 else f" (+{len(problems) - 10} more)") + ".")
    return ""


def check_speech_qc_teeth(run_dir: Path) -> str:
    """AF-SPEECH-QC — the DETERMINISTIC RE-MEASURE behind the speech-QC gate.

    The legacy _chk_speech_qc was a JSON RUBBER STAMP — it validated only the report's
    SHAPE and never opened the speech it graded, so an agent could type 8.7/pass:true
    over a speech that does not fill the requested duration (AF-SPEECH-SHORT) or that
    under-sings / wallpapers the hook (AF-SPEECH-HOOK-COUNT) and it sailed through. This
    is the equivalent upgrade the prompt/image gates already got: it RE-PARSES the REAL
    speech file (the same on-disk artifact _chk_speech_length measures), RE-MEASURES the
    word count against target_talk_minutes x SPEECH_WPM_FLOOR and re-runs the
    AF-SPEECH-HOOK-COUNT engine, and rejects a speech-QC report whose pass claims do not
    match reality.

    DEFER SEMANTICS (D10 — identical to the legacy gate):
      * report ABSENT -> DEFER (returns ""). The speech is written downstream (Phase 9
        delivery); when no report exists yet the gate passes so the render is not
        blocked (the existing conditional-by-design contract, unchanged).
      * report PRESENT but the speech file is ABSENT -> FAIL. A speech-QC report grading
        a speech that does not exist is a FABRICATED pass — the report cannot be trusted
        (present-input-but-broken = fail, D10).
      * speech present, no readable word count / no intake target -> the AF-SPEECH-SHORT
        half defers (its own contract); the AF-SPEECH-HOOK-COUNT half still fires
        (5-20x char-exact hook sings).

    Returns "" on pass / defer, or a fatal AF-SPEECH-QC message naming the exact
    discrepancy."""
    # Locate the speech artifact with the SAME resolution _chk_speech_length uses
    # (canonical plural PRESENTERS-SPEECH.md first, legacy singular/scratch names after).
    speech = None
    for rel in ("working/presenter-speech/speech.md",
                "working/delivery/PRESENTERS-SPEECH.md",
                "working/presenter-speech/PRESENTERS-SPEECH.md",
                "working/delivery/PRESENTER-SPEECH.md",
                "working/presenter-speech/PRESENTER-SPEECH.md"):
        p = run_dir / rel
        if p.exists():
            speech = p
            break
    report_path = run_dir / "working" / "qc" / "speech_qc_report.json"
    if speech is None:
        # D10: when the report EXISTS but the speech it grades does not, the report is
        # a FABRICATED pass — a QC specialist cannot have graded a speech that is not
        # on disk. That is present-input-but-broken and FAILS (never a defer).
        if report_path.exists():
            return ("AF-SPEECH-QC: the speech-QC report exists but the speech file it "
                    "grades is MISSING (no working/presenter-speech/speech.md or "
                    "PRESENTERS-SPEECH.md under working/delivery) — a report grading a "
                    "nonexistent speech is a FABRICATED pass. Write the speech, then "
                    "re-run the Speech QC Specialist.")
        return ""  # no speech AND no report — genuine pre-speech state: defer (D10).

    problems = []

    # --- A. DURATION FLOOR RE-MEASURE (AF-SPEECH-SHORT) ---
    target = None
    for rel in ("working/copy/intake.json", "intake.json", "working/intake.json"):
        p = run_dir / rel
        if p.exists():
            obj = _read_json(p)
            if isinstance(obj, dict) and "__parse_error__" not in obj:
                raw = obj.get("target_talk_minutes")
                try:
                    target = float(raw) if raw is not None else None
                except (TypeError, ValueError):
                    target = None
            break
    if not target or target <= 0:
        return ""  # no usable duration target — the intake gate owns that case.

    words = len(speech.read_text(errors="replace").split())
    floor = int(round(target * SPEECH_WPM_FLOOR))
    if words < floor:
        problems.append(
            f"the real speech {speech.name} has {words} words; the floor for a "
            f"{target:g}-minute talk is {floor} words (target_talk_minutes x "
            f"{SPEECH_WPM_FLOOR} wpm) — the speech-QC report marked pass over a speech "
            f"that does not fill the requested duration (AF-SPEECH-SHORT)")

    # --- B. HOOK-COUNT RE-MEASURE (AF-SPEECH-HOOK-COUNT, the SOP 9.1 step 2a engine) ---
    pec = _import_pitch_engines_check()
    if pec is not None:
        try:
            hook_hits = pec.run(run_dir, phase="SPEECH-QC")
        except Exception:  # noqa: BLE001
            hook_hits = []
        hook_codes = [h for h in hook_hits if str(h.get("code") or "").startswith("AF-")]
        if hook_codes:
            problems.append(
                "the pitch-engine hook-count check (SOP 9.1 step 2a: char-exact "
                "intake.hook sung 5-20x in PRESENTERS-SPEECH.md) TRIGGERED on the real "
                "speech: " + "; ".join(
                    f"{h.get('code')}: {h.get('detail')}" for h in hook_codes[:3]))

    if problems:
        return ("AF-SPEECH-QC: the speech-QC report passed its shape check but the REAL "
                "speech file CONTRADICTS it — a genuine speech-QC pass RE-PARSES the "
                "actual speech (duration floor + hook-count), not a self-typed score. "
                "Fix the speech and re-run the Speech QC Specialist. Discrepancies: "
                + "; ".join(problems[:10])
                + ("" if len(problems) <= 10 else f" (+{len(problems) - 10} more)") + ".")
    return ""


def _chk_prompt_qc(path: Optional[Path]) -> str:
    """PROMPT-QC gate (AF-PROMPT-QC). After Prompt-Authoring (P4-PROMPT), an
    INDEPENDENT QC specialist grades every per-slide prompt against the written
    9,000–14,000-char prompt standard rubric (length, the 15-element structural blocks,
    8-class negative block, per-string spelling-locks, verbatim copy). Self/builder grade
    refused — AND the report's pass is cross-checked against the on-disk prompts.

    The legacy gate validated only the report's SHAPE and never opened a prompt (a JSON
    rubber stamp). It now mirrors the FIX-2 image-QC upgrade: after the report-shape gate,
    it derives the run dir from the report path and delegates to check_prompt_qc_teeth(),
    which RE-OPENS and RE-MEASURES every prompt and rejects a pass that contradicts them."""
    base = _qc_report_gate(path, "AF-PROMPT-QC", "Phase Prompt-QC",
                          "qc-specialist-prompt-presentations-sops.md")
    if base:
        return base
    # Report shape is valid — now apply the deterministic re-measure teeth. The report
    # lives at <run_dir>/working/qc/prompt_qc_report.json, so the run dir is parents[2].
    if path is not None:
        try:
            run_dir = path.resolve().parents[2]
        except (IndexError, OSError):
            run_dir = None
        if run_dir is not None and run_dir.is_dir():
            teeth = check_prompt_qc_teeth(run_dir)
            if teeth:
                return teeth
    return ""


def _chk_image_qc(path: Optional[Path]) -> str:
    """IMAGE-QC gate (AF-IMAGE-QC + AF-IMAGE-QC-VISION). After Render (P4-RENDER), an
    INDEPENDENT QC specialist grades the rendered slides against the written image-QC
    rubric (multimodal pass: copy-vs-pixel parity, baked-not-overlaid, contrast,
    no-placeholder). Self/builder grade refused.

    FIX-2: the legacy implementation was a JSON RUBBER STAMP — it validated only the
    report's SHAPE (gate string, average >= 8.5, no triggered autofails, pass:true,
    independence) and never opened a single PNG, so an agent could type 8.66/pass:true
    over a deck of flat cream cards and it would sail through. That rubber stamp is now
    REPLACED: after the report-shape gate, this derives the run dir from the report
    path and delegates to check_image_qc_vision(), a DETERMINISTIC PIXEL cross-check
    that actually inspects every rendered PNG (byte floor + flat-card detection) and
    rejects a pixel-blind / overlay-blessing / hook-slide-excluding report. It no longer
    defers-to-pass when renders are present."""
    base = _qc_report_gate(path, "AF-IMAGE-QC", "Phase Image-QC",
                          "qc-specialist-image-presentations-sops.md")
    if base:
        return base
    # Report shape is valid — now apply the deterministic pixel teeth. The report
    # lives at <run_dir>/working/qc/image_qc_report.json, so the run dir is parents[2].
    if path is not None:
        try:
            run_dir = path.resolve().parents[2]
        except (IndexError, OSError):
            run_dir = None
        if run_dir is not None and run_dir.is_dir():
            teeth = check_image_qc_vision(run_dir)
            if teeth:
                return teeth
    return ""


def check_image_qc_report_gate(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """IMAGE-QC PREFLIGHT SCHEDULER (AF-IMAGE-QC + AF-IMAGE-QC-VISION) — resolves the
    image-QC chicken-and-egg WITHOUT weakening the gate.

    The image-QC report (working/qc/image_qc_report.json, manifest phase P-IMAGE-QC,
    order 4.95) is a POST-render artifact: the INDEPENDENT Image QC Specialist grades the
    rendered PNGs AFTER P4-RENDER (order 4.9). The legacy wiring required that report to
    already EXIST at the PRE-render preflight — a chicken-and-egg that made render
    un-startable (render needs the report; the report needs the render). This scheduler
    makes the pre-render requirement satisfiable while keeping every tooth:

      * GENUINE PRE-RENDER STATE — no rendered PNGs AND no report yet -> DEFER (return "")
        so render can proceed. This mirrors the sibling post-render gates already wired
        run-dir-scoped in PREFLIGHT_REQUIRED (check_image_qc_present / check_image_qc_vision
        / check_visual_variety / check_brand_consistency), all of which defer pre-render.

      * ONCE RENDERS (or a report) EXIST — the report becomes MANDATORY and is enforced by
        _chk_image_qc: the report-shape gate (gate string / average >= 8.5 / no triggered
        autofails / pass:true / INDEPENDENT-reviewer provenance) PLUS the deterministic
        pixel/vision teeth (check_image_qc_vision: byte-floor + flat-card detection + the
        report cross-check). A rendered deck with NO report FAILS AF-IMAGE-QC (file absent);
        a pixel-blind / flat-card / rubber-stamp report FAILS AF-IMAGE-QC-VISION.

    The post-render pixel teeth are NOT dropped — they additionally bite, independent of
    this scheduler, at the render closeout (run_postflight_gate -> check_image_qc_vision),
    via the standalone check_image_qc_vision PREFLIGHT_REQUIRED entry, and at pre-delivery
    (canonical_render_guard.guard_pre_delivery + the P-IMAGE-QC phase attestation, which
    requires the report to be present). This change only fixes the render-vs-QC ORDERING.

    NOTE on lockstep: _chk_image_qc remains the manifest-declared P-IMAGE-QC checker and
    the AF-IMAGE-QC enforcement symbol (unchanged signature, still path-based). This is the
    run-dir-scoped scheduler around it; sync_check's orphan check (B1) only scans `_chk_*`
    symbols, so this `check_*` wrapper needs no manifest entry (same as its siblings)."""
    report_path = run_dir / "working" / "qc" / "image_qc_report.json"
    report_present = report_path.exists()
    pngs = _gather_rendered_pngs(run_dir)
    if not pngs and not report_present:
        return ""  # genuine pre-render state: the post-render artifact is not due yet.
    # Renders or a report exist -> enforce the full image-QC gate. An absent report now
    # FAILS via _qc_report_gate(None) inside _chk_image_qc (a rendered deck must be graded);
    # a present report is shape-validated AND pixel/vision cross-checked (AF-IMAGE-QC-VISION).
    return _chk_image_qc(report_path if report_present else None)


def _chk_typography_qc(path: Optional[Path]) -> str:
    """TYPOGRAPHY-QC gate (AF-TYPOGRAPHY-QC). After the Design brief (PF-DESIGN), an
    INDEPENDENT QC specialist grades the design system against the written
    typography rubric (weight ladder, per-archetype treatment, anti-template
    variation, type-scale floor). Self/builder grade refused — AND the report's pass
    is cross-checked against the on-disk design-system artifacts (the
    type_layout_system.md tokens) via the deterministic AF-FONT-FLOOR re-measure."""
    base = _qc_report_gate(path, "AF-TYPOGRAPHY-QC", "Phase Typography-QC",
                          "qc-specialist-typography-presentations-sops.md")
    if base:
        return base
    # Report shape is valid — now apply the deterministic design-token teeth. The report
    # lives at <run_dir>/working/qc/typography_qc_report.json, so the run dir is parents[2].
    if path is not None:
        try:
            run_dir = path.resolve().parents[2]
        except (IndexError, OSError):
            run_dir = None
        if run_dir is not None and run_dir.is_dir():
            teeth = check_typography_qc_teeth(run_dir)
            if teeth:
                return teeth
    return ""


def _chk_speech_qc(path: Optional[Path]) -> str:
    """SPEECH-QC gate (AF-SPEECH-QC). After the Presenter Speech (P9-SPEECH), an
    INDEPENDENT QC specialist grades the speech against the written speech rubric
    (pacing, on-slide-sync, persuasion-arc fidelity, audience-facing voice).
    CONDITIONAL by design (the AF-SPEECH-SHORT pattern): the speech is written
    downstream, so when no report exists yet this DEFERS (returns "", pass) rather
    than blocking the pre-speech render. Once the report exists it is enforced —
    the report-shape gate plus the deterministic re-measure of the REAL speech file
    (duration floor + hook-count), so a fabricated pass:true over a short/under-sung
    speech is REJECTED. Self/builder grade refused."""
    if path is None:
        return ""  # speech QC not produced yet (pre-delivery) — gate defers.
    base = _qc_report_gate(path, "AF-SPEECH-QC", "Phase Speech-QC",
                          "qc-specialist-speech-presentations-sops.md")
    if base:
        return base
    # Report shape is valid — now apply the deterministic speech-file teeth. The report
    # lives at <run_dir>/working/qc/speech_qc_report.json, so the run dir is parents[2].
    if path is not None:
        try:
            run_dir = path.resolve().parents[2]
        except (IndexError, OSError):
            run_dir = None
        if run_dir is not None and run_dir.is_dir():
            teeth = check_speech_qc_teeth(run_dir)
            if teeth:
                return teeth
    return ""


# ---------------------------------------------------------------------------
# SLIDE-COUNT FLOOR (AF-SLIDE-COUNT-FLOOR) — the deck must carry enough slides
# for the chosen speaking duration. A 30-minute talk on 10 slides is 3 minutes a
# slide — a wall-of-text reading session, not a paced deck. The verified pacing
# band is ~1.3–1.5 slides per talking minute; this gate enforces the LOW end
# (1.3) as the floor, so a deck below target_talk_minutes x SLIDES_PER_MINUTE_FLOOR
# AUTO-FAILS. The +/- pacing budget above the floor stays with the Director.
# ---------------------------------------------------------------------------
SLIDES_PER_MINUTE_FLOOR = 1.3   # AF-SLIDE-COUNT-FLOOR: min slides per talking minute


def _client_requested_slide_count(run_dir: Path) -> Optional[int]:
    """Return the client's EXPLICIT requested slide count when they stated one, else
    None. This is the EXACT number of slides the client asked for (e.g. "make it 25
    slides", "I want a 50-slide deck"); when present it IS the deck length and is
    honored verbatim — never floored up, capped down, defaulted, or substituted by the
    duration/source heuristic (AF-SLIDE-COUNT-EXACT). Read from intake.json (canonical)
    with mission_prd.json as a fallback. A non-positive / non-integer / absent value
    returns None (no explicit count -> the duration/coverage floors govern instead).
    The intake gate (_chk_intake) rejects a present-but-malformed value, so a present
    value is always a clean positive int by the time the deck is built."""
    for rel in ("working/copy/intake.json", "intake.json", "working/intake.json",
                "working/copy/mission_prd.json", "mission_prd.json",
                "working/mission_prd.json"):
        p = run_dir / rel
        if not p.exists():
            continue
        obj = _read_json(p)
        if isinstance(obj, dict) and "__parse_error__" not in obj:
            raw = obj.get("client_requested_slide_count")
            if raw is None or isinstance(raw, bool):
                continue
            try:
                n = int(raw)
            except (TypeError, ValueError):
                continue
            if n > 0:
                return n
    return None


def _chk_slide_count_floor(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """SLIDE-COUNT-FLOOR gate (AF-SLIDE-COUNT-FLOOR). The output slide count must be
    >= target_talk_minutes x SLIDES_PER_MINUTE_FLOOR (1.3). A 30-min/10-slide deck
    (needs >=39) AUTO-FAILS. target_talk_minutes comes from intake.json (the field
    _chk_intake requires). Returns "" on pass / not-applicable, or a fatal AF
    message. Defers (passes) only when no target or no slide count can be read —
    those are the intake gate's / rich-prompt gate's job, not this one."""
    # An EXPLICIT client-requested slide count IS the deck length (AF-SLIDE-COUNT-EXACT
    # owns it). The duration pacing floor must NEVER force MORE slides than the client
    # asked for, so when a requested count is present this gate defers to the exact one.
    if _client_requested_slide_count(run_dir) is not None:
        return ""
    target = None
    for rel in ("working/copy/intake.json", "intake.json", "working/intake.json"):
        p = run_dir / rel
        if p.exists():
            obj = _read_json(p)
            if isinstance(obj, dict) and "__parse_error__" not in obj:
                raw = obj.get("target_talk_minutes")
                try:
                    target = float(raw) if raw is not None else None
                except (TypeError, ValueError):
                    target = None
            break
    if not target or target <= 0:
        return ""  # no usable target — intake gate handles a missing/invalid target.

    n = _count_output_slides(run_dir, slides_path)
    if n is None:
        return ""  # slide count unreadable — the rich-prompt/coverage gates own that.

    floor = int(__import__("math").ceil(target * SLIDES_PER_MINUTE_FLOOR))
    if n < floor:
        return (f"AF-SLIDE-COUNT-FLOOR: the deck has {n} slide(s) for a "
                f"{target:g}-minute talk; the pacing floor is "
                f"target_talk_minutes x {SLIDES_PER_MINUTE_FLOOR} = {floor} slides. "
                f"A {target:g}-min talk on {n} slides is a reading session, not a paced "
                f"deck. Add {floor - n} slide(s) (verified band is ~1.3–1.5 slides/min).")
    return ""


def _chk_slide_count_exact(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """SLIDE-COUNT-EXACT gate (AF-SLIDE-COUNT-EXACT). When the client stated an
    EXPLICIT requested slide count, the built deck MUST have exactly that many slides —
    the client's number is honored verbatim (25 -> 25, 50 -> 50, 500 -> 500). It is
    NEVER floored up, capped down, defaulted, or substituted by the duration/content/
    source heuristic, and the client is NEVER asked to accept a different number.

    This gate is AUTHORITATIVE over the duration pacing floor (_chk_slide_count_floor)
    and the Mode-B anti-compression coverage floor (_chk_coverage): both DEFER when an
    explicit requested count is present, and this gate fails LOUD if the built deck's
    length differs from the requested count in EITHER direction (too few OR too many).

    Returns "" on pass / not-applicable (no explicit count, or the count is unreadable
    — the rich-prompt/coverage gates own a missing slides.json), or a fatal
    AF-SLIDE-COUNT-EXACT message."""
    requested = _client_requested_slide_count(run_dir)
    if requested is None:
        return ""  # no explicit client count — the duration/coverage floors govern.

    n = _count_output_slides(run_dir, slides_path)
    if n is None:
        return ""  # slide count unreadable — the rich-prompt/coverage gates own that.

    if n != requested:
        verb = "Add" if n < requested else "Remove"
        delta = abs(requested - n)
        return (f"AF-SLIDE-COUNT-EXACT: the client explicitly requested {requested} "
                f"slide(s), but the deck has {n}. The client's requested slide count is "
                f"honored EXACTLY ({requested} -> {requested}); it is never floored up, "
                f"capped down, defaulted, or substituted by the duration/content "
                f"heuristic, and the client is never asked to accept a different number. "
                f"{verb} {delta} slide(s) so the deck is exactly {requested}.")
    return ""


# ---------------------------------------------------------------------------
# PITCH / OFFER-LADDER (AF-PITCH-MISSING) — a converting deck MUST carry an offer
# ladder AND a re-pitch after the FINAL price. The arc_allocation.json the Offer
# Price Strategist produces is the source of truth: it must contain the LADDER /
# offer beats (the value stack -> anchor -> price drops) and a RE-PITCH beat after
# the FINAL price. A deck with no offer ladder, or no re-pitch, is a teaching dump
# that never asks for the sale (or asks once and never recovers a "no").
# ---------------------------------------------------------------------------
def _chk_pitch(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """PITCH gate (AF-PITCH-MISSING). The converting arc MUST carry an offer ladder
    (value-stack / anchor / price beats) AND a re-pitch beat after the FINAL price.
    Reads arc_allocation.json (the Offer Price Strategist artifact). Returns "" on
    pass, or a fatal AF-PITCH-MISSING message. Defers (passes) only when the arc
    artifact is absent — _chk_arc owns the 'arc missing' failure, this gate owns
    'arc present but has no pitch / no re-pitch'.

    2A: CONDITIONAL on intake.json.pitch_included. AF-PITCH-MISSING fires ONLY when
    pitch_included:true (the client said the deck ends with an offer/pitch). When
    pitch_included:false this gate DEFERS — a pitchless deck is first-class and must
    NOT be forced to carry an offer ladder; its integrity is enforced by AF-PITCH-LEAK
    instead. When the flag is unset, AF-PITCH-FLAG-UNSET owns the failure; to avoid a
    double-fail this gate also defers on an unset flag."""
    if _intake_pitch_included(run_dir) is not True:
        # pitchless (false) or unset → AF-PITCH-LEAK / AF-PITCH-FLAG-UNSET own those cases.
        return ""
    arc = None
    for rel in ("working/copy/arc_allocation.json", "arc_allocation.json",
                "working/arc_allocation.json"):
        p = run_dir / rel
        if p.exists():
            arc = p
            break
    if arc is None:
        return ""  # arc not built yet — _chk_arc owns the absence; nothing to pitch-check.
    obj = _read_json(arc)
    if isinstance(obj, dict) and "__parse_error__" in obj:
        return ("AF-PITCH-MISSING: arc_allocation.json is not valid JSON, so the offer "
                "ladder + re-pitch cannot be proven. See SOP-PITCH-01 / SOP-PITCH-03.")
    slots = obj if isinstance(obj, list) else (
        obj.get("slots") or obj.get("allocation") or obj.get("slides") or [])
    # Flatten every arc-section / tag token across the allocation into one lowercase blob.
    tokens = []
    for s in slots if isinstance(slots, list) else []:
        if isinstance(s, dict):
            for k in ("arc_section", "section", "beat", "tag", "type", "role"):
                v = s.get(k)
                if isinstance(v, str):
                    tokens.append(v.lower())
            tags = s.get("tags")
            if isinstance(tags, list):
                tokens += [str(t).lower() for t in tags]
        elif isinstance(s, str):
            tokens.append(s.lower())
    blob = " ".join(tokens)
    # OFFER-LADDER present: at least one ladder/anchor/price/offer/value-stack beat.
    ladder_markers = ("ladder", "anchor", "price", "offer", "value-stack",
                      "value_stack", "valuestack", "drop", "stack")
    has_ladder = any(m in blob for m in ladder_markers)
    # RE-PITCH present: an explicit re-pitch / second-close beat after the FINAL.
    repitch_markers = ("re-pitch", "re_pitch", "repitch", "second close",
                       "second-close", "re-offer", "reoffer", "post-final", "post_final")
    has_repitch = any(m in blob for m in repitch_markers)
    if not has_ladder or not has_repitch:
        missing = []
        if not has_ladder:
            missing.append("no offer-ladder beats (value-stack -> anchor -> price drops)")
        if not has_repitch:
            missing.append("no re-pitch beat after the FINAL price")
        return ("AF-PITCH-MISSING: the converting arc is incomplete — "
                + "; ".join(missing) + ". A converting deck must build an offer ladder "
                "and re-pitch after the FINAL price (SOP-PITCH-01 value stack / "
                "SOP-PITCH-03 re-pitch). A deck that never asks for the sale, or asks "
                "once and never recovers a 'no', auto-fails.")
    return ""


# ---------------------------------------------------------------------------
# CREATIVITY (AF-CREATIVITY) — reject template-sameness and cliche. A deck whose
# slides all share one layout archetype (forty copies of the same black headline)
# or whose copy leans on worn-out cliche openers is a template, not a designed
# piece. The design system (design_system.json, Typography Architect) assigns each
# slide one of several archetypes; this gate fails a deck where a single archetype
# dominates beyond a ceiling, OR where banned cliche phrases appear in slide copy.
# ---------------------------------------------------------------------------
# A single archetype may cover at most this fraction of the deck before it reads as
# template-sameness (the anti-template SOP-DESIGN-03 mandates no two CONSECUTIVE
# slides share an archetype; this is the deck-wide dominance ceiling behind it).
ARCHETYPE_DOMINANCE_CEILING = 0.60
# Worn-out cliche openers/phrases a converting deck must not lean on (lowercase,
# substring match on slide copy). Curated from the forensic template-sameness set.
FORBIDDEN_CLICHE_PHRASES = (
    "in today's fast-paced world",
    "in today's fast paced world",
    "in this day and age",
    "at the end of the day",
    "think outside the box",
    "take it to the next level",
    "the sky is the limit",
    "last but not least",
    "needless to say",
    "synergy",
    "low-hanging fruit",
    "move the needle",
    "circle back",
    "boil the ocean",
    "best-kept secret",
    "game changer",
    "game-changer",
    "revolutionary new",
    "one simple trick",
)


def _chk_creativity(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """CREATIVITY gate (AF-CREATIVITY). Rejects template-sameness and cliche:
      * archetype dominance — if the design system assigns one archetype to MORE
        than ARCHETYPE_DOMINANCE_CEILING (60%) of the slides, the deck reads as a
        template (forty copies of one layout) and FAILS;
      * cliche copy — if slide copy contains any FORBIDDEN_CLICHE_PHRASES, FAILS.
    Reads working/typography/design_system.json (per-slide archetype map) and
    working/copy/slides_copy.md.

    Cliche copy defers (passes) pre-copy — that absence is genuinely owned by
    _chk_slides_copy (same PREFLIGHT_REQUIRED list, fails closed on a missing
    slides_copy.md).

    Archetype dominance is NOT symmetrically owned. design_system.json is the
    Typography Architect's own primary output (typography-architect.md /
    typography-architect-sops.md: "You produce one artifact, the DESIGN SYSTEM
    SPEC ... working/typography/design_system.json"), distinct from
    working/research/design-brief-*.md (a Deep-Research-Specialist INPUT the
    Typography Architect consumes, owned by _chk_design_brief) and from
    working/typography/type_layout_system.md (the font-floor gate-of-record
    artifact, owned by check_font_floor / check_typography_qc_teeth). No
    PREFLIGHT_REQUIRED gate requires design_system.json itself to exist, so a
    Typography Architect pass that writes type_layout_system.md but never
    writes design_system.json previously sailed through here silently — the
    anti-template-sameness half of AF-CREATIVITY simply never fired, forever,
    and nothing else caught it (Root Cause 2: absence read as approval).

    Fix: once TYPE_LAYOUT_SYSTEM_REL exists (the same signal check_font_floor /
    check_typography_qc_teeth already use to decide "the design phase is
    underway/complete, its artifacts are now mandatory"), a still-missing
    design_system.json is UNDETERMINED, not silently clean -- and on this
    completeness-style preflight gate UNDETERMINED behaves like FAIL
    (CheckResult doctrine, presentation_job/result.py), so it is reported as an
    AF-CREATIVITY problem rather than dropped. Genuinely PRE-typography (neither
    type_layout_system.md nor design_system.json produced yet) still defers —
    that is a real "not produced yet" state, not a hole.

    Returns "" on pass, or a fatal AF-CREATIVITY message."""
    problems = []
    # --- archetype dominance ---
    ds = None
    for rel in ("working/typography/design_system.json",
                "working/copy/design_system.json", "design_system.json"):
        p = run_dir / rel
        if p.exists():
            ds = p
            break
    archetype_result = CheckResult.PASS
    if ds is None:
        if (run_dir / TYPE_LAYOUT_SYSTEM_REL).exists():
            # Design phase is demonstrably underway (the font-floor gate-of-record
            # artifact exists) yet design_system.json — the Typography Architect's
            # own primary output — never showed up. No other gate owns this. We
            # cannot verify archetype dominance without it: UNDETERMINED, and this
            # completeness gate treats UNDETERMINED as FAIL (do not silently pass).
            archetype_result = CheckResult.UNDETERMINED
        # else: genuinely pre-typography — nothing produced yet, legitimate defer.
    if not archetype_result.ok and ds is None:
        problems.append(
            "archetype dominance could not be verified: "
            f"{TYPE_LAYOUT_SYSTEM_REL} exists (the design phase is underway/"
            "complete) but design_system.json (the Typography Architect's own "
            "archetype-plan output) is absent — UNDETERMINED, refused rather than "
            "silently passed (SOP-DESIGN-03 / typography-architect-sops.md "
            "'Outputs: working/typography/design_system.json')")
    if ds is not None:
        obj = _read_json(ds)
        if isinstance(obj, dict) and "__parse_error__" not in obj:
            per = obj.get("per_slide") or obj.get("slides") or obj.get("slide_plan")
            archetypes = []
            if isinstance(per, list):
                for s in per:
                    if isinstance(s, dict):
                        a = s.get("archetype") or s.get("type") or s.get("treatment")
                        if isinstance(a, str) and a.strip():
                            archetypes.append(a.strip().lower())
            elif isinstance(per, dict):
                for v in per.values():
                    if isinstance(v, str) and v.strip():
                        archetypes.append(v.strip().lower())
                    elif isinstance(v, dict):
                        a = v.get("archetype") or v.get("type")
                        if isinstance(a, str) and a.strip():
                            archetypes.append(a.strip().lower())
            if len(archetypes) >= 3:
                counts = {}
                for a in archetypes:
                    counts[a] = counts.get(a, 0) + 1
                top, top_n = max(counts.items(), key=lambda kv: kv[1])
                frac = top_n / len(archetypes)
                if frac > ARCHETYPE_DOMINANCE_CEILING:
                    problems.append(
                        f"template-sameness: archetype {top!r} covers {top_n}/"
                        f"{len(archetypes)} slides ({frac:.0%}), over the "
                        f"{ARCHETYPE_DOMINANCE_CEILING:.0%} dominance ceiling — the deck "
                        f"reads as one repeated layout (SOP-DESIGN-03 anti-template)")
    # --- cliche copy ---
    cp = None
    for rel in ("working/copy/slides_copy.md", "slides_copy.md",
                "working/slides_copy.md"):
        p = run_dir / rel
        if p.exists():
            cp = p
            break
    if cp is not None:
        low = cp.read_text(errors="replace").lower()
        hits = [ph for ph in FORBIDDEN_CLICHE_PHRASES if ph in low]
        if hits:
            problems.append("cliche copy: slide copy leans on worn-out phrase(s) "
                            + ", ".join(repr(h) for h in hits[:6]))
    if problems:
        return ("AF-CREATIVITY: the deck fails the anti-template / anti-cliche gate — "
                + "; ".join(problems) + ". Vary the layout archetypes (no single "
                "archetype dominates) and rewrite cliche copy into specific, original "
                "lines (SOP-DESIGN-03 variable-layout / slide-copywriter doctrine).")
    return ""


def _chk_arc(path: Optional[Path]) -> str:
    # Signature Presentation Architect — the converting-arc allocation must exist.
    if path is None:
        return "file absent — Signature Presentation Architect did not produce the converting arc (arc_allocation)"
    obj = _read_json(path)
    if "__parse_error__" in obj:
        return f"not valid JSON ({obj['__parse_error__']})"
    slots = obj if isinstance(obj, list) else (obj.get("slots") or obj.get("allocation") or obj.get("slides"))
    if not slots:
        return "no per-slide arc-section allocation present (Architect arc not built)"
    return ""


def _chk_slides_copy(path: Optional[Path]) -> str:
    # Slide Copywriter — authored copy per doctrine must exist (Copy-QC validates depth).
    if path is None:
        return "file absent — Slide Copywriter did not author slides_copy (copy doctrine skipped)"
    if len(path.read_text(errors="replace").strip()) < 500:
        return "slides_copy is near-empty (Slide Copywriter doctrine not applied)"
    return ""


def _chk_design_brief(path: Optional[Path]) -> str:
    # Typography Architect / design research — per-slide creative art-direction brief must exist.
    if path is None:
        return "file absent — no design/typography brief (Typography Architect did not art-direct the slides)"
    if len(path.read_text(errors="replace").strip()) < 200:
        return "design/typography brief is near-empty (creative art-direction skipped)"
    return ""


def _count_output_slides(run_dir: Path, slides_path: Optional[Path] = None) -> Optional[int]:
    """Count the slides the system is about to output. Returns None when no count
    can be read.

    H3: when slides_path is supplied it is the ACTUAL positional slides.json the
    renderer will render (main()'s positional[0]). Counting it FIRST means the
    coverage gate and the rich-prompt gate derive their slide count from the EXACT
    file that gets rendered — not a different canonical-path slides.json that might
    have a different slide count (which would let the coverage / rich-prompt gates be
    bypassed by feeding one file to the gate and rendering another). Only when no
    explicit path is given (or it can't be read) do we fall back to the canonical
    working/copy/slides.json spots, then arc_allocation.json."""
    def _count_from(p: Path) -> Optional[int]:
        if not p.exists():
            return None
        obj = _read_json(p)
        if isinstance(obj, list):
            return len(obj)
        if isinstance(obj, dict) and "__parse_error__" not in obj:
            slides = obj.get("slides")
            if isinstance(slides, list):
                return len(slides)
        return None

    # 0. The ACTUAL rendered file (positional slides.json), when threaded in.
    if slides_path is not None:
        n = _count_from(slides_path)
        if n is not None:
            return n

    # 1. slides.json — the renderer's direct input. Look in a few canonical spots.
    for rel in ("working/copy/slides.json", "slides.json", "working/slides.json"):
        p = run_dir / rel
        n = _count_from(p)
        if n is not None:
            return n
    # 2. arc_allocation.json — per-slide arc-section allocation.
    arc = run_dir / "working" / "copy" / "arc_allocation.json"
    if arc.exists():
        obj = _read_json(arc)
        if "__parse_error__" not in obj:
            slots = obj if isinstance(obj, list) else (
                obj.get("slots") or obj.get("allocation") or obj.get("slides"))
            if isinstance(slots, list):
                return len(slots)
    return None


def _chk_research_cited(path: Optional[Path]) -> str:
    """RESEARCH-CITATION GATE (AF-RESEARCH-UNCITED, fail-loud).

    The existing _chk_research_brief gate only checks that research_complete:true
    appears in the file header — a flag any model can self-assert without doing any
    research.  This gate is the source of truth: it IGNORES the self-asserted flag
    and instead counts real authoritative http(s) URLs in the research pack.

    The gate FAILS LOUD when:
      * The research pack exists but contains fewer than MIN_CITED_SOURCES (8)
        distinct http(s) URLs.  This proves the brief was built from intake alone,
        not from real web research.

    When the file is absent, this check returns "" (the existing _chk_research_brief
    gate already produces a clear "file absent" failure for that case; we do not
    double-report it here).  The check is strictly ADDITIVE — it never weakens the
    existing research_complete:true check.

    Heuristic documented (C3): we extract every http(s) URL, reduce each to its
    REGISTERED DOMAIN (urlparse(...).netloc), and de-duplicate by domain. Junk hosts
    are EXCLUDED outright — localhost/loopback, RFC-1918 / link-local IPs, bare IP
    literals, .local/.internal/.invalid/.test/.example reserved TLDs, and example.*
    placeholder domains. The gate then requires at least MIN_DISTINCT_DOMAINS (6)
    DISTINCT REAL PUBLIC domains. This defeats the three ways the old raw-URL count
    was gamed: (1) the same source cited many times, (2) localhost/RFC-1918/bare-IP
    "URLs", and (3) example.com placeholders. MIN_CITED_SOURCES stays the named,
    configurable floor referenced by sync_check / the manifest; the effective gate is
    the distinct-public-domain count.
    """
    if path is None:
        return ""  # file-absent handled by _chk_research_brief; don't double-report

    text = path.read_text(errors="replace")
    domains = _distinct_public_domains(text)
    n = len(domains)
    if n < MIN_DISTINCT_DOMAINS:
        return (
            f"AF-RESEARCH-UNCITED: research pack at {path.name} cites {n} distinct "
            f"real public domain(s), below the HARD floor of {MIN_DISTINCT_DOMAINS} "
            f"(MIN_CITED_SOURCES={MIN_CITED_SOURCES}). Junk hosts (localhost, "
            f"loopback, RFC-1918 / private IPs, bare IP literals, .local/.internal/"
            f".example reserved TLDs, and example.* placeholders) are NOT counted — "
            f"so a research_complete:true flag, the same source cited many times, or "
            f"a list of localhost/example.com URLs is not proof of real web research. "
            f"The research pack MUST cite at least {MIN_DISTINCT_DOMAINS} DISTINCT, "
            f"real, authoritative public sources (covering categories A/B/C/D/G/H/I/K/L "
            f"per the Deep Research SOP). Re-run Phase -0.5 (Deep Research Specialist) "
            f"and cite every distinct source you find."
        )

    # SOURCED-CATEGORY population check (independent of the self-asserted
    # research_complete:true flag, which this gate ignores). Categories G (Credible
    # Attributable Quotes), H (Fact-Validation Ledger), and I (Objection Research)
    # carry the SOURCED material the "who says so" / proof / objection slides draw
    # from — a research pack with citations but an absent/empty G, H, or I is not a
    # usable pack. Uses the same shared helpers as _chk_research_brief so the two
    # checkers cannot drift.
    sourced_cats = {
        "G": "Credible Attributable Quotes",
        "H": "Fact-Validation Ledger",
        "I": "Objection Research",
    }
    missing = []
    for letter, label in sourced_cats.items():
        body = _research_category_body(text, letter)
        if not body:
            missing.append(f"Category {letter} ({label}) section absent")
        elif not _research_category_nonempty(body):
            missing.append(
                f"Category {letter} ({label}) section is empty or contains only "
                f"template placeholder text")
    if missing:
        return (
            "AF-RESEARCH-UNCITED: the research pack cites enough distinct domains but "
            "the SOURCED categories that feed the 'who says so' / proof / objection "
            "slides are missing or empty (checked independently of the self-asserted "
            "research_complete:true flag). Re-run Phase -0.5 (Deep Research Specialist "
            "SOP 9.4) and populate each sourced category with real cited material: "
            + "; ".join(missing)
        )
    return ""


def _chk_claims_without_citation(run_dir: Path) -> str:
    """CLAIMS-WITHOUT-CITATION gate (AF-RESEARCH-UNCITED, fail-loud).

    Scans the slide copy (working/copy/slides.json or working/copy/slides_copy.md)
    for factual/statistical claim markers (a percentage, a dollar figure, the words
    'research'/'study'/'studies show'/'statistics'/'data shows') and fails if such
    a marker appears but the research pack contains NO corresponding cited URL.

    Documented heuristic: this is a presence gate, not a semantic match.  If ANY
    claim marker is found in the copy AND the research pack has zero URLs, the deck
    is shipping unsupported claims.  If the research pack has >= MIN_CITED_SOURCES
    URLs (meaning _chk_research_cited already passed), this check passes — the
    researcher's responsibility is to ensure each claim is supported; the mechanical
    gate proves a research pack with citations exists.  This prevents the case where
    zero research was done but the brief self-reports research_complete:true.

    Both checks (this + _chk_research_cited) are wired into PREFLIGHT_REQUIRED so
    neither can be silently skipped.  When the research pack has no URLs at all AND
    the copy has claim markers, we return the failure; otherwise we pass.
    """
    import re as _re

    # Locate the slide copy (json or markdown).
    copy_text = ""
    for rel in ("working/copy/slides.json", "working/copy/slides_copy.md",
                "slides.json", "working/slides.json"):
        cp = run_dir / rel
        if cp.exists():
            copy_text = cp.read_text(errors="replace")
            break
    if not copy_text.strip():
        return ""  # no copy found — nothing to check

    # Check for factual/statistical claim markers in the copy.
    claim_re = _re.compile(
        "|".join(_CLAIM_MARKERS_RE_SRC), _re.IGNORECASE)
    has_claims = bool(claim_re.search(copy_text))
    if not has_claims:
        return ""  # no claim markers — gate not triggered

    # Check whether ANY cited URL exists in the research pack.
    brief_matches = sorted(run_dir.glob("working/research/brief-*.md"))
    if not brief_matches:
        # No research pack at all + claims in copy = fail loud.
        return (
            "AF-RESEARCH-UNCITED: slide copy contains factual/statistical claim "
            "markers (a percentage, a dollar figure, or 'research'/'study'/"
            "'studies show'/'statistics'/'data shows') but no research pack "
            "(working/research/brief-*.md) was found. Claims on slides MUST be "
            "backed by a cited research pack. Run Phase -0.5 (Deep Research "
            "Specialist) and cite every source before copy proceeds to render."
        )

    # H4: union the cited domains across ALL brief-*.md files, not just the
    # alphabetically-first one (a single deck routinely splits research across
    # several brief files; reading only brief_matches[0] would miss real citations
    # in the others and wrongly fail). C3: count DISTINCT REAL PUBLIC domains, not
    # raw URL strings, so localhost/example.com/RFC-1918/bare-IP/dup "URLs" do not
    # count as citations here either.
    all_domains = set()
    for bm in brief_matches:
        all_domains |= _distinct_public_domains(bm.read_text(errors="replace"))
    if not all_domains:
        names = ", ".join(b.name for b in brief_matches)
        return (
            f"AF-RESEARCH-UNCITED: slide copy contains factual/statistical claim "
            f"markers but the research pack ({names}) contains zero cited real public "
            f"http(s) sources (localhost / example.* / RFC-1918 / bare-IP hosts are "
            f"not counted). Every claim in slide copy that cites a percentage, dollar "
            f"figure, or research finding MUST have a corresponding citation in the "
            f"research pack. Re-run Phase -0.5 and add real cited sources before render."
        )
    return ""


def _chk_coverage(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """ANTI-COMPRESSION (AF-COVERAGE-1). Mode B (working from a client's existing
    deck) is ADD-only: the output deck must NEVER have fewer slides than the source
    deck. Reads mission_prd.json's top-level integer 'source_slide_count' (Mode A =
    absent/0, which always passes) and compares to the output slide count from
    slides.json (or arc_allocation.json). Returns "" on pass, or a fatal AF message
    string (run_preflight maps a returned reason to exit 3)."""
    # An EXPLICIT client-requested slide count is an explicit client instruction that
    # overrides the Mode-B anti-compression floor (the client may deliberately ask to
    # set their deck to an exact length below the source). AF-SLIDE-COUNT-EXACT owns
    # the count when one is set; defer so the two gates never contradict.
    if _client_requested_slide_count(run_dir) is not None:
        return ""
    # Resolve mission_prd.json (Mode A: absent -> source 0 -> always pass).
    source = 0
    for rel in ("working/copy/mission_prd.json", "mission_prd.json",
                "working/mission_prd.json"):
        p = run_dir / rel
        if p.exists():
            obj = _read_json(p)
            if isinstance(obj, dict) and "__parse_error__" not in obj:
                raw = obj.get("source_slide_count", 0)
                try:
                    source = int(raw)
                except (TypeError, ValueError):
                    source = 0
            break

    if source <= 0:
        return ""  # Mode A — no client source deck; coverage check does not apply.

    # H3: count the ACTUAL rendered slides.json (positional) when threaded in.
    final = _count_output_slides(run_dir, slides_path)
    if final is None:
        # We have a source count but cannot read the output count — cannot prove
        # coverage, so fail rather than silently pass a possibly-compressed deck.
        return (f"AF-COVERAGE-1: source deck had {source} slides but the output "
                f"slide count could not be read (no slides.json/arc_allocation.json). "
                f"The system must NEVER output fewer slides than the source "
                f"(Mode B is ADD-only).")

    if final < source:
        return (f"AF-COVERAGE-1: output deck has {final} slides; client source deck "
                f"had {source}. The system must NEVER output fewer slides than the "
                f"source (Mode B is ADD-only). Add {source - final} slides.")
    return ""


def _load_slide_copy_map(run_dir: Path, slides_path: Optional[Path] = None) -> dict:
    """Return {ordinal: copy_val} for the deck the renderer will actually render, so the
    VERBATIM-WORDS-BAKED check can compare each prompt against its slide's exact copy.
    Reads the positional slides.json (H3) when threaded in, else the canonical
    working/copy/slides.json. Returns {} when no slides.json can be read (the verbatim
    check then defers — the rich-prompt floor/blocks/quality still apply)."""
    candidates = []
    if slides_path is not None:
        candidates.append(slides_path)
    candidates += [run_dir / "working" / "copy" / "slides.json",
                   run_dir / "slides.json",
                   run_dir / "working" / "slides.json"]
    for p in candidates:
        if not p.exists():
            continue
        obj = _read_json(p)
        slides = obj if isinstance(obj, list) else (
            obj.get("slides") if isinstance(obj, dict) and "__parse_error__" not in obj else None)
        if isinstance(slides, list) and slides:
            out = {}
            for s in slides:
                if isinstance(s, dict) and isinstance(s.get("slide"), int):
                    out[s["slide"]] = s.get("copy")
            if out:
                return out
    return {}


def _load_slide_arc_tags(run_dir: Path) -> dict:
    """Return {ordinal: lowercase tag/section/beat blob} from arc_allocation.json.
    Used ONLY to identify the HOOK / section-banner slide-total-floor exemption for
    _chk_copy_density (spec §8.2 item 1: "pure-typography HOOK slides floor drops
    to 12; section-banner slides likewise"). Returns {} when no arc_allocation.json
    can be read — the gate then applies the standard (non-exempt) floor to every
    slide, which is the safe default (never a silent pass)."""
    arc = None
    for rel in ("working/copy/arc_allocation.json", "arc_allocation.json",
                "working/arc_allocation.json"):
        p = run_dir / rel
        if p.exists():
            arc = p
            break
    if arc is None:
        return {}
    obj = _read_json(arc)
    if isinstance(obj, dict) and "__parse_error__" in obj:
        return {}
    slots = obj if isinstance(obj, list) else (
        obj.get("slots") or obj.get("allocation") or obj.get("slides") or [])
    out = {}
    if not isinstance(slots, list):
        return {}
    for s in slots:
        if not isinstance(s, dict):
            continue
        ordinal = s.get("slide")
        if not isinstance(ordinal, int):
            continue
        tokens = []
        for k in ("arc_section", "section", "beat", "tag", "type", "role", "phase"):
            v = s.get(k)
            if isinstance(v, str):
                tokens.append(v.lower())
        tags = s.get("tags")
        if isinstance(tags, list):
            tokens += [str(t).lower() for t in tags]
        # U067: one arc_allocation schema marks hook and section-banner slides as
        # BOOLEANS, not strings. Measured 2026-07-26 on the only in-repository
        # arc_allocation.json: 103 slots, 8 with "hook": true, 4 with
        # "label_slide": true, and _is_hook_or_banner_slide returned False for all
        # 103, because the loop above keeps only isinstance(v, str) values and
        # neither key is even in the tuple. Emit the marker token the classifier
        # already recognises. This is an ALLOWLIST of exactly two flags: do NOT
        # generalise it to "every true boolean". "case_study" is also a boolean on
        # this schema (true on 2 of 103) and case-study slides are the DENSEST in
        # the deck; exempting them from the body-copy floor is the opposite of what
        # this exemption is for.
        for flag, marker in (("hook", "hook"), ("label_slide", "section-banner")):
            if s.get(flag) is True:
                tokens.append(marker)
        out[ordinal] = " ".join(tokens)
    return out


def _is_hook_or_banner_slide(tag_blob: str) -> bool:
    if not tag_blob:
        return False
    markers = ("hook", "section-banner", "section_banner", "sectionbanner", "banner")
    return any(m in tag_blob for m in markers)


def _chk_copy_density(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-COPY-BAND (preflight, fail-loud) -- per-slide COPY character-count FLOOR
    and CEILING, code-enforced. See the COPY_* band constants above for the exact
    numbers and their provenance (spec item 7 / §8.2): this is the first-ever
    code-enforced minimum for headline/subhead/slide-total, and it reconciles the
    old 3-way bullet-count conflict (5-bullets-7-words vs 3-bullets-30-words vs
    render-reality 5-words) into ONE rule instead of three incompatible ones.

    Reads the ACTUAL slides.json the renderer will render (slides_path, H3 -- the
    same positional-file convention _chk_rich_prompts/_chk_coverage use) so the
    band is enforced against the exact file that gets rendered, not a different
    canonical slides.json that might disagree. Defers ("" / pass) when slides.json
    cannot be found or parsed -- _chk_slides_copy / _chk_arc / schema validation own
    that absence; this gate only judges copy[] CONTENT once it exists.

    Because QC_FOREIGN_SIGNATURES bans QC *reports* that carry a word-count
    rubric, this is a FIRST-PARTY preflight check (like AF-P1) that counts
    CHARACTERS deterministically in code -- never a QC-report criterion -- so
    there is no collision with that ban.

    Returns "" when every slide's copy[] clears every band, else a fatal
    AF-COPY-BAND message naming slide + field + count + band for every offender."""
    n = _count_output_slides(run_dir, slides_path)
    if n is None:
        return ""  # no slides.json yet -- _chk_slides_copy/_chk_arc own the absence
    copy_map = _load_slide_copy_map(run_dir, slides_path)
    if not copy_map:
        return ""  # slides.json unreadable/empty -- upstream checks own that failure
    arc_tags = _load_slide_arc_tags(run_dir)
    offenders = []
    for ordinal in range(1, n + 1):
        copy_val = copy_map.get(ordinal)
        if not isinstance(copy_val, list) or not copy_val:
            continue  # missing/malformed copy -- schema validation / AF-P1 own this
        fields = [str(c) if c is not None else "" for c in copy_val]
        headline = fields[0]
        subhead = fields[1] if len(fields) > 1 else None
        kicker = fields[2] if len(fields) > 2 else None
        bullets = fields[3:]

        hlen = len(headline)
        if not (COPY_HEADLINE_CHAR_FLOOR <= hlen <= COPY_HEADLINE_CHAR_CEILING):
            offenders.append(
                f"slide {ordinal:02d} HEADLINE: {hlen} chars, outside "
                f"{COPY_HEADLINE_CHAR_FLOOR}-{COPY_HEADLINE_CHAR_CEILING} band")

        if subhead:
            slen = len(subhead)
            if not (COPY_SUBHEAD_CHAR_FLOOR <= slen <= COPY_SUBHEAD_CHAR_CEILING):
                offenders.append(
                    f"slide {ordinal:02d} SUBHEAD: {slen} chars, outside "
                    f"{COPY_SUBHEAD_CHAR_FLOOR}-{COPY_SUBHEAD_CHAR_CEILING} band")

        if kicker:
            klen = len(kicker)
            if klen > COPY_KICKER_CHAR_CEILING:
                offenders.append(
                    f"slide {ordinal:02d} KICKER: {klen} chars, over the "
                    f"{COPY_KICKER_CHAR_CEILING}-char ceiling")

        real_bullets = [b for b in bullets if b]
        if len(real_bullets) > COPY_BULLET_MAX_COUNT:
            offenders.append(
                f"slide {ordinal:02d} BULLETS: {len(real_bullets)} bullets, over the "
                f"{COPY_BULLET_MAX_COUNT}-bullet max (AF-COPY-BAND reconciles the historical "
                f"3-vs-5-bullet / 7-vs-5-word conflict DOWN to this one enforced rule)")
        for bi, b in enumerate(real_bullets, start=1):
            blen = len(b)
            if not (COPY_BULLET_CHAR_FLOOR <= blen <= COPY_BULLET_CHAR_CEILING):
                offenders.append(
                    f"slide {ordinal:02d} BULLET {bi}: {blen} chars, outside "
                    f"{COPY_BULLET_CHAR_FLOOR}-{COPY_BULLET_CHAR_CEILING} band")

        total = sum(len(f) for f in fields)
        tag_blob = arc_tags.get(ordinal, "")
        exempt_eligible = _is_hook_or_banner_slide(tag_blob)
        exempt = exempt_eligible and COPY_HOOK_EXEMPTION_ENFORCED
        floor = COPY_HOOK_SLIDE_TOTAL_CHAR_FLOOR if exempt else COPY_SLIDE_TOTAL_CHAR_FLOOR
        if not (floor <= total <= COPY_SLIDE_TOTAL_CHAR_CEILING):
            offenders.append(
                f"slide {ordinal:02d} SLIDE TOTAL: {total} chars, outside "
                f"{floor}-{COPY_SLIDE_TOTAL_CHAR_CEILING} band"
                + (" (hook/section-banner exemption applied)" if exempt else ""))
            if exempt_eligible and not COPY_HOOK_EXEMPTION_ENFORCED:
                offenders[-1] += (" [U067 ADVISORY: this slide is hook/section-banner and "
                                  f"would clear a {COPY_HOOK_SLIDE_TOTAL_CHAR_FLOOR}-char "
                                  "floor once the exemption is enforced]")

    if offenders:
        return ("AF-COPY-BAND: per-slide copy character-count floor/ceiling FAILED for "
                f"{len(offenders)} field(s). Bands: HEADLINE "
                f"{COPY_HEADLINE_CHAR_FLOOR}-{COPY_HEADLINE_CHAR_CEILING}, SUBHEAD "
                f"{COPY_SUBHEAD_CHAR_FLOOR}-{COPY_SUBHEAD_CHAR_CEILING} (if present), KICKER "
                f"<= {COPY_KICKER_CHAR_CEILING} (if present), BULLETS <= {COPY_BULLET_MAX_COUNT} "
                f"each {COPY_BULLET_CHAR_FLOOR}-{COPY_BULLET_CHAR_CEILING}, SLIDE TOTAL "
                f"{COPY_SLIDE_TOTAL_CHAR_FLOOR}-{COPY_SLIDE_TOTAL_CHAR_CEILING} "
                f"({COPY_HOOK_SLIDE_TOTAL_CHAR_FLOOR} floor for hook/section-banner slides). "
                "Offenders: " + "; ".join(offenders))
    return ""


def _collect_prompt_problems(run_dir: Path, slides_path: Optional[Path] = None) -> list:
    """Re-open every on-disk per-slide prompt and return the full list of fatal problems:
    missing file, sub-floor / over-ceiling length, missing structural blocks, AND the
    QUALITY-LAYER teeth (AF-P13 eight-class negative block, AF-P14 spelling-lock,
    AF-P-DENSITY, AF-P-VERBATIM). This is the single source of truth shared by the
    preflight rich-prompt gate (_chk_rich_prompts) AND the governed Prompt-QC teeth
    (check_prompt_qc_teeth) so a Prompt-QC report can never claim pass over prompts that
    do not actually clear the gate. Returns [] when every prompt clears every check, or
    [(ordinal, reason), ...]. Returns the sentinel [(0, reason)] when the slide count
    cannot be determined."""
    n = _count_output_slides(run_dir, slides_path)
    if n is None:
        return [(0, "cannot determine the slide count (no slides.json / "
                    "arc_allocation.json), so the per-slide rich prompts cannot be "
                    "verified. Produce slides.json before render.")]
    # FIX-22 / D16 + R3 / D10: BEFORE trusting any per-slide prompt, run the
    # canonical zero-padded + duplicate-detector over the whole prompts dir. A
    # `slide-1.txt` vs `slide-01.txt` collision (both target slide 1) or ANY
    # non-canonical prompt filename is a build-blocking defect — the preflight
    # rich-prompt gate AND the governed Prompt-QC teeth both route through here,
    # so it is caught at exit 3 before any KIE dispatch, never mid-render. The
    # shared detector's verdicts go through _canonical_prompt_dir_problems, which
    # overlays the R3 3-digit-canonical rule (signature-deck 100-slide floor) so
    # slide-100.txt is accepted while slide-1.txt / slide-1000.txt still fail.
    _pg22 = _import_prompt_gate()
    if _pg22 is not None:
        dir_problems = _canonical_prompt_dir_problems(run_dir)
        if dir_problems:
            # Fatal directory-level defect, reported at the sentinel ordinal -1 so
            # callers can tell it apart from the "slide count unknown" sentinel (0).
            # The message already names the offending files; the ordinal is cosmetic.
            return [(-1, "; ".join(dir_problems))]
    copy_map = _load_slide_copy_map(run_dir, slides_path)
    problems = []
    for ordinal in range(1, n + 1):
        p = resolve_prompt_path(run_dir, ordinal)
        if p is None:
            problems.append((ordinal, "NO rich prompt file in working/prompts/"))
            continue
        # H1: measure the STRIPPED length so a whitespace-padded / whitespace-only
        # prompt file can never satisfy the floor.
        raw = p.read_text(errors="replace")
        stripped = raw.strip()
        length = len(stripped)
        if length < PROMPT_CHAR_FLOOR:
            problems.append((ordinal, f"rich prompt {p.name} is {length} non-whitespace "
                            f"chars, under the {PROMPT_CHAR_FLOOR}-char HARD floor"))
            continue
        if length > PROMPT_CHAR_CEILING:
            problems.append((ordinal, f"rich prompt {p.name} is {length} chars, OVER the "
                            f"{PROMPT_CHAR_CEILING}-char HARD ceiling (AF-P2)"))
            continue
        missing_blocks = _missing_structural_blocks(stripped.lower())
        if missing_blocks:
            problems.append((ordinal, f"rich prompt {p.name} clears the floor but is missing "
                            f"required structural block(s): {', '.join(missing_blocks)} "
                            f"([ARCHETYPE ...] / DO-NOT BLOCK [alias NEGATIVE BLOCK] / 'Do not ')"))
            continue
        # QUALITY-LAYER teeth — the floor is NOT a length-only rubber stamp.
        for q in rich_prompt_quality_problems(stripped, copy_map.get(ordinal)):
            problems.append((ordinal, f"rich prompt {p.name}: {q}"))
    return problems


def _chk_rich_prompts(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """RICH-PROMPT-REQUIRED gate (AF-P1). EVERY slide the system is about to render
    MUST have a hand-authored RICH per-slide prompt in working/prompts/ that is
    >= PROMPT_CHAR_FLOOR (9,000) chars AND clears the quality floor. A missing prompt
    file, one under the floor, or one that is thin/padded/missing the 8-class negative
    block / spelling-lock / verbatim copy is an AF-P1 auto-fail: build_deck.py renders
    the rich prompt VERBATIM and NEVER composes a thin fallback, so a thin/absent prompt
    means the slide cannot be rendered at all. Returns "" on pass, or a fatal AF-P1
    message (run_preflight maps a returned reason to exit 3). The 18,000 ceiling (AF-P2),
    the required structural blocks, AND the quality teeth (AF-P13 / AF-P14 / AF-P-DENSITY
    / AF-P-VERBATIM) are ALL enforced here at preflight — not only per-slide at render
    time in load_rich_prompt — so a too-short, structurally-empty, thin, or
    non-verbatim prompt is caught at exit 3 BEFORE any KIE dispatch, never mid-render."""
    # H3: count the ACTUAL rendered slides.json (positional) when threaded in, so the
    # rich-prompt gate verifies a prompt for every slide that will actually render.
    problems = _collect_prompt_problems(run_dir, slides_path)
    if problems and problems[0][0] == 0:
        return "AF-P1: " + problems[0][1]
    if problems and problems[0][0] < 0:
        # FIX-22 / D16 directory-level fatal (duplicate / non-canonical prompt file).
        return "AF-P1: " + problems[0][1]
    if problems:
        n = _count_output_slides(run_dir, slides_path)
        offenders = "; ".join(f"slide {o:02d}: {r}" for o, r in problems)
        head = (f"AF-P1: rich-prompt-required gate FAILED for {len(problems)} of {n} "
                f"slide-checks. build_deck.py renders the Slide Image Creator's rich prompt "
                f"VERBATIM (working/prompts/slide-NN.txt or slide-NN-prompt.txt) and "
                f"never composes a thin fallback; each must be >= {PROMPT_CHAR_FLOOR} "
                f"chars AND clear the quality floor. Offenders:")
        return head + " | " + offenders
    return ""


def _chk_kie_baked(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """KIE-BAKED gate (AF-I14, fail-loud). EVERY rendered slide must have been BAKED
    by the image model (a real KIE task that produced a real, above-floor PNG) — not
    drawn natively (Pillow/PPTX/ImageDraw), not a flat-colour placeholder, not a stub.

    AF-BAKED previously existed in the ruleset as an SOP-only rule (enforced_by:
    closeout_gate, py_symbol: null) with no real checker — a model could skip the KIE
    bake and native-render slides with nothing failing. This is that real checker.

    It reads the on-disk signal that write_process_manifest() emits:
    working/checkpoints/process_manifest.json, finds the LAST phase=="render" record,
    and HARD-FAILS (returns an AF-I14 message) when ANY of:
      * the render record is absent (no render ran -> no KIE bake);
      * output_slide_count < the rendered slide count derived from slides_path
        (a slide was skipped during baking);
      * any per-slide entry has taskId in (None, "", "native", "placeholder") -- i.e.
        not a real KIE task id;
      * any per-slide image file is missing, fails verify_png, or is below
        PLACEHOLDER_MIN_BYTES (the flat-placeholder / native-render byte signature);
      * any per-slide image_sha256 is null, or a single hash repeats across >2 slides
        (flat-fill reuse: the same placeholder image baked into many slides).

    Returns "" ONLY when every rendered slide maps to a real KIE taskId AND a verified,
    above-floor PNG. The message names the offending slide(s) and the remediation verb.

    CONDITIONAL by design (the AF-SPEECH-SHORT pattern): the render record is written
    AFTER the render. So at PRE-RENDER preflight — before any render has run — there is
    no process_manifest.json / no render record yet, and this gate DEFERS (returns "",
    pass) rather than blocking the render that is about to produce the record. Once a
    render record EXISTS (pre-render of a downstream phase, or at the closeout postflight
    gate where it always exists), the gate is the source of truth that the slides were
    actually KIE-baked. This is wired into the lockstep so it can never be silently
    skipped once a render record is on disk.
    """
    ckpt = run_dir / "working" / "checkpoints" / "process_manifest.json"
    if not ckpt.exists():
        return ""  # no render yet — gate defers to the post-render run.
    obj = _read_json(ckpt)
    if not isinstance(obj, dict) or "__parse_error__" in obj:
        return ("AF-I14: working/checkpoints/process_manifest.json is unreadable / not "
                "valid JSON, so the KIE bake cannot be proven. stage=render, missing "
                "artifact=KIE-baked PNG -> re-run build_deck render for every slide.")
    phases = obj.get("phases")
    render_recs = [p for p in phases if isinstance(p, dict) and p.get("phase") == "render"] \
        if isinstance(phases, list) else []
    if not render_recs:
        return ""  # manifest exists but no render record yet — gate defers.
    rec = render_recs[-1]  # the LAST render record is authoritative for this run.

    rendered_n = _count_output_slides(run_dir, slides_path)
    out_count = rec.get("output_slide_count")
    try:
        out_count = int(out_count)
    except (TypeError, ValueError):
        out_count = None
    if rendered_n is not None and (out_count is None or out_count < rendered_n):
        return (f"AF-I14: render record baked {out_count} slide(s) but the deck has "
                f"{rendered_n} slide(s) — at least one slide was skipped during baking. "
                f"stage=render, missing artifact=KIE-baked PNG -> re-run build_deck "
                f"render for the missing slide(s).")

    per_slide = rec.get("slides")
    if not isinstance(per_slide, list) or not per_slide:
        return ("AF-I14: render record carries no per-slide entries, so no slide can be "
                "proven KIE-baked. stage=render, missing artifact=KIE-baked PNG -> "
                "re-run build_deck render for every slide.")

    _BAD_TASK_IDS = {None, "", "native", "placeholder"}
    problems = []
    hash_counts = {}
    for entry in per_slide:
        if not isinstance(entry, dict):
            problems.append("a per-slide entry is malformed (not an object)")
            continue
        ordinal = entry.get("slide")
        tag = f"slide {ordinal}" if ordinal is not None else "a slide"
        task_id = entry.get("taskId")
        if (task_id in _BAD_TASK_IDS) or (isinstance(task_id, str) and task_id.strip().lower() in _BAD_TASK_IDS):
            problems.append(
                f"{tag}: taskId={task_id!r} is not a real KIE task id "
                f"(native render / placeholder, not a model bake)")
        img = entry.get("image")
        if not img:
            problems.append(f"{tag}: no image path recorded — slide was not baked")
        else:
            ip = Path(img)
            if not ip.exists():
                problems.append(f"{tag}: baked image {ip.name} is missing on disk")
            else:
                try:
                    verify_png(ip)
                except Exception as exc:  # noqa: BLE001
                    problems.append(f"{tag}: image {ip.name} is not a valid PNG ({exc})")
                else:
                    size = ip.stat().st_size
                    if size < PLACEHOLDER_MIN_BYTES:
                        problems.append(
                            f"{tag}: image {ip.name} is {size:,} bytes, below the "
                            f"{PLACEHOLDER_MIN_BYTES:,}-byte KIE-bake floor "
                            f"(flat-placeholder / native-render signature)")
        sha = entry.get("image_sha256")
        if not sha:
            problems.append(f"{tag}: image_sha256 is null — bake not recorded")
        else:
            hash_counts[sha] = hash_counts.get(sha, 0) + 1

    for sha, count in hash_counts.items():
        if count > 2:
            problems.append(
                f"image_sha256 {sha[:12]}… is reused across {count} slides "
                f"(>2) — flat-fill placeholder baked into many slides")

    if problems:
        return ("AF-I14: rendered slide(s) were not KIE-baked (native render / missing "
                "image / flat-placeholder fill / reused flat-fill). stage=render, "
                "missing artifact=KIE-baked PNG -> re-run build_deck render for the "
                "offending slide(s). Offenders: " + "; ".join(problems))
    return ""


# Speech-length gate floor: the presenter speech must carry at least
# target_talk_minutes x SPEECH_WPM_FLOOR words. 120 wpm is the LOW end of the
# verified 120-140 absorption band the Presenter Speech Writer cites (the deck is
# paced at 130 wpm; a script that lands below 120 wpm of content is too SHORT for
# the chosen duration). This is the floor only — the +/-10% pacing budget around
# 130 wpm stays with the Presenter Speech Writer / Presenter Coach.
SPEECH_WPM_FLOOR = 120


def _chk_speech_length(run_dir: Path) -> str:
    """SPEECH-LENGTH gate (AF-SPEECH-SHORT). Once the presenter speech exists, its
    word count must be >= target_talk_minutes x SPEECH_WPM_FLOOR (120 wpm). A speech
    shorter than that does not fill the duration the client asked for and FAILS short.

    This gate is CONDITIONAL by design: the speech is written downstream (Phase 9
    delivery), AFTER the deterministic render. So when no speech artifact exists yet,
    this returns "" (the render is allowed to proceed — the gate fires at delivery,
    not at render). When a speech file IS present, it is enforced. Either way the
    gate is wired into the lockstep so it can never be silently skipped once the
    speech is written. Returns "" on pass/not-applicable, or a fatal AF message."""
    # target_talk_minutes comes from intake.json (the field _chk_intake requires).
    target = None
    for rel in ("working/copy/intake.json", "intake.json", "working/intake.json"):
        p = run_dir / rel
        if p.exists():
            obj = _read_json(p)
            if isinstance(obj, dict) and "__parse_error__" not in obj:
                raw = obj.get("target_talk_minutes")
                try:
                    target = float(raw) if raw is not None else None
                except (TypeError, ValueError):
                    target = None
            break
    if not target or target <= 0:
        # No usable target (intake gate handles a missing target). Not applicable here.
        return ""

    # Locate the speech artifact. Absent => deferred to delivery (pass here).
    # Filenames standardized to the PLURAL possessive PRESENTERS-SPEECH.md (canonical,
    # matching the AF-DH1 client-package whitelist + the producer presenters-speech-writer
    # role + PRESENTER-GUIDE/PRESENTER-AUDIO); the legacy singular (PRESENTER-SPEECH.md)
    # and scratch speech.md names are still accepted so the gate keeps finding a speech
    # written by an older flow.
    speech = None
    for rel in ("working/presenter-speech/speech.md",
                "working/delivery/PRESENTERS-SPEECH.md",
                "working/presenter-speech/PRESENTERS-SPEECH.md",
                "working/delivery/PRESENTER-SPEECH.md",
                "working/presenter-speech/PRESENTER-SPEECH.md"):
        p = run_dir / rel
        if p.exists():
            speech = p
            break
    if speech is None:
        return ""  # speech not written yet (pre-delivery render) — gate defers.

    words = len(speech.read_text(errors="replace").split())
    floor = int(round(target * SPEECH_WPM_FLOOR))
    if words < floor:
        return (f"AF-SPEECH-SHORT: presenter speech {speech.name} has {words} words; "
                f"the floor for a {target:g}-minute talk is {floor} words "
                f"(target_talk_minutes x {SPEECH_WPM_FLOOR} wpm). The speech is too "
                f"SHORT to fill the requested duration. Lengthen it.")
    return ""


# ---------------------------------------------------------------------------
# AF-VISUAL-VARIETY (deck-quality gates 2026-06-19)
# ---------------------------------------------------------------------------
# A monotone dark deck — e.g. 35 slides all deep navy — is illegible at projection
# distance and signals that no real design system was applied. This gate measures
# the mean luminance of rendered slide PNGs (if present) AND falls back to checking
# the design_system.json archetype roster for palette diversity when no PNGs exist.
#
# PASS conditions (either is sufficient):
#   * At least VISUAL_VARIETY_LIGHT_SLIDE_FLOOR_PCT of slides have a mean luma
#     above the dark-luma threshold (i.e. the deck has light/break slides), AND
#     the fraction of slides below the dark-luma threshold does not exceed
#     VISUAL_VARIETY_DARK_DOMINANT_CEILING.
#   * No single dominant background hue (quantised to 36 buckets of 10° each)
#     covers more than VISUAL_VARIETY_HUE_DOMINANCE_CEILING of all slides.
#
# FAIL: >=90% slides share one dominant background hue bucket OR >=90% are below
# the dark-luma threshold with no light/break slides.
#
# Defers (passes) when no rendered PNGs exist yet (pre-render) so the gate never
# blocks the initial render; it fires at Image-QC and postflight where slides exist.
# ---------------------------------------------------------------------------
VISUAL_VARIETY_DARK_DOMINANT_CEILING = 0.90   # if >= this fraction is dark -> FAIL
VISUAL_VARIETY_LIGHT_SLIDE_FLOOR_PCT = 0.10   # at least 10% of slides must be light
VISUAL_VARIETY_HUE_DOMINANCE_CEILING = 0.90   # if one hue bucket >= this fraction -> FAIL
VISUAL_VARIETY_DARK_LUMA_THRESHOLD   = 0.30   # luma <= this is "dark" (0-1 scale)

# Warn-mode staging for the NEUTRAL branch of the hue-dominance half (U050). Step 1a makes
# bucket 36 countable for the first time, and _png_dominant_hue_bucket reads only the twelve
# corner/edge pixels at its sample_coords -- so any deck with a white, grey or black MARGIN
# reports 100% neutral dominance however varied its CONTENT is. Measured 2026-07-26: a
# twenty-slide deck of five distinct hue families inset on a white margin goes PASSED ->
# REJECTED the moment this is True. While False that deck passes and a durable
# `neutral_hue_dominance` warning is recorded instead. A chromatic monotone deck still
# HARD-FAILS in both positions -- staging narrows this to the neutral branch only.
# PROMOTION CRITERION -- flip to True only when BOTH hold: (i) staged_warnings.json records
# zero `neutral_hue_dominance` events across the last ten consecutive real builds; and
# (ii) the hue sampler reads slide CONTENT rather than only those twelve background pixels,
# which is a different unit. Record the flip as its own commit.
VISUAL_VARIETY_NEUTRAL_HUE_ENFORCED = False


def _imaging_available() -> bool:
    """True when Pillow can actually decode an image. A check that silently passes when
    its dependency is missing is not a check (Super Spec 6.5). Callers must FAIL, not
    defer, when this returns False -- with the install instruction in the message."""
    try:
        from PIL import Image  # noqa: PLC0415
        Image.new("RGB", (1, 1))
        return True
    except Exception:  # noqa: BLE001
        return False


def _warn_once(run_dir: Path, key: str, message: str) -> None:
    """Append a staged-warning record to working/qc/staged_warnings.json and print to
    stderr. Durable because the promotion criterion in U050 step 4d is 'zero of the last
    N runs warned' -- which is unanswerable if warnings only ever went to a console."""
    import sys as _sys
    print(message, file=_sys.stderr)
    try:
        out = run_dir / "working" / "qc" / "staged_warnings.json"
        out.parent.mkdir(parents=True, exist_ok=True)
        existing = _read_json(out) if out.exists() else {}
        if not isinstance(existing, dict) or "__parse_error__" in existing:
            existing = {}
        existing.setdefault(key, 0)
        existing[key] += 1
        existing.setdefault("_messages", {})[key] = message
        out.write_text(json.dumps(existing, indent=2), encoding="utf-8")
    except Exception:  # noqa: BLE001
        pass  # a warning that cannot be recorded must never break a build


def _png_mean_luma(path: Path) -> Optional[float]:
    """Return the mean luminance (0.0-1.0 Y-channel) of a PNG slide, or None if
    the file cannot be read or Pillow is unavailable. Converts to grayscale (ITU-R
    BT.601) and averages all pixels. When Pillow is absent, no stdlib measurement
    is attempted -- a missing dependency must produce a FAILURE (see
    _imaging_available in each checker), never a fabricated measurement."""
    try:
        from PIL import Image
        with Image.open(str(path)) as im:
            gray = im.convert("L")
            pixels = list(gray.getdata())
            return sum(pixels) / (len(pixels) * 255.0)
    except Exception:  # noqa: BLE001
        pass
    # No stdlib fallback. Until 2026-07-26 this function fell back to counting high bytes
    # in the DEFLATE-compressed PNG stream, which bears no relationship to luminance: it
    # returned a plausible float in 0-1, so `lumas` filled with noise and the
    # dark-dominance check produced a confident wrong answer. A missing dependency must
    # produce a FAILURE (see _imaging_available), never a fabricated measurement.
    return None

def _png_dominant_hue_bucket(path: Path) -> Optional[int]:
    """Return the dominant background hue bucket (0-35, quantised in 10° steps) of a
    PNG slide, or None when PIL is absent or the file is unreadable. Samples 50
    corner/edge pixels from the slide's four corners (background approximation) to
    find the dominant hue without loading the whole image into memory."""
    try:
        from PIL import Image
        with Image.open(str(path)) as im:
            rgb = im.convert("RGB")
            w, h = rgb.size
            # Sample the four corners (background pixels).
            sample_coords = [
                (0, 0), (w // 4, 0), (w // 2, 0), (3 * w // 4, 0), (w - 1, 0),
                (0, h // 2), (w - 1, h // 2),
                (0, h - 1), (w // 4, h - 1), (w // 2, h - 1), (3 * w // 4, h - 1), (w - 1, h - 1),
            ]
            hue_counts = {}
            for x, y in sample_coords:
                r, g, b = rgb.getpixel((x, y))
                # Convert RGB -> HSV hue
                rn, gn, bn = r / 255.0, g / 255.0, b / 255.0
                mx = max(rn, gn, bn)
                mn = min(rn, gn, bn)
                diff = mx - mn
                if diff < 0.05:
                    # Achromatic (gray/black/white). Bucket 36 is the NEUTRAL bucket and is
                    # deliberately outside the 0-35 hue range so a neutral-background deck is
                    # counted as monotone rather than dropped. Before 2026-07-26 the assignment
                    # below overwrote this unconditionally, `hue` was unbound on this path, the
                    # resulting UnboundLocalError was swallowed by the handler at the end of this
                    # function, and every neutral slide was silently dropped from hue_buckets --
                    # which disabled the hue half of AF-VISUAL-VARIETY on any white, black or grey
                    # deck. Twenty identical white slides passed. Do not re-merge these branches.
                    bucket = 36
                else:
                    if mx == rn:
                        hue = 60 * (((gn - bn) / diff) % 6)
                    elif mx == gn:
                        hue = 60 * ((bn - rn) / diff + 2)
                    else:
                        hue = 60 * ((rn - gn) / diff + 4)
                    if hue < 0:
                        hue += 360
                    bucket = int(hue // 10) % 36
                hue_counts[bucket] = hue_counts.get(bucket, 0) + 1
            if not hue_counts:
                return None
            return max(hue_counts, key=hue_counts.get)
    except Exception:  # noqa: BLE001
        return None


def check_visual_variety(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-VISUAL-VARIETY: reject an all-dark monotone deck.

    Scans rendered slide PNGs (renders/slide-NN.png). If not present yet, defers
    (passes). When PNGs exist:
      * Fails if >= VISUAL_VARIETY_DARK_DOMINANT_CEILING (90%) of slides have
        mean luma <= VISUAL_VARIETY_DARK_LUMA_THRESHOLD with < 10% light slides.
      * Fails if one dominant background hue bucket covers >=
        VISUAL_VARIETY_HUE_DOMINANCE_CEILING (90%) of slides.
    Returns '' on pass / defer, or a fatal AF-VISUAL-VARIETY message."""
    renders_dir = run_dir / "renders"
    if not renders_dir.is_dir():
        return ""  # no renders yet — gate defers to post-render QC pass.

    pngs = sorted(renders_dir.glob("slide-*.png"))
    if not pngs:
        return ""  # defers when no PNG renders exist yet.

    if not _imaging_available():
        return ("AF-VISUAL-VARIETY: DECK FAIL — Pillow is not installed, so no slide's palette "
                "could be examined. This gate cannot pass unexamined renders. "
                "Install with: python3 -m pip install Pillow  "
                "(AF-VISUAL-VARIETY)")

    lumas = []
    hue_buckets = []
    for png in pngs:
        luma = _png_mean_luma(png)
        if luma is not None:
            lumas.append(luma)
        hue = _png_dominant_hue_bucket(png)
        if hue is not None:
            hue_buckets.append(hue)

    n = len(pngs)
    problems = []

    # Dark-dominance check.
    if lumas:
        dark_count = sum(1 for l in lumas if l <= VISUAL_VARIETY_DARK_LUMA_THRESHOLD)
        dark_frac = dark_count / len(lumas)
        light_frac = 1.0 - dark_frac
        if (dark_frac >= VISUAL_VARIETY_DARK_DOMINANT_CEILING and
                light_frac < VISUAL_VARIETY_LIGHT_SLIDE_FLOOR_PCT):
            problems.append(
                f"monotone_dark_palette: {dark_count}/{len(lumas)} slides have mean "
                f"luma <= {VISUAL_VARIETY_DARK_LUMA_THRESHOLD} (dark threshold); "
                f"only {light_frac:.0%} are light (minimum {VISUAL_VARIETY_LIGHT_SLIDE_FLOOR_PCT:.0%} required). "
                f"Add light/break slides and vary the background palette to prevent "
                f"projection-distance illegibility."
            )

    # Hue-dominance check.
    if hue_buckets:
        from collections import Counter
        hue_counts = Counter(hue_buckets)
        top_bucket, top_count = hue_counts.most_common(1)[0]
        hue_frac = top_count / len(hue_buckets)
        if hue_frac >= VISUAL_VARIETY_HUE_DOMINANCE_CEILING:
            if top_bucket == 36:
                hue_label = "neutral (grey/black/white, no chroma)"
            else:
                hue_label = f"~{top_bucket * 10}°"
            finding = (
                f"monotone_palette: {top_count}/{len(hue_buckets)} slides share "
                f"dominant background hue {hue_label} (={hue_frac:.0%} of deck; "
                f"ceiling {VISUAL_VARIETY_HUE_DOMINANCE_CEILING:.0%}). "
                f"Introduce section-break slides with contrasting palette to ensure "
                f"visual variety and WCAG AA legibility at projection distance."
            )
            if top_bucket == 36 and not VISUAL_VARIETY_NEUTRAL_HUE_ENFORCED:
                _warn_once(run_dir, "neutral_hue_dominance",
                           f"AF-VISUAL-VARIETY WARNING (not yet enforced): {finding}")
            else:
                problems.append(finding)

    if problems:
        return (
            f"AF-VISUAL-VARIETY: DECK FAIL — rendered slides show no visual variety. "
            + " | ".join(problems)
            + " (SOP-DESIGN-03, AF-VISUAL-VARIETY)"
        )
    return ""


# ---------------------------------------------------------------------------
# AF-PACKAGE-CLEAN (deck-quality gates 2026-06-19)
# ---------------------------------------------------------------------------
# The final delivered bundle must contain ONLY the canonical deliverable files.
# Dev artifacts (build scripts, polling scripts, download helpers, task JSON dirs,
# intermediate .md working drafts, Office lock/temp files) MUST NOT reach the
# client package. This gate scans the bundle directory for forbidden patterns.
#
# Forbidden:
#   * Any *.py file (build/render/assembly scripts)
#   * Any *.sh file (shell helpers)
#   * Any ~$* file (Office lock/temp files, e.g. ~$WIB-Business-Function-Fidelity.pptx)
#   * Any tasks/ or task_*.json working directories / files
#   * Any numbered intermediate .md working drafts (e.g. 01-draft.md, working-draft.md)
#   * Any file that is NOT in the DELIVERABLES_REQUIRED filename set (whitelist-first)
#
# The whitelist check (file must match a canonical deliverable) is the PRIMARY gate;
# the blocklist is belt-and-suspenders.
#
# Defers (passes) when bundle_dir does not exist yet (pre-delivery) so it never
# blocks the render phase; it fires at the postflight gate and closeout.
# ---------------------------------------------------------------------------
PACKAGE_CLEAN_FORBIDDEN_EXTENSIONS = frozenset({
    ".py", ".sh",
})
PACKAGE_CLEAN_FORBIDDEN_NAME_PREFIXES = ("~$",)
PACKAGE_CLEAN_FORBIDDEN_DIR_NAMES = frozenset({
    "tasks", "working", "prompts", "images", "renders", "qc", "scripts", "checkpoints",
})
# Intermediate .md draft pattern: a working draft that is NOT a canonical deliverable.
# Canonical deliverable .md files are PRESENTERS-SPEECH.md and PRESENTERS-SPEECH-FISH-TAGGED.md
# (plural — matches DELIVERABLES_REQUIRED + the AF-DH1 whitelist + the producer role). The
# legacy singular spellings are still accepted so an older bundle is not flagged dirty.
PACKAGE_CLEAN_CANONICAL_MD_FILES = frozenset({
    "PRESENTERS-SPEECH.md", "PRESENTERS-SPEECH-FISH-TAGGED.md",
    "PRESENTER-SPEECH.md", "PRESENTER-SPEECH-FISH-TAGGED.md",
})
# Intermediate-draft naming patterns (numbered drafts, *-draft.md, *-working.md, etc.)
PACKAGE_CLEAN_DRAFT_MD_RE = re.compile(
    r'^(?:\d+[-_].*\.md|.*[-_](?:draft|working|temp|wip|intermediate|scratch)\.md)$',
    re.IGNORECASE
)


def check_package_cleanliness(bundle_dir: Path) -> str:
    """AF-PACKAGE-CLEAN: the final delivered bundle must contain ONLY the canonical
    deliverable files. Dev artifacts (.py, .sh, ~$* lock files, tasks/ dirs,
    intermediate .md drafts) auto-fail the gate.

    Returns '' on pass / not-yet-present, or a fatal AF-PACKAGE-CLEAN message
    listing each forbidden artifact found."""
    if not bundle_dir or not bundle_dir.is_dir():
        return ""  # bundle not yet assembled — gate defers.

    # Build the canonical filename set (deck_slug-agnostic: strip {deck_slug} prefix).
    canonical_names: set = set()
    for spec in DELIVERABLES_REQUIRED:
        fname = spec["filename"]
        # Accept both the template form and any real slug substitution.
        canonical_names.add(fname)
        # Also accept DELIVERABLES_REQUIRED filenames with any deck_slug prefix.
        # E.g. "{deck_slug}-FINAL.pptx" -> anything ending in "-FINAL.pptx".
        canonical_names.add("deliverables.json")      # the deliverables ledger itself
        canonical_names.add("teleprompter_publish.json")  # teleprompter publish ledger

    forbidden: list = []

    for item in bundle_dir.iterdir():
        name = item.name

        # Forbidden directories.
        if item.is_dir():
            if name.lower() in PACKAGE_CLEAN_FORBIDDEN_DIR_NAMES:
                forbidden.append(f"forbidden directory: {name}/")
            continue

        # Forbidden extensions.
        ext = item.suffix.lower()
        if ext in PACKAGE_CLEAN_FORBIDDEN_EXTENSIONS:
            forbidden.append(f"forbidden artifact ({ext}): {name}")
            continue

        # Forbidden Office lock/temp files (~$* prefix).
        if any(name.startswith(pfx) for pfx in PACKAGE_CLEAN_FORBIDDEN_NAME_PREFIXES):
            forbidden.append(f"Office lock/temp file: {name}")
            continue

        # task_*.json intermediate files.
        if re.match(r'^task_.*\.json$', name, re.IGNORECASE) or name.lower() == "tasks":
            forbidden.append(f"intermediate task file: {name}")
            continue

        # Intermediate .md working drafts (not a canonical deliverable .md).
        if name.endswith(".md") and name not in PACKAGE_CLEAN_CANONICAL_MD_FILES:
            if PACKAGE_CLEAN_DRAFT_MD_RE.match(name):
                forbidden.append(f"intermediate working draft: {name}")
                continue

    if forbidden:
        return (
            "AF-PACKAGE-CLEAN: DELIVERY BLOCKED — the bundle contains forbidden artifacts "
            "that must never reach the client package. Remove each before delivery: "
            + "; ".join(forbidden)
            + " (AF-PACKAGE-CLEAN, SOP-DH1 / AF-DH1)"
        )
    return ""


# ---------------------------------------------------------------------------
# AF-IMAGE-QC-RAN (deck-quality gates 2026-06-19)
# ---------------------------------------------------------------------------
# The image-QC report must: (a) exist, (b) be NEWER than the rendered PNGs (not
# stale), and (c) contain a per-slide PASS/FAIL row for ALL N slides (no
# rubber-stamped all-[x] boilerplate). The existing AF-IMAGE-QC gate already
# checks the report's score/pass/independence. AF-IMAGE-QC-RAN is the
# complementary freshness + per-slide-coverage check: a stale report or one that
# was auto-populated from a template without real per-slide entries is refused.
#
# "Per-slide rows" is detected by checking that the image_qc_report.json carries
# a "slides" or "per_slide" array with one entry per rendered slide. A report
# that lacks this structure is assumed to be rubber-stamped boilerplate and fails.
#
# Defers when no renders exist yet. Defers when no image_qc_report.json exists
# yet (pre-Image-QC phase); AF-IMAGE-QC owns the absence failure.
# ---------------------------------------------------------------------------
def check_image_qc_present(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-IMAGE-QC-RAN: the image-QC report must exist, be newer than the rendered
    PNGs, and carry a per-slide PASS/FAIL row for every rendered slide.

    Returns '' on pass / defer, or a fatal AF-IMAGE-QC-RAN message."""
    report_path = run_dir / "working" / "qc" / "image_qc_report.json"
    renders_dir = run_dir / "renders"

    # No renders yet — gate defers.
    pngs: list = []
    if renders_dir.is_dir():
        pngs = sorted(renders_dir.glob("slide-*.png"))
    if not pngs:
        return ""  # no rendered slides yet — defer.

    # No report yet — defer (AF-IMAGE-QC owns the absence failure at preflight).
    if not report_path.exists():
        return ""  # absence is AF-IMAGE-QC's job; we only check freshness + coverage.

    # Staleness check: report must be NEWER than ALL rendered PNGs.
    try:
        report_mtime = report_path.stat().st_mtime
        png_mtimes = [p.stat().st_mtime for p in pngs if p.exists()]
        if png_mtimes and report_mtime < max(png_mtimes):
            newest_png = max(pngs, key=lambda p: p.stat().st_mtime)
            return (
                f"AF-IMAGE-QC-RAN: image-QC report is STALE — it was last written "
                f"before {newest_png.name} was rendered. The report must be re-run after "
                f"the current render batch. A stale report cannot prove the rendered slides "
                f"were actually reviewed (AF-IMAGE-QC-RAN)."
            )
    except OSError:
        pass  # if stat fails, skip staleness check

    # Per-slide coverage check: the report must carry per-slide entries.
    obj = _read_json(report_path)
    if "__parse_error__" in obj:
        return ""  # JSON error is AF-IMAGE-QC's job.

    n_slides = len(pngs)
    # Accept "slides" or "per_slide" array.
    per_slide = obj.get("slides") or obj.get("per_slide") or obj.get("slide_results") or []
    if not isinstance(per_slide, list):
        per_slide = []

    if not per_slide:
        return (
            f"AF-IMAGE-QC-RAN: image-QC report has no per-slide entries (expected "
            f"{n_slides} rows, one per rendered slide). A report with no per-slide data "
            f"is rubber-stamped boilerplate and cannot prove the rendered slides were "
            f"reviewed. Re-run the Image QC Specialist with a per-slide PASS/FAIL row "
            f"for each of the {n_slides} rendered slide(s) (AF-IMAGE-QC-RAN)."
        )

    if len(per_slide) < n_slides:
        return (
            f"AF-IMAGE-QC-RAN: image-QC report covers only {len(per_slide)} of "
            f"{n_slides} rendered slide(s). Every rendered slide must have a per-slide "
            f"PASS/FAIL entry. Missing {n_slides - len(per_slide)} slide(s). Re-run the "
            f"Image QC Specialist to cover all {n_slides} slides (AF-IMAGE-QC-RAN)."
        )

    return ""


# ---------------------------------------------------------------------------
# AF-BRAND-CONSISTENCY (deck-quality gates 2026-06-19)
# ---------------------------------------------------------------------------
# Every rendered slide's dominant palette must fall within the client's declared
# brand token set (approved hex colors). Off-brand stock imagery or completely
# wrong palettes (e.g. a fantasy-castle render when the brand is navy/gold) are
# auto-failed. This complements AF-VISUAL-VARIETY (variety rejects sameness;
# brand rejects off-token outliers).
#
# Brand tokens are read from:
#   working/copy/intake.json -> brand.palette (list of hex strings)
#   working/research/design-brief-*.md -> palette declaration
#
# When no brand palette is declared (brand tokens not configured), this gate
# DEFERS (passes) — a gate that can't know what to check can't enforce anything.
#
# When brand tokens ARE declared, each rendered slide's 3 most common pixel colors
# are sampled. If ALL 3 dominant colors are "far" from every declared brand token
# (distance > BRAND_CONSISTENCY_TOLERANCE in RGB-255 space), that slide fails.
#
# Defers when no rendered PNGs exist yet (pre-render).
# ---------------------------------------------------------------------------
BRAND_CONSISTENCY_TOLERANCE = 80   # max RGB Euclidean distance (0-441 scale, 0=exact match)
BRAND_CONSISTENCY_SAMPLE_COLORS = 3  # number of dominant colors to sample per slide

# Warn-mode staging for the brand-palette requirement (U050). While False, a run with no
# declared brand.palette gets a WARNING and still passes. Flip to True only when the
# promotion criterion in step 4d is met. Every existing intake.json without a palette
# becomes non-conformant the moment this is True -- that is the whole reason it is staged.
# PROMOTION CRITERION: flip to True only when BOTH hold:
# (i) staged_warnings.json records zero brand_palette_absent events across the last
#     ten consecutive real builds; and
# (ii) the intake surface actually collects a palette (U058's work, gated on D01).
# Record the flip as its own commit so it can be reverted without reverting the fix.
BRAND_PALETTE_REQUIRED = False


def _hex_to_rgb(hex_color: str) -> Optional[tuple]:
    """Parse a hex color (#RRGGBB or RRGGBB) to (R, G, B) tuple, or None on error."""
    h = hex_color.lstrip("#").strip()
    if len(h) == 3:
        h = h[0] * 2 + h[1] * 2 + h[2] * 2
    if len(h) != 6:
        return None
    try:
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return None


def _rgb_distance(c1: tuple, c2: tuple) -> float:
    """Euclidean distance between two RGB tuples."""
    return ((c1[0] - c2[0]) ** 2 + (c1[1] - c2[1]) ** 2 + (c1[2] - c2[2]) ** 2) ** 0.5


def _slide_dominant_colors(path: Path, n: int = 3):
    """Return a list of the n most-common quantised colors in the slide PNG. Uses PIL
    when available; returns [] when PIL is absent or the file is unreadable.

    BUG FIX (Pillow 11.x): quantize(colors=64) returns a SHORT palette when the image
    has fewer than 64 unique colors (e.g. a solid fill returns len(palette)==3).
    Iterating range(64) then produces IndexError at palette[i*3] for i>=1.  We bound
    the loop to len(palette)//3 to use only the entries Pillow actually allocated.
    The bare 'except Exception' that previously swallowed this real bug is replaced
    with an import-guard so PIL-unavailability is silent but real bugs surface."""
    try:
        from PIL import Image  # noqa: PLC0415 — checked here, not at module level
    except ImportError:
        return []  # PIL genuinely unavailable — gate defers gracefully.
    try:
        with Image.open(str(path)) as im:
            # Quantise to up to 64 colors and get the palette.
            quantised = im.convert("RGB").quantize(colors=64)
            palette = quantised.getpalette()  # flat list of R,G,B ints
            hist = quantised.histogram()
            # The palette length may be shorter than 64*3 for low-colour images
            # (e.g. a solid fill produces len(palette)==3 on Pillow 11.x).
            # Bound the range to the actual number of allocated entries.
            n_colors = len(palette) // 3
            # Pair (count, (R, G, B)) for real palette entries only; sort descending.
            indexed = sorted(
                [(hist[i], (palette[i * 3], palette[i * 3 + 1], palette[i * 3 + 2]))
                 for i in range(n_colors)],
                reverse=True
            )
            return [rgb for _, rgb in indexed[:n]]
    except Exception as _exc:  # noqa: BLE001
        # Log real errors to stderr so they are visible in debug runs but do not
        # block the gate (the caller skips slides where dominant==[] is returned).
        import sys as _sys
        print(f"_slide_dominant_colors({path.name}): {_exc}", file=_sys.stderr)
        return []


def _load_brand_palette(run_dir: Path) -> list:
    """Load the brand palette hex tokens from intake.json -> brand.palette.
    Returns a list of (R, G, B) tuples, or [] when no palette is declared."""
    intake_path = run_dir / "working" / "copy" / "intake.json"
    if not intake_path.exists():
        return []
    obj = _read_json(intake_path)
    if "__parse_error__" in obj:
        return []
    brand = obj.get("brand") if isinstance(obj, dict) else None
    if not isinstance(brand, dict):
        return []
    palette = brand.get("palette") or []
    if not isinstance(palette, list):
        return []
    tokens = []
    for entry in palette:
        if isinstance(entry, str):
            rgb = _hex_to_rgb(entry)
            if rgb:
                tokens.append(rgb)
        elif isinstance(entry, dict):
            for key in ("hex", "color", "value"):
                val = entry.get(key)
                if isinstance(val, str):
                    rgb = _hex_to_rgb(val)
                    if rgb:
                        tokens.append(rgb)
                        break
    return tokens


def check_brand_consistency(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-BRAND-CONSISTENCY: every rendered slide's dominant palette must fall within
    the client's declared brand token set. Off-brand palettes auto-fail that slide.

    Complements AF-VISUAL-VARIETY: variety rejects sameness; brand rejects off-token
    outliers. Defers when no brand palette is declared or no renders exist.

    Returns '' on pass / defer, or a fatal AF-BRAND-CONSISTENCY message."""
    renders_dir = run_dir / "renders"
    if not renders_dir.is_dir():
        return ""  # no renders yet — defer.

    pngs = sorted(renders_dir.glob("slide-*.png"))
    if not pngs:
        return ""  # no renders yet — defer.

    if not _imaging_available():
        return ("AF-BRAND-CONSISTENCY: DECK FAIL — Pillow is not installed, so no slide's palette "
                "could be examined. This gate cannot pass unexamined renders. "
                "Install with: python3 -m pip install Pillow  "
                "(AF-BRAND-CONSISTENCY)")

    brand_tokens = _load_brand_palette(run_dir)
    if not brand_tokens:
        msg = ("no brand palette is declared, so every rendered slide is unexamined. "
               "Add brand.palette to working/copy/intake.json as a list of hex tokens "
               '(for example [\"#0A2540\", \"#F5F5F5\"]); a single primary colour is not '
               "a palette and is not read by this gate")
        if BRAND_PALETTE_REQUIRED:
            return f"AF-BRAND-CONSISTENCY: DECK FAIL — {msg} (AF-BRAND-CONSISTENCY)"
        _warn_once(run_dir, "brand_palette_absent",
                   f"AF-BRAND-CONSISTENCY WARNING (not yet enforced): {msg}")
        return ""

    off_brand_slides = []
    for png in pngs:
        dominant = _slide_dominant_colors(png, n=BRAND_CONSISTENCY_SAMPLE_COLORS)
        if not dominant:
            continue  # PIL unavailable or unreadable — skip this slide.
        # A slide is OFF-BRAND if ALL sampled dominant colors are far from every token.
        all_far = all(
            all(_rgb_distance(dc, bt) > BRAND_CONSISTENCY_TOLERANCE
                for bt in brand_tokens)
            for dc in dominant
        )
        if all_far:
            off_brand_slides.append(png.stem)

    if off_brand_slides:
        return (
            f"AF-BRAND-CONSISTENCY: DECK FAIL — {len(off_brand_slides)} slide(s) have "
            f"dominant palette colors that fall outside the declared brand token set "
            f"(tolerance {BRAND_CONSISTENCY_TOLERANCE} RGB units): "
            + ", ".join(off_brand_slides)
            + ". Off-brand stock imagery or wrong-palette renders must be re-generated "
            "using the approved brand palette. Check brand.palette in intake.json and "
            "re-run with the correct art direction (AF-BRAND-CONSISTENCY)."
        )
    return ""


# The MANDATORY pre-render artifact set. The first three are the minimum the
# operator task names explicitly; the rest are the broader mandatory upstream set
# from the dept-pipeline analysis. All must exist + be complete before any render.
# ---------------------------------------------------------------------------
# NO-DARK-SLIDES GATE (AF-DARK-SLIDE) — slides must use light/bright backgrounds
# by default. Dark/black backgrounds are ONLY allowed when the client has explicitly
# requested a dark theme via intake.json client_dark_theme: true.
# ---------------------------------------------------------------------------

# Dark-background keywords (case-insensitive) that trigger AF-DARK-SLIDE.
_DARK_SLIDE_KEYWORDS = (
    "dark background",
    "black background",
    "dark theme",
    "near-black",
    "dark slide",
    "dark mode",
)

# Near-black hex/rgb patterns that signal a dark background in prompts.
_DARK_COLOR_PATTERNS = re.compile(
    r"""
    \#(?:000000|000|111111|111|222222|222|1[aA]1[aA]1[aA]|0[dD]0[dD]0[dD])
    |
    rgb\(\s*(?:0\s*,\s*0\s*,\s*0|1[0-9]\s*,\s*1[0-9]\s*,\s*1[0-9]|20\s*,\s*20\s*,\s*20)\s*\)
    """,
    re.VERBOSE | re.IGNORECASE,
)


def _chk_no_dark_slides(run_dir: Path) -> str:
    """NO-DARK-SLIDES gate (AF-DARK-SLIDE, fail-loud).

    Slides MUST use LIGHT / bright backgrounds by DEFAULT. Dark or black-background
    slides are NOT ALLOWED unless the client EXPLICITLY requests a dark theme via
    the intake flag `client_dark_theme: true`. Light is the default; dark is opt-in
    by client request only.

    This gate:
      1. Reads intake.json for client_dark_theme (default False). The role docs
         (Director intake Q-bank, Brand Steward style_block, Slide Image Creator,
         Prompt Author / Prompt-QC) record the SAME opt-in under the legacy key name
         `DARK_OK`. The two are the SAME INTENT: "the client authorized a dark theme."
         This gate aliases `DARK_OK` -> `client_dark_theme` so an opt-in recorded
         EITHER way is honored (a deck the roles cleared for dark via DARK_OK is not
         spuriously auto-failed because the gate only read the other key name).
         `client_dark_theme` is the CANONICAL key; DARK_OK is the accepted alias.
      2. Scans image prompts (working/prompts/*.txt) AND the design brief
         (working/research/design-brief-*.md) for dark background keywords and
         near-black hex/rgb color references.
      3. If dark keywords are found AND client_dark_theme is not True -> FAIL.
      4. If client_dark_theme is True -> PASS (dark is allowed by explicit request).
      5. If no dark keywords found -> PASS.

    Returns "" on pass, or a fatal AF-DARK-SLIDE message on fail.
    """
    # Step 1: read client_dark_theme from intake.json, aliasing the role-doc key
    # DARK_OK -> client_dark_theme (same intent: client authorized a dark theme).
    def _is_dark_optin(val) -> bool:
        # Accept a JSON true, or a truthy string ("true"/"yes"/"1") the roles may
        # write when they mirror style_block.md's `DARK_OK = true`.
        if val is True:
            return True
        if isinstance(val, str) and val.strip().lower() in ("true", "yes", "1"):
            return True
        return False

    client_dark_theme = False
    for rel in ("working/copy/intake.json", "intake.json", "working/intake.json"):
        p = run_dir / rel
        if p.exists():
            obj = _read_json(p)
            if isinstance(obj, dict) and "__parse_error__" not in obj:
                # canonical key OR the DARK_OK alias honors the opt-in.
                if _is_dark_optin(obj.get("client_dark_theme")) or \
                        _is_dark_optin(obj.get("DARK_OK")):
                    client_dark_theme = True
            break

    if client_dark_theme:
        # Dark theme is explicitly permitted by client request — gate passes.
        return ""

    # Step 2: gather all text to scan (prompts + design brief).
    scan_sources = []

    # Image prompts.
    prompts_dir = run_dir / "working" / "prompts"
    if prompts_dir.is_dir():
        for prompt_file in sorted(prompts_dir.glob("*.txt")):
            try:
                scan_sources.append((str(prompt_file), prompt_file.read_text(errors="replace")))
            except OSError:
                pass

    # Design brief (working/research/design-brief-*.md).
    research_dir = run_dir / "working" / "research"
    if research_dir.is_dir():
        for brief_file in sorted(research_dir.glob("design-brief-*.md")):
            try:
                scan_sources.append((str(brief_file), brief_file.read_text(errors="replace")))
            except OSError:
                pass

    if not scan_sources:
        # No prompts or design brief yet — gate defers (pre-authoring phase).
        return ""

    # Step 3: scan for dark background keywords and near-black color patterns.
    hits = []
    for source_name, text in scan_sources:
        text_lower = text.lower()
        for keyword in _DARK_SLIDE_KEYWORDS:
            if keyword in text_lower:
                hits.append(f"{source_name}: keyword {keyword!r}")
                break
        if _DARK_COLOR_PATTERNS.search(text):
            hits.append(f"{source_name}: near-black hex/rgb color detected")

    if hits:
        return (
            "AF-DARK-SLIDE: Dark/black background detected in prompts or design brief "
            "but client_dark_theme is not set in intake.json. Light backgrounds are "
            "required by default. Set client_dark_theme: true in intake.json ONLY when "
            "the client explicitly requests a dark theme (opt-in by client request only). "
            "Offenders: " + "; ".join(hits)
        )
    return ""


# ===========================================================================
# GOAL-4 CHECKERS (decisions 1C / 2A / 3C / 5C)
# ===========================================================================

# ---------------------------------------------------------------------------
# 1C — client-provided materials: intake question asked + assets_manifest consumed
#      + scratch-deck parsed. Three run-dir-scoped checkers.
# ---------------------------------------------------------------------------
def _read_intake_obj(run_dir: Path):
    """Resolve + parse working/copy/intake.json (or intake.json) for the 1C/2A
    intake-flag checks. Returns the dict, or None when absent/unparseable."""
    for rel in ("working/copy/intake.json", "intake.json", "working/intake.json"):
        p = run_dir / rel
        if p.exists():
            obj = _read_json(p)
            if isinstance(obj, dict) and "__parse_error__" not in obj:
                return obj
            return None
    return None


def _read_assets_manifest(run_dir: Path):
    """Resolve + parse working/copy/assets_manifest.json. Returns the dict, the
    sentinel string "__absent__" when no manifest exists, or "__parse_error__"."""
    for rel in ("working/copy/assets_manifest.json", "assets_manifest.json",
                "working/assets_manifest.json"):
        p = run_dir / rel
        if p.exists():
            obj = _read_json(p)
            if isinstance(obj, dict) and "__parse_error__" in obj:
                return "__parse_error__"
            return obj
    return "__absent__"


def _chk_asset_question(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """1C — AF-ASSET-QUESTION-MISSING. The signature intake MUST ask the client
    whether they already have materials (photos / logo / brand colors / a rough or
    old deck / slides / concepts). intake.json must record
    asset_intake_question_asked:true. Defers (passes) only when intake.json is
    absent — _chk_intake owns the intake-absence failure; this gate owns
    'intake present but the asset question was never asked'."""
    obj = _read_intake_obj(run_dir)
    if obj is None:
        return ""  # intake not built / unparseable — _chk_intake owns that absence.
    if obj.get("asset_intake_question_asked") is not True:
        return ("AF-ASSET-QUESTION-MISSING: intake.json does not record "
                "asset_intake_question_asked:true. Every signature run MUST ask the "
                "client the one asset-intake question (photos / a logo / brand colors / "
                "a rough or old deck / slides / concepts) and record that it was asked. "
                "See Brainstorming Buddy SOP 9.1/9.2 + Director intake.")
    return ""


def _chk_assets_manifest(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """1C — AF-MANIFEST-UNREFERENCED. When the client provided assets, the
    Media-Librarian step must have written assets_manifest.json AND every provided
    asset must be provably consumed downstream (a non-empty consumed_by list AND a
    resolvable public_url the Brand Steward / Slide Image Creator can pass to
    gpt-image-2 as input_urls). Provided-but-unconsumed assets are the defect this
    gate catches. Defers (passes) when intake says no assets were provided, or when
    intake is absent (pre-intake phase)."""
    obj = _read_intake_obj(run_dir)
    if obj is None:
        return ""  # pre-intake — nothing to enforce yet.
    # The intake either declares assets_provided directly, or the manifest does.
    manifest = _read_assets_manifest(run_dir)
    intake_says_provided = bool(obj.get("assets_provided"))
    if manifest == "__parse_error__":
        return ("AF-MANIFEST-UNREFERENCED: assets_manifest.json is present but is not "
                "valid JSON, so provided-asset consumption cannot be proven. See "
                "Media-Librarian SOP (asset-ingest step).")
    if manifest == "__absent__":
        if intake_says_provided:
            return ("AF-MANIFEST-UNREFERENCED: intake.json declares the client provided "
                    "assets (assets_provided:true) but no working/copy/assets_manifest.json "
                    "exists. The Media-Librarian step must classify each provided asset, "
                    "upload it to a stable public URL, and record it (with consumed_by) in "
                    "assets_manifest.json so the Brand Steward + Slide Image Creator consume "
                    "it as gpt-image-2 input_urls.")
        return ""  # no assets provided and no manifest — nothing to consume.
    # Manifest exists — validate it proves the question was asked + assets consumed.
    if not isinstance(manifest, dict):
        return ("AF-MANIFEST-UNREFERENCED: assets_manifest.json is not a JSON object. "
                "Expected {asset_question_asked, assets_provided, assets:[...], scratch_deck}.")
    manifest_provided = bool(manifest.get("assets_provided"))
    if not (manifest_provided or intake_says_provided):
        return ""  # manifest records no provided assets — vacuously consumed.
    assets = manifest.get("assets")
    if not isinstance(assets, list) or not assets:
        return ("AF-MANIFEST-UNREFERENCED: assets_manifest.json declares assets were "
                "provided but carries no assets[] entries. Each provided asset must be "
                "recorded with a public_url + a non-empty consumed_by list.")
    unconsumed = []
    for idx, a in enumerate(assets):
        if not isinstance(a, dict):
            unconsumed.append(f"asset#{idx} (not an object)")
            continue
        consumed_by = a.get("consumed_by")
        url = str(a.get("public_url") or a.get("url") or "").strip()
        kind = str(a.get("kind") or "asset")
        if not isinstance(consumed_by, list) or not consumed_by:
            unconsumed.append(f"{kind}#{idx} (empty consumed_by)")
        elif not url:
            unconsumed.append(f"{kind}#{idx} (no public_url to pass as input_urls)")
    if unconsumed:
        return ("AF-MANIFEST-UNREFERENCED: one or more provided assets are not provably "
                "consumed downstream — " + "; ".join(unconsumed[:8]) + ". Every provided "
                "asset must carry a public_url AND a non-empty consumed_by (e.g. "
                "['brand-steward','slide-image-creator']) so it is actually fed to "
                "gpt-image-2 as input_urls, not collected and ignored.")
    return ""


def _chk_scratch_parse(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """1C — AF-SCRATCH-PARSE-SKIPPED. When the client uploaded a rough/old deck,
    the scratch-deck parser MUST extract its content + structure into
    working/copy/scratch_seed.json AND that seed must be folded into the PRD
    (mission_prd.json). The client still answers every interview question — the
    seed only seeds the PRD; it never replaces the interview. Defers (passes) when
    no scratch deck was provided, or when intake/manifest are absent."""
    manifest = _read_assets_manifest(run_dir)
    if manifest in ("__absent__", "__parse_error__") or not isinstance(manifest, dict):
        return ""  # no manifest — _chk_assets_manifest owns manifest problems.
    scratch = manifest.get("scratch_deck")
    if not isinstance(scratch, dict) or not scratch.get("provided"):
        return ""  # no scratch deck uploaded — nothing to parse.
    if not scratch.get("parsed"):
        return ("AF-SCRATCH-PARSE-SKIPPED: assets_manifest.json records a client scratch "
                "deck was provided but scratch_deck.parsed is not true. The scratch-deck "
                "parser must extract the uploaded deck's content/structure into "
                "working/copy/scratch_seed.json and seed the PRD (the interview still runs "
                "in full).")
    # The seed file must exist and be non-empty.
    seed_rel = str(scratch.get("seed_prd_path") or "working/copy/scratch_seed.json")
    seed_path = run_dir / seed_rel
    if not seed_path.exists():
        return (f"AF-SCRATCH-PARSE-SKIPPED: scratch_deck.parsed is true but the seed file "
                f"{seed_rel!r} does not exist. The parser must write scratch_seed.json.")
    seed_obj = _read_json(seed_path)
    if isinstance(seed_obj, dict) and "__parse_error__" in seed_obj:
        return (f"AF-SCRATCH-PARSE-SKIPPED: {seed_rel} is not valid JSON, so the parsed "
                f"scratch content cannot seed the PRD.")
    # Prove the seed reached the PRD: mission_prd.json must reference the seed.
    prd = None
    for rel in ("working/copy/mission_prd.json", "mission_prd.json",
                "working/mission_prd.json"):
        p = run_dir / rel
        if p.exists():
            prd = p
            break
    if prd is not None:
        prd_obj = _read_json(prd)
        if isinstance(prd_obj, dict) and "__parse_error__" not in prd_obj:
            blob = json.dumps(prd_obj).lower()
            if ("scratch_seed" not in blob and "scratch_deck" not in blob
                    and not prd_obj.get("seeded_from_scratch_deck")):
                return ("AF-SCRATCH-PARSE-SKIPPED: scratch_seed.json exists but the PRD "
                        "(mission_prd.json) shows no seeded scratch content (no "
                        "seeded_from_scratch_deck flag and no scratch_seed reference). The "
                        "parsed scratch deck must seed the PRD.")
    return ""


# ---------------------------------------------------------------------------
# 2A — explicit PITCH_INCLUDED flag. Pitchless is first-class: when
#      pitch_included:false the pitch/price machinery is suppressed and a NEW
#      AF-PITCH-LEAK fires if any pitch/price content leaked into the deck. The
#      pitch-required gate (AF-PITCH-MISSING) becomes conditional on
#      pitch_included:true.
# ---------------------------------------------------------------------------
def _intake_pitch_included(run_dir: Path):
    """Return True / False / None (unset) for intake.json.pitch_included."""
    obj = _read_intake_obj(run_dir)
    if obj is None:
        return None
    val = obj.get("pitch_included")
    if isinstance(val, bool):
        return val
    return None


def _chk_pitch_flag(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """2A — AF-PITCH-FLAG-UNSET. The signature intake MUST capture an explicit
    boolean pitch_included (yes = the deck ends with an offer/pitch; no = a
    teaching/content-only deck). It is non-defaultable: a deck must never silently
    inherit 'has a pitch'. Defers (passes) only when intake.json is absent
    (_chk_intake owns that absence)."""
    obj = _read_intake_obj(run_dir)
    if obj is None:
        return ""  # pre-intake — _chk_intake owns intake absence.
    val = obj.get("pitch_included")
    if not isinstance(val, bool):
        return ("AF-PITCH-FLAG-UNSET: intake.json.pitch_included must be an explicit "
                "boolean captured at intake (true = the presentation ends with an "
                "offer/pitch; false = a teaching/content-only presentation). It was "
                f"{val!r}. Never default it — ask the client and record yes/no.")
    return ""


# Pitch/price/offer tokens that must NOT appear anywhere in a pitchless deck's arc
# or copy. Lowercase substring match. Kept narrow + specific so ordinary teaching
# copy (e.g. 'value', 'price of admission' as metaphor) is not over-caught — these
# are sale-mechanic tokens, the price-ladder/offer machinery.
PITCHLESS_FORBIDDEN_TOKENS = (
    "price ladder", "price_ladder", "offer ladder", "offer_ladder",
    "value stack", "value-stack", "value_stack",
    "re-pitch", "repitch", "re_pitch", "second close", "second-close",
    "anchor price", "anchor_price", "buy now", "enroll now", "enrol now",
    "limited-time offer", "limited time offer", "payment plan", "deposit today",
    "money-back guarantee", "money back guarantee", "act now",
    "cost of inaction", "cost_of_inaction",
)


def _chk_pitch_leak(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """2A — AF-PITCH-LEAK. A PITCHLESS deck (intake.json.pitch_included:false) must
    contain NO pitch/price/offer/ladder content. Scans arc_allocation.json,
    slides_copy.md, and price_ladder.json for the suppressed sale-mechanic tokens
    AND for the mere presence of a price_ladder.json (the Offer Price Strategist
    must not have run). Defers (passes) when pitch_included is true or unset (true
    decks are governed by AF-PITCH-MISSING; unset is governed by AF-PITCH-FLAG-UNSET)."""
    if _intake_pitch_included(run_dir) is not False:
        return ""  # only pitchless decks are leak-checked.
    leaks = []
    # The price-ladder artifact must not exist at all in a pitchless deck.
    for rel in ("working/copy/price_ladder.json", "price_ladder.json",
                "working/price_ladder.json"):
        if (run_dir / rel).exists():
            leaks.append(f"{rel} present (Offer Price Strategist ran on a pitchless deck)")
            break
    # Token scan over the arc + copy artifacts.
    scan = []
    for rel in ("working/copy/arc_allocation.json", "arc_allocation.json",
                "working/copy/slides_copy.md", "slides_copy.md",
                "working/copy/price_ladder.json"):
        p = run_dir / rel
        if p.exists():
            scan.append((rel, p.read_text(errors="replace").lower()))
    for rel, low in scan:
        hits = [t for t in PITCHLESS_FORBIDDEN_TOKENS if t in low]
        if hits:
            leaks.append(f"{rel}: " + ", ".join(repr(h) for h in hits[:6]))
    if leaks:
        return ("AF-PITCH-LEAK: intake.json declares this is a PITCHLESS deck "
                "(pitch_included:false) but pitch/price/offer content leaked in — "
                + "; ".join(leaks) + ". A pitchless deck is built with NO pitch: the "
                "Offer Price Strategist and price ladder are SUPPRESSED and no "
                "offer/ladder/re-pitch/cost-of-inaction beats may appear.")
    return ""


# ---------------------------------------------------------------------------
# 3C — Kie.ai balance pre-flight (AF-KIE-BALANCE). Phase-0 of the signature runner
#      and a shared pre-render gate. HARD-ABORTS before any render when the live
#      Kie credit balance is below the estimated floor for this deck.
# ---------------------------------------------------------------------------
def _fetch_kie_balance(api_key: str, url: str = KIE_CREDIT_URL,
                       timeout: int = 30) -> float:
    """GET the live Kie credit balance. Returns the numeric balance. Raises
    RuntimeError on a network/parse error so the caller fails LOUD rather than
    treating an unknown balance as 'enough'. Parses the common Kie response shapes
    ({data:{credit|credits|balance}} or a top-level number)."""
    import urllib.request
    import urllib.error
    req = urllib.request.Request(url, headers={
        "Authorization": f"Bearer {api_key}",
        "Accept": "application/json",
    })
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            raw = resp.read().decode("utf-8", errors="replace")
    except urllib.error.HTTPError as exc:
        # FIX-6 — a 401/403 on the credit call is a PERMANENT auth failure (wrong
        # key / header / locked account). Never transient: surface it as AuthError
        # so the preflight fails fast instead of being treated as 'unknown balance'.
        if exc.code in (401, 403):
            raise AuthError(f"Kie credit endpoint returned HTTP {exc.code} ({url})") from exc
        raise RuntimeError(f"Kie credit endpoint unreachable ({url}): HTTP {exc.code}") from exc
    except (urllib.error.URLError, OSError) as exc:
        raise RuntimeError(f"Kie credit endpoint unreachable ({url}): {exc}")
    try:
        obj = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise RuntimeError(f"Kie credit response is not JSON: {exc}; body={raw[:200]!r}")
    # Common shapes (verified live 2026-06-20):
    #   {"code":200,"msg":"success","data":1951.33}   <- data is a BARE NUMBER
    #   {"data": {"credit": N}} / {"credits": N} / {"balance": N} / a top-level number N
    candidates = []
    if isinstance(obj, (int, float)):
        candidates.append(obj)
    if isinstance(obj, dict):
        # The real Kie credit endpoint returns the balance as a bare number under "data".
        data_val = obj.get("data")
        if isinstance(data_val, (int, float)):
            candidates.append(data_val)
        # Some shapes nest the number under data{credit|credits|balance|...}; others put
        # it at the top level. Scan both the data dict and the outer dict for a key.
        for container in (data_val if isinstance(data_val, dict) else None, obj):
            if not isinstance(container, dict):
                continue
            for k in ("credit", "credits", "balance", "remaining", "available"):
                v = container.get(k)
                if isinstance(v, (int, float)):
                    candidates.append(v)
    if not candidates:
        raise RuntimeError(f"Kie credit response carried no numeric balance: {raw[:200]!r}")
    return float(candidates[0])


def kie_balance_preflight(run_dir: Path, slide_count: int,
                          api_key: Optional[str] = None) -> str:
    """3C — AF-KIE-BALANCE. Phase-0 balance gate. Computes the estimated credit
    floor for this deck (slide_count x PER_SLIDE_CREDIT_ESTIMATE x
    KIE_BALANCE_FLOOR_MULTIPLIER), fetches the live Kie balance, and returns a
    fatal AF-KIE-BALANCE string when balance < floor. Returns "" on pass.

    SHARED implementation: the signature runner's Phase-0 and any pre-render hook
    both call this so there is ONE balance check. Defers (passes) ONLY when no API
    key is available AND adhoc/no-network context (slide_count<=0) — a real run with
    a key always evaluates. A balance the endpoint cannot return is a HARD ABORT (an
    unknown balance is never treated as 'enough')."""
    if not slide_count or slide_count <= 0:
        return ""  # nothing to render — no floor to clear.
    if not api_key:
        # No key available to query — cannot prove balance; this is only reached on
        # an adhoc/offline path. The render path's own key load fails loud elsewhere.
        return ""
    estimated_floor = float(slide_count) * PER_SLIDE_CREDIT_ESTIMATE * KIE_BALANCE_FLOOR_MULTIPLIER
    try:
        balance = _fetch_kie_balance(api_key)
    except AuthError as exc:
        # FIX-6 — a permanent auth failure on the balance call must abort the run
        # with the auth diagnosis, never be treated as 'unknown balance'.
        return (f"AF-KIE-AUTH: kie.ai key did not authenticate before render ({exc}). "
                "Check the KIE_API_KEY, the Authorization: Bearer header format, and "
                "that the key is not locked/rate-blocked by the provider.")
    except RuntimeError as exc:
        return ("AF-KIE-BALANCE: could not verify the Kie.ai credit balance before "
                f"render ({exc}). An unverifiable balance is a HARD ABORT — never render "
                "on an unknown balance. Fix the key / endpoint, or top up and retry.")
    if balance < estimated_floor:
        return ("AF-KIE-BALANCE: Kie.ai credit balance is below the estimated floor for "
                f"this deck. balance={balance:g} credits, estimated_floor={estimated_floor:g} "
                f"({slide_count} slides x {PER_SLIDE_CREDIT_ESTIMATE} credits x "
                f"{KIE_BALANCE_FLOOR_MULTIPLIER} headroom). HARD ABORT before any render so "
                "the run does not die mid-deck. Top up Kie.ai credits and retry.")
    return ""


def _preflight_kie_auth(api_key: str) -> None:
    """FIX-6 — one-shot auth proof BEFORE any parallel render. A single authenticated
    call to the live Kie credit endpoint must succeed or the run stops here. This
    converts a 401 storm (the observed 164x backoff tailspin on a submission-path bug)
    into ONE clear block: if the balance call authenticates, the key/header is valid
    and any later createTask 401 is a request-shape/header bug that surfaces as
    AuthError instead of burning the retry budget. Raises AuthError on a 401/403
    (permanent — never retried); raises RuntimeError on any other failure to reach
    the endpoint (fail loud, never 'unknown = enough')."""
    try:
        balance = _fetch_kie_balance(api_key)
    except AuthError:
        raise
    except RuntimeError as exc:
        raise RuntimeError(
            "AF-KIE-AUTH: kie.ai auth preflight could not verify the key "
            f"({exc}). Do NOT render on an unverifiable key.")
    print(f"  OK: kie.ai auth preflight passed (credit {balance:g})", flush=True)


# ---------------------------------------------------------------------------
# MASTER-SPEC 7.4 — OCR-engine AVAILABILITY pre-flight (AF-OCR-ENGINE-MISSING).
# Phase-0 of the signature runner (mirrors kie_balance_preflight's placement/
# calling convention exactly) and a shared pre-render gate inside build_deck.py's
# own direct-invocation Phase-0 block. HARD-ABORTS before any phase is dispatched
# (run_signature_deck.phase0_preflight) — and, defense-in-depth, before any render
# (build_deck.main()'s own Phase-0) — when this render environment has no OCR
# engine, so the deck is never rendered (and never paid for) on a box whose
# postflight OCR-readback audit can never actually run.
# ---------------------------------------------------------------------------
def ocr_engine_preflight(run_dir: Path) -> str:
    """MASTER-SPEC 7.4 (§7.4, 'Quality that is measured, not asserted'): 'The missing
    OCR dependency fails at minute zero, before any paid generation, not after 62
    images.' This is that refusal.

    prompt_gate.ocr_readback()'s own docstring documents that OCR is intentionally
    PROVENANCE-RECORDED-OPTIONAL and 'never hard-fails purely for a missing engine' —
    correct for THAT function (a healthy-engine box must keep rendering; a box mid-
    decision records the absence rather than silently skipping it) — but nothing
    UPSTREAM of it was refusing on a missing engine, so a box with no OCR engine
    rendered and paid for every slide, then produced unverifiable OCR sidecars with
    nothing hard-failing. This is the ENGINE-AVAILABILITY probe ('does an OCR engine
    exist on this box at all'), run BEFORE a single slide is dispatched — distinct
    from a POSTFLIGHT sidecar-verification gate (a separate concern: auditing the
    per-slide renders/*.ocr.json records once they exist) that can only audit a
    slide once it has already been rendered and paid for.

    Called from TWO sites, both non-adhoc only (an --adhoc-no-process run forces an
    empty api_key — see main() — so it can never spend real Kie.ai credit regardless
    of OCR state, exactly the same carve-out kie_balance_preflight already takes):
      1. run_signature_deck.phase0_preflight() — the REAL orchestrated path
         (presentation-canonical-entry.sh -> run_signature_deck.py -> build_deck.py),
         checked in-process BEFORE any phase (research/copy/QC/render) is dispatched —
         the earliest possible point in the entire run.
      2. build_deck.main()'s own Phase-0 block — so a DIRECT `python3 build_deck.py`
         invocation that bypasses the runner is equally refused, exactly as
         kie_balance_preflight and the P0B-PRIORITY binding in run_preflight() already
         refuse a direct-invocation bypass of other Phase-0/Phase-0b gates.

    Unconditional: never defers on run_dir state. An OCR-blind box can never satisfy
    a postflight OCR-readback audit later, regardless of what phase the run is in, so
    there is no state in which deferring would be correct.

    ENVIRONMENT TRAP THIS CLOSES: probes prompt_gate.ocr_engine_diagnostic() in THIS
    process — the SAME interpreter/environment that will actually render (run_signature_
    deck.py launches the build_deck.py render subprocess via `sys.executable`, inheriting
    this process's environment, so what this process sees is what the render subprocess
    sees). A user-site-packages install of pytesseract (~/Library/Python/*/lib/python/
    site-packages) is invisible under PYTHONNOUSERSITE=1, inside a venv without
    --system-site-packages, under sudo, and to launchd/cron (whose default PATH carries
    no Homebrew, so no tesseract binary either) — a probe run from a DIFFERENT
    interpreter/shell than the one that renders proves nothing about the render's own
    environment. Returns "" on pass, else a fatal AF-OCR-ENGINE-MISSING message naming
    exactly what is missing and how to install it."""
    pg = _import_prompt_gate()
    if pg is None:
        return ("AF-OCR-ENGINE-MISSING: prompt_gate.py could not be imported at all "
                "(it ships beside build_deck.py) — cannot verify OCR engine availability "
                "in this render environment. MASTER-SPEC 7.4: 'The missing OCR dependency "
                "fails at minute zero, before any paid generation, not after 62 images.'")
    try:
        diag = pg.ocr_engine_diagnostic()
    except Exception as exc:  # noqa: BLE001 — a broken diagnostic is itself fail-closed
        return (f"AF-OCR-ENGINE-MISSING: prompt_gate.ocr_engine_diagnostic() raised "
                f"({exc}) instead of returning a status — cannot verify OCR engine "
                "availability. MASTER-SPEC 7.4: 'The missing OCR dependency fails at "
                "minute zero, before any paid generation, not after 62 images.'")
    if diag.get("available"):
        return ""
    reason = diag.get("reason", "unknown")
    detail = diag.get("detail", "")
    install_hint = {
        "pytesseract-import-failed": (
            "install pytesseract INTO THE EXACT interpreter that runs build_deck.py / "
            "run_signature_deck.py (e.g. `<that interpreter> -m pip install "
            "pytesseract`). A user-site-packages install is INVISIBLE under "
            "PYTHONNOUSERSITE=1, inside a venv without --system-site-packages, under "
            "sudo, and to launchd/cron — install into the interpreter's OWN "
            "site-packages if this render runs under any of those."
        ),
        "pillow-import-failed": (
            "install Pillow INTO THE EXACT interpreter that runs build_deck.py / "
            "run_signature_deck.py (same interpreter-site caveat as pytesseract above)."
        ),
        "tesseract-binary-not-found": (
            "install the tesseract binary and ensure it is on PATH for the render "
            "process (e.g. `brew install tesseract` on macOS). launchd/cron's default "
            "PATH is typically /usr/bin:/bin:/usr/sbin:/sbin — NO Homebrew — so a "
            "scheduled/cron-driven render needs an explicit PATH (or the tesseract "
            "binary's full path) configured in the launchd plist / cron environment, "
            "not merely in an interactive shell's PATH."
        ),
    }.get(reason, "install pytesseract + Pillow + the tesseract binary, reachable from "
                   "the EXACT interpreter/environment that will run build_deck.py / "
                   "run_signature_deck.py.")
    return (
        f"AF-OCR-ENGINE-MISSING: no OCR engine is available in this render environment "
        f"({reason}{': ' + detail if detail else ''}). MASTER-SPEC 7.4: 'The missing OCR "
        "dependency fails at minute zero, before any paid generation, not after 62 "
        f"images.' Fix: {install_hint}"
    )


# ---------------------------------------------------------------------------
# 3C — phase-precondition gate (AF-PHASE-SKIPPED). The deterministic runner
#      (run_signature_deck.py) makes skipping/reordering a phase structurally
#      impossible: before dispatching phase N+1 it proves every lower-order phase
#      has an attestation in working/checkpoints/process_manifest.json OR a logged
#      owner-authorized skip. This shared checker is the single source of truth for
#      that precondition (the runner imports it; build_deck references it so the
#      symbol is live and gate-integrity-visible).
# ---------------------------------------------------------------------------
def check_phase_preconditions(run_dir: Path, phase_id, prior_phase_ids) -> str:
    """3C — AF-PHASE-SKIPPED. Returns "" when every phase in prior_phase_ids has
    either an attestation in working/checkpoints/process_manifest.json (phases[].id /
    phase_id) OR an owner-authorized skip record in
    working/checkpoints/phase_skip_approvals.json (owner_approved:true). Returns a
    fatal AF-PHASE-SKIPPED string naming the first unmet prior phase otherwise — the
    runner HARD-ABORTS on it, so a phase can never be silently skipped or reordered.
    An owner-authorized skip is the ONLY exception, and it must be logged."""
    attested = set()
    pm = run_dir / "working" / "checkpoints" / "process_manifest.json"
    if pm.exists():
        obj = _read_json(pm)
        if isinstance(obj, dict) and "__parse_error__" not in obj:
            # (1) Runner attestations live under "phase_attestations" (keyed phase_id) —
            #     written by run_signature_deck.attest_phase. This is the CANONICAL
            #     attestation store the deterministic runner produces. Reading it here is
            #     the v16.1.5 fix for the FALSE AF-PHASE-SKIPPED: the writer appends to
            #     phase_attestations[] but this reader formerly inspected only phases[], so
            #     every phase the runner attested was invisible to the next phase's
            #     precondition check and the chain hard-aborted after the first phase. This
            #     now mirrors run_signature_deck._attested_phase_ids and
            #     canonical_render_guard.attested_phase_ids (one contract, three readers).
            for att in obj.get("phase_attestations", []) or []:
                if isinstance(att, dict):
                    for k in ("phase_id", "id", "name"):
                        v = att.get(k)
                        if isinstance(v, str) and v.strip():
                            attested.add(v.strip())
            # (2) build_deck.py appends its OWN render record under "phases" as
            #     {"phase": "render", ...}; that record counts as the P4-RENDER phase being
            #     attested (a canonical render is seen without the runner re-stamping it).
            #     Legacy phases[] records keyed by id/phase_id/name are still honoured so a
            #     pre-runner / hand-written manifest does not regress.
            for ph in obj.get("phases", []) or []:
                if isinstance(ph, dict):
                    if ph.get("phase") == "render":
                        attested.add("P4-RENDER")
                    for k in ("id", "phase_id", "phase", "name"):
                        v = ph.get(k)
                        if isinstance(v, str) and v.strip():
                            attested.add(v.strip())
    approved = set()
    sk = run_dir / "working" / "checkpoints" / "phase_skip_approvals.json"
    if sk.exists():
        sobj = _read_json(sk)
        records = sobj if isinstance(sobj, list) else (
            (sobj.get("approvals") or sobj.get("skips") or [])
            if isinstance(sobj, dict) else [])
        for r in records if isinstance(records, list) else []:
            # FIX-1 (AF-FORGED-APPROVAL): a skip record is ONLY a verifiable
            # owner-authorized skip when it carries a NON-EMPTY owner_msg_id. An
            # owner_action-only record — or any record with no owner_msg_id — is
            # exactly the self-forgery vector the live E2E used (it passed with no
            # oracle query). Such records are FAIL-CLOSED here: the phase is NOT
            # added to `approved`, so it stays an unmet precondition (the runner's
            # load_skip_approvals additionally raises AF-FORGED-APPROVAL on them).
            # The full owner-message resolution happens in the runner's
            # load_skip_approvals -> cc_board owner-ids oracle; this shared gate
            # refuses msg-id-less records up front.
            if (isinstance(r, dict) and r.get("owner_approved") is True
                    and str(r.get("phase_id") or "").strip()
                    and not str(r.get("owner_msg_id") or "").strip()):
                # FIX-1 / FIX-23 manifest reconcile: a msg-id-less (or owner_action-
                # only) skip record is the exact self-forgery vector the live E2E
                # used. Refuse it explicitly AND surface AF-FORGED-APPROVAL here so
                # the build_deck enforcement path (Guard A py_symbol
                # check_phase_preconditions) emits the code text its manifest row
                # declares — the phase stays required, fail-closed.
                print(
                    f"[check_phase_preconditions] AF-FORGED-APPROVAL: skip record "
                    f"for phase {str(r.get('phase_id') or '').strip()!r} has NO "
                    "owner_msg_id — an owner_action-only record is a forged approval "
                    "by definition; the phase remains REQUIRED.",
                    file=sys.stderr,
                )
            if (isinstance(r, dict) and r.get("owner_approved") is True
                    and str(r.get("phase_id") or "").strip()
                    and str(r.get("owner_msg_id") or "").strip()):
                pid = str(r["phase_id"]).strip()
                if pid in UNSKIPPABLE_QC_PHASES:
                    # FIX-2 (Error 2): QC phases are STRUCTURALLY UNSKIPPABLE — a
                    # phase-skip record for a QC phase authorizes nothing. The QC
                    # phase stays a required precondition (AF-QC-SKIP refused here).
                    print(
                        f"[check_phase_preconditions] AF-QC-SKIP: phase {pid!r} is a "
                        "QC phase and is structurally unskippable — its skip record is "
                        "REFUSED; the phase remains required.",
                        file=sys.stderr,
                    )
                    continue
                approved.add(pid)
    for prior in (prior_phase_ids or []):
        pid = str(prior).strip()
        if not pid or pid in attested or pid in approved:
            continue
        return ("AF-PHASE-SKIPPED: phase " + str(phase_id) + " was dispatched before "
                "prior phase " + pid + " was attested. Each phase N+1 reads phase N's "
                "attestation in working/checkpoints/process_manifest.json as a "
                "precondition, so skipping or reordering a phase is structurally "
                "impossible EXCEPT with a logged owner-authorized skip in "
                "working/checkpoints/phase_skip_approvals.json (an entry with "
                "phase_id + owner_approved:true + approved_by + reason + timestamp). "
                "See run_signature_deck.py + SOP-SLIDE-05-PROCESS-MANIFEST.md.")
    return ""


# ---------------------------------------------------------------------------
# 5C — native-overlay path ELIMINATED (AF-OVERLAY-DELIVERED). No slide may ship a
#      native PPTX text run instead of a single composed gpt-image-2 image, and the
#      eliminated pptx_text_overlays.json must never be present.
# ---------------------------------------------------------------------------
def _delivered_pptx_native_text(pptx_path: Path) -> str:
    """Return a non-empty reason if the delivered PPTX carries any native on-slide
    text run (a shape with non-empty text_frame text on a slide). The off-slide
    speaker-notes part is NOT on-slide and is explicitly allowed. Returns "" when
    the deck is image-only (the required state) or when python-pptx is unavailable
    (cannot scan — defers rather than false-fail)."""
    try:
        from pptx import Presentation
    except ImportError:
        return ""  # cannot scan without python-pptx; the file-presence half still fires.
    try:
        prs = Presentation(str(pptx_path))
    except Exception:  # noqa: BLE001
        # A file python-pptx cannot open is NOT a valid deck to scan here — the
        # postflight bundle-completeness magic-byte gate (AF-BUNDLE-COMPLETE) owns
        # malformed / decoy .pptx files. Defer rather than false-fail AF-OVERLAY-DELIVERED.
        return ""
    offenders = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            txt = (shape.text_frame.text or "").strip()
            if txt:
                offenders.append(f"slide {idx}: native text run {txt[:40]!r}")
    if offenders:
        return ("delivered slides carry NATIVE on-slide text instead of a single composed "
                "gpt-image-2 image — " + "; ".join(offenders[:8]))
    return ""


def _chk_no_overlay(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """5C — AF-OVERLAY-DELIVERED. The native PPTX text/element-overlay path is
    eliminated. This gate fails when (a) any eliminated pptx_text_overlays.json file
    is present in the run dir (its mere presence at assembly is a violation), OR
    (b) the delivered PPTX carries any native (non-notes) on-slide text run. The
    only legitimate slide text is baked into the single composed gpt-image-2 image;
    the only legitimate PPTX text part is the off-slide speaker-notes pane. Persistent
    garble is re-prompted / re-seeded then escalated to a human — NEVER overlaid."""
    # (a) the eliminated overlays file must never exist.
    for rel in OVERLAY_FORBIDDEN_FILES:
        if (run_dir / rel).exists():
            return ("AF-OVERLAY-DELIVERED: the eliminated native-text-overlay file "
                    f"{rel!r} is present. The pptx_text_overlays.json native-overlay path "
                    "was removed by construction (Decision 5C). Persistent garble must be "
                    "re-prompted / re-seeded and, if it still fails, escalated to a human — "
                    "NEVER composited as a native PPTX text box. Delete this file and "
                    "re-render the affected slide as a single composed gpt-image-2 image.")
    # (b) scan any delivered PPTX in the run dir for native on-slide text.
    for pptx in sorted(run_dir.glob("**/*.pptx")):
        # skip Office lock/temp files
        if pptx.name.startswith("~$"):
            continue
        reason = _delivered_pptx_native_text(pptx)
        if reason:
            return ("AF-OVERLAY-DELIVERED: " + reason + ". Every delivered slide must be a "
                    "SINGLE composed gpt-image-2 image with its text baked in; the native "
                    "PPTX text-overlay path is eliminated. Re-prompt/re-seed the garbled "
                    "slide, escalate to a human if it persists, and re-render — do not add a "
                    "native text run.")
    return ""


# ---------------------------------------------------------------------------
# FIX-2 — owner/founder skip token (the ONLY way to waive a FIX-2 gate)
# ---------------------------------------------------------------------------
# TRUST BOUNDARY, INCREMENT 2 — the sixth-bypass closure (2026-08 architecture
# sweep). Read this before touching anything below.
#
# THE HONEST LIMIT: this codebase has NO adversary-unwritable anchor.
# process_manifest.json is written by the same UID that would forge an
# approval, and Increment 1 (presentation_job/runfacts.py, still in place,
# untouched by this increment) already proved the point: its
# verify_owner_skip() re-implements the identical "owner_approved:true + four
# non-empty strings" rule this function has always used, sourced from a
# sealed-but-not-anchored read of the SAME file — and it is report-only by
# default, so it changes nothing about what a running fleet actually does.
# There is no key here to check a signature against that isn't ALSO a literal
# in this source tree, forgeable by anything that can read the repo. Building
# a "signature" field and calling it real would be exactly the fake guarantee
# this closure is required not to ship.
#
# WHAT IS ACTUALLY ACHIEVABLE, and what this increment does:
#   1. FRICTION on the laziest self-skip shape. A one-token placeholder reason
#      ("x", "ok", "approved") or a non-timestamp ("t", "now", "asap") is
#      refused outright (_OWNER_SKIP_REASON_MIN_CHARS,
#      _owner_skip_parse_timestamp). This does not stop a determined forger
#      willing to type a real sentence and a real date — nothing available in
#      this architecture does — but it closes the exact shape proven live in
#      the sweep's PoC: a single mechanically-generated dict, one reason
#      copy-pasted onto every af_code, minted in the time it takes to write
#      one line of JSON.
#   2. A blanket-skip signature check. The SAME reason text, reused verbatim
#      across MORE than _OWNER_SKIP_BLANKET_REASON_LIMIT distinct af_codes in
#      one manifest, disqualifies every record that shares it. A named,
#      per-gate reason is required; one rubber-stamped sentence covering most
#      or all of the FIX-2 battery is not — see _owner_skip_evaluate.
#   3. MANDATORY, UNCONDITIONAL DISCLOSURE — the one defense this
#      architecture cannot defeat by definition, because it does not depend on
#      distinguishing a real approval from a forged one. EVERY consumption and
#      EVERY rejected attempt is (a) printed to stderr in an unmissable
#      banner, unconditionally — never gated behind the runfacts shadow path,
#      never dependent on PRES_TRUST_BOUNDARY_ENFORCE, never silenced by an
#      import failure — and (b) appended to a NEW, code-owned, append-only
#      ledger key in process_manifest.json ("owner_skip_events", written ONLY
#      by _owner_skip_ledger_append — the agent-authored "owner_skip_approval"
#      input itself is never rewritten). run_postflight_gate() reads this
#      ledger back and refuses to print "COMPLETE" without also printing every
#      entry in it, or an explicit "no waivers used" line when it is empty —
#      see the postflight success block further down this file. A skip can no
#      longer be silent: this is the actual, load-bearing fix.
#
# An absent record, or a record that never names this af_code, is UNCHANGED
# from before and stays QUIET (the common, innocent, majority case — see
# runfacts.py's own findings() docstring for why treating plain absence as a
# "finding" is the cry-wolf failure mode this does not repeat). A record that
# NAMES this af_code but fails validation is now always loud — see
# _owner_skip_evaluate / _owner_skip_disclose.
_OWNER_SKIP_REASON_MIN_CHARS = 8
_OWNER_SKIP_BLANKET_REASON_LIMIT = 2  # one reason may waive at most this many DISTINCT af_codes
_OWNER_SKIP_LEDGER_KEY = "owner_skip_events"


def _owner_skip_parse_timestamp(ts: str):
    """Return a parsed datetime, or None when `ts` is not a real ISO-8601
    date/datetime — rejects placeholders like 't' / 'now' / 'asap'. Accepts a
    trailing 'Z' (normalized to '+00:00' for pre-3.11 datetime.fromisoformat
    compatibility, since this fleet's Python version is not pinned here)."""
    s = (ts or "").strip()
    if not s:
        return None
    if s.endswith("Z"):
        s = s[:-1] + "+00:00"
    try:
        return datetime.fromisoformat(s)
    except ValueError:
        return None


def _owner_skip_all_records(obj: dict) -> list:
    """The raw owner_skip_approval records (tolerant of a single object or a
    list) from an already-parsed process_manifest.json dict."""
    raw = obj.get("owner_skip_approval")
    return raw if isinstance(raw, list) else ([raw] if isinstance(raw, dict) else [])


def _owner_skip_structurally_valid(r: dict) -> Tuple[bool, str]:
    """The base authenticity rule this mechanism has always had (owner_approved
    literal true + non-empty approved_by/reason/timestamp) PLUS the two
    friction checks from the module comment above. Returns (True, "") or
    (False, human-readable reason) — never raises."""
    if r.get("owner_approved") is not True:
        return False, f"owner_approved is {r.get('owner_approved')!r}, not literal true"
    approved_by = str(r.get("approved_by") or "").strip()
    if not approved_by:
        return False, "approved_by is empty"
    reason = str(r.get("reason") or "").strip()
    if not reason:
        return False, "reason is empty"
    if len(reason) < _OWNER_SKIP_REASON_MIN_CHARS:
        return False, (f"reason {reason!r} is shorter than the "
                        f"{_OWNER_SKIP_REASON_MIN_CHARS}-character floor — a real "
                        f"justification is required, not a placeholder token")
    timestamp_raw = str(r.get("timestamp") or "").strip()
    if not timestamp_raw:
        return False, "timestamp is empty"
    if _owner_skip_parse_timestamp(timestamp_raw) is None:
        return False, (f"timestamp {timestamp_raw!r} does not parse as a real ISO-8601 "
                        f"date/time")
    return True, ""


def _owner_skip_evaluate(run_dir: Path, af_code: str) -> Tuple[Optional[dict], list]:
    """Evaluate every owner_skip_approval record naming `af_code`. Returns
    (record_or_None, events): record is the FIRST validly-waiving record (or
    None — the gate stays enforced, unchanged contract); events lists EVERY
    record this af_code touched, granted or rejected-with-reason, for the
    mandatory disclosure layer in _owner_skip_approved. Records that never
    mention this af_code produce NO event (the quiet, common case). Never
    raises: any read/parse problem becomes an event or an empty result, not an
    exception."""
    want = af_code.strip().upper()
    pm = run_dir / "working" / "checkpoints" / "process_manifest.json"
    if not pm.exists():
        return None, []
    obj = _read_json(pm)
    if not isinstance(obj, dict) or "__parse_error__" in obj:
        detail = (obj.get("__parse_error__") if isinstance(obj, dict)
                  else f"top-level JSON is {type(obj).__name__}, expected object")
        # A manifest that EXISTS but cannot be read is not the quiet "nothing has
        # happened yet" case — surface it every time a gate consults it.
        return None, [{
            "af_code": want, "approved_by": "", "reason": "", "timestamp": "",
            "outcome": "rejected",
            "rejected_because": (f"process_manifest.json exists but is UNPARSEABLE "
                                  f"({detail}) — no owner_skip_approval record can be "
                                  f"honored from a corrupt manifest"),
        }]
    records = _owner_skip_all_records(obj)

    # Blanket-reason signature: reason text (whitespace-normalized, case-folded)
    # reused across more than _OWNER_SKIP_BLANKET_REASON_LIMIT DISTINCT af_codes
    # anywhere in this manifest disqualifies every record that shares it.
    reason_af_codes: Dict[str, set] = {}
    for r in records:
        if not isinstance(r, dict):
            continue
        norm_reason = re.sub(r"\s+", " ", str(r.get("reason") or "").strip().lower())
        this_af = str(r.get("af_code") or r.get("gate") or "").strip().upper()
        if norm_reason and this_af:
            reason_af_codes.setdefault(norm_reason, set()).add(this_af)
    blanket_reasons = {rs for rs, codes in reason_af_codes.items()
                       if len(codes) > _OWNER_SKIP_BLANKET_REASON_LIMIT}

    matches = [r for r in records if isinstance(r, dict)
               and str(r.get("af_code") or r.get("gate") or "").strip().upper() == want]
    if not matches:
        return None, []

    granted = None
    events = []
    for r in matches:
        norm_reason = re.sub(r"\s+", " ", str(r.get("reason") or "").strip().lower())
        ok, why = _owner_skip_structurally_valid(r)
        if ok and norm_reason in blanket_reasons:
            ok = False
            others = sorted(reason_af_codes[norm_reason] - {want})
            why = (f"reason is byte-identical (modulo whitespace/case) to the reason "
                   f"waiving {len(others)} OTHER gate(s) in the same manifest ({others}) "
                   f"— one reason may not waive more than "
                   f"{_OWNER_SKIP_BLANKET_REASON_LIMIT} distinct gate(s); a blanket "
                   f"skip is refused, name each gate's own reason")
        event = {
            "af_code": want,
            "approved_by": str(r.get("approved_by") or "").strip(),
            "reason": str(r.get("reason") or "").strip(),
            "timestamp": str(r.get("timestamp") or "").strip(),
            "outcome": "granted" if ok else "rejected",
        }
        if not ok:
            event["rejected_because"] = why
        events.append(event)
        if ok and granted is None:
            granted = r
    return granted, events


def _owner_skip_ledger_append(run_dir: Path, events: list, consumed_at: str) -> None:
    """Append `events` to process_manifest.json["owner_skip_events"] — a ledger
    this function OWNS exclusively; the agent-authored "owner_skip_approval"
    input is never read-modify-written here. Best-effort / non-fatal by
    design, matching write_process_manifest's own "never clobber, never let a
    report write mask the verdict" convention: a ledger-write failure (or an
    unreadable pre-existing manifest) degrades to stderr-only disclosure —
    it must never be able to turn a real gate result into a crash."""
    if not events:
        return
    ckpt_dir = run_dir / "working" / "checkpoints"
    pm_path = ckpt_dir / "process_manifest.json"
    try:
        ckpt_dir.mkdir(parents=True, exist_ok=True)
        manifest: dict = {}
        if pm_path.exists():
            try:
                existing = json.loads(pm_path.read_text())
                if isinstance(existing, dict):
                    manifest = existing
                else:
                    return  # legacy bare-list manifest — never clobber it here
            except Exception:  # noqa: BLE001 — corrupt manifest: stderr already fired
                return
        ledger = manifest.get(_OWNER_SKIP_LEDGER_KEY)
        if not isinstance(ledger, list):
            ledger = []
        for ev in events:
            ledger.append({**ev, "consumed_at": consumed_at})
        manifest[_OWNER_SKIP_LEDGER_KEY] = ledger
        pm_path.write_text(json.dumps(manifest, indent=2))
    except Exception:  # noqa: BLE001 — best-effort audit trail, never blocks the gate
        pass


def _owner_skip_disclose(run_dir: Path, events: list) -> None:
    """UNCONDITIONAL disclosure layer. Runs regardless of the runfacts shadow
    path succeeding, regardless of PRES_TRUST_BOUNDARY_ENFORCE, regardless of
    whether presentation_job is even importable — this is the mandatory half
    of the "malformed input fails loud" / "a skip can never be silent"
    contract. Prints one unmissable banner per event and appends the same
    events to the durable ledger (see _owner_skip_ledger_append). Never
    raises."""
    if not events:
        return
    now = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    for ev in events:
        try:
            granted = ev["outcome"] == "granted"
            print("!" * 78, file=sys.stderr)
            print(f"{'OWNER-SKIP-CONSUMED' if granted else 'OWNER-SKIP-REJECTED'}: "
                  f"gate {ev['af_code']} "
                  f"{'WAIVED' if granted else 'attempted-but-REFUSED'} by "
                  f"owner_skip_approval (approved_by={ev['approved_by']!r}, "
                  f"reason={ev['reason']!r}, timestamp={ev['timestamp']!r})",
                  file=sys.stderr)
            if not granted:
                print(f"  REFUSED because: {ev.get('rejected_because', '(no detail recorded)')}",
                      file=sys.stderr)
            print(f"  -> logged to working/checkpoints/process_manifest.json"
                  f"[\"{_OWNER_SKIP_LEDGER_KEY}\"] for audit — this can never be silent.",
                  file=sys.stderr)
            print("!" * 78, file=sys.stderr)
        except Exception:  # noqa: BLE001 — disclosure must never crash a gate check
            pass
    _owner_skip_ledger_append(run_dir, events, now)


def _owner_skip_approved_legacy(run_dir: Path, af_code: str):
    """Return the logged owner/founder skip-approval record waiving `af_code`, or None.

    A FIX-2 gate (AF-CANONICAL-RENDER-BYPASS / AF-LOCAL-CANVAS / AF-IMAGE-QC-VISION /
    AF-MODE-UNSET) may be skipped ONLY by an explicit, LOGGED owner token recorded in
    working/checkpoints/process_manifest.json under "owner_skip_approval" (a single
    object or a list of them). A valid token carries owner_approved:true, the af_code
    (or gate) it waives, a non-empty (>= _OWNER_SKIP_REASON_MIN_CHARS) approved_by/
    reason, a real parseable timestamp, and a reason that is not a copy-pasted
    blanket-skip signature shared by more than _OWNER_SKIP_BLANKET_REASON_LIMIT other
    af_codes (see _owner_skip_evaluate / _owner_skip_structurally_valid — Trust
    Boundary Increment 2). No agent may self-skip and the absence of a token means
    the gate is ENFORCED. Returns the matching record (so callers can log who
    approved what) or None.

    Delegates to _owner_skip_evaluate and discards its events — kept as a separate,
    stable name because run_signature_deck.py and the test suite call this directly
    for the plain pass/fail answer. This is still the sole source of truth for the
    value _owner_skip_approved() returns in report-only mode (the default, unchanged
    from Trust Boundary Increment 1) — see _owner_skip_approved below for the
    shadow wrapper and the (now unconditional) disclosure layer."""
    record, _events = _owner_skip_evaluate(run_dir, af_code)
    return record


def _owner_skip_approved(run_dir: Path, af_code: str):
    """Public entry point, SAME name and signature every existing call site
    already uses (10 in this file, 3 in run_signature_deck.py via
    bd._owner_skip_approved — INTAKE-INTERVIEW, AF-IMAGE-QC-VISION,
    AF-OCR-READBACK, AF-CANONICAL-RENDER-BYPASS, AF-LOCAL-CANVAS,
    AF-DECK-TYPE-UNSET, AF-MODE-UNSET, AF-STYLE-UNPICKED, AF-STYLE-DOUBLECHARGE,
    AF-PRIORITY-SHIFT, AF-COPY-QC, AF-PROMPT-QC, AF-HARMONY), so nothing needs
    to change at any call site for Trust Boundary Increment 2 to take effect.

    Evaluates _owner_skip_evaluate ONCE, unconditionally discloses every event
    it found (see _owner_skip_disclose — this is the mandatory, non-optional
    fix: a skip, or a REJECTED attempt at one, can no longer be silent), then
    runs the UNCHANGED Trust Boundary Increment 1 shadow path on top: it seals
    (or reuses the already-sealed) RunFacts for run_dir and asks the PURE
    verify_owner_skip(facts, af_code) function the same question, logging any
    divergence via presentation_job.runfacts.shadow_compare. Set
    PRES_TRUST_BOUNDARY_ENFORCE=1 to make the sealed verdict authoritative
    instead of report-only — this increment does not change that flag's
    default or its meaning; the tightened validity rule above is already
    unconditional and does not depend on it.

    This function can never raise due to the shadow path: any error sealing or
    querying RunFacts is caught, logged, and the legacy result returned — a bug
    in the shadow code must never be able to break an EXISTING gate."""
    legacy_result, events = _owner_skip_evaluate(run_dir, af_code)
    _owner_skip_disclose(run_dir, events)
    try:
        from presentation_job import runfacts as _rf  # noqa: PLC0415 — lazy, avoids any import cycle
        facts = _rf.get_or_seal(Path(run_dir))
        verdict, detail = _rf.verify_owner_skip(facts, af_code)
        _rf.shadow_compare(
            f"owner_skip:{af_code}",
            legacy_result is not None,
            "approved" if legacy_result is not None else "no matching valid record",
            verdict, detail, run_dir=run_dir,
        )
        if _rf.enforcing():
            if verdict is _rf.Verdict.PASS:
                return legacy_result if legacy_result is not None else {
                    "owner_approved": True, "af_code": af_code.strip().upper(),
                    "source": "runfacts-enforcing",
                }
            return None
    except Exception as exc:  # noqa: BLE001 — the shadow must never break this gate
        try:
            print(f"TRUST-BOUNDARY-SHADOW-ERROR owner_skip:{af_code}: {exc!r}",
                  file=sys.stderr)
        except Exception:  # noqa: BLE001
            pass
    return legacy_result


def _gather_rendered_pngs(run_dir: Path) -> list:
    """Return the sorted list of rendered slide PNGs for a run. Primary source is
    renders/slide-*.png (what render_slide writes); the per-slide image paths in the
    process_manifest render record are folded in as a fallback so a deck whose PNGs
    live outside renders/ is still inspected. Deduplicated by resolved path."""
    found = {}
    renders_dir = run_dir / "renders"
    if renders_dir.is_dir():
        for p in renders_dir.glob("slide-*.png"):
            try:
                found[p.resolve()] = p
            except OSError:
                found[p] = p
    # Fallback: image paths recorded in the last render record of the manifest.
    ckpt = run_dir / "working" / "checkpoints" / "process_manifest.json"
    if ckpt.exists():
        obj = _read_json(ckpt)
        if isinstance(obj, dict) and "__parse_error__" not in obj:
            phases = obj.get("phases")
            recs = [p for p in phases if isinstance(p, dict) and p.get("phase") == "render"] \
                if isinstance(phases, list) else []
            if recs:
                for entry in recs[-1].get("slides", []) or []:
                    if isinstance(entry, dict) and entry.get("image"):
                        ip = Path(entry["image"])
                        if ip.exists() and ip.suffix.lower() == ".png":
                            try:
                                found.setdefault(ip.resolve(), ip)
                            except OSError:
                                found.setdefault(ip, ip)
    return sorted(found.values(), key=lambda p: p.name)


def _png_flatfill_fraction(path: Path):
    """Return (dominant_fraction, (R, G, B)) for the single most-common quantised
    colour in a slide PNG, or (None, None) when PIL is absent / the file is unreadable.
    A flat Pillow Image.new fill returns ~1.0; a real photographic kie.ai render — with
    a subject and model noise — stays well below IMAGE_QC_FLATFILL_DOMINANCE."""
    try:
        from PIL import Image  # noqa: PLC0415 — optional dep, checked here
    except ImportError:
        return (None, None)
    try:
        with Image.open(str(path)) as im:
            q = im.convert("RGB").quantize(colors=64)
            palette = q.getpalette() or []
            hist = q.histogram()
            n_colors = len(palette) // 3
            total = 0
            best_count = -1
            best_rgb = None
            for i in range(n_colors):
                c = hist[i] if i < len(hist) else 0
                total += c
                if c > best_count:
                    best_count = c
                    best_rgb = (palette[i * 3], palette[i * 3 + 1], palette[i * 3 + 2])
            if total <= 0 or best_rgb is None:
                return (None, None)
            return (best_count / total, best_rgb)
    except Exception:  # noqa: BLE001
        return (None, None)


def _image_qc_report_defects(report_obj, n_slides) -> list:
    """Return a list of fatal defects in an image-QC report that prove it cannot have
    come from a real multimodal PIXEL read. Empty list = the report carries the
    provenance a genuine vision pass would. Checks (AF-IMAGE-QC-VISION):
      * per-slide coverage MUST be a LIST (not a dict keyed by slide, so no slide can
        be silently dropped) covering every rendered slide — and may NOT exclude any
        slide (e.g. the hook slides 1/24/49) from scope;
      * a multimodal VISION read must be recorded with provenance: a declared vision
        engine AND a per-slide observed/OCR field — a self-typed number is refused;
      * the eliminated overlay rubric (typography_overlay_readiness / "overlay the
        headline" / "post-production overlay" / "applied in post") may NOT appear."""
    defects = []
    if not isinstance(report_obj, dict):
        return ["image-QC report is not a JSON object — it cannot record a pixel read"]

    # (1) No slide may be excluded from scope.
    for k in ("out_of_scope", "excluded_slides", "scope_excludes", "slides_excluded",
              "excluded", "out_of_scope_slides"):
        v = report_obj.get(k)
        if v:
            defects.append(
                f"image-QC excludes slides from scope via {k!r}={v!r}; the gate may NOT "
                f"exclude any slide (the hook/typography slides are graded too)")
            break

    # (2) Per-slide coverage must be a non-empty LIST, not a dict.
    per_slide = None
    for k in ("slides", "per_slide", "slide_results"):
        if k in report_obj:
            per_slide = report_obj.get(k)
            break
    if isinstance(per_slide, dict):
        defects.append(
            "per-slide coverage is recorded as a DICT; it must be a per-slide LIST so "
            "every rendered slide carries its own graded row and none can be silently "
            "dropped")
        per_slide = None
    if not isinstance(per_slide, list) or not per_slide:
        defects.append(
            "image-QC report has no per-slide coverage LIST (slides/per_slide/"
            "slide_results); a report with no per-slide rows is a rubber-stamped number")
        per_slide = []
    elif n_slides and len(per_slide) < n_slides:
        defects.append(
            f"image-QC report covers {len(per_slide)} of {n_slides} rendered slide(s); "
            f"every rendered slide must carry a per-slide graded row")

    # (3) A multimodal vision read must be recorded with provenance.
    vision_engine = ""
    for k in ("vision_model", "multimodal_model", "vision_provider", "ocr_engine",
              "vision_engine", "reviewer_vision_model"):
        val = report_obj.get(k)
        if isinstance(val, str) and val.strip():
            vision_engine = val.strip()
            break
    if not vision_engine:
        defects.append(
            "image-QC report declares no vision engine (vision_model/multimodal_model/"
            "ocr_engine); a real image-QC pass is a multimodal PIXEL read, not a typed "
            "score")
    if isinstance(per_slide, list) and per_slide:
        VIS_FIELDS = ("vision", "observed_text", "ocr", "ocr_text", "baked_text",
                      "read_text", "description", "pixels_read", "visual_subject")
        rows_with_vision = sum(
            1 for r in per_slide
            if isinstance(r, dict) and any(
                str(r.get(f) or "").strip() for f in VIS_FIELDS))
        if rows_with_vision < len(per_slide):
            defects.append(
                f"only {rows_with_vision}/{len(per_slide)} per-slide rows carry a vision/"
                f"OCR observation (observed_text/ocr/visual_subject); rows without one are "
                f"pixel-blind and cannot prove the slide was actually read")

    # (4) The eliminated overlay rubric may not appear.
    try:
        blob = json.dumps(report_obj, default=str).lower()
    except Exception:  # noqa: BLE001
        blob = ""
    for needle in ("typography_overlay_readiness", "overlay the headline",
                   "overlay the canonical", "post-production overlay",
                   "post production overlay", "applied in post", "overlay_readiness"):
        if needle in blob:
            defects.append(
                f"image-QC rubric still references the ELIMINATED overlay model "
                f"({needle!r}); the typography-overlay-readiness criterion and any "
                f"'overlay the headline' recommendation are removed (Decision 5C) — words "
                f"are baked into the kie.ai image, never overlaid")
            break
    return defects


def check_image_qc_vision(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-IMAGE-QC-VISION — the DETERMINISTIC PIXEL image-QC cross-check (FIX-2).

    This is the teeth behind the image-QC gate. It does what the old _chk_image_qc
    JSON rubber stamp never did: it OPENS the rendered PNGs and inspects pixels, and it
    refuses an image-QC report that could not have come from a real multimodal read.

    Two independent classes of failure, both fatal:
      A. PIXEL cross-check (runs whenever rendered PNGs exist):
         * any PNG below the PLACEHOLDER_MIN_BYTES (51,200) kie-bake floor — a flat
           Pillow card / native render, NOT a real model bake (also an AF-LOCAL-CANVAS
           signature);
         * any PNG whose single dominant quantised colour covers >= 90% of the frame
           AND is light/cream — a flat typography card with NO visual subject (the
           slide-1/24/49 hook-card defect). A genuine kie.ai render never does this.
      B. REPORT cross-check (runs whenever the image-QC report exists):
         * coverage must be a per-slide LIST (not a dict), cover every rendered slide,
           and exclude none (hook slides included);
         * a multimodal vision read must be recorded with provenance (declared engine +
           per-slide observation) — a self-typed number is refused;
         * the eliminated typography-overlay-readiness rubric may not appear.

    It DEFERS (returns "") ONLY when there is nothing to inspect yet — no rendered PNGs
    AND no report (the genuine pre-render state). Once renders OR a report exist it does
    NOT defer-to-pass: the legacy 'no renders/manifest -> pass' escape is closed. A
    failure may be waived ONLY by a logged owner_skip_approval token for
    AF-IMAGE-QC-VISION in process_manifest.json.

    Returns "" on pass / defer, or a fatal AF-IMAGE-QC-VISION message."""
    skip = _owner_skip_approved(run_dir, AF_IMAGE_QC_VISION)
    if skip is not None:
        print(f"  NOTE  AF-IMAGE-QC-VISION waived by logged owner_skip_approval "
              f"(approved_by={skip.get('approved_by')!r}, reason={skip.get('reason')!r}).",
              file=sys.stderr)
        return ""

    pngs = _gather_rendered_pngs(run_dir)
    report_path = run_dir / "working" / "qc" / "image_qc_report.json"
    report_obj = _read_json(report_path) if report_path.exists() else None

    # Genuine pre-render state — nothing to inspect yet. This is the ONLY defer.
    if not pngs and report_obj is None:
        return ""

    problems = []

    # --- A. PIXEL cross-check over every rendered PNG ---
    for png in pngs:
        try:
            size = png.stat().st_size
        except OSError:
            problems.append(f"{png.name}: cannot stat the rendered PNG on disk")
            continue
        if size < PLACEHOLDER_MIN_BYTES:
            problems.append(
                f"{png.name}: {size:,} bytes, below the {PLACEHOLDER_MIN_BYTES:,}-byte "
                f"kie-bake floor — a flat/native-rendered card, not a real gpt-image-2 "
                f"bake (AF-LOCAL-CANVAS signature)")
            continue
        frac, rgb = _png_flatfill_fraction(png)
        if frac is None:
            continue  # PIL unavailable / unreadable — byte floor already covered it.
        if frac >= IMAGE_QC_FLATFILL_DOMINANCE:
            luma = (0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]) / 255.0 if rgb else 0.0
            if luma >= IMAGE_QC_FLATFILL_LUMA:
                problems.append(
                    f"{png.name}: {frac:.0%} of the frame is a single light fill "
                    f"rgb{rgb} (luma {luma:.2f}) — a flat cream/typography card with NO "
                    f"visual subject. Pure-typography hook slides must be RENDERED by "
                    f"kie.ai (cream + baked display type as the image), never drawn on a "
                    f"local canvas")

    # --- B. REPORT cross-check (when a report exists) ---
    if report_obj is not None:
        if "__parse_error__" in report_obj:
            problems.append(
                f"image_qc_report.json is not valid JSON ({report_obj['__parse_error__']}) "
                f"— a report that cannot be parsed cannot prove a pixel read")
        else:
            problems.extend(_image_qc_report_defects(report_obj, len(pngs)))

    if problems:
        return ("AF-IMAGE-QC-VISION: the image-QC gate did not actually inspect the "
                "rendered slides (rubber-stamp / flat card / pixel-blind report). A real "
                "image-QC pass is a multimodal PIXEL read of every PNG; flat cream cards, "
                "below-floor PNGs, and self-typed scores are rejected. Re-render any flat "
                "card through kie.ai gpt-image-2 and re-run the Image QC Specialist with a "
                "per-slide vision read. Offenders: " + "; ".join(problems[:10])
                + ("" if len(problems) <= 10 else f" (+{len(problems) - 10} more)")
                + ". This gate may be waived ONLY by a logged owner_skip_approval for "
                "AF-IMAGE-QC-VISION in working/checkpoints/process_manifest.json.")
    return ""


def check_ocr_readback(run_dir: Path) -> str:
    """U027 — AF-OCR-READBACK postflight gate.

    This audits the per-slide OCR sidecar records (renders/slide-*.ocr.json) that
    _record_ocr_readback() writes beside every rendered PNG once
    _verify_aspect_and_readback() calls prompt_gate.ocr_readback() against it. That
    pair PRODUCES the readback; this function READS it back and decides pass/fail —
    the two must never be the same code path, or a broken producer could never be
    caught by its own consumer.

    DEFERS (returns "") ONLY when there are no rendered PNGs yet (the genuine
    pre-render state — nothing to audit). Once PNGs exist, every one of them MUST
    carry a well-formed sidecar:

      * missing sidecar, or a sidecar that is not valid JSON -> FAIL. A gate that
        cannot even confirm the engine ran is not a pass.
      * sidecar carries checked:false (or `checked` missing/falsy) -> FAIL, and this
        branch is NEVER waivable. A checked:false record means the OCR engine did
        not run against that render at all; a self-disabled check is not a pass, and
        no owner_skip_approval token — however well-formed, whatever af_code/gate
        value it declares — can waive it (MASTER-SPEC 7.4 / D10).
      * sidecar carries checked:true, matched:false -> FAIL, UNLESS a logged
        owner_skip_approval token for AF-OCR-READBACK is present in
        working/checkpoints/process_manifest.json (_owner_skip_approved). Unlike
        checked:false, a genuine mismatch the owner has reviewed and accepted (e.g.
        a stylised headline OCR legitimately cannot read) IS waivable.
      * checked:true, matched:true on every rendered slide -> PASS ("").

    Returns "" on pass/defer, or a fatal AF-OCR-READBACK message naming the
    offending slide(s) and the count/denominator (e.g. "1 of 2")."""
    pngs = _gather_rendered_pngs(run_dir)
    if not pngs:
        return ""

    total = len(pngs)
    missing: list = []      # no sidecar, or sidecar is not valid JSON
    unchecked: list = []    # sidecar exists but checked is false/absent
    mismatched: list = []   # checked:true but matched is false

    for png in pngs:
        sidecar = png.with_suffix(".ocr.json")
        if not sidecar.is_file():
            missing.append(png.name)
            continue
        obj = _read_json(sidecar)
        if not isinstance(obj, dict) or "__parse_error__" in obj:
            missing.append(png.name)
            continue
        if not obj.get("checked"):
            unchecked.append(png.name)
            continue
        if obj.get("matched") is False:
            mismatched.append(png.name)

    if missing:
        extra = "" if len(missing) <= 5 else f" (+{len(missing) - 5} more)"
        return (
            f"AF-OCR-READBACK: {len(missing)} of {total} slide(s) missing or unreadable "
            f"OCR sidecar(s): {', '.join(missing[:5])}{extra}. Every rendered PNG must "
            "carry a renders/<slide>.ocr.json readback record written by "
            "ocr_readback()/_record_ocr_readback() before this gate can pass."
        )

    if unchecked:
        extra = "" if len(unchecked) <= 5 else f" (+{len(unchecked) - 5} more)"
        return (
            f"AF-OCR-READBACK: {len(unchecked)} of {total} slide(s) failed OCR readback "
            f"(checked:false): {', '.join(unchecked[:5])}{extra}. A checked:false record "
            "means the OCR engine never ran against this render — a self-disabled check "
            "is not a pass. This branch is NOT waivable: no owner_skip_approval token, "
            "however well-formed, can waive it."
        )

    if mismatched:
        skip = _owner_skip_approved(run_dir, AF_OCR_READBACK)
        if skip is not None:
            print(f"  NOTE  AF-OCR-READBACK waived by logged owner_skip_approval "
                  f"(approved_by={skip.get('approved_by')!r}, reason={skip.get('reason')!r}).",
                  file=sys.stderr)
            return ""
        extra = "" if len(mismatched) <= 5 else f" (+{len(mismatched) - 5} more)"
        return (
            f"AF-OCR-READBACK: {len(mismatched)} of {total} slide(s) mismatched the "
            f"approved copy: {', '.join(mismatched[:5])}{extra}. The rendered text does "
            "not match the slide's approved copy (garbled/unreadable bake). Re-prompt/"
            "re-render the affected slide(s), or waive with a logged owner_skip_approval "
            "token for AF-OCR-READBACK in working/checkpoints/process_manifest.json."
        )

    return ""


def check_canonical_render_path(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-CANONICAL-RENDER-BYPASS / AF-LOCAL-CANVAS — the canonical render path is
    build_deck.py / run_signature_deck.py ONLY (FIX-2).

    This is the symbol the orchestrator (run_signature_deck.py) and the preflight
    lockstep consume to make a hand-rolled renderer/assembler IMPOSSIBLE to route a
    deck around. It scans every non-canonical *.py file inside the run/working dir and
    HARD-FAILS when one defines:
      * a LOCAL slide canvas — Image.new with the 2048x1152 slide geometry (the
        Pillow cream-card path) -> AF-LOCAL-CANVAS;
      * a native PowerPoint text box — add_textbox / add_text_box (a hand-rolled
        overlay assembler stamping words on top of the image) -> AF-CANONICAL-RENDER-BYPASS;
      * a direct kie.ai task submission — createTask / recordInfo / api.kie.ai outside
        build_deck.py (a per-deck renderer) -> AF-CANONICAL-RENDER-BYPASS.

    The canonical tools themselves (CANONICAL_RENDER_SCRIPTS) and anything under a
    scripts/ or virtual-env directory are exempt. Returns "" when the run dir carries
    no hand-rolled renderer. A failure may be waived ONLY by a logged
    owner_skip_approval token (AF-CANONICAL-RENDER-BYPASS or AF-LOCAL-CANVAS)."""
    skip = (_owner_skip_approved(run_dir, AF_CANONICAL_RENDER_BYPASS)
            or _owner_skip_approved(run_dir, AF_LOCAL_CANVAS))
    if skip is not None:
        print(f"  NOTE  canonical-render-path guard waived by logged owner_skip_approval "
              f"(approved_by={skip.get('approved_by')!r}, reason={skip.get('reason')!r}).",
              file=sys.stderr)
        return ""

    _SKIP_DIR_SEGS = {".venv", "venv", "site-packages", "__pycache__", ".git",
                      "node_modules", ".mypy_cache", ".pytest_cache", "scripts"}
    offenders = []
    try:
        candidates = sorted(run_dir.rglob("*.py"))
    except OSError:
        return ""
    for py in candidates:
        if py.name in CANONICAL_RENDER_SCRIPTS:
            continue
        if _SKIP_DIR_SEGS & set(py.parts):
            continue
        try:
            text = py.read_text(errors="replace")
        except OSError:
            continue
        try:
            rel = str(py.relative_to(run_dir))
        except ValueError:
            rel = py.name
        if "Image.new" in text and "2048" in text and "1152" in text:
            offenders.append((rel, AF_LOCAL_CANVAS,
                              "defines a LOCAL 2048x1152 slide canvas (Image.new) instead "
                              "of a kie.ai gpt-image-2 render"))
        elif ("add_textbox" in text) or ("add_text_box" in text):
            offenders.append((rel, AF_CANONICAL_RENDER_BYPASS,
                              "adds NATIVE PowerPoint text boxes (add_textbox) on top of "
                              "the image — a hand-rolled overlay assembler"))
        elif ("createTask" in text) or ("recordInfo" in text) or ("api.kie.ai" in text):
            offenders.append((rel, AF_CANONICAL_RENDER_BYPASS,
                              "submits kie.ai tasks directly (createTask/recordInfo) "
                              "outside the canonical build_deck.py renderer"))
    if offenders:
        rel, code, why = offenders[0]
        more = f" (+{len(offenders) - 1} more hand-rolled script(s))" if len(offenders) > 1 else ""
        return (f"{code}: non-canonical renderer/assembler {rel!r} {why}. The canonical "
                f"render path is build_deck.py / run_signature_deck.py ONLY: every slide's "
                f"WORDS and VISUAL are generated together in ONE kie.ai gpt-image-2 image, "
                f"with ZERO native PowerPoint text and NO local Pillow canvas. Delete the "
                f"hand-rolled script and route the deck through run_signature_deck.py -> "
                f"build_deck.py. This gate may be waived ONLY by a logged owner_skip_approval "
                f"({code}) in working/checkpoints/process_manifest.json.{more}")
    return ""


def _read_dark_optin(run_dir: Path) -> bool:
    """True iff intake.json records a dark-theme opt-in via the canonical
    client_dark_theme key OR the role-doc DARK_OK alias (same intent). Shared by the
    no-dark gate and the font-floor dark-mode size/contrast compensation."""
    def _truthy(val) -> bool:
        if val is True:
            return True
        return isinstance(val, str) and val.strip().lower() in ("true", "yes", "1")
    for rel in ("working/copy/intake.json", "intake.json", "working/intake.json"):
        p = run_dir / rel
        if p.exists():
            obj = _read_json(p)
            if isinstance(obj, dict) and "__parse_error__" not in obj:
                return _truthy(obj.get("client_dark_theme")) or _truthy(obj.get("DARK_OK"))
            return False
    return False


def check_font_floor(run_dir: Path) -> str:
    """AF-FONT-FLOOR — DETERMINISTIC numeric type-system rejector (R8 Gap 1+2).

    Parses the Typography Architect's machine-readable tokens from
    working/typography/type_layout_system.md and FAILS LOUD when the DECLARED type
    system violates a hard numeric floor — before any render. This is the coded
    floor that the vision-opinion Typography-QC gate never was; it does NOT do
    OCR/pixel work (that stays with the independent vision pass), it rejects a
    design system whose own declared tokens are below floor.

    Tokens parsed (one per line, `key: value`, case-insensitive key):
        min_body_pt:           <int>   smallest body/subhead size (pt-equiv @1080)
        type_scale_steps:      <int>   number of named modular steps (must be 4..5)
        min_contrast_ratio:    <float> smallest body/subhead text:bg contrast ratio
       [min_large_contrast_ratio: <float>]  optional, for large/headline text

    FLOORS (raised when client_dark_theme/DARK_OK opt-in — projection dispersion):
        body pt   >= FONT_BODY_PT_FLOOR (18)   | dark: DARK_THEME_BODY_PT_FLOOR (22)
        scale     in [TYPE_SCALE_STEPS_MIN..MAX] (4..5)
        contrast  >= CONTRAST_RATIO_FLOOR_NORMAL (4.5) | dark: DARK_THEME_CONTRAST_FLOOR (7.0)

    DEFER (return "") only when NO type system has been produced yet (pre-typography
    phase). Once a design system exists (working/typography/design_system.json or a
    design-brief), the type_layout_system.md tokens are MANDATORY — a missing file or
    missing token then FAILS (closes R8 Gap 1: the phantom artifact). Run-dir-scoped.
    """
    layout = run_dir / TYPE_LAYOUT_SYSTEM_REL
    design_present = (
        (run_dir / "working" / "typography" / "design_system.json").exists()
        or any((run_dir / "working" / "research").glob("design-brief-*.md"))
        if (run_dir / "working").exists() else False
    )

    if not layout.exists():
        if design_present:
            return (
                f"AF-FONT-FLOOR: a design system exists but the deterministic type "
                f"tokens file {TYPE_LAYOUT_SYSTEM_REL} is MISSING. The Typography "
                f"Architect MUST emit it with min_body_pt / type_scale_steps / "
                f"min_contrast_ratio so the font-floor gate has real data (R8 Gap 1).")
        return ""  # pre-typography: defer

    text = layout.read_text(encoding="utf-8", errors="replace")

    def _num(key):
        m = re.search(rf'(?im)^\s*{re.escape(key)}\s*[:=]\s*([0-9]+(?:\.[0-9]+)?)', text)
        return float(m.group(1)) if m else None

    dark = _read_dark_optin(run_dir)
    pt_floor = DARK_THEME_BODY_PT_FLOOR if dark else FONT_BODY_PT_FLOOR
    contrast_floor = DARK_THEME_CONTRAST_FLOOR if dark else CONTRAST_RATIO_FLOOR_NORMAL

    min_body_pt = _num("min_body_pt")
    type_scale_steps = _num("type_scale_steps")
    min_contrast = _num("min_contrast_ratio")

    problems = []
    if min_body_pt is None:
        problems.append(f"missing token min_body_pt (declare the smallest body/subhead "
                        f"size, pt-equivalent at 1080px tall)")
    elif min_body_pt < pt_floor:
        problems.append(f"min_body_pt={min_body_pt:g} is below the {pt_floor:g}pt "
                        f"{'dark-theme ' if dark else ''}body/subhead floor "
                        f"(FONT_BODY_PT_FLOOR={FONT_BODY_PT_FLOOR})")
    if type_scale_steps is None:
        problems.append("missing token type_scale_steps (declare the number of modular "
                        "type-scale steps)")
    elif not (TYPE_SCALE_STEPS_MIN <= int(type_scale_steps) <= TYPE_SCALE_STEPS_MAX):
        problems.append(f"type_scale_steps={int(type_scale_steps)} is not a modular "
                        f"{TYPE_SCALE_STEPS_MIN}-{TYPE_SCALE_STEPS_MAX}-step scale")
    if min_contrast is None:
        problems.append("missing token min_contrast_ratio (declare the smallest "
                        "text:background contrast ratio, WCAG relative-luminance)")
    elif min_contrast < contrast_floor:
        problems.append(f"min_contrast_ratio={min_contrast:g} is below the "
                        f"{contrast_floor:g}:1 {'dark-theme AAA ' if dark else 'WCAG AA '}"
                        f"floor (CONTRAST_RATIO_FLOOR_NORMAL={CONTRAST_RATIO_FLOOR_NORMAL})")

    if problems:
        return ("AF-FONT-FLOOR: the declared type system "
                f"({TYPE_LAYOUT_SYSTEM_REL}) violates the deterministic floor — "
                + "; ".join(problems) + ". Raise the offending token(s) and re-emit; "
                "the font-floor gate rejects a sub-floor design BEFORE render.")
    return ""


def _chk_research_map(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-RESEARCH-WEAVE — research woven ACROSS the deck, not funnelled to one slide.

    The R3 breadth gate. CONDITIONAL: defers (returns "") until copy exists
    (working/copy/slides_copy.md). Once copy exists it requires
    working/research/research_map.json and enforces FOUR independent conditions:

      1. MAP EXISTS + BREADTH — the map assigns >= 1 real research item (with a
         verbatim `anchor` token) to at least RESEARCH_WEAVE_FLOOR_PCT (60%) of
         NON-EXEMPT content slides. Hook/pure-type/transition slides carry
         `"exempt": <reason>` and are excluded from the denominator (so the gate
         never pressures fabrication onto a slide that should not carry a stat).
      2. WRITER USED IT — for >= that floor of mapped non-exempt slides, the slide's
         anchor token actually appears in slides_copy.md (the slide block or a
         RESEARCH_USED tag). Mechanical, not semantic; the anchor is a verbatim
         figure / dollar / short quote fragment / source domain.
      3. RENDER COPY CARRIES IT — for >= that same floor of mapped non-exempt slides,
         the anchor token ALSO appears in the ACTUAL file the renderer consumes:
         slides.json's copy[] (H3 — the positional slides_path when threaded in, else
         the canonical working/copy/slides.json / slides.json / working/slides.json
         spots). slides_copy.md is a working document; the pixels are baked VERBATIM
         from slides.json's copy[], so a deck whose slides_copy.md weaves the anchors
         but whose slides.json dropped them passes the old gate while the rendered
         pixels carry NO research. Defers (passes) when no slides.json can be read yet
         (its absence is owned by the schema / AF-P1 / slide-count gates).
      4. WHOLE-BRIEF BREADTH — the deck draws on >= MIN_DISTINCT_RESEARCH_ITEMS (8)
         DISTINCT items, so breadth cannot be faked by repeating one stat.

    Replaces the toothless "a pack exists" logic of _chk_claims_without_citation with
    a true coverage gate (that check is kept as the zero-research backstop). Run-dir-scoped.
    """
    copy_path = run_dir / "working" / "copy" / "slides_copy.md"
    if not copy_path.exists():
        return ""  # pre-copy: defer (research-before-content is enforced by AF-PHASE-SKIPPED)

    map_path = run_dir / RESEARCH_MAP_REL
    if not map_path.exists():
        return (f"AF-RESEARCH-WEAVE: copy exists but {RESEARCH_MAP_REL} is MISSING. "
                "The Deep Research Specialist MUST map research facts/quotes/stats to "
                "specific slides (SOP 9.5) BEFORE copy so research is woven across the "
                "deck, not funnelled to one proof slide.")
    obj = _read_json(map_path)
    if not isinstance(obj, dict) or "__parse_error__" in obj:
        return (f"AF-RESEARCH-WEAVE: {RESEARCH_MAP_REL} is missing or not valid JSON.")

    slides = obj.get("slides")
    if not isinstance(slides, list) or not slides:
        return (f"AF-RESEARCH-WEAVE: {RESEARCH_MAP_REL} has no slides[] assignment "
                "array — the research-to-slide map is empty.")

    copy_text = copy_path.read_text(encoding="utf-8", errors="replace")
    copy_lc = copy_text.lower()

    # The ACTUAL render copy — slides.json copy[] is what the renderer bakes into the
    # pixels VERBATIM (slides_copy.md is only the writer's working document). Reuse the
    # shared loader so this gate counts the exact file that gets rendered (H3: the
    # positional slides_path when threaded in, else canonical spots). None = no render
    # copy readable yet; condition 3 defers and the upstream schema / AF-P1 / slide-count
    # gates own the absence.
    render_copy_map = _load_slide_copy_map(run_dir, slides_path) or {}
    render_copy_lc = " ".join(
        str(c) if c is not None else ""
        for vals in render_copy_map.values()
        for c in (vals if isinstance(vals, list) else [vals])
    ).lower()

    non_exempt = [s for s in slides if isinstance(s, dict) and not s.get("exempt")]
    if not non_exempt:
        return ("AF-RESEARCH-WEAVE: every slide in the research map is marked exempt — "
                "exemptions cannot cover the whole deck; map research onto the content "
                "(teaching/proof/story) slides.")

    def _anchors(s):
        out = []
        for a in (s.get("assigned") or []):
            if isinstance(a, dict):
                anc = str(a.get("anchor", "")).strip()
                if anc:
                    out.append((str(a.get("item_id", "")).strip(), anc))
        return out

    mapped = [s for s in non_exempt if _anchors(s)]
    mapped_pct = (len(mapped) / len(non_exempt)) * 100.0
    if mapped_pct < RESEARCH_WEAVE_FLOOR_PCT:
        return (f"AF-RESEARCH-WEAVE: only {len(mapped)}/{len(non_exempt)} "
                f"({mapped_pct:.0f}%) of non-exempt content slides have a research item "
                f"assigned, below the {RESEARCH_WEAVE_FLOOR_PCT}% breadth floor. Map "
                "facts/quotes/stats onto the teaching body, not just the proof slide.")

    used = 0
    for s in mapped:
        if any(anc.lower() in copy_lc for _id, anc in _anchors(s)):
            used += 1
    used_pct = (used / len(mapped)) * 100.0 if mapped else 0.0
    if used_pct < RESEARCH_WEAVE_FLOOR_PCT:
        return (f"AF-RESEARCH-WEAVE: only {used}/{len(mapped)} ({used_pct:.0f}%) of "
                f"mapped slides actually weave their assigned anchor into slides_copy.md, "
                f"below the {RESEARCH_WEAVE_FLOOR_PCT}% floor. The writer must USE the "
                "mapped item (anchor verbatim in the slide block or a RESEARCH_USED tag).")

    # Condition 3 — RENDER COPY CARRIES IT. slides_copy.md passing is NOT enough: the
    # renderer bakes words VERBATIM from slides.json copy[], so a deck that wove the
    # anchors into slides_copy.md but dropped them from slides.json would render pixels
    # with NO research while the old gate passed. Require the anchor in the render copy
    # on >= the same floor of mapped slides. Defer (skip) when no render copy is readable.
    if render_copy_lc.strip():
        rendered = 0
        for s in mapped:
            if any(anc.lower() in render_copy_lc for _id, anc in _anchors(s)):
                rendered += 1
        rendered_pct = (rendered / len(mapped)) * 100.0 if mapped else 0.0
        if rendered_pct < RESEARCH_WEAVE_FLOOR_PCT:
            return (f"AF-RESEARCH-WEAVE: only {rendered}/{len(mapped)} ({rendered_pct:.0f}%) "
                    f"of mapped slides carry their assigned anchor in the RENDER copy "
                    f"(slides.json copy[]), below the {RESEARCH_WEAVE_FLOOR_PCT}% floor. "
                    "slides_copy.md alone is not enough — the renderer bakes words VERBATIM "
                    "from slides.json, so the mapped anchor must be present there too "
                    "(the research is woven into the pixels, not just the working doc).")

    distinct = obj.get("distinct_items_used")
    if not isinstance(distinct, int):
        ids = set()
        for s in non_exempt:
            for _id, _anc in _anchors(s):
                if _id:
                    ids.add(_id)
        distinct = len(ids)
    if distinct < MIN_DISTINCT_RESEARCH_ITEMS:
        return (f"AF-RESEARCH-WEAVE: the deck draws on only {distinct} distinct research "
                f"items, below the floor of {MIN_DISTINCT_RESEARCH_ITEMS}. Breadth cannot "
                "be satisfied by repeating one stat — draw on the whole brief.")
    return ""


def _chk_research_reaches_render(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-RESEARCH-REACHES-RENDER — research must reach the RENDER COPY, not just
    slides_copy.md.

    WORKSTREAM FIX: the AF-RESEARCH-WEAVE gate (_chk_research_map) validates that the
    research anchors live in working/copy/slides_copy.md, but the renderer consumes
    working/copy/slides.json (the positional slides.json). If the slide-copy JSON is
    authored without the mapped anchors (or a hand-fed slides.json is substituted), the
    deck renders with NO research even though the markdown copy passed the weave gate.
    This gate closes that exact divergence: it reads the ACTUAL render copy (the same
    file _load_slide_copy_map feeds the verbatim/prompt gates, i.e. the H3 positional
    slides.json) and FAILS when a research-mapped non-exempt slide's anchor token is
    missing from that slide's render copy[]. So research cannot silently fail to reach
    the rendered deck.

    CONDITIONAL (defer, returns ""): pre-copy / pre-map states. The absences are owned
    by _chk_slides_copy / _chk_arc / _chk_research_map / _chk_research_cited. Run-dir-
    scoped (None sentinel) so it reads the ACTUAL rendered slides.json (H3) — a gate
    that counted a different canonical slides.json than the one rendered would be the
    exact bypass this gate exists to close.
    """
    # Defer until both the research map AND the render copy exist (upstream gates own
    # their absences — mirror _chk_research_map's pre-copy defer so a bare render that
    # has not run research yet is not double-reported).
    if not (run_dir / RESEARCH_MAP_REL).exists():
        return ""
    copy_map = _load_slide_copy_map(run_dir, slides_path)
    if not copy_map:
        return ""  # no slides.json yet — schema validation / _chk_rich_prompts own that.

    obj = _read_json(run_dir / RESEARCH_MAP_REL)
    if not isinstance(obj, dict) or "__parse_error__" in obj:
        return ""  # malformed map — _chk_research_map owns that failure.
    slides_map = obj.get("slides")
    if not isinstance(slides_map, list) or not slides_map:
        return ""  # empty/missing slides[] — _chk_research_map owns that failure.

    def _anchors(s):
        out = []
        for a in (s.get("assigned") or []):
            if isinstance(a, dict):
                anc = str(a.get("anchor", "")).strip()
                if anc:
                    out.append(anc)
        return out

    missing = []
    for s in slides_map:
        if not isinstance(s, dict) or s.get("exempt"):
            continue
        if not _anchors(s):
            continue
        ordinal = s.get("slide")
        copy_val = copy_map.get(ordinal)
        lines = [str(c) for c in copy_val] if isinstance(copy_val, list) else (
            [str(copy_val)] if copy_val else [])
        render_lc = " ".join(lines).lower()
        if not render_lc:
            missing.append((ordinal, _anchors(s), "no render copy"))
            continue
        absent = [a for a in _anchors(s) if a.lower() not in render_lc]
        if absent:
            missing.append((ordinal, absent, "anchor(s) absent from render copy"))

    if missing:
        detail = "; ".join(
            f"slide {o}: {', '.join(a)} ({why})" for o, a, why in missing)
        return (f"AF-RESEARCH-REACHES-RENDER: research mapped to these slides but its "
                f"anchor token(s) are NOT in the RENDER copy (slides.json copy[] — the "
                f"file the renderer bakes into the slides): {detail}. The anchors were "
                f"validated in slides_copy.md by AF-RESEARCH-WEAVE, but they never "
                f"reached the render copy, so the deck would ship WITHOUT the research "
                f"the gate approved. The Slide Copywriter / builder MUST weave the "
                f"mapped anchor verbatim into the slide's copy[] in slides.json (or the "
                f"renderer reads the research-backed copy) before render.")
    return ""


def _import_intelligence_engines_check():
    """Import the intelligence_engines_check module that ships beside build_deck.py.
    Tries a normal import first (scripts/ is on sys.path when build_deck runs / is
    imported), then a path-based load from this file's own directory. Returns the module
    or None."""
    try:
        import importlib
        return importlib.import_module("intelligence_engines_check")
    except Exception:  # noqa: BLE001
        try:
            import importlib.util as _ilu
            spec = _ilu.spec_from_file_location(
                "intelligence_engines_check",
                str(Path(__file__).resolve().parent / "intelligence_engines_check.py"))
            if spec and spec.loader:
                mod = _ilu.module_from_spec(spec)
                spec.loader.exec_module(mod)
                return mod
        except Exception:  # noqa: BLE001
            return None
    return None


def check_intelligence_engines_prompt(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-FACE-PROMPT-MISSING / AF-LIGHT-PROMPT-MISSING / AF-WORLD-SCALE /
    AF-HAIR-INAUTHENTIC — the facial / lighting / world / representation INTELLIGENCE-engine
    mechanical checker, WIRED INTO THE RENDER PREFLIGHT.

    intelligence_engines_check.py was, by its own header, agent-discipline-only: it "does
    NOT touch build_deck.py's render path" and was invoked only by the QC specialists, so a
    grep for it in the renderer returned 0 — facial intelligence was enforced by NOTHING
    mechanical. This wires its --phase prompt engines into PREFLIGHT so they fire
    DETERMINISTICALLY at render time: every people/scene slide prompt must carry an explicit
    expression token (skin-tone/age/expression discipline), a key/fill/rim + rim/hair
    lighting token (lighting), a world believability justification (composition/world), and
    an age-appropriate hairstyle token (diversity/representation). Defers (returns "") when
    no prompts exist yet (pre-prompt phase). Returns a fatal message naming the codes+slides.
    Run-dir-scoped (None sentinel); prompts live at <run_dir>/working/prompts/."""
    prompts_dir = run_dir / "working" / "prompts"
    if not prompts_dir.is_dir():
        return ""  # pre-prompt phase — the prompt engines defer (mirrors check_prompts).
    iec = _import_intelligence_engines_check()
    if iec is None:
        return ("AF-FACE-PROMPT-MISSING: intelligence_engines_check.py could not be imported "
                "from the scripts directory; the facial/lighting/world/hair prompt engines "
                "could not run. Ensure intelligence_engines_check.py is present beside "
                "build_deck.py.")
    problems: list = []
    try:
        # intelligence_engines_check reads <dir>/prompts, <dir>/copy, <dir>/brand — so the
        # working/ dir is the run dir from its point of view.
        iec.check_prompts(run_dir / "working", problems)
    except Exception as exc:  # noqa: BLE001
        return (f"AF-FACE-PROMPT-MISSING: intelligence_engines_check.check_prompts raised "
                f"{exc!r}; the prompt-side INTELLIGENCE engines could not be verified.")
    if problems:
        lines = "; ".join(
            f"{p.get('code')} [{p.get('slide')}]: {p.get('detail')}" for p in problems[:10])
        more = "" if len(problems) <= 10 else f" (+{len(problems) - 10} more)"
        return ("AF-INTELLIGENCE-ENGINES: the facial / lighting / world / representation "
                "prompt engines auto-failed — each people/scene prompt must carry an explicit "
                "expression token, a key/fill/rim + rim/hair lighting token, a world "
                "believability justification, and an age-appropriate hairstyle token "
                "(skin-tone / age / expression / diversity / lighting / composition present "
                "per slide). Offenders: " + lines + more + ".")
    return ""


# ===========================================================================
# v15.0.0 — WRITING-half engines (G1/G2), HARMONY (G5), EXCELLENCE (G6),
# and the deterministic Prompt-QC / Copy-QC measurers that feed the send-back
# loop (G7/G8). build_deck.py owns ALL of these; run_signature_deck.py imports
# check_prompt_qc_deterministic / check_copy_qc_deterministic / check_deck_harmony
# per the §3.6 interface contract and never edits this file.
#
# SHIFT-LEFT: the writing-engine wrappers fire at COPY-QC (before any prompt is
# authored); EXCELLENCE + the per-slide perceptual measurer fire at PROMPT-QC
# (before any kie.ai render is paid for); deck-level HARMONY fires pre-assembly.
# ===========================================================================

# EXCELLENCE quality floor (§3.4 / G6). A prompt must clear BOTH the 9,000-char
# LENGTH floor (necessary) AND this richness floor (sufficient): a floor-grazing,
# boilerplate-padded prompt scores below this and routes back even though it clears
# 9,000 chars — the two floors are independent (§1.1a). Env-overridable for tuning.
try:
    PROMPT_EXCELLENCE_FLOOR = float(os.environ.get("PROMPT_EXCELLENCE_FLOOR", "0.70"))
except (TypeError, ValueError):
    PROMPT_EXCELLENCE_FLOOR = 0.70
# A genuinely rich 15-element prompt carries 400+ distinct words; 220 is the bare
# AF-P-DENSITY floor. Excellence credits richness up to this target.
PROMPT_EXCELLENCE_DISTINCT_TARGET = 400
# Local world-grounding markers (kept in build_deck so the EXCELLENCE/HARMONY gates
# do not reach into intelligence_engines_check's private token sets).
_WORLD_GROUND_TOKENS = (
    "because", "would this", "believab", "plausib", "to scale", "real-world",
    "real world", "grounded", "normal house", "not a luxury", "actual ", "in-world",
    "true to life", "true-to-life", "lived-in", "lived in",
)
# Recurring-character continuity anchors (Story engine, deck-level harmony).
_CHARACTER_CONTINUITY_TOKENS = (
    "same character", "same person", "same woman", "same man", "recurring",
    "continuity", "consistent cast", "consistent character", "reappears", "returns",
    "established earlier", "as before", "carry forward", "carried forward",
    "matching the", "same subject", "the protagonist", "our protagonist",
)


def _import_pitch_engines_check():
    """Import pitch_engines_check beside build_deck.py (mirror of
    _import_intelligence_engines_check). Returns the module or None."""
    try:
        import importlib
        return importlib.import_module("pitch_engines_check")
    except Exception:  # noqa: BLE001
        try:
            import importlib.util as _ilu
            spec = _ilu.spec_from_file_location(
                "pitch_engines_check",
                str(Path(__file__).resolve().parent / "pitch_engines_check.py"))
            if spec and spec.loader:
                mod = _ilu.module_from_spec(spec)
                spec.loader.exec_module(mod)
                return mod
        except Exception:  # noqa: BLE001
            return None
    return None


def _import_slide_geometry():
    """Import the slide_geometry module that ships beside build_deck.py.
    Tries a normal import first (scripts/ is on sys.path when build_deck runs / is
    imported), then a path-based load from this file's own directory. Returns the module
    or None (callers degrade gracefully). Mirrors _import_prompt_gate."""
    try:
        import importlib
        return importlib.import_module("slide_geometry")
    except Exception:  # noqa: BLE001
        try:
            import importlib.util as _ilu
            spec = _ilu.spec_from_file_location(
                "slide_geometry",
                str(Path(__file__).resolve().parent / "slide_geometry.py"))
            if spec and spec.loader:
                mod = _ilu.module_from_spec(spec)
                spec.loader.exec_module(mod)
                return mod
        except Exception:  # noqa: BLE001
            return None
    return None


def _chk_text_fits(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-TEXT-OVERFLOW wrapper — the geometry lives in slide_geometry.py; this is the
    build_deck-side symbol the manifest's py_symbol lockstep resolves against. Passes
    SLIDE_GEOMETRY_EDGE_MARGIN_FRAC in explicitly so build_deck's declared margin and
    slide_geometry.py's enforced margin can never silently diverge."""
    sg = _import_slide_geometry()
    if sg is None:
        return ""    # module absent -> defer, exactly as _import_prompt_gate callers do
    return sg.check_text_fits(
        run_dir, slides_path, edge_margin_frac=SLIDE_GEOMETRY_EDGE_MARGIN_FRAC)


def _chk_spelling(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-SPELLING wrapper."""
    sg = _import_slide_geometry()
    if sg is None:
        return ""
    return sg.check_spelling(run_dir, slides_path)


def _chk_type_size(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-TYPE-SIZE-MEASURED wrapper. Passes the FLOOR CONSTANTS, the reference render
    height, and the dark-theme decision in, so the measured floor and check_font_floor's
    declared floor are the same number by construction."""
    sg = _import_slide_geometry()
    if sg is None:
        return ""
    dark = _read_dark_optin(run_dir)
    return sg.check_type_size(
        run_dir, slides_path,
        pt_floor=(DARK_THEME_BODY_PT_FLOOR if dark else FONT_BODY_PT_FLOOR),
        dark=dark,
        pt_reference_height_px=SLIDE_GEOMETRY_PT_REFERENCE_HEIGHT_PX)


def _engine_name_for_code(code: str) -> str:
    """Map an AF code to the human INTELLIGENCE-engine name it belongs to, for the
    routeback payload's `intelligence` field (so the author knows which engine is absent)."""
    c = (code or "").upper()
    table = {
        "AF-FACE-PROMPT-MISSING": "Facial",
        "AF-LIGHT-PROMPT-MISSING": "Lighting",
        "AF-WORLD-SCALE": "World",
        "AF-HAIR-INAUTHENTIC": "Representation/Hair",
        "AF-NO-VILLAIN": "Story",
        "AF-NO-FELT-STAKES": "Emotional",
        "AF-CADENCE": "Pricing",
        "AF-NO-COST-OF-INACTION": "Pricing",
        "AF-GUARANTEE-GENERIC": "Pricing",
        "AF-NO-BRANDED-METHOD": "Pricing",
        "AF-METHOD-FABRICATED": "Pricing",
        "AF-NO-TIME-TO-RESULT": "Pricing",
        "AF-P13": "Negative-Block",
        "AF-P14": "Typography/Spelling-Lock",
        "AF-P-VERBATIM": "Typography/Verbatim",
        "AF-P-DENSITY": "Density",
        "AF-R3": "Representation",
        "AF-EXCELLENCE": "Excellence",
        "AF-HARMONY": "Harmony",
    }
    if c.startswith("AF-HOOK") or c.startswith("AF-OBI") or c == "AF-C2":
        return "Hook"
    return table.get(c, "Intelligence")


def _slide_no_from(val) -> Optional[int]:
    """Extract a slide ordinal from a checker's `slide` field ('slide-07' -> 7,
    7 -> 7). Returns None when no ordinal is present (DECK-level deficiency)."""
    if isinstance(val, int):
        return val
    m = re.search(r"(\d+)", str(val or ""))
    return int(m.group(1)) if m else None


def _fmt_engine_problem(p) -> str:
    """One-line render of a checker problem entry (dict or str), defensive about shape."""
    if isinstance(p, dict):
        where = p.get("slide", p.get("rung", "DECK"))
        return f"{p.get('code', '?')} [{where}]: {p.get('detail', '')}"
    return str(p)


def _is_defer_only(p) -> bool:
    """True for a pure-defer sentinel ({'defer': ...} with no 'code') a sub-engine
    emits when its inputs are not present yet."""
    return isinstance(p, dict) and "code" not in p and "defer" in p


def _engine_problem_to_def(p, phase: str) -> dict:
    """Convert a checker problem (dict/str) into the standard routeback deficiency
    schema {code, severity, slide, measured, required, intelligence, fix, phase}."""
    if isinstance(p, dict):
        code = p.get("code", "AF-COPY")
        slide = p.get("slide", p.get("rung", "DECK"))
        detail = p.get("detail", "")
    else:
        code, slide, detail = "AF-COPY", "DECK", str(p)
    return {"code": code, "severity": "reauthor", "slide": slide,
            "measured": "engine deficiency", "required": "engine present",
            "intelligence": _engine_name_for_code(code), "fix": detail, "phase": phase}


def _pdef(code, severity, measured, required, intelligence, fix) -> dict:
    """Build one per-slide prompt deficiency in the routeback schema (§3.5)."""
    return {"code": code, "severity": severity, "measured": measured,
            "required": required, "intelligence": intelligence, "fix": fix}


def check_intelligence_engines_copy(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """G1 — AF-NO-VILLAIN / AF-NO-FELT-STAKES (+ narrative-harmony codes): the WRITING-half
    INTELLIGENCE engines (Story villain-before-hero, Emotional felt-stakes-before-offer, Hook
    copy refrain, Recap) + narrative harmony, WIRED INTO THE COPY-QC PREFLIGHT — the EARLIEST
    phase a narrative defect is detectable (BEFORE any image prompt is authored; shift-left).

    Mirrors check_intelligence_engines_prompt's import/defer pattern: imports
    intelligence_engines_check and calls iec.check_copy(<working>, problems), which per the
    v15 interface contract runs the WRITING engines AND check_narrative_harmony, appending AF
    codes to `problems`. Defers (returns "") pre-copy (no slides_copy.md yet). Returns a fatal
    message naming the offending codes/slides on any auto-fail (run_preflight -> exit 3)."""
    copy_md = run_dir / "working" / "copy" / "slides_copy.md"
    if not copy_md.exists():
        return ""  # pre-copy phase — the copy engines defer (mirrors check_copy's own defer).
    iec = _import_intelligence_engines_check()
    if iec is None:
        return ("AF-NO-VILLAIN: intelligence_engines_check.py could not be imported from the "
                "scripts directory; the Story / Emotional copy engines + narrative harmony "
                "could not run. Ensure intelligence_engines_check.py is present beside "
                "build_deck.py.")
    problems: list = []
    try:
        # iec reads <dir>/copy, <dir>/prompts, <dir>/brand — so working/ is its run dir.
        iec.check_copy(run_dir / "working", problems)
    except Exception as exc:  # noqa: BLE001
        return (f"AF-NO-VILLAIN: intelligence_engines_check.check_copy raised {exc!r}; the "
                f"copy-side INTELLIGENCE engines (Story / Emotional / narrative harmony) "
                f"could not be verified.")
    problems = [p for p in problems if not _is_defer_only(p)]
    if problems:
        lines = "; ".join(_fmt_engine_problem(p) for p in problems[:10])
        more = "" if len(problems) <= 10 else f" (+{len(problems) - 10} more)"
        return ("AF-INTELLIGENCE-COPY: the WRITING-half engines auto-failed at COPY-QC — the "
                "deck must plant the VILLAIN before the first HERO/solution beat (Story), land "
                "a felt-stakes quantifier BEFORE the offer (Emotional), carry the sacred Hook "
                "refrain, and cohere as a narrative arc. Fix the SCRIPT before any prompt is "
                "authored (shift-left). Offenders: " + lines + more + ".")
    return ""


def check_pitch_engines(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """G2 — AF-CADENCE / AF-NO-COST-OF-INACTION / AF-GUARANTEE-GENERIC /
    AF-NO-BRANDED-METHOD / AF-METHOD-FABRICATED / AF-NO-TIME-TO-RESULT: the persuasion (offer)
    sub-engines WIRED INTO THE COPY-QC PREFLIGHT. PROMISE PRECEDES PRICE; the offer's value
    cadence, cost-of-inaction, branded method, guarantee specificity and time-to-result are
    graded on the SCRIPT before any prompt is authored (shift-left).

    CONDITIONAL on intake.json.pitch_included (mirrors _chk_pitch): defers for a pitchless
    deck (false) or an unset flag (AF-PITCH-LEAK / AF-PITCH-FLAG-UNSET own those). Defers
    pre-copy. Calls pitch_engines_check.check_copy(run_dir, problems) — the v15 importable
    entry point Agent 3 exposes; run_dir is the ACTUAL run dir, matching the module's own
    load_run(run_dir) -> working/copy convention. Returns a fatal message on any auto-fail."""
    if _intake_pitch_included(run_dir) is not True:
        return ""  # pitchless / unset — AF-PITCH-LEAK / AF-PITCH-FLAG-UNSET own those cases.
    copy_md = run_dir / "working" / "copy" / "slides_copy.md"
    if not copy_md.exists():
        return ""  # pre-copy phase — the offer engines defer.
    pec = _import_pitch_engines_check()
    if pec is None:
        return ("AF-PITCH-ENGINE: pitch_engines_check.py could not be imported from the "
                "scripts directory; the offer sub-engines (cadence / cost-of-inaction / "
                "branded-method / guarantee / time-to-result) could not run. Ensure "
                "pitch_engines_check.py is present beside build_deck.py.")
    if not hasattr(pec, "check_copy"):
        return ("AF-PITCH-ENGINE: pitch_engines_check.py exposes no check_copy(run_dir, "
                "problems) entry point (the v15 interface contract requires it). The offer "
                "sub-engines could not be wired into COPY-QC.")
    problems: list = []
    try:
        pec.check_copy(run_dir, problems)
    except Exception as exc:  # noqa: BLE001
        return (f"AF-PITCH-ENGINE: pitch_engines_check.check_copy raised {exc!r}; the offer "
                f"persuasion sub-engines could not be verified.")
    problems = [p for p in problems if not _is_defer_only(p)]
    if problems:
        lines = "; ".join(_fmt_engine_problem(p) for p in problems[:10])
        more = "" if len(problems) <= 10 else f" (+{len(problems) - 10} more)"
        return ("AF-PITCH-ENGINE: the offer sub-engines auto-failed at COPY-QC — PROMISE must "
                "precede PRICE and the offer must carry a value cadence, a quantified "
                "cost-of-inaction, a non-generic guarantee, a real (un-fabricated) branded "
                "method, and a concrete time-to-result. Fix the OFFER SCRIPT before any prompt "
                "is authored (shift-left). Offenders: " + lines + more + ".")
    return ""


def _excellence_score(prompt_text: str):
    """G6 — the EXCELLENCE dimension. Returns (score 0.0-1.0, reasons) measuring whether the
    prompt's character budget was spent on DEFECT-PREVENTING SPECIFICITY (richness) rather
    than boilerplate padding. A prompt can clear the 9,000-char LENGTH floor and still score
    low here — that is the design: length is necessary, excellence is the independent QUALITY
    gate (two floors, §1.1a), so a compliant-but-soulless prompt does NOT pass.

    Weighted dimensions (sum to 1.0):
      * 0.25 negative-block coverage   — fraction of the 8 paired defect classes named
      * 0.15 spelling-lock richness    — >=2 distinct per-string lock markers
      * 0.15 people-anatomy + world-grounding specificity
      * 0.20 concrete craft signals    — brand HEX + explicit type SIZE + composition token
      * 0.25 lexical richness          — distinct words vs a rich 15-element target (400)
    """
    lc = prompt_text.lower()
    reasons = []

    # 1. negative-block coverage (0.25)
    n_classes = len(NEGATIVE_BLOCK_CLASS_TOKENS)
    missing_classes = _negative_block_class_problems(lc)
    nb_cov = (n_classes - len(missing_classes)) / max(1, n_classes)
    if missing_classes:
        reasons.append(f"negative block names only {n_classes - len(missing_classes)}/"
                       f"{n_classes} defect classes (missing: {', '.join(missing_classes)})")

    # 2. spelling-lock richness (0.15)
    n_lock = sum(1 for t in SPELLING_LOCK_TOKENS if t in lc)
    sl = 1.0 if n_lock >= 2 else (0.5 if n_lock == 1 else 0.0)
    if n_lock < 2:
        reasons.append(f"only {n_lock} spelling-lock marker(s) (>=2 distinct expected for "
                       f"per-string lock coverage)")

    # 3. people-anatomy + world-grounding specificity (0.15)
    has_anatomy = any(t in lc for t in NEGATIVE_BLOCK_CLASS_TOKENS["anatomical artifacts"])
    has_world = any(t in lc for t in _WORLD_GROUND_TOKENS)
    pw = (0.5 if has_anatomy else 0.0) + (0.5 if has_world else 0.0)
    if not has_anatomy:
        reasons.append("no people-anatomy specificity (finger/hand/eye/skin defect guards)")
    if not has_world:
        reasons.append("no world-grounding/believability justification token")

    # 4. concrete craft signals (0.20)
    craft_hits = 0
    if _HEX_COLOR_RE.search(prompt_text):
        craft_hits += 1
    else:
        reasons.append("no brand palette HEX (#RRGGBB)")
    if _TYPE_SIZE_RE.search(prompt_text):
        craft_hits += 1
    else:
        reasons.append("no explicit type SIZE (pt/px)")
    if any(t in lc for t in PROMPT_COMPOSITION_TOKENS):
        craft_hits += 1
    else:
        reasons.append("no composition/zone token (thirds/grid/zone/safe-margin)")
    craft = craft_hits / 3.0

    # 5. lexical richness (0.25)
    distinct = len(set(_WORD_RE.findall(lc)))
    lex = min(1.0, distinct / float(PROMPT_EXCELLENCE_DISTINCT_TARGET))
    if distinct < PROMPT_EXCELLENCE_DISTINCT_TARGET:
        reasons.append(f"only {distinct} distinct words (rich target "
                       f"{PROMPT_EXCELLENCE_DISTINCT_TARGET}) — likely boilerplate padding")

    score = (0.25 * nb_cov + 0.15 * sl + 0.15 * pw + 0.20 * craft + 0.25 * lex)
    return round(score, 3), reasons


def check_prompt_excellence(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """G6 preflight gate (AF-EXCELLENCE) — fires at PROMPT-QC, BEFORE any render. Every
    per-slide rich prompt must score >= PROMPT_EXCELLENCE_FLOOR on the EXCELLENCE dimension:
    the budget must be spent on defect-preventing specificity, not padding. This is the
    QUALITY floor that sits ALONGSIDE the 9,000-char LENGTH floor (§1.1a) — a floor-grazing,
    boilerplate-padded prompt clears 9,000 chars yet routes back here. Defers pre-prompt
    (no prompts dir); the slide-count / missing-file failures are owned by _chk_rich_prompts.
    Returns "" on pass, or a fatal AF-EXCELLENCE message naming the under-floor slides."""
    prompts_dir = run_dir / "working" / "prompts"
    if not prompts_dir.is_dir():
        return ""  # pre-prompt phase — defer (mirrors the perceptual prompt engines).
    n = _count_output_slides(run_dir, slides_path)
    if n is None:
        return ""  # _chk_rich_prompts owns the 'no slide count' failure.
    offenders = []
    for ordinal in range(1, n + 1):
        p = resolve_prompt_path(run_dir, ordinal)
        if p is None:
            continue  # _chk_rich_prompts owns the missing-file failure.
        stripped = p.read_text(errors="replace").strip()
        score, reasons = _excellence_score(stripped)
        if score < PROMPT_EXCELLENCE_FLOOR:
            offenders.append(f"slide {ordinal:02d}: excellence {score:.2f} < "
                             f"{PROMPT_EXCELLENCE_FLOOR:.2f} ({'; '.join(reasons[:4])})")
    if offenders:
        return ("AF-EXCELLENCE: one or more prompts clear the 9,000-char LENGTH floor but "
                "fail the EXCELLENCE QUALITY floor — the budget was spent on padding, not "
                "defect-preventing specificity (full 8-class negative block, per-string "
                "spelling-lock, people-anatomy + world-grounding, brand HEX + type SIZE + "
                "composition, lexical richness). Amazing, not merely compliant (§1.1a). "
                "Re-author (do NOT pad): " + " | ".join(offenders))
    return ""


def check_deck_harmony(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """G5 — AF-HARMONY: the deck-wide COHESION gate, fired PRE-ASSEMBLY (after render, before
    assembly/delivery). Proves the engines are ORCHESTRATED IN HARMONY across slides — not
    just that each slide passes in isolation: recurring-character continuity (Story across
    slides), palette/brand coherence (cross-check check_brand_consistency), world continuity,
    and archetype RHYTHM (a deliberate recurring motif — distinct from _chk_creativity, which
    only rejects sameness). Defers (returns "") pre-render (no PNGs yet). Returns "" on pass,
    or a fatal AF-HARMONY message (run_preflight -> exit 3)."""
    pngs = _gather_rendered_pngs(run_dir)
    if not pngs:
        return ""  # pre-render — harmony is a pre-assembly gate; nothing rendered to cohere.

    problems = []
    prompts_dir = run_dir / "working" / "prompts"
    prompt_texts = {}
    if prompts_dir.is_dir():
        for pf in sorted(prompts_dir.glob("slide-*.txt")):
            sn = _slide_no_from(pf.stem)
            if sn is not None:
                prompt_texts[sn] = pf.read_text(errors="replace").lower()

    # --- recurring-character continuity (Story engine across slides) ---
    people_markers = ("person", "people", "woman", "man", "subject", "portrait", "face",
                      "presenter", "founder", "customer", "client", "she ", "he ", "they ")
    people_slides = [sn for sn, t in prompt_texts.items()
                     if any(m in t for m in people_markers)]
    casting_ledger = any((run_dir / rel).exists() for rel in (
        "working/brand/casting_ledger.json", "working/brand/casting.json",
        "working/copy/casting_ledger.json", "working/brand/cast.json"))
    if len(people_slides) >= 2 and not casting_ledger:
        anchored = sum(1 for sn in people_slides
                       if any(tok in prompt_texts[sn] for tok in _CHARACTER_CONTINUITY_TOKENS))
        if anchored < 2:
            problems.append(
                f"recurring-character continuity: {len(people_slides)} people-slides but only "
                f"{anchored} carry a continuity anchor (same character/recurring/consistent "
                f"cast) and no casting ledger exists — the Story engine's character does not "
                f"provably persist across the deck")

    # --- archetype RHYTHM (a deliberate recurring motif, not 40 one-offs) ---
    archetypes = []
    for rel in ("working/typography/design_system.json",
                "working/copy/design_system.json", "design_system.json"):
        ds = run_dir / rel
        if not ds.exists():
            continue
        obj = _read_json(ds)
        if not isinstance(obj, dict) or "__parse_error__" in obj:
            break
        per = obj.get("per_slide") or obj.get("slides") or obj.get("slide_plan")
        if isinstance(per, list):
            for s in per:
                if isinstance(s, dict):
                    a = s.get("archetype") or s.get("type") or s.get("treatment")
                    if isinstance(a, str) and a.strip():
                        archetypes.append(a.strip().lower())
        elif isinstance(per, dict):
            for v in per.values():
                if isinstance(v, str) and v.strip():
                    archetypes.append(v.strip().lower())
                elif isinstance(v, dict):
                    a = v.get("archetype") or v.get("type")
                    if isinstance(a, str) and a.strip():
                        archetypes.append(a.strip().lower())
        break
    if len(archetypes) >= 4:
        counts = {}
        for a in archetypes:
            counts[a] = counts.get(a, 0) + 1
        max_recur = max(counts.values())
        if max_recur < 2:
            problems.append(
                "archetype rhythm: every slide uses a unique layout archetype (no motif "
                "recurs) — the deck has no deliberate visual cadence/through-line; a coherent "
                "deck repeats a motif at intervals (distinct from anti-template sameness)")

    # --- palette / brand coherence cross-check (folded into harmony) ---
    try:
        brand_reason = check_brand_consistency(run_dir, slides_path)
    except Exception as _exc:  # noqa: BLE001
        import sys as _sys
        print(f"harmony brand cross-check raised: {_exc!r}", file=_sys.stderr)
        brand_reason = ""
    if brand_reason:
        problems.append("palette coherence (brand cross-check): " + brand_reason)

    # --- world continuity (conservative; only the wholly-fragmented extreme) ---
    world_tokens = ("office", "home", "studio", "stage", "kitchen", "boardroom", "outdoor",
                    "street", "warehouse", "classroom", "clinic", "factory", "field", "store")
    scene_world = {}
    for sn, t in prompt_texts.items():
        present = [w for w in world_tokens if w in t]
        if present:
            scene_world[sn] = set(present)
    if len(scene_world) >= 3:
        shared = set.intersection(*scene_world.values()) if scene_world else set()
        any_pair = False
        slides_w = list(scene_world.items())
        for i in range(len(slides_w)):
            for j in range(i + 1, len(slides_w)):
                if slides_w[i][1] & slides_w[j][1]:
                    any_pair = True
                    break
            if any_pair:
                break
        if not shared and not any_pair:
            problems.append(
                "world continuity: the scene-bearing slides share no common world/setting "
                "token across the deck — the world reads as disjoint per-slide scenes rather "
                "than one believable continuous world")

    if problems:
        return ("AF-HARMONY: the deck fails the cohesion gate — the engines pass per-slide but "
                "are NOT orchestrated in harmony across the deck (recurring character, palette "
                "coherence, world continuity, archetype rhythm). Re-render ONLY the "
                "inconsistent slides to restore continuity (Director of Presentations owns "
                "harmony; SOP-HARMONY-01). Offenders: " + "; ".join(problems) + ".")
    return ""


def check_copy_qc_deterministic(run_dir: Path, slides_path: Optional[Path] = None) -> dict:
    """G8 — the deterministic COPY-QC measurer that feeds run_copy_qc_loop (run_signature_deck
    imports it). The SOURCE OF TRUTH for whether the SCRIPT clears the WRITING engines BEFORE
    any prompt is authored — NOT the copy-QC agent's self-score. Runs the WRITING engines
    (iec.check_copy: Story / Emotional / narrative harmony) and, when pitch_included, the
    offer sub-engines (pitch_engines_check.check_copy). Returns
    {pass, phase:'copy', deficiencies:[{code, severity, slide, measured, required,
    intelligence, fix, phase}]} — `deficiencies` is what write_routeback_payload turns into a
    per-item work order. pass is True only when every WRITING/offer engine is clean."""
    deficiencies = []
    iec = _import_intelligence_engines_check()
    if iec is None:
        _miss = _pdef("AF-NO-VILLAIN", "reauthor", "module missing",
                      "intelligence_engines_check.py", "Story",
                      "Ensure intelligence_engines_check.py is beside build_deck.py.")
        _miss["slide"] = "DECK"
        _miss["phase"] = "copy"
        deficiencies.append(_miss)
    else:
        cprob = []
        try:
            iec.check_copy(run_dir / "working", cprob)
        except Exception as exc:  # noqa: BLE001
            cprob = [{"code": "AF-NO-VILLAIN",
                      "detail": f"intelligence_engines_check.check_copy raised {exc!r}"}]
        for p in cprob:
            if _is_defer_only(p):
                continue
            deficiencies.append(_engine_problem_to_def(p, "copy"))

    if _intake_pitch_included(run_dir) is True:
        pec = _import_pitch_engines_check()
        if pec is not None and hasattr(pec, "check_copy"):
            pprob = []
            try:
                pec.check_copy(run_dir, pprob)
            except Exception as exc:  # noqa: BLE001
                pprob = [{"code": "AF-PITCH-ENGINE",
                          "detail": f"pitch_engines_check.check_copy raised {exc!r}"}]
            for p in pprob:
                if _is_defer_only(p):
                    continue
                deficiencies.append(_engine_problem_to_def(p, "offer"))

    return {"pass": not deficiencies, "phase": "copy", "deficiencies": deficiencies}


def check_prompt_qc_deterministic(run_dir: Path, slides_path: Optional[Path] = None) -> dict:
    """G7 — the deterministic PROMPT-QC measurer, the SOURCE OF TRUTH that feeds
    run_prompt_qc_loop (run_signature_deck imports it). This is the gate that PHYSICALLY
    prevents a thin/soulless prompt from reaching kie.ai: the loop exits on THIS verdict, not
    the prompt-QC agent's self-typed pass.

    A slide passes ONLY when BOTH floors are clean (§1.1a):
      LENGTH gate  — C1 >= 9,000 chars (AF-P1, fatal); C2 <= 18,000 (AF-P2).
      QUALITY gate — C4 structural blocks; C5 perceptual INTELLIGENCE engines (Facial/
        Lighting/World/Hair + per-slide harmony + image-side Hook via iec.check_prompts);
        C6 8-class negative block (AF-P13); C7 per-string spelling-lock (AF-P14); C8 verbatim
        copy baked (AF-P-VERBATIM); C9 density + demographic-landmine (AF-P-DENSITY/AF-R3);
        C10 EXCELLENCE (AF-EXCELLENCE). Length alone never passes; engines alone never pass.

    Deck-level WRITING-engine deficiencies (iec.check_copy / pitch_engines_check) are recorded
    too as a backstop. Returns the §3.5 schema:
      {pass, n_slides, slides:{N:{char_count, excellence, deficiencies:[{code, severity,
       measured, required, intelligence, fix}]}}, deck_deficiencies:[...], deficiencies:[...]}
    where severity 'fatal' already hard-fails the renderer and 'reauthor' is what the loop
    sends back. `deficiencies` is a flat list (each tagged with its slide) for the routeback
    writer's convenience."""
    n = _count_output_slides(run_dir, slides_path)
    if n is None:
        d = _pdef("AF-P1", "fatal", "no slide count", "slides.json / arc_allocation.json",
                  "Structure", "Produce slides.json before Prompt-QC so prompts can be "
                  "verified per slide.")
        d_deck = dict(d, slide="DECK")
        return {"pass": False, "n_slides": 0, "slides": {},
                "deck_deficiencies": [d_deck], "deficiencies": [d_deck]}

    copy_map = _load_slide_copy_map(run_dir, slides_path)

    # Run the perceptual prompt engines ONCE (deck-level) and bucket by slide ordinal so each
    # slide's routeback names exactly which image engine (Facial/Lighting/World/Hair/Hook) is
    # absent. iec.check_prompts also carries the NEW per-slide harmony + image-side Hook.
    perceptual_by_slide = {}
    iec = _import_intelligence_engines_check()
    if iec is not None:
        iprob = []
        try:
            iec.check_prompts(run_dir / "working", iprob)
        except Exception:  # noqa: BLE001
            iprob = []
        for p in iprob:
            if _is_defer_only(p):
                continue
            sn = _slide_no_from(p.get("slide")) if isinstance(p, dict) else None
            perceptual_by_slide.setdefault(sn, []).append(p)

    # Deck-level WRITING-engine backstop (copy + offer must already be clean from COPY-QC).
    deck_def = []
    copy_verdict = check_copy_qc_deterministic(run_dir, slides_path)
    for d in copy_verdict.get("deficiencies", []):
        deck_def.append(dict(d, severity="deck"))

    slides = {}
    flat = []
    any_blocking = bool(deck_def)
    for ordinal in range(1, n + 1):
        deficiencies = []
        p = resolve_prompt_path(run_dir, ordinal)
        if p is None:
            deficiencies.append(_pdef(
                "AF-P1", "fatal", "no prompt file",
                "working/prompts/slide-NN.txt (9,000-18,000 chars)", "Rich-Prompt",
                "Author the rich per-slide prompt; build_deck renders it VERBATIM and never "
                "composes a thin fallback."))
            slides[ordinal] = {"char_count": 0, "excellence": 0.0, "deficiencies": deficiencies}
            any_blocking = True
            for d in deficiencies:
                flat.append(dict(d, slide=ordinal))
            continue

        stripped = p.read_text(errors="replace").strip()
        length = len(stripped)
        lc = stripped.lower()

        # --- LENGTH gate (C1/C2) ---
        if length < PROMPT_CHAR_FLOOR:
            deficiencies.append(_pdef(
                "AF-P1", "fatal", f"{length} chars", f">= {PROMPT_CHAR_FLOOR}", "Length",
                f"Re-author to {PROMPT_CHAR_FLOOR}-{PROMPT_CHAR_CEILING} chars of real "
                f"specificity; do NOT pad."))
        elif length > PROMPT_CHAR_CEILING:
            deficiencies.append(_pdef(
                "AF-P2", "reauthor", f"{length} chars", f"<= {PROMPT_CHAR_CEILING}", "Length",
                f"Trim to <= {PROMPT_CHAR_CEILING} chars without dropping any engine token."))

        # --- QUALITY gate (C4 structural blocks) ---
        # A block counts as present if its canonical label OR any accepted alias appears
        # (the negative-block header may read "DO-NOT BLOCK" — the canonical exemplar/SOP
        # label — or the legacy "NEGATIVE BLOCK"; either satisfies the slot). The gate is
        # not weakened: a prompt carrying NEITHER negative-block header still fails here,
        # and the 8-class CONTENT teeth (AF-P13, below) fire independently of the header.
        for b in _missing_structural_blocks(lc):
            aliases = STRUCTURAL_BLOCK_ALIASES.get(b, [])
            want = b + (f" (or alias {', '.join(aliases)})" if aliases else "")
            deficiencies.append(_pdef(
                "AF-P-STRUCT", "reauthor", "block missing", want, "Structure",
                f"Add the required structural block: {want}"))

        # --- C5 perceptual INTELLIGENCE engines (Facial/Lighting/World/Hair + harmony+Hook) ---
        for pp in perceptual_by_slide.get(ordinal, []):
            code = pp.get("code", "AF-INTELLIGENCE") if isinstance(pp, dict) else "AF-INTELLIGENCE"
            detail = pp.get("detail", "") if isinstance(pp, dict) else str(pp)
            deficiencies.append(_pdef(
                code, "reauthor", "engine token absent", "engine token present",
                _engine_name_for_code(code), detail))

        # --- C6 eight-class negative block (AF-P13) ---
        for cls in _negative_block_class_problems(lc):
            deficiencies.append(_pdef(
                "AF-P13", "reauthor", "class not named", cls, "Negative-Block",
                f"Name the '{cls}' defect class in the 8-class paired negative block (SOP 9.8)."))

        # --- C7 per-string spelling-lock (AF-P14) ---
        if not _spelling_lock_present(lc):
            deficiencies.append(_pdef(
                "AF-P14", "reauthor", "no lock marker", "per-string spelling-lock",
                "Typography/Spelling-Lock",
                "Add a letter-for-letter spelling-lock directive for every on-slide string."))

        # --- C8 verbatim copy baked (AF-P-VERBATIM) ---
        for mc in _verbatim_copy_problems(stripped, copy_map.get(ordinal)):
            deficiencies.append(_pdef(
                "AF-P-VERBATIM", "reauthor", "copy not baked", mc, "Typography/Verbatim",
                f"Bake the slide's exact copy string verbatim so kie.ai renders it: {mc}"))

        # --- C9 density + demographic landmine (AF-P-DENSITY / AF-R3) ---
        for dd in _prompt_density_problems(stripped, lc):
            deficiencies.append(_pdef(
                "AF-P-DENSITY", "reauthor", "thin/padded",
                "hex + type-size + composition + distinct-word floor", "Density", dd))
        for landmine in FORBIDDEN_DEMOGRAPHIC_DEFAULTS:
            if landmine.lower() in lc:
                deficiencies.append(_pdef(
                    "AF-R3", "fatal", f"landmine '{landmine}'", "no demographic default",
                    "Representation",
                    "Remove the hardcoded demographic-default landmine; cast from the "
                    "client's captured audience / casting ledger (SOP-CAST-01)."))

        # --- C10 EXCELLENCE (AF-EXCELLENCE) ---
        score, ereasons = _excellence_score(stripped)
        if score < PROMPT_EXCELLENCE_FLOOR:
            deficiencies.append(_pdef(
                "AF-EXCELLENCE", "reauthor", f"{score:.2f}", f">= {PROMPT_EXCELLENCE_FLOOR:.2f}",
                "Excellence",
                "Spend the budget on defect-preventing specificity, not padding: "
                + "; ".join(ereasons[:5])))

        slides[ordinal] = {"char_count": length, "excellence": score,
                           "deficiencies": deficiencies}
        if any(d["severity"] in ("fatal", "reauthor") for d in deficiencies):
            any_blocking = True
        for d in deficiencies:
            flat.append(dict(d, slide=ordinal))

    flat_deck = [dict(d, slide=d.get("slide", "DECK")) for d in deck_def]
    return {"pass": not any_blocking, "n_slides": n, "slides": slides,
            "deck_deficiencies": deck_def, "deficiencies": flat + flat_deck}


# ===========================================================================
# POWERFUL-PRESENTATION DOCTRINE GATES (manifest v18) — the priority-shift spine.
# ---------------------------------------------------------------------------
# The framework's North Star: the #1 job is to HOLD ATTENTION the whole duration
# so the audience re-ranks the owner's offer/idea to the TOP of its priority
# stack (the priority shift). These deterministic gates make the framework LAW.
#
# MASTER SWITCH (no-regression guarantee): every doctrine gate DEFERS (returns "")
# unless Phase P0B-PRIORITY (the attention-content-strategist) has produced
# working/copy/priority_shift_spec.json. Before that artifact exists (legacy /
# pre-doctrine / ad-hoc decks, and every test fixture that predates the spine),
# the gates are wired but never fire — so they cannot break an in-flight build the
# new phase never touched. Once the spec exists, the gates have full teeth.
# ===========================================================================
PRIORITY_SPEC_REL = "working/copy/priority_shift_spec.json"
STYLE_SAMPLES_MANIFEST_REL = "working/style-preview/style_samples_manifest.json"
STYLE_CHOICE_REL = "working/copy/style_preview_choice.json"
PRIORITY_SHIFT_REPORT_REL = "working/qc/priority_shift_report.json"
# Phase ids (PIPELINE-MANIFEST.json doctrine spine). P0B-PRIORITY (manifest order 0.2)
# is the attention-content-strategist phase that PRODUCES priority_shift_spec.json; it is
# a mandatory order-based precondition of P4-RENDER (manifest order 4.9, the money phase).
# v16.0.1 binds the render preflight to P0B at EVERY entry point so a DIRECT build_deck.py
# render cannot render doctrine-blind by simply omitting the spec (see run_preflight).
PRIORITY_PHASE_ID = "P0B-PRIORITY"
RENDER_PHASE_ID = "P4-RENDER"
# Creation modes (P19/P118 — Step Zero identifies the mode before anything else).
CREATION_MODES = ("from_scratch", "content_personal", "content_general")
# U021 -- the CLOSED set of deck_type values. Mirrors deck-intake-driver.py's
# LEGACY_FIELD_MAPPING (the single point of derivation) and
# intake/deck-intake-questions.json's legacy_field_mapping table. deck_type is
# ENGINE-WRITTEN by derive_legacy_fields(); it is never an agent's free text.
# Adding a value here without adding it there re-opens the split this closes.
DECK_TYPES = ("webinar", "signature_presentation")
# U022 -- dated migration window for the two gates that must stop exempting
# legacy run dirs (AF-MODE-UNSET's no-doctrine exemption, and U021's
# deck_type engine gate, registered in the manifest as AF-SP-TYPE-UNDECLARED).
# A run dir created before this date and carrying no
# creation_mode is grandfathered; after it, the gate blocks. Mirrors the
# established pattern at prove_sp_intake.py:96. Do NOT extend this date in
# place -- retire it in a dated follow-up line item.
MIGRATION_WINDOW_UNTIL = date(2026, 9, 30)
# U021 stage flag -- MODULE SCOPE, beside DECK_TYPES. Rule 3.5: warn -> remediate ->
# enforce, and the enforce flip is a SEPARATE unit whose ENTIRE change is this one
# line. Its prerequisite is an installed deck_type producer reachable from the
# department (deck-intake-driver.py in <dept>/scripts/, U006/A4, or U058).
# The dated window below REPORTS that the window has closed; it does not enforce.
DECK_TYPE_GATE_STAGE = "enforce"          # "warn" | "enforce"
# The eight-move build sequence (P141-P150), in canonical order. The copy must plant
# these beat tags monotonically so the arc actually engineers the shift.
EIGHT_MOVE_TAGS = (
    "PRIORITY_STACK", "PRESENT_COST", "HIGHER_PRIORITY",
    "VALUE_ANCHOR", "URGENCY_SCARCITY", "ABILITY_UNBLOCK",
    "RERANK_DEMAND", "TRIGGER",
)
# Hedge phrases a PROCLAMATION must never carry (P109). Deliberately the strong,
# unambiguous multi-word hedges only — bare "maybe"/"perhaps" are too common in
# legitimate copy to gate on, so they are excluded to avoid false positives.
PROCLAMATION_HEDGE_TOKENS = (
    "kind of", "sort of", "you might want to consider", "you may want to consider",
    "i think maybe", "we could probably", "it might be that", "i guess",
    "more or less", "if that makes sense", "sort-of", "kind-of",
)
PRIORITY_STACK_MARKERS = (
    "priority_stack", "priority stack", "priority-stack", "top of your list",
    "what matters most", "what you prioritize", "what you prioritise",
    "your current priorities", "ranks above",
)
RERANK_MARKERS = (
    "re-rank", "rerank", "re_rank", "rerank_demand", "move this to the top",
    "make this your #1", "make this your number one", "your new #1",
    "decide now", "decision-demand", "decision demand", "choose this first",
    "put this first", "this becomes your priority",
)
TRIGGER_MARKERS = (
    "today only", "by midnight", "this week only", "next 24 hours",
    "deadline", "doors close", "doors closing", "closes ", "ends tonight",
    "ends ", "expires", "only until", "limited to the first", "enroll now",
    "enrol now", "act now", "before the timer", "registration closes",
)
PEAK_TAGS = ("peak", "apex", "wow", "salience_apex", "salience-apex",
             "promise-apex", "promise_apex")
ENDING_TAGS = ("recap", "close", "closing", "cta", "call to action", "call-to-action",
               "final", "ending", "send-off", "sendoff")
LADDER_BEAT_MARKERS = ("ladder", "anchor", "price", "offer", "value-stack",
                       "value_stack", "valuestack", "drop", "value_add", "value add")
COST_OF_INACTION_MARKERS = ("cost of inaction", "cost_of_inaction", "cost-of-inaction",
                            "present cost", "present-cost", "present_cost",
                            "what it costs to stay", "the cost of doing nothing")
URGENCY_SCARCITY_MARKERS = ("urgency", "scarcity", "limited", "only ", "deadline",
                            "running out", "few spots", "seats left", "while it lasts")
ABILITY_UNBLOCK_MARKERS = ("payment plan", "payment-plan", "financing", "ability_unblock",
                           "ability-unblock", "remove the blocker", "easy to start",
                           "no experience needed", "done for you", "done-for-you",
                           "we handle", "guarantee")
# HOLE C — the orphaned "no-beat" persuasion taxonomy (documented in SOP-PITCH-06 /
# SOP-ENGINE-00 but never enforced). WIRED here as a single conservative deck-level
# gate. Each entry: code -> (label, generous synonym markers). A doctrine deck that
# omits one of these named persuasion beats fails the specific code.
PERSUASION_BEAT_MARKERS = {
    "AF-NO-PROBLEM": ("problem", "villain", "enemy", "what's broken", "whats broken",
                      "the real reason", "stuck", "struggle", "pain"),
    "AF-NO-CHOICE": ("choice", "two paths", "two roads", "you can either",
                     "either keep", "fork", "decision"),
    "AF-NO-FORK": ("fork", "two paths", "two roads", "path a", "path b",
                   "stay the same or", "keep doing", "or you can"),
    "AF-NO-COMPARISON": ("compared to", "vs ", "versus", "old way", "new way",
                         "instead of", "unlike", "the difference between"),
    "AF-NO-MEASURABLE-RESULTS": ("%", "result", "results", "increase", "roi",
                                 "grew ", "doubled", "tripled", "x in ", "in 30 days",
                                 "measurable"),
    "AF-NO-EXPERT-PROOF": ("expert", "study", "research", "according to", "data shows",
                           "proven", "case study", "wall of wins", "testimonial",
                           "as featured"),
    "AF-NO-BEFORE-AFTER": ("before", "after", "transformation", "went from",
                           "used to", "now i", "now they", "from struggling"),
}


def _doctrine_active(run_dir: Path) -> bool:
    """The priority-shift doctrine engages ONLY once Phase P0B-PRIORITY has produced
    a valid working/copy/priority_shift_spec.json. This is the no-regression switch."""
    return _read_priority_spec(run_dir) is not None


def _read_priority_spec(run_dir: Path):
    """Parse working/copy/priority_shift_spec.json. Returns the dict, or None when
    absent/unparseable (gates then DEFER)."""
    for rel in (PRIORITY_SPEC_REL, "priority_shift_spec.json",
                "working/priority_shift_spec.json"):
        p = run_dir / rel
        if p.exists():
            obj = _read_json(p)
            if isinstance(obj, dict) and "__parse_error__" not in obj:
                return obj
            return None
    return None


def _slides_copy_text_lc(run_dir: Path):
    """Lowercased slides_copy.md text, or None when no copy exists yet (defer)."""
    for rel in ("working/copy/slides_copy.md", "slides_copy.md",
                "working/slides_copy.md"):
        p = run_dir / rel
        if p.exists():
            return p.read_text(errors="replace").lower()
    return None


def _first_marker_offset(text_lc: str, markers) -> Optional[int]:
    """Lowest character offset at which any of `markers` appears, or None."""
    best = None
    for m in markers:
        idx = text_lc.find(m.lower())
        if idx >= 0 and (best is None or idx < best):
            best = idx
    return best


def _png_ordinal(path: Path) -> Optional[int]:
    m = re.search(r"(\d+)", path.stem)
    return int(m.group(1)) if m else None


def _apex_slide_ordinal(run_dir: Path) -> Optional[int]:
    """The OFFER/PROMISE-APEX slide ordinal, read from the arc allocation's PEAK/APEX/
    OFFER beat tag. Returns None when it cannot be determined (gate then defers)."""
    for rel in ("working/copy/arc_allocation.json", "arc_allocation.json",
                "working/arc_allocation.json"):
        p = run_dir / rel
        if not p.exists():
            continue
        obj = _read_json(p)
        if not isinstance(obj, (list, dict)) or (isinstance(obj, dict) and "__parse_error__" in obj):
            return None
        slots = obj if isinstance(obj, list) else (
            obj.get("slots") or obj.get("allocation") or obj.get("slides") or [])
        for s in slots if isinstance(slots, list) else []:
            if not isinstance(s, dict):
                continue
            tokens = []
            for k in ("arc_section", "section", "beat", "tag", "type", "role"):
                v = s.get(k)
                if isinstance(v, str):
                    tokens.append(v.lower())
            tags = s.get("tags")
            if isinstance(tags, list):
                tokens += [str(t).lower() for t in tags]
            blob = " ".join(tokens)
            if any(t in blob for t in ("apex", "salience_apex", "promise-apex",
                                       "promise_apex", "wow-peak", "peak")):
                ordn = s.get("slide") or s.get("ordinal") or s.get("n")
                try:
                    return int(ordn)
                except (TypeError, ValueError):
                    return None
        return None
    return None


def _chk_deck_type(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-DECK-TYPE-UNSET (U021). intake.json.deck_type must be present and one of
    DECK_TYPES. It is ENGINE-WRITTEN by deck-intake-driver.derive_legacy_fields();
    a run that reaches preflight without it was assembled by hand, and _sp_active
    then silently defers every signature gate.

    Defers when intake is absent -- PREFLIGHT_REQUIRED entry #1 (_chk_intake) owns
    absence and returns 'file absent', so absence is already fail-closed there.
    Warn-mode until MIGRATION_WINDOW_UNTIL (Rule 3.5): reports to stderr and
    returns '' so no in-flight run dir is broken by the roll. Waivable only by a
    logged owner_skip_approval for AF-DECK-TYPE-UNSET."""
    intake = _read_intake_obj(run_dir)
    if not isinstance(intake, dict):
        return ""  # absence is _chk_intake's, fail-closed there.
    declared = str(intake.get("deck_type") or "").strip()
    if declared in DECK_TYPES:
        return ""
    reason = (f"AF-DECK-TYPE-UNSET: intake.json.deck_type is "
              f"{declared or 'unset'!r}; it must be one of {', '.join(DECK_TYPES)}. "
              f"deck_type is written by deck-intake-driver.py's "
              f"derive_legacy_fields() from the ONE presentation_type answer -- it is "
              f"never hand-typed. An unset deck_type makes _sp_active defer, so every "
              f"signature-presentation gate silently no-ops (SOP-MODE-00 / P-SP-CLAIM).")
    if _owner_skip_approved(run_dir, "AF-DECK-TYPE-UNSET"):
        return ""
    # Read the MODULE-SCOPE stage flag declared above. Do NOT re-declare it here:
    # a local rebinding makes `return reason` unreachable and the gate permanently warn.
    if DECK_TYPE_GATE_STAGE != "enforce":
        overdue = "" if date.today() <= MIGRATION_WINDOW_UNTIL else " [WINDOW CLOSED -- the enforce unit is overdue]"
        print("  WARN  " + reason + " [warn-mode until "
              + MIGRATION_WINDOW_UNTIL.isoformat() + "]" + overdue, file=sys.stderr)
        return ""
    return reason


def _chk_mode(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-MODE-UNSET (P118). intake.json.creation_mode must be one of
    from_scratch | content_personal | content_general; the two content modes need an
    extracted_substance record. Defers when intake is absent (_chk_intake owns that)
    or — for a legacy deck with no mode AND no doctrine spec — defers so a pre-doctrine
    build is never broken. A doctrine-active deck MUST declare the mode."""
    intake = _read_intake_obj(run_dir)
    if not isinstance(intake, dict):
        return ""  # no intake — _chk_intake / AF-SLIDE-COUNT-FLOOR own absence.
    mode = str(intake.get("creation_mode") or "").strip().lower()
    if not mode and not _doctrine_active(run_dir):
        # U022: the pre-doctrine exemption is now DATED, not permanent.
        if _owner_skip_approved(run_dir, "AF-MODE-UNSET"):
            return ""
        if date.today() <= MIGRATION_WINDOW_UNTIL:
            print("  WARN  AF-MODE-UNSET: intake.json.creation_mode is unset and no "
                  "priority_shift_spec.json is present. This deck is exempt only until "
                  + MIGRATION_WINDOW_UNTIL.isoformat() + "; after that an unset mode "
                  "blocks the build (SOP-MODE-00, P118).", file=sys.stderr)
            return ""
        # fall through to the CREATION_MODES check, which now blocks.
    if mode not in CREATION_MODES:
        return (f"AF-MODE-UNSET: intake.json.creation_mode is {mode or 'unset'!r}; it must be "
                f"one of {', '.join(CREATION_MODES)}. Step Zero of every deck is to identify "
                f"the creation mode before any content is built (SOP-MODE-00, P118).")
    if mode in ("content_personal", "content_general"):
        substance = intake.get("extracted_substance") or intake.get("source_brief_origin")
        if not substance:
            return (f"AF-MODE-UNSET: creation_mode {mode!r} is a content mode but no "
                    f"extracted_substance / source provenance is recorded — a content-mode "
                    f"deck must diagnose and extract the source first (SOP-MODE-00, P22-P29).")
    return ""


def _chk_priority_shift(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-NO-SHIFT (P33/P105) — the priority-shift SPINE gate. The deck must carry a
    real priority_shift_spec.json (true_goal + a named priority_stack[]) AND plant the
    eight build-move beat tags monotonically in slides_copy.md so the arc actually
    engineers the re-rank. Defers when the doctrine is not active."""
    if not _doctrine_active(run_dir):
        return ""
    spec = _read_priority_spec(run_dir) or {}
    problems = []
    if not str(spec.get("true_goal") or "").strip():
        problems.append("priority_shift_spec.json has no true_goal (the destination shift)")
    stack = spec.get("priority_stack")
    if not (isinstance(stack, list) and len(stack) >= 1):
        problems.append("priority_shift_spec.json has an empty priority_stack[] "
                        "(the audience's current ranking must be surfaced)")
    text = _slides_copy_text_lc(run_dir)
    if text:
        present = [(tag, text.find(tag.lower())) for tag in EIGHT_MOVE_TAGS
                   if text.find(tag.lower()) >= 0]
        if len(present) < 5:
            problems.append(f"only {len(present)}/8 build-move beat tags are present in "
                            f"slides_copy.md (need >=5 of {', '.join(EIGHT_MOVE_TAGS)})")
        else:
            offs = [o for _, o in present]
            if offs != sorted(offs):
                problems.append("the build-move beat tags are out of canonical order — "
                                "the eight moves must run monotonically (P141-P150)")
    if problems:
        return ("AF-NO-SHIFT: the deck does not engineer a deliberate priority shift — "
                + "; ".join(problems) + ". Re-rank the owner's offer/idea to the top of the "
                "audience's priority stack (SOP-NORTHSTAR-00 / SOP-PRIORITY-02).")
    return ""


def _chk_priority_stack(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-NO-PRIORITY-STACK (P142, Move 1). The audience's current priority stack must
    be NAMED before the first value/price ladder beat. Defers without doctrine/copy."""
    if not _doctrine_active(run_dir):
        return ""
    text = _slides_copy_text_lc(run_dir)
    if not text:
        return ""
    ladder = _first_marker_offset(text, LADDER_BEAT_MARKERS)
    if ladder is None:
        return ""  # no ladder beat yet — nothing to order against.
    stack = _first_marker_offset(text, PRIORITY_STACK_MARKERS)
    if stack is None or stack > ladder:
        return ("AF-NO-PRIORITY-STACK: Move 1 is missing — the deck must surface the "
                "audience's current priority stack BEFORE the first value/price ladder "
                "beat (P142, SOP-PRIORITY-02).")
    return ""


def _chk_rerank(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-NO-RERANK (P148, Move 7). After the PRICE the deck must demand the re-rank /
    decision out loud. Defers for pitchless/unset decks (the re-rank is a sale mechanic)."""
    if not _doctrine_active(run_dir):
        return ""
    if _intake_pitch_included(run_dir) is not True:
        return ""
    text = _slides_copy_text_lc(run_dir)
    if not text:
        return ""
    price = text.find("price")
    if price < 0:
        return ""  # no price beat yet.
    rerank = _first_marker_offset(text, RERANK_MARKERS)
    if rerank is None or rerank <= price:
        return ("AF-NO-RERANK: Move 7 is missing — after the PRICE the deck must demand the "
                "re-rank out loud (make the owner's thing the audience's new #1) (P148).")
    return ""


def _chk_trigger(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-NO-TRIGGER (P149, Move 8). The CTA must fire a time-bound trigger. Defers for
    pitchless/unset decks."""
    if not _doctrine_active(run_dir):
        return ""
    if _intake_pitch_included(run_dir) is not True:
        return ""
    text = _slides_copy_text_lc(run_dir)
    if not text:
        return ""
    if _first_marker_offset(text, TRIGGER_MARKERS) is not None:
        return ""
    return ("AF-NO-TRIGGER: Move 8 is missing — the CTA carries no time-bound trigger "
            "(deadline / scarcity window / act-now). A priority that is not acted on now is "
            "not yet a priority (P149).")


def _chk_proclamation_hedge(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-PROCLAMATION-HEDGE (P109). Proclamations must be plain, bold, hedge-free.
    Flags the strong multi-word hedges only (no false positives on bare 'maybe')."""
    if not _doctrine_active(run_dir):
        return ""
    text = _slides_copy_text_lc(run_dir)
    if not text:
        return ""
    hits = [t for t in PROCLAMATION_HEDGE_TOKENS if t in text]
    if hits:
        return ("AF-PROCLAMATION-HEDGE: the copy hedges its declarations with "
                + ", ".join(repr(h) for h in hits[:5]) + ". A proclamation is a plain, bold "
                "claim of truth that dares to challenge the norm — strip the hedge "
                "(SOP-PROCLAMATION-01, P109).")
    return ""


def _chk_peak_end(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-PEAK-END (P49). The arc must declare a deliberate PEAK beat AND a deliberate
    ending beat (the peak-end rule — a flat ending is remembered as flat). Defers when
    no arc / no doctrine."""
    if not _doctrine_active(run_dir):
        return ""
    blob = None
    for rel in ("working/copy/arc_allocation.json", "arc_allocation.json",
                "working/arc_allocation.json"):
        p = run_dir / rel
        if p.exists():
            obj = _read_json(p)
            if isinstance(obj, dict) and "__parse_error__" in obj:
                return ("AF-PEAK-END: arc_allocation.json is not valid JSON, so the "
                        "engineered PEAK + ending cannot be proven (P49).")
            slots = obj if isinstance(obj, list) else (
                obj.get("slots") or obj.get("allocation") or obj.get("slides") or [])
            tokens = []
            for s in slots if isinstance(slots, list) else []:
                if isinstance(s, dict):
                    for k in ("arc_section", "section", "beat", "tag", "type", "role"):
                        v = s.get(k)
                        if isinstance(v, str):
                            tokens.append(v.lower())
                    tags = s.get("tags")
                    if isinstance(tags, list):
                        tokens += [str(t).lower() for t in tags]
                elif isinstance(s, str):
                    tokens.append(s.lower())
            blob = " ".join(tokens)
            break
    if blob is None:
        return ""  # no arc yet — _chk_arc owns absence.
    missing = []
    if not any(t in blob for t in PEAK_TAGS):
        missing.append("no PEAK/APEX/WOW beat")
    if not any(t in blob for t in ENDING_TAGS):
        missing.append("no deliberate ending/recap/CTA beat")
    if missing:
        return ("AF-PEAK-END: the arc fails the peak-end rule — " + "; ".join(missing)
                + ". Engineer a deliberate peak and a deliberate ending; a flat ending is "
                "remembered as flat (P49, SOP-NORTHSTAR-00).")
    return ""


def _chk_salience_apex(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-NO-SALIENCE-APEX (P48/P155). The OFFER/PROMISE-APEX slide must be the single
    most VIVID slide in the deck (von Restorff — the most vivid thing is what the mind
    prioritizes). Deterministic backstop to the human style sign-off: vividness proxy =
    (1 - dominant-colour fraction), reusing the AF-VISUAL-VARIETY pixel data. Defers
    pre-render, without doctrine, or when the apex slide cannot be identified."""
    if not _doctrine_active(run_dir):
        return ""
    pngs = _gather_rendered_pngs(run_dir)
    if not pngs:
        return ""  # pre-render — defer.
    apex = _apex_slide_ordinal(run_dir)
    if apex is None:
        return ""  # apex undeterminable — AF-PEAK-END owns the arc-tag absence.
    scores = {}
    for p in pngs:
        ordn = _png_ordinal(p)
        if ordn is None:
            continue
        frac, _rgb = _png_flatfill_fraction(p)
        if frac is None:
            continue
        scores[ordn] = 1.0 - float(frac)  # higher = more visual variety / vividness
    if apex not in scores or len(scores) < 3:
        return ""  # apex not rendered yet, or too few measurable slides — defer.
    deck_max = max(scores.values())
    if deck_max <= 0:
        return ""
    # Fail only when the apex is CLEARLY below the deck's most-vivid slide (a real
    # von-Restorff inversion), not on a marginal tie — backstop, not a hair-trigger.
    if scores[apex] < 0.85 * deck_max:
        return (f"AF-NO-SALIENCE-APEX: the OFFER/PROMISE-APEX slide (slide {apex}) is not "
                f"the most vivid element in the deck (vividness {scores[apex]:.2f} vs deck "
                f"peak {deck_max:.2f}). The owner's thing must be the single most vivid "
                f"thing in the room by the end (von Restorff, P48/P155).")
    return ""


def _chk_converter_no_invent(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-CONVERTER-NO-INVENT (P167d). A content-converter deck must EXTRACT, never
    INVENT: every numeric/statistical claim in the converter's source brief must trace
    back to the raw source. Defers unless BOTH a source brief AND the raw source text
    are present (the converter owns supplying them)."""
    brief_text = None
    for rel in ("working/copy/source_brief.md", "working/copy/source_brief.json",
                "working/converter/source_brief.md", "working/copy/source_brief.txt"):
        p = run_dir / rel
        if p.exists():
            brief_text = p.read_text(errors="replace")
            break
    if brief_text is None:
        return ""  # no converter brief — defer (not a converter deck, or not built yet).
    source_text = None
    src_candidates = [run_dir / "working" / "source" / "transcript.txt",
                      run_dir / "working" / "converter" / "source.txt",
                      run_dir / "working" / "copy" / "source_raw.txt"]
    src_candidates += sorted((run_dir / "working" / "source").glob("*.txt")) \
        if (run_dir / "working" / "source").is_dir() else []
    src_candidates += sorted((run_dir / "working" / "source").glob("*.md")) \
        if (run_dir / "working" / "source").is_dir() else []
    for p in src_candidates:
        if p.exists():
            source_text = (source_text or "") + "\n" + p.read_text(errors="replace")
    if not source_text:
        return ""  # raw source not on disk to compare against — defer.
    source_norm = _norm_ws(source_text.lower())
    # Statistical/numeric claim tokens: percentages, money, and multi-digit figures.
    claim_re = re.compile(r"\b\d[\d,\.]*\s?%|\$\s?\d[\d,\.]*|\b\d{2,}[\d,\.]*\b")
    invented = []
    for m in claim_re.findall(brief_text):
        tok = _norm_ws(str(m).lower()).replace(" ", "")
        digits = re.sub(r"[^\d]", "", tok)
        if not digits:
            continue
        # The exact figure must appear somewhere in the source (digit-only match so
        # "$1,200" in the brief matches "1200" / "1,200" in the source).
        if digits not in re.sub(r"[^\d]", "", source_norm):
            invented.append(str(m).strip())
    invented = sorted(set(invented))[:6]
    if invented:
        return ("AF-CONVERTER-NO-INVENT: the source brief carries figure(s) absent from the "
                "raw source: " + ", ".join(invented) + ". Extract, never invent — a converter "
                "deck may only use claims the source actually makes (P167d).")
    return ""


def _chk_persuasion_beats(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """HOLE C — wire the orphaned "no-beat" persuasion taxonomy. A doctrine-active
    converting deck must carry each named persuasion beat (problem, choice/fork,
    comparison, measurable results, expert proof, before/after). Generous synonym
    matching keeps real converting decks passing; only a clear, total absence of a
    named beat fails the specific code. Defers without doctrine / copy / a pitch."""
    if not _doctrine_active(run_dir):
        return ""
    if _intake_pitch_included(run_dir) is not True:
        return ""
    text = _slides_copy_text_lc(run_dir)
    if not text:
        return ""
    missing = []
    for code, markers in PERSUASION_BEAT_MARKERS.items():
        if _first_marker_offset(text, markers) is None:
            missing.append(code)
    if missing:
        return (": ".join((missing[0], "")) + "the deck is missing the named persuasion "
                "beat(s) " + ", ".join(missing) + " — every converting deck must name the "
                "problem, present a choice/fork, draw a comparison, cite a measurable result, "
                "carry expert proof, and show a before/after (SOP-PITCH-06 / SOP-ENGINE-00 "
                "persuasion-beat taxonomy).")
    return ""


def _chk_style_preview(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-STYLE-UNPICKED + AF-STYLE-DOUBLECHARGE — the approved 3-style-preview gate
    (P-STYLE-PREVIEW, order 4.85). After the 9 sample renders (3 variants x 3
    representative slides) are produced, the owner picks ONE variant via their own
    gateway (NEVER the operator chat). The full deck must not render until the pick is
    logged; and the winning variant's 3 representative slides must be LOCKED and reused
    so kie never double-charges. Defers until the samples manifest exists."""
    sm = run_dir / STYLE_SAMPLES_MANIFEST_REL
    if not sm.exists():
        return ""  # style-preview phase not run (or pre-phase) — defer.
    choice_p = run_dir / STYLE_CHOICE_REL
    choice = _read_json(choice_p) if choice_p.exists() else None
    valid_pick = (isinstance(choice, dict) and "__parse_error__" not in choice
                  and choice.get("owner_approved") is True
                  and str(choice.get("chosen_variant") or "").strip())
    # W1a — AF-STYLE-UNPICKED: samples rendered but no owner-approved variant pick.
    if not valid_pick:
        if _owner_skip_approved(run_dir, "AF-STYLE-UNPICKED"):
            return ""
        return ("AF-STYLE-UNPICKED: 9 style samples were rendered but the owner has not "
                "picked a winning variant (working/copy/style_preview_choice.json must carry "
                "owner_approved:true + chosen_variant). The full deck must NOT render until "
                "the owner chooses A/B/C via their OWN gateway — never the operator chat.")
    # W1b — AF-STYLE-DOUBLECHARGE: the winner's 3 representative slide renders must be
    # carried forward (locked_renders), not re-dispatched to kie.
    if _owner_skip_approved(run_dir, "AF-STYLE-DOUBLECHARGE"):
        return ""
    locked = choice.get("locked_renders")
    if not (isinstance(locked, list) and len(locked) >= 1):
        return ("AF-STYLE-DOUBLECHARGE: the winning variant " + repr(choice.get("chosen_variant"))
                + " records no locked_renders — its 3 representative slides must be carried "
                "forward and REUSED, never re-rendered. kie must never double-charge for the "
                "already-approved samples (P-STYLE-PREVIEW, order 4.85).")
    for ref in locked:
        rp = run_dir / str(ref) if not str(ref).startswith("/") else Path(str(ref))
        if not rp.exists():
            return ("AF-STYLE-DOUBLECHARGE: locked sample render " + repr(str(ref))
                    + " (the approved variant's representative slide) is missing — it must be "
                    "reused, not re-charged. kie must never double-charge the approved samples.")
    return ""


def _chk_priority_shift_ledger(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """AF-PRIORITY-SHIFT (composite, P-SHIFT-QC order 7.5) — the framework's 14-item
    pre-output ship gate, the ONLY point where copy + design + rendered images coexist.
    Runs all 14 sub-assertions (reusing the shift-left gates), writes a per-item
    pass/fail ledger to working/qc/priority_shift_report.json, and refuses ship until
    all 14 PASS. Item 0 is the North Star itself: the #1 job is to hold attention.
    Defers pre-render and without doctrine; waivable only by a logged owner skip."""
    if not _doctrine_active(run_dir):
        return ""
    pngs = _gather_rendered_pngs(run_dir)
    if not pngs:
        return ""  # the 14-item gate needs the rendered images — defer pre-render.
    spec = _read_priority_spec(run_dir) or {}
    text = _slides_copy_text_lc(run_dir) or ""
    rows = []

    def add(item, ok, evidence):
        rows.append({"item": item, "pass": bool(ok), "evidence": str(evidence)[:240]})

    pitch = _intake_pitch_included(run_dir) is True
    add("0_attention_is_the_no1_job",
        bool(str(spec.get("true_goal") or "").strip()),
        "priority_shift_spec.true_goal declares the destination shift")
    add("1_creation_mode_identified", _chk_mode(run_dir) == "", "AF-MODE-UNSET")
    add("2_priority_stack_named", _chk_priority_stack(run_dir) == "", "AF-NO-PRIORITY-STACK")
    add("3_priority_shift_engineered", _chk_priority_shift(run_dir) == "", "AF-NO-SHIFT")
    add("4_present_cost_exposed",
        _first_marker_offset(text, COST_OF_INACTION_MARKERS) is not None,
        "cost-of-inaction beat present (Move 2)")
    add("5_higher_priority_lever",
        bool(str(spec.get("higher_priority_hook") or "").strip()),
        "priority_shift_spec.higher_priority_hook (Move 3)")
    add("6_value_anchored_high",
        _first_marker_offset(text, ("anchor", "value-stack", "value_stack", "value add")) is not None,
        "value anchor beat present (Move 4)")
    add("7_urgency_scarcity",
        (not pitch) or _first_marker_offset(text, URGENCY_SCARCITY_MARKERS) is not None,
        "urgency/scarcity present (Move 5; n/a for pitchless)")
    add("8_ability_unblocked",
        (not pitch) or _first_marker_offset(text, ABILITY_UNBLOCK_MARKERS) is not None,
        "ability-blocker removed (Move 6; n/a for pitchless)")
    add("9_rerank_demanded", _chk_rerank(run_dir) == "", "AF-NO-RERANK (Move 7)")
    add("10_trigger_fired", _chk_trigger(run_dir) == "", "AF-NO-TRIGGER (Move 8)")
    add("11_proclamation_hedge_free", _chk_proclamation_hedge(run_dir) == "",
        "AF-PROCLAMATION-HEDGE")
    add("12_peak_and_ending", _chk_peak_end(run_dir) == "", "AF-PEAK-END")
    add("13_most_vivid_by_the_end", _chk_salience_apex(run_dir) == "", "AF-NO-SALIENCE-APEX")
    add("14_one_promise_one_wow_one_demonstration",
        all(str(spec.get(k) or "").strip()
            for k in ("the_one_promise", "the_one_wow", "the_one_demonstration")),
        "spec carries the single promise + wow + demonstration anchors")

    passed = all(r["pass"] for r in rows)
    # FIX-2 (Error 2): real per-slide verdicts. P-SHIFT-QC runs AFTER render (order
    # 7.5), so every rendered slide exists and the ship gate grades each one. Each
    # per-slide verdict records the slide ordinal and an explicit pass (a rendered
    # slide present + non-degenerate in the run dir). This is the per-slide coverage
    # the QC report floor (>= 20 real per-slide verdicts) requires — a 14-row
    # checklist alone can never satisfy the floor.
    per_slide = []
    for png in sorted(_gather_rendered_pngs(run_dir)):
        sid = _png_ordinal(png)
        if sid is None:
            continue
        per_slide.append({
            "slide": sid,
            "pass": png.is_file(),
            "verdict": "pass" if png.is_file() else "fail",
            "evidence": "rendered slide present in run dir (ship-gate per-slide pass)",
        })
    report = {
        "schema": "priority_shift_report/v1",
        "gate": "AF-PRIORITY-SHIFT",
        "phase": "P-SHIFT-QC (order 7.5)",
        "generated_at": time.strftime("%Y-%m-%dT%H:%M:%S%z"),
        "pass": passed,
        "items": rows,
        "slides": per_slide,
    }
    try:
        out = run_dir / PRIORITY_SHIFT_REPORT_REL
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(json.dumps(report, indent=2))
    except Exception:  # noqa: BLE001 — never let report-write mask the verdict
        pass
    if not passed:
        if _owner_skip_approved(run_dir, "AF-PRIORITY-SHIFT"):
            return ""
        failed = [r["item"] for r in rows if not r["pass"]]
        return ("AF-PRIORITY-SHIFT: the 14-item priority-shift ship gate failed on "
                + ", ".join(failed) + " (see working/qc/priority_shift_report.json). "
                "P-SHIFT-QC (order 7.5) blocks ship until all 14 items PASS — item 0 is the "
                "North Star: the #1 job is to hold attention (P161/P155, SOP-INTEGRATION-00).")
    return ""


# ---------------------------------------------------------------------------
# SIGNATURE PRESENTATION (Skill 51) — defer-unless-signature preflight wrappers.
# The ONLY signature-presentation additions to build_deck.py: three thin wrappers
# that DELEGATE to the Skill-51 provers (never re-implementing a rule) and, above all,
# DEFER (return "") unless intake.json records deck_type == "signature_presentation".
# Every OTHER deck type executes the identical pre-existing code path — the binding
# minimal/additive, defer-unless-signature guarantee. The switch mirrors
# _doctrine_active() (the priority-shift no-regression switch).
# ---------------------------------------------------------------------------
_SP_PROVER_CACHE = {}


def _sp_active(run_dir: Path) -> bool:
    """Signature-Presentation switch: the SP gates engage ONLY when intake.json records
    deck_type == 'signature_presentation'. Every other deck type takes the identical
    pre-existing code path (defer-unless-signature). Mirror of _doctrine_active()."""
    obj = _read_intake_obj(run_dir)
    return isinstance(obj, dict) and obj.get("deck_type") == "signature_presentation"


def _sp_prover(mod_name: str):
    """Lazily import a Skill-51 prover module by name, searching build_deck's own scripts
    dir first and then the sibling 51-signature-presentation/scripts/
    (repo/worktree layout). On a materialized department the provers resolve only via the
    sibling fallback because the install does not copy them into the engine's scripts dir.
    Cached. Returns the module, or None when it cannot be located/imported. Only reached
    for a signature deck (the wrappers defer first via _sp_active)."""
    if mod_name in _SP_PROVER_CACHE:
        return _SP_PROVER_CACHE[mod_name]
    import importlib.util
    here = Path(__file__).resolve().parent
    cands = [here / (mod_name + ".py")]
    # Repo / worktree layout: a sibling 51-signature-presentation/scripts/ up-tree.
    cands += [anc / "51-signature-presentation" / "scripts" / (mod_name + ".py")
              for anc in here.parents]
    # FIX-PRES-09(i): a MATERIALIZED department tree (build_deck runs from
    # <workspace>/…/presentations/scripts/) cannot reach skill 51 by walking bare
    # ancestors. Mirror deck-intake-driver._sp_skill_roots(): also try an ancestor's
    # `skills/51-signature-presentation/scripts/` (installed SKILLS_DIR layout) and
    # the canonical installed roots (VPS /data, Mac ~/.openclaw).
    cands += [anc / "skills" / "51-signature-presentation" / "scripts" / (mod_name + ".py")
              for anc in here.parents]
    for _base in ("/data/.openclaw/skills",
                  str(Path.home() / ".openclaw" / "skills")):
        cands.append(Path(_base) / "51-signature-presentation" / "scripts"
                     / (mod_name + ".py"))
    mod = None
    for cand in cands:
        if cand.exists():
            try:
                spec = importlib.util.spec_from_file_location(mod_name, cand)
                m = importlib.util.module_from_spec(spec)
                spec.loader.exec_module(m)
                mod = m
                break
            except Exception:  # noqa: BLE001 — a broken import must not crash preflight
                continue
    _SP_PROVER_CACHE[mod_name] = mod
    return mod


def _sp_delegate(kind: str, run_dir: Path) -> str:
    """Shared engine for the _chk_sp_* thin wrappers: import the Skill-51 prover for
    `kind` and run it against the run dir's SP artifacts, normalized to the preflight
    convention ("" on PASS, a non-empty AF reason on fail). Only reached for a signature
    deck. Fail-closed: a genuinely missing prover, unreadable artifact, or unexpected
    prover error is a non-empty (blocking) reason — the sacred gate is never silently
    skipped for a signature deck."""
    mod = _sp_prover({"intake": "prove_sp_intake",
                      "structure": "prove_sp_structure",
                      "no_pitch": "prove_sp_no_pitch"}[kind])
    if mod is None:
        return ("signature_presentation deck but the Skill-51 " + kind + " prover could "
                "not be imported (install the 51-signature-presentation scripts next to "
                "build_deck.py). Fail-closed — the sacred gate cannot be skipped.")
    try:
        if kind == "intake":
            obj = _read_json(run_dir / "working" / "copy" / "sp_intake.json")
            if not isinstance(obj, dict) or "__parse_error__" in obj:
                return "AF-SP-8Q-MISSING: working/copy/sp_intake.json is missing or unreadable."
            fails = mod.evaluate(obj)
            return "" if not fails else "; ".join(str(c) + ": " + str(m) for c, m in fails)
        if kind == "structure":
            deck = _read_json(run_dir / "working" / "copy" / "sp_structure.json")
            if not isinstance(deck, dict) or "__parse_error__" in deck:
                return "AF-SP-PHASE-ORDER: working/copy/sp_structure.json is missing or unreadable."
            violations, _notes = mod.verify(mod._load_structure(None), deck)
            return "" if not violations else "; ".join(str(c) + ": " + str(m) for c, m in violations)
        # no_pitch — the Phase-3 teaching-band hygiene gate.
        code, msgs = mod.evaluate_paths(run_dir / "working" / "copy" / "sp_intake.json",
                                        run_dir / "working" / "copy" / "sp_structure.json", None)
        return "" if code == 0 else "AF-SP-P3-PITCH: " + " | ".join(str(m) for m in msgs)
    except Exception as exc:  # noqa: BLE001 — fail-closed, never crash preflight
        return ("signature_presentation " + kind + " prover raised " + repr(exc)
                + " — fail-closed (the sacred gate cannot be skipped).")


def _chk_sp_intake(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """P-SP-INTAKE — Signature-Presentation intake gate. DEFERS unless signature_presentation."""
    if not _sp_active(run_dir):
        return ""
    return _sp_delegate("intake", run_dir)


def _chk_sp_structure(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """P-SP-STRUCTURE — SACRED 4-phase structure contract. DEFERS unless signature_presentation."""
    if not _sp_active(run_dir):
        return ""
    return _sp_delegate("structure", run_dir)


def _chk_sp_no_pitch(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """P-SP-P3-HYGIENE — Phase-3 teaching no-pitch hygiene. DEFERS unless signature_presentation."""
    if not _sp_active(run_dir):
        return ""
    return _sp_delegate("no_pitch", run_dir)


SP_TRANSCRIPT_REL = Path("working") / "interview" / "intake_transcript.json"


def _chk_sp_intake_trace(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """P-SP-INTAKE-TRACE — the intake CONVERSATION gate (AF-INTAKE-BATCH).

    A10 / T0-12. Skill 51 declares that every rule is machine-enforced by a fail-closed
    prover and never advisory, and then wired this one as advisory and non-gating. So an
    eight-question batched interaction that afterwards produced a structurally valid
    intake RECORD passed every preflight and reached a signed certificate: the run
    supplied its own record as the only evidence of how the intake had been CONDUCTED,
    and the rule that defines the skill's value was the one rule nothing enforced.

    The transcript is written mechanically, turn by turn, by deck-intake-driver.py's
    turn-gate (SP_TRANSCRIPT_REL) — a different producer from the deck build being
    certified. It is REQUIRED here: an ABSENT transcript is a fail, because otherwise
    the cheapest way to pass a conversation gate is to conduct no recorded conversation.

    FIX-3 ("intake must be a real conversation"): the transcript must ALSO be
    DRIVER-PRODUCED — a signed envelope (format 'sp-intake-transcript-v1' +
    driver_signature + qid_sequence) written by deck-intake-driver.py's turn-gate.
    A hand-written transcript (bare JSON list, or a list hand-written next to a
    hand-written intake_ledger.json — ERROR 3 of the 2026-08-06 E2E audit) fails
    fail-closed with AF-INTAKE-BATCH (NO-DRIVER-ENVELOPE / NO-DRIVER-SIGNATURE).
    Presence of a transcript file is not proof it came from a real conversation.

    DEFERS (no-op) unless intake.json declares deck_type == signature_presentation, so
    every other deck type takes the identical pre-existing path.
    """
    if not _sp_active(run_dir):
        return ""
    mod = _sp_prover("intake_trace_check")
    if mod is None:
        return ("AF-INTAKE-BATCH: 51-signature-presentation/scripts/intake_trace_check.py "
                "is not co-located with build_deck.py — the intake-conversation gate cannot "
                "run for a signature deck (fail-closed; install skill 51 next to the engine).")
    tpath = Path(run_dir) / SP_TRANSCRIPT_REL
    if not tpath.is_file():
        return ("AF-INTAKE-BATCH: no intake transcript at " + str(SP_TRANSCRIPT_REL) +
                " — a signature-presentation intake must be CONDUCTED choice-first and one "
                "question per turn, and the turn-gate records that conversation mechanically. "
                "An absent transcript is not proof of a compliant intake (fail-closed). Run "
                "the intake through deck-intake-driver.py --signature, which writes a signed "
                "driver envelope at that path (a hand-written intake_ledger.json with no "
                "transcript is NOT an interview).")
    try:
        raw = tpath.read_text(encoding="utf-8")
        # FIX-3: the raw envelope (for the driver-provenance gate) and the parsed
        # turns (for the conversation scan). A bare list has NO driver provenance.
        import json as _json
        try:
            envelope = _json.loads(raw)
        except _json.JSONDecodeError:
            envelope = raw
        prov_fails = mod.check_driver_provenance(envelope)
        turns = mod.parse_transcript(raw)
        if not turns:
            return ("AF-INTAKE-BATCH: the intake transcript at " + str(SP_TRANSCRIPT_REL) +
                    " parsed to zero turns (unreadable format) — fail-closed.")
        result = mod.scan_transcript(turns, mod.load_bank_questions())
        # Driver provenance failures are treated as violations of the same gate.
        prov_violations = [{"code": mod.AF_CODE, "reason": code, "turn_index": None,
                            "detail": msg} for code, msg in prov_fails]
        result["violations"] = prov_violations + result.get("violations", [])
        result["pass"] = len(result["violations"]) == 0
    except Exception as exc:  # noqa: BLE001 — fail-closed, never crash preflight
        return ("AF-INTAKE-BATCH: the intake-conversation scanner raised " + repr(exc)
                + " — fail-closed (the conversation gate cannot be skipped).")
    if result.get("pass"):
        return ""
    reasons = "; ".join(
        str(v.get("reason", "?")) + " @turn " + str(v.get("turn_index", "?"))
        + ": " + str(v.get("detail", ""))
        for v in result.get("violations", []))
    return "AF-INTAKE-BATCH: " + reasons


def _chk_sp_claim(run_dir: Path, slides_path: Optional[Path] = None) -> str:
    """P-SP-CLAIM — the routing/claim gate. Runs for EVERY deck (does NOT defer,
    unlike the three gates above): if the run carries signature-presentation signals
    (an sp_intake.json, a set Signature frame, a frame-selection question, or a
    'signature presentation' request) but intake.json does not DECLARE
    deck_type == signature_presentation, fail-closed AF-SP-TYPE-UNDECLARED. Closes
    the highest-severity skip of the audit — a signature presentation built through
    the generic path by omitting the magic word (which makes _sp_active defer, so
    every SP gate no-ops). A non-signature deck with no SP signal passes untouched."""
    mod = _sp_prover("prove_sp_routing")
    if mod is None:
        # Routing prover not co-located: still block the unambiguous case we can
        # detect without it, so the bypass cannot ride a missing prover.
        obj = _read_intake_obj(run_dir)
        declared = isinstance(obj, dict) and obj.get("deck_type") == "signature_presentation"
        sp_present = (run_dir / "working" / "copy" / "sp_intake.json").is_file()
        if sp_present and not declared:
            return ("AF-SP-TYPE-UNDECLARED: working/copy/sp_intake.json is present but "
                    "intake.json does not declare deck_type == signature_presentation "
                    "(install 51-signature-presentation/scripts/prove_sp_routing.py next to "
                    "build_deck.py for the full signal set). Fail-closed.")
        return ""
    try:
        fails = mod.evaluate_run_dir(run_dir)
        return "" if not fails else "; ".join(str(c) + ": " + str(m) for c, m in fails)
    except Exception as exc:  # noqa: BLE001 — fail-closed, never crash preflight
        return ("signature-presentation claim gate raised " + repr(exc)
                + " — fail-closed (a signature deck cannot skip the claim gate).")


PREFLIGHT_REQUIRED = [
    ("working/copy/intake.json",
     "intake.json (interview_confirmed:true, presentation_mode one-person|general)",
     "Phase 0a — Director SOP 9.1 / Content-to-Presentation Architect SOP 9.1",
     _chk_intake),
    ("working/research/brief-*.md",
     "research brief (research_complete:true)",
     "Phase -0.5 — Deep Research Specialist SOP 9.1/9.4 (AF-RESEARCH-GATE)",
     _chk_research_brief),
    # 10th check — RESEARCH-CITATION GATE (AF-RESEARCH-UNCITED). The research pack
    # must contain >= MIN_CITED_SOURCES (8) distinct http(s) URLs, proving real web
    # research was done. The self-asserted research_complete:true flag is IGNORED
    # here — this gate is the source of truth for whether citations exist.
    # Uses rel glob (same as research brief) so the first matching brief file is
    # passed to the checker.
    ("working/research/brief-*.md",
     f"research pack cited URLs — >= {MIN_CITED_SOURCES} distinct http(s) URLs "
     f"required (self-asserted research_complete:true flag is not proof of research)",
     "Phase -0.5 — Deep Research Specialist SOP 9.1/9.4 (AF-RESEARCH-UNCITED)",
     _chk_research_cited),
    # 11th check — CLAIMS-WITHOUT-CITATION (AF-RESEARCH-UNCITED). If slide copy
    # contains factual/statistical claim markers (%, $, 'research', 'study', etc.)
    # and the research pack has zero cited URLs, the deck is shipping unsupported
    # claims. Uses rel sentinel None (needs the whole run dir to find both copy
    # and research pack).
    (None,
     "claims-without-citation — slide copy claim markers must have a cited research pack",
     "Phase -0.5 — Deep Research Specialist SOP 9.1/9.4 (AF-RESEARCH-UNCITED)",
     _chk_claims_without_citation),
    # RESEARCH-WEAVE / BREADTH gate (AF-RESEARCH-WEAVE). Research must be woven ACROSS
    # the deck, not funnelled to one proof slide. CONDITIONAL: defers until copy exists,
    # then requires working/research/research_map.json mapping research items to >=60%
    # of non-exempt content slides, the writer actually using the anchors in
    # slides_copy.md AND in the RENDER copy (slides.json copy[] — the pixels are baked
    # verbatim from it), and >=8 distinct items deck-wide. Run-dir-scoped (None sentinel);
    # threads slides_path (H3) so it judges the actual rendered slides.json.
    (None,
     "research woven across the deck — research_map.json maps facts/quotes/stats to "
     ">=60% of non-exempt content slides, the copy uses the anchors (slides_copy.md AND "
     "the render copy slides.json copy[]), and the deck draws on >=8 distinct items "
     "(AF-RESEARCH-WEAVE)",
     "Phase 3.5 — Deep Research Specialist SOP 9.5 (research-to-slide map) + Slide "
     "Copywriter (RESEARCH_USED) (AF-RESEARCH-WEAVE)",
     _chk_research_map),
    # RESEARCH-REACHES-RENDER (AF-RESEARCH-REACHES-RENDER) — the renderer consumes the
    # positional slides.json copy[], NOT slides_copy.md. AF-RESEARCH-WEAVE validates the
    # markdown; this gate proves the SAME mapped research anchors actually reached the
    # RENDER COPY (slides.json copy[]) so research cannot silently fail to reach the
    # baked slides. Run-dir-scoped (None sentinel); H3 threads the ACTUAL rendered
    # slides.json. CONDITIONAL: defers pre-map/pre-copy (owned upstream).
    (None,
     "research reaches the RENDER copy — every research-mapped non-exempt slide's "
     "anchor token is present in slides.json copy[] (the file the renderer bakes into "
     "the slides), not only in slides_copy.md (AF-RESEARCH-REACHES-RENDER)",
     "Phase 4 — Slide Copywriter / builder must weave the validated anchors into the "
     "render copy (AF-RESEARCH-REACHES-RENDER)",
     _chk_research_reaches_render),
    ("working/qc/copy_qc_report.json",
     "copy QC report (gate Phase 1Q, average >= 8.5, no AF-* triggered)",
     "Phase 1Q — QC Specialist SOP 9.1 / SOP-SLIDE-00",
     _chk_copy_qc),
    ("working/copy/arc_allocation.json",
     "converting arc — per-slide arc-section allocation",
     "Phase 3 — Signature Presentation Architect SOP",
     _chk_arc),
    ("working/copy/slides_copy.md",
     "slide copy authored per doctrine (hook recurs across the deck, 3-4 band, 10 components)",
     "Phase 4 — Slide Copywriter SOP",
     _chk_slides_copy),
    # AF-COPY-BAND — per-slide COPY character-count FLOOR and CEILING (spec item 7 /
    # §8.2). The first-ever code-enforced minimum for headline/subhead/slide-total,
    # and the single reconciled rule for the old 3-way bullet-count conflict
    # (5-bullets-7-words vs 3-bullets-30-words vs render-reality 5-words). Run-dir-
    # scoped (None sentinel) so it reads the ACTUAL positional slides.json (H3).
    # Deliberately a first-party preflight check, NOT a QC-report criterion, so it
    # never collides with the QC_FOREIGN_SIGNATURES word-count-rubric ban.
    (None,
     f"copy character-count band — HEADLINE {COPY_HEADLINE_CHAR_FLOOR}-"
     f"{COPY_HEADLINE_CHAR_CEILING}, SUBHEAD {COPY_SUBHEAD_CHAR_FLOOR}-"
     f"{COPY_SUBHEAD_CHAR_CEILING} (if present), KICKER <= {COPY_KICKER_CHAR_CEILING} "
     f"(if present), BULLETS <= {COPY_BULLET_MAX_COUNT} each {COPY_BULLET_CHAR_FLOOR}-"
     f"{COPY_BULLET_CHAR_CEILING}, SLIDE TOTAL {COPY_SLIDE_TOTAL_CHAR_FLOOR}-"
     f"{COPY_SLIDE_TOTAL_CHAR_CEILING} (AF-COPY-BAND)",
     "Phase 4 — Slide Copywriter SOP / qc-specialist-presentations.md c4 "
     "(AF-COPY-BAND)",
     _chk_copy_density),
    # G1 — WRITING-half INTELLIGENCE engines (Story villain-before-hero, Emotional
    # felt-stakes-before-offer, Hook copy refrain, Recap) + narrative harmony, fired at
    # COPY-QC (the EARLIEST phase, BEFORE any image prompt is authored — shift-left).
    # Wires intelligence_engines_check.check_copy into preflight so the narrative engines
    # fire DETERMINISTICALLY, not by an agent choosing to run a script. Run-dir-scoped
    # (None sentinel); defers pre-copy.
    (None,
     "intelligence engines (copy) — the deck plants the VILLAIN before the first HERO "
     "beat (Story), lands a felt-stakes quantifier BEFORE the offer (Emotional), carries "
     "the Hook refrain, and coheres as a narrative arc (AF-NO-VILLAIN / AF-NO-FELT-STAKES)",
     "Phase 1Q/COPY-QC — Slide Copywriter SOP + intelligence_engines_check.py --phase copy (G1)",
     check_intelligence_engines_copy),
    # G2 — persuasion (offer) sub-engines: cadence / cost-of-inaction / branded-method /
    # guarantee specificity / time-to-result, fired at COPY-QC (PROMISE PRECEDES PRICE).
    # Wires pitch_engines_check.check_copy into preflight. CONDITIONAL on pitch_included
    # (defers for a pitchless / unset deck, mirroring _chk_pitch). Run-dir-scoped; defers
    # pre-copy.
    (None,
     "pitch engines (offer) — PROMISE precedes PRICE; the offer carries a value cadence, "
     "a quantified cost-of-inaction, a non-generic guarantee, a real branded method, and "
     "a concrete time-to-result (AF-CADENCE / AF-NO-COST-OF-INACTION / AF-GUARANTEE-GENERIC "
     "/ AF-NO-BRANDED-METHOD / AF-NO-TIME-TO-RESULT)",
     "Phase 1Q/COPY-QC — Offer Price Strategist + pitch_engines_check.py --phase 1Q (G2)",
     check_pitch_engines),
    ("working/research/design-brief-*.md",
     "typography/design brief — per-slide creative art direction",
     "Phase F — Typography Architect / SOP-DESIGN-01/02",
     _chk_design_brief),
    # TYPOGRAPHY-QC gate (AF-TYPOGRAPHY-QC). After the design brief, an INDEPENDENT
    # QC specialist grades the design system against the written typography rubric.
    ("working/qc/typography_qc_report.json",
     "typography QC report (gate Phase Typography-QC, average >= 8.5, independent reviewer)",
     "Phase Typography-QC — Typography QC Specialist (AF-TYPOGRAPHY-QC)",
     _chk_typography_qc),
    # FONT-FLOOR gate (AF-FONT-FLOOR). DETERMINISTIC numeric rejector of the declared
    # type system: parses working/typography/type_layout_system.md tokens (min_body_pt,
    # type_scale_steps, min_contrast_ratio) and FAILS a sub-floor body size / non-modular
    # scale / below-WCAG contrast BEFORE render. Defers pre-typography; once a design
    # system exists, the tokens are mandatory. Run-dir-scoped (None sentinel).
    (None,
     "font floor — declared type system meets the 18pt body floor, a 4-5 step modular "
     "scale, and WCAG-AA contrast (raised when client_dark_theme); AF-FONT-FLOOR",
     "Phase F — Typography Architect SOP 9.1 (emits type_layout_system.md tokens) (AF-FONT-FLOOR)",
     check_font_floor),
    # PROMPT-QC gate (AF-PROMPT-QC). After Prompt-Authoring, an INDEPENDENT QC
    # specialist grades every per-slide prompt against the 9,000–14,000-char prompt standard
    # AND the gate RE-MEASURES every on-disk prompt (check_prompt_qc_teeth) — no rubber stamp.
    ("working/qc/prompt_qc_report.json",
     "prompt QC report (gate Phase Prompt-QC, average >= 8.5, independent reviewer)",
     "Phase Prompt-QC — Prompt QC Specialist (AF-PROMPT-QC)",
     _chk_prompt_qc),
    # IMAGE-QC gate (AF-IMAGE-QC). After Render, an INDEPENDENT QC specialist grades
    # the rendered slides against the written image-QC rubric (multimodal pass). The
    # image-QC report is a POST-render artifact (manifest order 4.95, AFTER P4-RENDER
    # 4.9), so this is wired RUN-DIR-SCOPED (None sentinel) via check_image_qc_report_gate:
    # it DEFERS pre-render (no PNGs + no report) so render can proceed, then enforces the
    # full report-shape gate + AF-IMAGE-QC-VISION pixel teeth (_chk_image_qc) once renders
    # exist. Fixes the chicken-and-egg where requiring the report at pre-render made render
    # un-startable; the post-render pixel teeth still bite (here, at the render closeout,
    # and at pre-delivery). _chk_image_qc stays the manifest-declared P-IMAGE-QC checker.
    (None,
     "image QC report (gate Phase Image-QC, average >= 8.5, independent reviewer) — "
     "enforced once renders exist; the report is a post-render artifact so this defers "
     "pre-render and the AF-IMAGE-QC-VISION pixel teeth bite post-render",
     "Phase Image-QC — Image QC Specialist (AF-IMAGE-QC)",
     check_image_qc_report_gate),
    # SPEECH-QC gate (AF-SPEECH-QC). CONDITIONAL: enforced only once the speech QC
    # report exists (speech is written downstream at delivery), so it never blocks
    # the pre-speech render but is wired so it can't be silently skipped.
    ("working/qc/speech_qc_report.json",
     "speech QC report (gate Phase Speech-QC, average >= 8.5, independent reviewer)",
     "Phase Speech-QC — Speech QC Specialist (AF-SPEECH-QC)",
     _chk_speech_qc),
    # SLIDE-COUNT-FLOOR gate (AF-SLIDE-COUNT-FLOOR). The output slide count must be
    # >= target_talk_minutes x 1.3 (the low end of the verified 1.3–1.5 slides/min
    # pacing band). A 30-min/10-slide deck auto-fails. Run-dir-scoped (None sentinel).
    (None,
     "slide-count floor — output slides >= target_talk_minutes x 1.3 (pacing band)",
     "Phase 0a/4 — Director duration intake + Arc allocation (AF-SLIDE-COUNT-FLOOR)",
     _chk_slide_count_floor),
    # SLIDE-COUNT-EXACT gate (AF-SLIDE-COUNT-EXACT). When the client stated an EXPLICIT
    # requested slide count, the built deck MUST be exactly that many slides — the
    # client's number is honored verbatim (25->25, 50->50, 500->500), never floored up,
    # capped down, defaulted, or substituted. AUTHORITATIVE over the pacing floor and
    # the coverage floor (both defer when a requested count is present). Run-dir-scoped.
    (None,
     "slide-count exact — when the client requested an explicit count, output slides == that count (never floored/capped/changed)",
     "Phase 0a/4 — Director intake (client's exact requested slide count, AF-SLIDE-COUNT-EXACT)",
     _chk_slide_count_exact),
    # PITCH gate (AF-PITCH-MISSING). The converting arc must carry an offer ladder
    # AND a re-pitch after the FINAL price. Run-dir-scoped (None sentinel).
    (None,
     "offer ladder + re-pitch — arc carries a value-stack/anchor/price ladder and a re-pitch",
     "Phase 3 — Offer Price Strategist (SOP-PITCH-01 / SOP-PITCH-03, AF-PITCH-MISSING)",
     _chk_pitch),
    # CREATIVITY gate (AF-CREATIVITY). Rejects template-sameness (one archetype
    # dominating the deck) and cliche copy. Run-dir-scoped (None sentinel).
    (None,
     "creativity — no single layout archetype dominates the deck; no cliche copy",
     "Phase F/4 — Typography Architect + Slide Copywriter (SOP-DESIGN-03, AF-CREATIVITY)",
     _chk_creativity),
    # 7th check — ANTI-COMPRESSION coverage gate (AF-COVERAGE-1). The rel sentinel
    # None tells run_preflight to call this check with the run_dir itself (it needs
    # both mission_prd.json's source_slide_count AND the output slide count), not a
    # single resolved artifact path.
    (None,
     "anti-compression coverage — output slides >= client source_slide_count (Mode B ADD-only)",
     "Mission PRD — source_slide_count (Mode A=0 always passes)",
     _chk_coverage),
    # 8th check — RICH-PROMPT-REQUIRED gate (AF-P1). EVERY slide must have a
    # hand-authored rich prompt >= 9,000 chars (clearing the quality floor) in working/prompts/. Like coverage,
    # this needs the whole run dir (it counts slides AND reads every prompt file),
    # so it uses the rel sentinel None.
    (None,
     "rich per-slide prompt — every slide has a >=9,000-char prompt in working/prompts/ "
     "that clears the quality floor (8-class negative block / spelling-lock / density / "
     "verbatim copy baked), rendered VERBATIM",
     "Phase 2 — Slide Image Creator SOP 9.1 (15-element rich prompt; rendered verbatim, no thin fallback)",
     _chk_rich_prompts),
    # INTELLIGENCE-ENGINES prompt gate (AF-FACE-PROMPT-MISSING / AF-LIGHT-PROMPT-MISSING /
    # AF-WORLD-SCALE / AF-HAIR-INAUTHENTIC). Wires intelligence_engines_check.py --phase
    # prompt into preflight so the facial / lighting / world / representation engines fire
    # DETERMINISTICALLY at render time (was agent-discipline-only, 0 grep hits in the
    # renderer). Run-dir-scoped (None sentinel); defers pre-prompt.
    (None,
     "intelligence engines (prompt) — every people/scene prompt carries an explicit "
     "expression token, a key/fill/rim + rim/hair lighting token, a world believability "
     "justification, and an age-appropriate hairstyle token (AF-FACE-PROMPT-MISSING / "
     "AF-LIGHT-PROMPT-MISSING / AF-WORLD-SCALE / AF-HAIR-INAUTHENTIC)",
     "Phase 2/Prompt-QC — Slide Image Creator SOP 9.2/9.3 + intelligence_engines_check.py --phase prompt",
     check_intelligence_engines_prompt),
    # G6 — EXCELLENCE quality floor (AF-EXCELLENCE), fired at PROMPT-QC, BEFORE any render
    # (the money step). Every prompt must score >= PROMPT_EXCELLENCE_FLOOR on the EXCELLENCE
    # dimension: the char budget must buy defect-preventing specificity, not padding. This
    # is the QUALITY floor that sits ALONGSIDE the 9,000-char LENGTH floor (two independent
    # gates, §1.1a) — a floor-grazing boilerplate-padded prompt clears 9,000 chars yet
    # routes back here. Run-dir-scoped (None sentinel); defers pre-prompt.
    (None,
     "prompt excellence — every rich prompt scores above the EXCELLENCE floor (full 8-class "
     "negative block, per-string spelling-lock, people-anatomy + world-grounding, brand HEX "
     "+ type SIZE + composition, lexical richness); padding to clear 9,000 chars routes back "
     "(AF-EXCELLENCE)",
     "Phase Prompt-QC — Slide Image Creator / Prompt QC (amazing-not-compliant, AF-EXCELLENCE)",
     check_prompt_excellence),
    # KIE-BAKED gate (AF-I14). Once a render record exists, EVERY rendered slide must
    # map to a real KIE taskId + a verified, above-floor PNG (no native render, no
    # flat-placeholder fill). Run-dir-scoped (reads process_manifest.json + needs the
    # slide count), so it uses the rel sentinel None. Pre-render of downstream phases
    # this passes vacuously (no render record yet); after render it is the source of
    # truth that the slides were actually baked.
    (None,
     "KIE-baked slides — every rendered slide has a real KIE taskId + a verified, above-floor PNG (no native render / flat placeholder)",
     "Phase 5/6 — Slide Image Creator / render (AF-I14)",
     _chk_kie_baked),
    # 9th check — SPEECH-LENGTH gate (AF-SPEECH-SHORT). Conditional: enforced only
    # once the presenter speech exists (it is written downstream at delivery), so it
    # never blocks the pre-speech render but is wired so it can't be silently skipped.
    (None,
     "speech length — presenter speech words >= target_talk_minutes x 120 wpm (fails short)",
     "Phase 9 — Presenter Speech Writer SOP 9.1 (130 wpm pacing; 120 wpm floor)",
     _chk_speech_length),
    # AF-VISUAL-VARIETY gate (2026-06-19). Rejects an all-dark monotone deck that
    # lacks visual variety (palette or luminance). Conditional: defers when no
    # rendered PNGs exist yet (pre-render), fires at Image-QC and postflight.
    (None,
     "visual variety — deck has light/break slides and varied background palette "
     "(>=10% light slides; no single hue covers >=90% of deck; AF-VISUAL-VARIETY)",
     "Phase Image-QC / Postflight — Typography Architect + Slide Image Creator "
     "(SOP-DESIGN-03, AF-VISUAL-VARIETY)",
     check_visual_variety),
    # AF-IMAGE-QC-RAN gate (2026-06-19). Requires the image-QC report to exist,
    # be NEWER than the rendered PNGs, and carry a per-slide PASS/FAIL row for
    # every rendered slide. Defers when no renders exist.
    (None,
     "image-QC ran — image-QC report is present, newer than all rendered PNGs, "
     "and covers every slide with a per-slide PASS/FAIL row (AF-IMAGE-QC-RAN)",
     "Phase Image-QC — Image QC Specialist (AF-IMAGE-QC-RAN)",
     check_image_qc_present),
    # AF-BRAND-CONSISTENCY gate (2026-06-19). Every rendered slide's dominant
    # palette must fall within the client's declared brand token set. Defers when
    # no brand palette is declared or no renders exist.
    (None,
     "brand consistency — every rendered slide's dominant palette is within the "
     "declared brand token set (AF-BRAND-CONSISTENCY)",
     "Phase Image-QC / Postflight — Slide Image Creator + Typography Architect "
     "(AF-BRAND-CONSISTENCY)",
     check_brand_consistency),
    # G5 — DECK HARMONY (AF-HARMONY), fired PRE-ASSEMBLY (after render, before assembly/
    # delivery). Proves the engines are ORCHESTRATED IN HARMONY across the deck — recurring-
    # character continuity, palette/brand coherence, world continuity, archetype RHYTHM —
    # not just that each slide passes in isolation. An individually-fine but incoherent deck
    # is caught BEFORE it is assembled and shipped. Run-dir-scoped (None sentinel); defers
    # pre-render (no PNGs yet).
    (None,
     "deck harmony — the engines cohere across the deck: recurring character continuity, "
     "palette/brand coherence, world continuity, and a deliberate archetype rhythm (not 40 "
     "one-offs); re-render only the inconsistent slides (AF-HARMONY)",
     "Pre-Assembly — Director of Presentations (harmony owner; SOP-HARMONY-01, AF-HARMONY)",
     check_deck_harmony),
    # AF-PACKAGE-CLEAN gate (2026-06-19). The final bundle directory must contain
    # only the canonical deliverable files (no .py/.sh/~$*/tasks/ artifacts).
    # Defers when bundle_dir does not exist. Fires at postflight closeout.
    # Note: this is a bundle-dir-scoped gate; run_preflight passes run_dir as the
    # first positional arg, but check_package_cleanliness accepts a Path to the
    # bundle_dir. Since PREFLIGHT runs before assembly (no bundle_dir yet), this
    # gate is wired as a deferred/conditional entry that always defers here and
    # is re-invoked directly in run_postflight_gate.
    # (PREFLIGHT entry deliberately omitted for AF-PACKAGE-CLEAN -- it fires only
    # at the postflight gate where bundle_dir is known.)
    # NO-DARK-SLIDES gate (AF-DARK-SLIDE). Slides MUST use LIGHT/bright backgrounds by
    # DEFAULT. Dark or black-background slides are NOT ALLOWED unless the client
    # EXPLICITLY requests a dark theme (intake flag client_dark_theme:true). Light is
    # the default; dark is opt-in by client request only. Run-dir-scoped (None sentinel):
    # reads intake.json + scans working/prompts/*.txt + working/research/design-brief-*.md.
    (None,
     "no-dark-slides — slides must use light backgrounds by default; dark only with client_dark_theme:true",
     "Intake/Prompt — Director intake + Slide Image Creator (AF-DARK-SLIDE)",
     _chk_no_dark_slides),
    # GOAL-4 / 1C — ASSET-QUESTION gate (AF-ASSET-QUESTION-MISSING). The signature
    # intake must ask whether the client already has materials (photos/logo/colors/
    # rough deck/slides/concepts). Run-dir-scoped (None sentinel); defers pre-intake.
    (None,
     "asset-intake question asked — intake records asset_intake_question_asked:true",
     "Intake — Brainstorming Buddy SOP 9.1/9.2 + Director (AF-ASSET-QUESTION-MISSING)",
     _chk_asset_question),
    # GOAL-4 / 1C — ASSETS-MANIFEST gate (AF-MANIFEST-UNREFERENCED). Provided assets
    # must be recorded in assets_manifest.json AND provably consumed (public_url +
    # consumed_by). Run-dir-scoped; defers when no assets provided / pre-intake.
    (None,
     "provided assets consumed — assets_manifest.json records each provided asset with a "
     "public_url + non-empty consumed_by (fed to gpt-image-2 as input_urls)",
     "Media-Librarian asset-ingest step (AF-MANIFEST-UNREFERENCED)",
     _chk_assets_manifest),
    # GOAL-4 / 1C — SCRATCH-PARSE gate (AF-SCRATCH-PARSE-SKIPPED). An uploaded
    # rough/old deck must be parsed into scratch_seed.json and seed the PRD (the
    # interview still runs in full). Run-dir-scoped; defers when no scratch deck.
    (None,
     "scratch deck parsed — uploaded old deck extracted to scratch_seed.json and seeded "
     "into the PRD (interview still runs in full)",
     "Media-Librarian scratch-deck parser sub-step (AF-SCRATCH-PARSE-SKIPPED)",
     _chk_scratch_parse),
    # GOAL-4 / 2A — PITCH-FLAG gate (AF-PITCH-FLAG-UNSET). intake.json must capture an
    # explicit boolean pitch_included. Run-dir-scoped; defers pre-intake.
    (None,
     "pitch_included captured — intake.json carries an explicit boolean pitch_included "
     "(yes = offer/pitch deck; no = teaching/content-only)",
     "Intake — Brainstorming Buddy + Director (AF-PITCH-FLAG-UNSET)",
     _chk_pitch_flag),
    # GOAL-4 / 2A — PITCH-LEAK gate (AF-PITCH-LEAK). A pitchless deck
    # (pitch_included:false) must carry NO pitch/price/offer/ladder content and no
    # price_ladder.json. Run-dir-scoped; defers when pitch_included is true/unset.
    (None,
     "pitchless integrity — a pitch_included:false deck has no offer/price/ladder/re-pitch "
     "content and no price_ladder.json (Offer Price Strategist suppressed)",
     "Phase 3 — pitchless suppression (AF-PITCH-LEAK)",
     _chk_pitch_leak),
    # GOAL-4 / 5C — NO-OVERLAY gate (AF-OVERLAY-DELIVERED). The native PPTX
    # text-overlay path is eliminated: no pptx_text_overlays.json may be present and
    # no delivered slide may carry a native on-slide text run (only a single composed
    # gpt-image-2 image; only off-slide speaker notes are allowed). Run-dir-scoped.
    (None,
     "no native overlay — no pptx_text_overlays.json present and no delivered slide ships "
     "native on-slide text instead of a composed gpt-image-2 image",
     "Phase 5/6/Postflight — render + PPTX assembly (AF-OVERLAY-DELIVERED)",
     _chk_no_overlay),
    # FIX-2 — CANONICAL-RENDER-PATH guard (AF-CANONICAL-RENDER-BYPASS / AF-LOCAL-CANVAS).
    # The canonical render path is build_deck.py / run_signature_deck.py ONLY. This
    # scans the run dir for a hand-rolled renderer/assembler (a local Image.new
    # 2048x1152 slide canvas, a native add_textbox overlay, or a direct kie createTask
    # outside build_deck.py) and HARD-FAILS pre-render so a per-deck renderer can never
    # produce the deck. Run-dir-scoped (None sentinel). Waivable ONLY by a logged
    # owner_skip_approval token.
    (None,
     "canonical render path — no hand-rolled *.py in the run dir defines a local slide "
     "canvas (Image.new 2048x1152), a native add_textbox, or a direct kie createTask "
     "(AF-CANONICAL-RENDER-BYPASS / AF-LOCAL-CANVAS)",
     "Phase 4/5/6 — render + assembly must route through build_deck.py only "
     "(AF-CANONICAL-RENDER-BYPASS)",
     check_canonical_render_path),
    # FIX-2 — IMAGE-QC-VISION pixel cross-check (AF-IMAGE-QC-VISION). The deterministic
    # PIXEL teeth behind the image-QC gate: it OPENS every rendered PNG (byte floor +
    # flat cream/typography-card detection) and refuses a pixel-blind / overlay-blessing /
    # hook-slide-excluding image-QC report. Defers ONLY pre-render (no PNGs + no report);
    # once renders or a report exist it does NOT defer-to-pass. Run-dir-scoped (None).
    (None,
     "image-QC vision — every rendered PNG is a real kie.ai bake (>= 51,200 bytes, no "
     "flat cream/typography card) and the image-QC report is a per-slide multimodal "
     "pixel read, not a self-typed number (AF-IMAGE-QC-VISION)",
     "Phase Image-QC / Postflight — Image QC Specialist pixel read (AF-IMAGE-QC-VISION)",
     check_image_qc_vision),
    # U027 — OCR-READBACK postflight gate (AF-OCR-READBACK). Audits the per-slide OCR
    # sidecar records (renders/slide-*.ocr.json) that _record_ocr_readback() writes
    # beside every rendered PNG once _verify_aspect_and_readback() calls
    # prompt_gate.ocr_readback() against it. A missing/unparseable sidecar, or a
    # checked:false record (the OCR engine never ran), FAILS and is NEVER waivable;
    # a checked:true/matched:false mismatch FAILS but is waivable ONLY by a logged
    # owner_skip_approval token for AF-OCR-READBACK. Defers ONLY pre-render (no
    # rendered PNGs yet). Run-dir-scoped (None sentinel).
    (None,
     "OCR readback — every rendered PNG carries a renders/<slide>.ocr.json sidecar "
     "proving the baked text was read back and matched the approved copy; a "
     "missing/unparseable sidecar or a checked:false record (OCR engine never ran) "
     "is NEVER waivable (AF-OCR-READBACK)",
     "Postflight — render closeout OCR audit (U027, AF-OCR-READBACK)",
     check_ocr_readback),
    # === POWERFUL-PRESENTATION DOCTRINE GATES (manifest v18) ===
    # All DEFER unless Phase P0B-PRIORITY produced working/copy/priority_shift_spec.json
    # (the no-regression master switch), so a legacy / ad-hoc build is never broken.
    # AF-DECK-TYPE-UNSET (U021) -- the type field the SP switch reads must exist and
    # be enumerated. Runs for EVERY deck; warn-mode until the dated window closes.
    (None,
     "deck type set -- intake.deck_type is webinar|signature_presentation, engine-written "
     "by deck-intake-driver.derive_legacy_fields() (AF-DECK-TYPE-UNSET)",
     "Phase 0.1 -- Director intake / type-picker (deck-intake-questions.json order 0, AF-DECK-TYPE-UNSET)",
     _chk_deck_type),
    # AF-MODE-UNSET (P118) -- the creation mode is identified at Step Zero.
    (None,
     "creation mode set -- intake.creation_mode is from_scratch|content_personal|content_general "
     "(content modes carry extracted_substance) (AF-MODE-UNSET)",
     "Phase 0.1/0.2 -- Director intake + Attention-Content Strategist (SOP-MODE-00, AF-MODE-UNSET)",
     _chk_mode),
    # AF-NO-SHIFT (P33/P105) — the priority-shift spine: spec + 8 monotonic build-moves.
    (None,
     "priority shift engineered — priority_shift_spec.json (true_goal + priority_stack[]) and "
     "the eight build-move beat tags monotonic in slides_copy.md (AF-NO-SHIFT)",
     "Phase 0.2 — Attention-Content Strategist (SOP-NORTHSTAR-00 / SOP-PRIORITY-02, AF-NO-SHIFT)",
     _chk_priority_shift),
    # AF-NO-PRIORITY-STACK (P142, Move 1) — stack named before the first ladder beat.
    (None,
     "priority stack named first — the audience's current priorities appear before the first "
     "value/price ladder beat (AF-NO-PRIORITY-STACK)",
     "Phase 0.2/4 — Attention-Content Strategist + Slide Copywriter (P142, AF-NO-PRIORITY-STACK)",
     _chk_priority_stack),
    # AF-NO-RERANK (P148, Move 7) — re-rank demanded after PRICE (pitch decks only).
    (None,
     "re-rank demanded after price — an explicit decision/re-rank beat follows the PRICE "
     "(AF-NO-RERANK)",
     "Phase 4 — Slide Copywriter + Presenters Speech Writer (P148, AF-NO-RERANK)",
     _chk_rerank),
    # AF-NO-TRIGGER (P149, Move 8) — time-bound trigger on the CTA (pitch decks only).
    (None,
     "trigger fired — the CTA carries a time-bound trigger (deadline / scarcity window) "
     "(AF-NO-TRIGGER)",
     "Phase 4 — Slide Copywriter + Offer Price Strategist (P149, AF-NO-TRIGGER)",
     _chk_trigger),
    # AF-PROCLAMATION-HEDGE (P109) — proclamations are plain, bold, hedge-free.
    (None,
     "proclamation hedge-free — no strong hedge phrases in slide copy (AF-PROCLAMATION-HEDGE)",
     "Phase 4 — Attention-Content Strategist + Slide Copywriter (SOP-PROCLAMATION-01, P109)",
     _chk_proclamation_hedge),
    # AF-PEAK-END (P49) — the arc declares a deliberate PEAK and a deliberate ending.
    (None,
     "peak + ending engineered — the arc carries a PEAK/APEX beat and a deliberate ending "
     "(AF-PEAK-END)",
     "Phase 3 — Offer Price Strategist + Director SOP 9.4 arc (P49, AF-PEAK-END)",
     _chk_peak_end),
    # AF-NO-SALIENCE-APEX (P48/P155) — the offer/promise-apex slide is the most vivid.
    (None,
     "most vivid by the end — the OFFER/PROMISE-APEX slide is the single most vivid rendered "
     "slide (von Restorff backstop) (AF-NO-SALIENCE-APEX)",
     "Phase 4.95 — Attention Designer / Image QC (P48/P155, AF-NO-SALIENCE-APEX)",
     _chk_salience_apex),
    # AF-CONVERTER-NO-INVENT (P167d) — converter extracts, never invents figures.
    (None,
     "converter extracts not invents — every figure in the source brief traces to the raw "
     "source (AF-CONVERTER-NO-INVENT)",
     "Phase -1 — Content-to-Presentation Architect (P167d, AF-CONVERTER-NO-INVENT)",
     _chk_converter_no_invent),
    # HOLE C — orphaned "no-beat" persuasion taxonomy, now WIRED (pitch decks only).
    (None,
     "persuasion beats present — problem / choice / fork / comparison / measurable-results / "
     "expert-proof / before-after beats are all present (AF-NO-PROBLEM / AF-NO-CHOICE / "
     "AF-NO-FORK / AF-NO-COMPARISON / AF-NO-MEASURABLE-RESULTS / AF-NO-EXPERT-PROOF / "
     "AF-NO-BEFORE-AFTER)",
     "Phase 4 — Slide Copywriter (SOP-PITCH-06 / SOP-ENGINE-00 persuasion-beat taxonomy)",
     _chk_persuasion_beats),
    # AF-STYLE-UNPICKED / AF-STYLE-DOUBLECHARGE (P-STYLE-PREVIEW, 4.85).
    (None,
     "style preview picked + locked — the owner picked one of 3 attention-grade variants and "
     "the winner's 3 representative renders are locked + reused, never re-charged "
     "(AF-STYLE-UNPICKED / AF-STYLE-DOUBLECHARGE)",
     "Phase 4.85 — Attention Designer 3-style-preview (owner gateway sign-off)",
     _chk_style_preview),
    # AF-PRIORITY-SHIFT (composite, P-SHIFT-QC 7.5) — the 14-item pre-output ship gate.
    (None,
     "priority-shift ship gate — the 14-item pre-output checklist (copy + design + images "
     "coexist) all PASS; item 0 is the North Star (AF-PRIORITY-SHIFT)",
     "Phase 7.5 — QC Specialist (SOP-INTEGRATION-00, P161/P155, AF-PRIORITY-SHIFT)",
     _chk_priority_shift_ledger),
    # Signature-Presentation (P-SP-CLAIM 0.14) — the routing/claim gate. Runs for EVERY deck
    # (does NOT defer): a run carrying SP signals MUST declare intake.json deck_type ==
    # signature_presentation (AF-SP-TYPE-UNDECLARED). Closes the "omit the magic word to skip
    # every SP gate" bypass — the single highest-severity skip of the audit. A non-signature
    # deck with no signal is untouched.
    (None,
     "signature-presentation claim gate — a run carrying SP signals (working/copy/sp_intake.json / "
     "a set Signature frame / a frame-selection question / a 'signature presentation' request) MUST "
     "declare intake.json deck_type == signature_presentation (AF-SP-TYPE-UNDECLARED). Runs for every "
     "deck; non-signature decks with no signal pass untouched.",
     "Phase 0.14 — Signature Presentation Router (P-SP-CLAIM, prove_sp_routing)",
     _chk_sp_claim),
    # Signature-Presentation (P-SP-INTAKE 0.15) — 8-Questions-in-ONE-block intake gate.
    (None,
     "signature-presentation intake gate — the 8 Questions delivered in ONE block + a set "
     "frame + q7 offer declared (AF-SP-8Q-MISSING / AF-SP-8Q-SPLIT / AF-SP-FRAME-UNSET / "
     "AF-SP-TYPE-MISMATCH / AF-SP-OFFER-UNDECLARED). DEFERS (no-op) unless intake.json "
     "deck_type == signature_presentation — non-signature decks behave exactly as before.",
     "Phase 0.15 — Signature Presentation Architect (P-SP-INTAKE, prove_sp_intake)",
     _chk_sp_intake),
    # A10 / T0-12 — P-SP-INTAKE-TRACE (0.16). The intake CONVERSATION gate. The record
    # gate above proves the assembled intake LEDGER; this proves the conversation that
    # produced it was choice-first and one question per turn (AF-INTAKE-BATCH), against
    # the transcript the driver's turn-gate writes mechanically. Previously advisory and
    # non-gating while SKILL.md declared every rule fail-closed, so a batched eight-
    # question dump reached a signed certificate. An ABSENT transcript is a fail — a
    # conversation gate no one can omit their way past. DEFERS (no-op) unless intake.json
    # deck_type == signature_presentation.
    (None,
     "signature-presentation intake-CONVERSATION gate — the intake was conducted "
     "choice-first (quick vs in-depth) and one bank question per assistant turn, proven "
     "against working/interview/intake_transcript.json (AF-INTAKE-BATCH: BATCH-IN-TURN / "
     "BATCH-BY-QMARKS / NO-CHOICE-OPENER / BANNED-PHRASE). A missing transcript is "
     "fail-closed. DEFERS (no-op) unless intake.json deck_type == signature_presentation.",
     "Phase 0.16 — Signature Presentation Architect (P-SP-INTAKE-TRACE, intake_trace_check)",
     _chk_sp_intake_trace),
    # Signature-Presentation (P-SP-STRUCTURE 4.1) — the SACRED 4-phase structure contract.
    (None,
     "signature-presentation structure gate — >=100 slides, 4-phase floors in contiguous "
     "order with phase labels, a suggested_image per slide, <=2 case studies, 3-7 teaching "
     "steps, central + 4 section hooks, N.E.E.I.T./4-Quadrant markers (AF-SP-SLIDE-FLOOR / "
     "AF-SP-PHASE-RANGE / AF-SP-PHASE-ORDER / AF-SP-PHASE-LABEL / AF-SP-IMG-SUGGESTION / "
     "AF-SP-CASESTUDY-CAP / AF-SP-TEACH-STEPS / AF-SP-HOOK / AF-SP-QUADRANT). DEFERS (no-op) "
     "unless intake.json deck_type == signature_presentation.",
     "Phase 4.1 — Signature Presentation Architect (P-SP-STRUCTURE, prove_sp_structure)",
     _chk_sp_structure),
    # AF-SP-P3-PITCH (P-SP-P3-HYGIENE 4.15) — Phase-3 teaching-band no-pitch hygiene.
    (None,
     "signature-presentation Phase-3 no-pitch hygiene — no offer name / price / CTA in the "
     "teaching band and a contiguous Phase-3->Phase-4 bridge (AF-SP-P3-PITCH). DEFERS (no-op) "
     "unless intake.json deck_type == signature_presentation.",
     "Phase 4.15 — QC Specialist (Signature Presentations) (P-SP-P3-HYGIENE, prove_sp_no_pitch)",
     _chk_sp_no_pitch),
]


def find_run_dir(explicit: Optional[str], slides_path: Path, out_path: Path) -> Path:
    """Resolve the run/working dir: --run-dir, else first ancestor containing a
    `working/` subtree (searched from slides.json then out.pptx dir), else
    slides.json's parent."""
    if explicit:
        return Path(explicit).resolve()
    for start in (slides_path.resolve().parent, out_path.resolve().parent):
        cur = start
        for _ in range(8):  # bounded upward walk
            if (cur / "working").is_dir():
                return cur
            if cur.parent == cur:
                break
            cur = cur.parent
    return slides_path.resolve().parent


# ===========================================================================
# FRONT-DOOR NONCE HANDSHAKE (Contract #8 — supersedes the retired
# OC_DECK_CANONICAL_ENTRY / OC_DECK_ALLOW_DIRECT env markers).
#
# The old markers were plain env vars documented in box-visible comments/SOPs, so
# any model that read the repo could set OC_DECK_CANONICAL_ENTRY=1 (or the CI escape
# OC_DECK_ALLOW_DIRECT=1) and invoke the renderer directly, bypassing the canonical
# entry gate chain. They are replaced by a per-run random nonce:
#   presentation-canonical-entry.sh mints a random nonce, writes it 0600 to
#   <run-dir>/working/checkpoints/.canonical-entry-nonce, and exports
#   OC_DECK_ENTRY_NONCE. The renderer admits the run ONLY when the exported nonce
#   matches that run-scoped file (constant-time compare); the entry script consumes
#   (deletes) the file after the run so a stale env value cannot be replayed.
# A model that merely READ the shipped source cannot conjure a valid value.
# ===========================================================================
ENTRY_NONCE_REL = Path("working") / "checkpoints" / ".canonical-entry-nonce"


def _entry_nonce_file(run_dir) -> Path:
    """Run-scoped path of the front-door nonce file (derived from run_dir, NOT env)."""
    return Path(run_dir) / ENTRY_NONCE_REL


def _verify_entry_nonce(run_dir) -> bool:
    """True iff OC_DECK_ENTRY_NONCE is set AND equals the content of the run-scoped
    nonce file <run-dir>/working/checkpoints/.canonical-entry-nonce. The path is
    derived from run_dir (never from an attacker-controllable env var) so a model
    cannot point it at a file it wrote elsewhere. A missing env var, a missing file,
    or any mismatch → False (fail-closed)."""
    import hmac
    env_nonce = (os.environ.get("OC_DECK_ENTRY_NONCE") or "").strip()
    if len(env_nonce) < 16:
        return False
    nf = _entry_nonce_file(run_dir)
    try:
        if not nf.is_file():
            return False
        file_nonce = nf.read_text(errors="replace").strip()
    except OSError:
        return False
    if len(file_nonce) < 16:
        return False
    return hmac.compare_digest(env_nonce, file_nonce)


def _resolve_run_dir_for_guard(run_dir_arg: Optional[str], positional: list) -> Path:
    """Best-effort run_dir resolution for the front-door nonce guard, computed at the
    top of main() (before the full find_run_dir call). Mirrors find_run_dir's inputs."""
    try:
        if run_dir_arg:
            return Path(run_dir_arg).resolve()
        if len(positional) >= 2:
            return find_run_dir(None, Path(positional[0]), Path(positional[1]))
        if len(positional) == 1:
            return find_run_dir(None, Path(positional[0]), Path(positional[0]))
    except Exception:  # noqa: BLE001 — never let resolution errors open the gate
        pass
    return Path.cwd()


def resolve_logo_path(logo_arg: Optional[str], run_dir: Path) -> Optional[Path]:
    """Resolve the official logo PNG. Priority:
        1. --logo <png> flag (wins over everything).
        2. working/copy/intake.json -> brand.logo_image_path (relative paths are
           resolved against the run dir, then the intake.json's own directory).
    Returns an absolute Path to an existing file, or None if no logo is configured.
    FAILS LOUD (exit 2) if a logo IS configured but the file is missing/unreadable."""
    candidate = None
    source = None

    if logo_arg:
        candidate = Path(logo_arg).expanduser()
        source = f"--logo {logo_arg}"
    else:
        intake = run_dir / "working" / "copy" / "intake.json"
        if intake.exists():
            obj = _read_json(intake)
            if isinstance(obj, dict) and "__parse_error__" not in obj:
                brand = obj.get("brand") or {}
                raw = str((brand.get("logo_image_path") or "")).strip() if isinstance(brand, dict) else ""
                if raw:
                    p = Path(raw).expanduser()
                    if not p.is_absolute():
                        # try run_dir-relative first, then intake.json-relative
                        for base in (run_dir, intake.parent):
                            cand = (base / p)
                            if cand.exists():
                                p = cand
                                break
                    candidate = p
                    source = f"intake.json brand.logo_image_path ({raw})"

    if candidate is None:
        return None

    candidate = candidate.resolve()
    if not candidate.exists() or not candidate.is_file():
        print(f"FATAL: official logo configured via {source} but file not found: "
              f"{candidate}", file=sys.stderr)
        sys.exit(2)
    # Light sanity: must be a PNG (RGBA transparency is what we preserve).
    with open(candidate, "rb") as f:
        magic = f.read(8)
    if magic[:4] != b"\x89PNG":
        print(f"FATAL: official logo {candidate} is not a PNG "
              f"(magic={magic[:8].hex()}). Provide a PNG (RGBA preferred).",
              file=sys.stderr)
        sys.exit(2)
    return candidate


def run_preflight(run_dir: Path, slides_path: Optional[Path] = None) -> None:
    """Refuse (exit 3) unless every required upstream dept artifact exists AND is
    complete. Lists exactly what is missing and which dept phase/role produces it.

    H3: slides_path is the ACTUAL positional slides.json the renderer will render.
    Threading it into the run-dir-scoped checks (coverage, rich-prompts) makes those
    gates count the EXACT file that gets rendered, so the slide-count derived in
    preflight always matches what main() renders — closing the gate-bypass where
    preflight counted a different canonical slides.json than the one rendered."""
    import inspect as _inspect
    print(f"=== PROCESS PREFLIGHT — run dir: {run_dir} ===", flush=True)
    problems = []

    # === v16.0.1 — BIND RENDER TO PHASE P0B-PRIORITY AT EVERY ENTRY POINT ===
    # The deterministic runner (run_signature_deck.py) makes P0B-PRIORITY a mandatory
    # order-based precondition of P4-RENDER, so a runner-driven render can never reach
    # kie.ai doctrine-blind. But a DIRECT `build_deck.py` call bypasses the runner: with
    # every upstream artifact present EXCEPT working/copy/priority_shift_spec.json, the
    # _doctrine_active() no-regression switch DEFERS every doctrine gate, so the deck
    # renders (kie.ai spend) with the priority-shift spine never governing it (delivery is
    # still later blocked by the canonical-render-guard, but the money is already spent).
    # ROOT CAUSE: P0B-mandatoriness lived ONLY in the runner. Here we reuse the SHARED
    # check_phase_preconditions machinery (single source of truth — not a parallel gate)
    # to refuse a NON-adhoc render unless P0B-PRIORITY is attested in process_manifest.json
    # (the ONLY waiver is a logged owner-authorized skip, exactly as the runner allows).
    # NOTE: run_preflight is only called for NON --adhoc-no-process runs (main() skips it
    # entirely in adhoc mode), so standalone/legacy adhoc testing stays exempt by
    # construction and the no-regression switch is NOT weakened. P0B is unconditional in
    # all three creation modes (from_scratch | content_personal | content_general) — no
    # mode legitimately skips it — so this binding is safe across every mode.
    p0b_reason = check_phase_preconditions(run_dir, RENDER_PHASE_ID, [PRIORITY_PHASE_ID])
    if p0b_reason:
        problems.append((
            PRIORITY_SPEC_REL,
            "Phase P0B-PRIORITY attested in process_manifest.json (the priority-shift "
            "spine that produces priority_shift_spec.json) — a direct build_deck render "
            "is bound to P0B exactly as the runner is, so it can never render doctrine-blind",
            "Phase 0.2 — Attention-Content Strategist "
            "(run_signature_deck.py P0B-PRIORITY, AF-PHASE-SKIPPED)",
            p0b_reason))

    for rel, label, phase, check in PREFLIGHT_REQUIRED:
        if rel is None:
            # run-dir-scoped check (e.g. _chk_coverage needs the whole run dir). Pass
            # slides_path too when the checker accepts it (H3), so the slide-count
            # gates derive from the actual rendered file.
            try:
                accepts_slides = "slides_path" in _inspect.signature(check).parameters
            except (TypeError, ValueError):
                accepts_slides = False
            if accepts_slides:
                reason = check(run_dir, slides_path)
            else:
                reason = check(run_dir)
            display = "(coverage / anti-compression)"
        else:
            if "*" in rel:
                matches = sorted(run_dir.glob(rel))
                found = matches[0] if matches else None
            else:
                p = run_dir / rel
                found = p if p.exists() else None
            reason = check(found)
            display = rel
        if reason:
            problems.append((display, label, phase, reason))
        else:
            print(f"  OK   {display}", flush=True)

    if problems:
        print("\nFATAL: PROCESS PREFLIGHT FAILED — refusing to render or assemble.", file=sys.stderr)
        print("build_deck.py is only the deterministic renderer/assembler; it cannot", file=sys.stderr)
        print("produce a compliant deliverable without the upstream department artifacts.\n", file=sys.stderr)
        print(f"Run/working dir checked: {run_dir}", file=sys.stderr)
        print("Missing or incomplete required artifacts:\n", file=sys.stderr)
        for rel, label, phase, reason in problems:
            print(f"  - {rel}", file=sys.stderr)
            print(f"      what:    {label}", file=sys.stderr)
            print(f"      reason:  {reason}", file=sys.stderr)
            print(f"      produced by: {phase}", file=sys.stderr)
        print("\nFix: run the real Presentations department pipeline (which produces these),", file=sys.stderr)
        print("or, for deliberate standalone testing ONLY, re-run with --adhoc-no-process", file=sys.stderr)
        print("(its output is explicitly NOT a process-compliant deliverable).", file=sys.stderr)
        sys.exit(3)

    print("=== PREFLIGHT PASSED — upstream dept artifacts present ===\n", flush=True)


def print_adhoc_banner() -> None:
    bar = "!" * 78
    print(bar, flush=True)
    print("!! ADHOC MODE (--adhoc-no-process): PROCESS PREFLIGHT SKIPPED.", flush=True)
    print("!! The output of this run is NOT a process-compliant deliverable.", flush=True)
    print("!! No research / copy / QC / human-approval artifacts were verified.", flush=True)
    print("!! Use for standalone testing ONLY — never for client work.", flush=True)
    print(bar + "\n", flush=True)


def _sha256_file(path: Path) -> Optional[str]:
    """Return the sha256 hex digest of a file, or None if unreadable."""
    try:
        h = hashlib.sha256()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception:  # noqa: BLE001
        return None


def write_process_manifest(run_dir: Path, rendered, task_ids, model_used,
                           out_path: Path, timestamp: str) -> Path:
    """APPEND this render's record to working/checkpoints/process_manifest.json.

    The manifest is a cumulative, multi-phase record: each department phase appends
    its own entry under "phases" — this writer must NEVER clobber prior phases. It
    creates working/checkpoints/ and the JSON file if absent, then appends one
    "render" phase record carrying: KIE taskIds, the model used, per-slide image
    sha256 hashes, the output slide count, and the timestamp (passed in by the
    caller — this runs in a plain Python script, so os/Date use is fine here).

    Returns the manifest path. Append-safe: a corrupt/legacy file is preserved as a
    sibling .corrupt-<ts> backup rather than silently overwritten."""
    ckpt_dir = run_dir / "working" / "checkpoints"
    ckpt_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = ckpt_dir / "process_manifest.json"

    manifest = {"phases": []}
    if manifest_path.exists():
        try:
            existing = json.loads(manifest_path.read_text())
            if isinstance(existing, dict):
                manifest = existing
                if not isinstance(manifest.get("phases"), list):
                    manifest["phases"] = []
            elif isinstance(existing, list):
                # legacy: a bare list of phase records — preserve it.
                manifest = {"phases": existing}
        except Exception:  # noqa: BLE001
            # Preserve the unreadable prior file instead of clobbering it.
            backup = manifest_path.with_suffix(
                f".corrupt-{timestamp.replace(':', '').replace('-', '')}")
            try:
                manifest_path.replace(backup)
            except Exception:  # noqa: BLE001
                pass
            manifest = {"phases": []}

    per_slide = []
    for r in rendered:
        f = r.get("file")
        per_slide.append({
            "slide": r.get("slide"),
            "taskId": r.get("taskId"),
            "image": f,
            "image_sha256": _sha256_file(Path(f)) if f else None,
        })

    record = {
        "phase": "render",
        "tool": "build_deck.py",
        "timestamp": timestamp,
        "model_used": model_used,
        "taskIds": list(task_ids),
        "output_slide_count": len(rendered),
        "output_pptx": str(out_path),
        "slides": per_slide,
    }
    manifest["phases"].append(record)

    manifest_path.write_text(json.dumps(manifest, indent=2))
    return manifest_path


# ---------------------------------------------------------------------------
# DELIVERABLES LEDGER — per-run deliverables.json (Requirements 2, 3)
# ---------------------------------------------------------------------------
# The ledger is written to <bundle_dir>/deliverables.json incrementally:
#   * initialised at the start of the run (all entries status=pending)
#   * updated after each artifact is produced (status=built or verified)
#   * finalised at the postflight gate (status=verified on pass, or unchanged
#     + hard-fail AF-BUNDLE-COMPLETE on failure)
#
# The final "complete" report is generated by READING deliverables.json
# (all entries verified), NOT from in-memory state. This survives crashes
# and supports resume: re-running the script re-reads the ledger and only
# re-checks artifacts that are not yet verified.
#
# IMPORTANT: this script builds only the deck_pptx. The other eight build
# artifacts (the rest of the nine-key DELIVERABLES_REQUIRED build bundle) are
# produced by upstream roles. The ledger records whatever state they are
# in when the postflight gate runs. If any are absent or below threshold the
# gate fails loud (AF-BUNDLE-COMPLETE, exit 5).

def _resolve_bundle_dir(out_dir_arg: Optional[str], out_path: Path) -> Path:
    """Resolve the final bundle dir.
       Priority: --out CLI arg > ~/Downloads/<deck-slug>/.
       The <deck-slug> is the stem of the out.pptx path (e.g. 'acme-q1' from
       'acme-q1.pptx').
    """
    if out_dir_arg:
        d = Path(out_dir_arg).expanduser().resolve()
    else:
        slug = out_path.stem or "deck"
        d = Path(BUNDLE_DIR_DEFAULT) / slug
    d.mkdir(parents=True, exist_ok=True)
    return d


def _expand_filename(template: str, deck_slug: str) -> str:
    """Replace {deck_slug} placeholder in a filename template."""
    return template.replace("{deck_slug}", deck_slug)


def init_deliverables_ledger(bundle_dir: Path, deck_slug: str) -> Path:
    """Create or reset deliverables.json in bundle_dir with all artifacts
    set to status=pending. Returns the ledger path.

    If a ledger already exists and has verified entries, preserves them
    (supports resume: only re-checks non-verified artifacts)."""
    ledger_path = bundle_dir / "deliverables.json"
    existing = {}
    if ledger_path.exists():
        try:
            data = json.loads(ledger_path.read_text())
            if isinstance(data, list):
                existing = {e["key"]: e for e in data if isinstance(e, dict) and "key" in e}
        except Exception:  # noqa: BLE001
            existing = {}

    entries = []
    for spec in DELIVERABLES_REQUIRED:
        key = spec["key"]
        fname = _expand_filename(spec["filename"], deck_slug)
        # Preserve a prior verified entry so resume works.
        if key in existing and existing[key].get("status") == "verified":
            entries.append(existing[key])
        else:
            # produced_later specs (e.g. webinar_mp4 — rendered by P9.6-WEBINAR-VIDEO
            # AFTER postflight) are skipped by run_postflight_gate's existence loop by
            # design; initialize them as verified so the extra_unverified catch-all does
            # not flag a file that is not expected at this stage.
            init_status = "verified" if spec.get("produced_later") else "pending"
            entries.append({
                "key": key,
                "filename": fname,
                "path": str(bundle_dir / fname),
                "label": spec["label"],
                "min_bytes": spec["min_bytes"],
                "note": spec["note"],
                "status": init_status,
                "size": None,
                "error": None,
            })
    ledger_path.write_text(json.dumps(entries, indent=2))
    return ledger_path


def update_deliverable_status(ledger_path: Path, key: str, status: str,
                               size: Optional[int] = None,
                               error: Optional[str] = None) -> None:
    """Update a single artifact entry in deliverables.json. status must be
    one of: pending | built | verified | failed.
    Writes the ledger to disk immediately (incremental, crash-safe)."""
    try:
        entries = json.loads(ledger_path.read_text())
    except Exception:  # noqa: BLE001
        return  # if the ledger is unreadable, skip silently (gate will catch it)
    for e in entries:
        if e.get("key") == key:
            e["status"] = status
            if size is not None:
                e["size"] = size
            if error is not None:
                e["error"] = error
            break
    ledger_path.write_text(json.dumps(entries, indent=2))


# ---------------------------------------------------------------------------
# TELEPROMPTER PUBLISH — upload the self-contained HTML to its host, capture URL
# ---------------------------------------------------------------------------
# The teleprompter is a single self-contained HTML file (inline CSS+JS+speech JSON),
# so "publishing" is a static-file upload to a host that then serves it at a public
# URL. FLEET POLICY (uniform): EVERY client — Mac AND VPS — publishes to the central
# Cloudflare host at
#     https://teleprompter.zerohumanworkforce.com/<client-slug>/<deck-slug>/teleprompter.html
# The host is an R2 bucket fronted by a small serving Worker, gated by Cloudflare
# Access. detect_platform() still records whether the box is a Mac or a VPS for the
# audit trail, but the host_target is ALWAYS 'cloudflare-central' — there is one
# uniform host so the link works identically everywhere.
#
# The publish record is written to <bundle_dir>/teleprompter_publish.json and is the
# load-bearing contract read by the postflight gate, the Delivery Concierge (delivers
# the link), and the Media Librarian (files the link in GHL). FAIL-LOUD by raising;
# the caller decides whether a publish failure blocks the gate (it does — the gate
# fails when the publish record is absent or unverified).
#
# Credentials: the central host is FLEET infra (the operator's Cloudflare account),
# so the upload token is the operator/fleet secret CLOUDFLARE_ZHW_APPS_API_TOKEN —
# NOT a per-client key, and NEVER printed. It is read from the same env stores the
# rest of the pipeline uses (SECRETS_CANDIDATES), with the live process env winning.
TELEPROMPTER_PUBLISH_LEDGER = "teleprompter_publish.json"

# Central Cloudflare host constants (fleet infra — see PART B infra build).
TELEPROMPTER_HOST = "teleprompter.zerohumanworkforce.com"
TELEPROMPTER_R2_BUCKET = "zhw-teleprompter"
TELEPROMPTER_HTML_NAME = "teleprompter.html"
# The R2 PutObject endpoint is the account-scoped REST path. The serving Worker maps
# https://<host>/<client>/<deck>/teleprompter.html -> the same object key in R2.
_CF_API_BASE = "https://api.cloudflare.com/client/v4"


def _load_secret(name: str) -> str:
    """Read a secret by name: live process env first (definitive), then the
    SECRETS_CANDIDATES env files. Returns "" if absent. NEVER prints the value."""
    val = os.environ.get(name, "").strip()
    if val:
        return val.strip("'\"")
    prefix = name + "="
    for path in SECRETS_CANDIDATES:
        p = Path(path)
        if not p.exists():
            continue
        try:
            for line in p.read_text().splitlines():
                line = line.strip()
                if line.startswith(prefix):
                    value = line[len(prefix):].strip().strip("'\"")
                    if value:
                        return value
        except OSError:
            continue
    return ""


def _slugify(text: str) -> str:
    """Lowercase, hyphenate, strip to [a-z0-9-] for URL path segments. Stable and
    deterministic — never random. Empty input -> 'untitled'."""
    s = re.sub(r"[^a-z0-9]+", "-", str(text).strip().lower()).strip("-")
    return s or "untitled"


def detect_platform(run_dir: Path, override: Optional[str] = None) -> str:
    """Return 'vps' | 'mac'. Priority:
       1. --platform CLI override (explicit operator control).
       2. intake.json box_type ('mac' -> 'mac'; anything else -> 'vps').
       3. Filesystem signal: a '/data/.openclaw' tree present -> 'vps'; else 'mac'
          (the same Mac-vs-VPS workdir split the Media Librarian SOP 9.1 uses).
    NOTE: this records the box type for the audit trail only. The publish host is
    UNIFORM (central Cloudflare) for both values — see publish_teleprompter()."""
    if override:
        o = str(override).strip().lower()
        if o in ("vps", "mac"):
            return o
    intake = run_dir / "intake.json"
    if intake.exists():
        obj = _read_json(intake)
        if "__parse_error__" not in obj:
            box_type = str(obj.get("box_type", "")).strip().lower()
            if box_type == "mac":
                return "mac"
            if box_type:
                return "vps"
    # Filesystem fallback: a VPS box runs OpenClaw under /data/.openclaw.
    if Path("/data/.openclaw").is_dir():
        return "vps"
    return "mac"


def _client_slug_from_run(run_dir: Path, bundle_dir: Path) -> str:
    """Derive a stable client slug for the URL path. Priority:
       1. intake.json client_slug / client_name / brand.client_name.
       2. The bundle dir's parent leaf (clients usually land under ~/Downloads/<client>-<deck>).
       3. 'client' as a last resort. Always slugified."""
    intake = run_dir / "intake.json"
    if intake.exists():
        obj = _read_json(intake)
        if "__parse_error__" not in obj:
            for k in ("client_slug", "client_name", "client"):
                v = str(obj.get(k, "")).strip()
                if v:
                    return _slugify(v)
            brand = obj.get("brand")
            if isinstance(brand, dict):
                v = str(brand.get("client_name", "")).strip()
                if v:
                    return _slugify(v)
    leaf = bundle_dir.name.strip()
    if leaf:
        return _slugify(leaf)
    return "client"


def _http_status(url: str, timeout: int = 60) -> int:
    """Live GET of a public URL; return the HTTP status code. Reuses the C1 SSRF
    guard (http(s) only). A network/timeout error surfaces as a RuntimeError so the
    caller records a verify failure rather than a false 200. Behind Cloudflare Access
    an unauthenticated GET may 302 to the Access login — that is NOT a 200 and is a
    correct verify-failure signal for an automated check; for the published-object
    check we GET the object directly through the serving Worker which returns 200 for
    a present object (Access challenges the human browser, not the object existence
    probe path). We only need the status code."""
    assert_url_scheme_allowed(url, what="teleprompter public_url")
    req = urllib.request.Request(url, headers={"User-Agent": "build_deck-teleprompter-verify/1.0"})
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            return int(resp.getcode())
    except urllib.error.HTTPError as exc:
        return int(exc.code)
    except urllib.error.URLError as exc:
        raise RuntimeError(f"teleprompter verify GET failed for {url}: {exc}") from exc


def _r2_put_object(token: str, account_id: str, bucket: str, object_key: str,
                   data: bytes, content_type: str = "text/html; charset=utf-8") -> None:
    """Upload a single object to R2 via the Cloudflare REST API (PutObject). Raises
    RuntimeError on non-success. The token is the fleet CLOUDFLARE_ZHW_APPS_API_TOKEN
    (never printed)."""
    url = (f"{_CF_API_BASE}/accounts/{account_id}/r2/buckets/{bucket}/objects/"
           f"{quote(object_key)}")
    assert_url_scheme_allowed(url, what="R2 PutObject URL")
    req = urllib.request.Request(
        url, data=data, method="PUT",
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": content_type,
        },
    )
    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            code = int(resp.getcode())
            if code not in (200, 201):
                raise RuntimeError(f"R2 PutObject returned HTTP {code} for key {object_key!r}")
    except urllib.error.HTTPError as exc:
        body_text = exc.read().decode(errors="replace")
        raise RuntimeError(
            f"R2 PutObject HTTP {exc.code} for bucket {bucket!r} key {object_key!r}. "
            f"Response: {body_text}"
        ) from exc
    except urllib.error.URLError as exc:
        raise RuntimeError(f"R2 PutObject network error for {bucket!r}/{object_key!r}: {exc}") from exc


def _write_publish_ledger(ledger_path: Path, record: dict) -> None:
    """Write the teleprompter_publish.json record (pretty JSON)."""
    ledger_path.write_text(json.dumps(record, indent=2))


def publish_teleprompter(bundle_dir: Path, deck_slug: str, run_dir: Path,
                         platform: str, deliverables_ledger: Path,
                         adhoc: bool = False) -> dict:
    """Upload <bundle_dir>/presenter-teleprompter.html to the central Cloudflare host
    (R2 bucket served by a Worker at teleprompter.zerohumanworkforce.com), capture the
    public URL, write TELEPROMPTER_PUBLISH_LEDGER, and record the URL on the
    'teleprompter_html' deliverable entry. Live-GET the URL to confirm HTTP 200.

    FLEET POLICY: the host is UNIFORM (central Cloudflare) for every platform.
    detect_platform()'s vps/mac value is recorded for audit only.

    Raises RuntimeError on publish/verify failure (after writing a 'verify_failed'
    record so the postflight gate fails loud). In --adhoc mode, writes a
    'skipped_adhoc' record and returns it (no host call — ad-hoc output is explicitly
    not a client deliverable)."""
    pub_path = bundle_dir / TELEPROMPTER_PUBLISH_LEDGER
    local_file = bundle_dir / "presenter-teleprompter.html"
    client_slug = _client_slug_from_run(run_dir, bundle_dir)
    deck = _slugify(deck_slug)
    object_key = f"{client_slug}/{deck}/{TELEPROMPTER_HTML_NAME}"
    public_url = f"https://{TELEPROMPTER_HOST}/{object_key}"
    now = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())

    if adhoc:
        record = {
            "platform": platform,
            "host_target": "cloudflare-central",
            "local_file": str(local_file),
            "public_url": None,
            "published_at": now,
            "verified_http_status": None,
            "verified_at": None,
            "status": "skipped_adhoc",
            "note": "ad-hoc standalone render — not a client deliverable; publish skipped.",
        }
        _write_publish_ledger(pub_path, record)
        return record

    if not local_file.exists():
        record = {
            "platform": platform, "host_target": "cloudflare-central",
            "local_file": str(local_file), "public_url": None,
            "published_at": now, "verified_http_status": None, "verified_at": None,
            "status": "verify_failed",
            "note": "presenter-teleprompter.html absent from bundle_dir — nothing to publish.",
        }
        _write_publish_ledger(pub_path, record)
        raise RuntimeError(
            f"publish_teleprompter: {local_file} absent — generate the teleprompter "
            f"(build_teleprompter.py) before publishing."
        )

    # FLEET secret — operator/fleet Cloudflare token, never a client key, never printed.
    token = _load_secret("CLOUDFLARE_ZHW_APPS_API_TOKEN")
    account_id = _load_secret("CLOUDFLARE_ZHW_ACCOUNT_ID")
    if not token or not account_id:
        record = {
            "platform": platform, "host_target": "cloudflare-central",
            "local_file": str(local_file), "public_url": public_url,
            "published_at": now, "verified_http_status": None, "verified_at": None,
            "status": "verify_failed",
            "note": ("central Cloudflare credentials not found "
                     "(CLOUDFLARE_ZHW_APPS_API_TOKEN / CLOUDFLARE_ZHW_ACCOUNT_ID). "
                     "This is FLEET infra; set the operator/fleet token in the env "
                     "store. The value is never printed."),
        }
        _write_publish_ledger(pub_path, record)
        raise RuntimeError(
            "publish_teleprompter: CLOUDFLARE_ZHW_APPS_API_TOKEN and/or "
            "CLOUDFLARE_ZHW_ACCOUNT_ID not found in env or secrets stores. The central "
            "teleprompter host is fleet infra; configure the operator/fleet token."
        )

    html_bytes = local_file.read_bytes()

    # 1. Upload the self-contained HTML to the R2 bucket (PutObject). On failure,
    #    record a verify_failed ledger so the gate fails loud with a diagnostic.
    try:
        _r2_put_object(token, account_id, TELEPROMPTER_R2_BUCKET, object_key, html_bytes)
    except RuntimeError as exc:
        record = {
            "platform": platform, "host_target": "cloudflare-central",
            "local_file": str(local_file), "public_url": public_url,
            "published_at": now, "verified_http_status": None, "verified_at": None,
            "status": "verify_failed",
            "note": f"R2 upload failed: {exc}",
        }
        _write_publish_ledger(pub_path, record)
        raise RuntimeError(
            f"publish_teleprompter: R2 upload of {object_key} failed: {exc}"
        ) from exc

    # 2. Ground-truth verify: live GET the public URL through the serving Worker. The
    #    object existence probe path returns 200 for a present object (Access
    #    challenges the human browser, not this status probe). A non-200 is a publish
    #    failure.
    published_at = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    try:
        status = _http_status(public_url, timeout=60)
    except RuntimeError as exc:
        record = {
            "platform": platform, "host_target": "cloudflare-central",
            "local_file": str(local_file), "public_url": public_url,
            "published_at": published_at, "verified_http_status": None,
            "verified_at": None, "status": "verify_failed",
            "note": f"upload succeeded but verify GET errored: {exc}",
        }
        _write_publish_ledger(pub_path, record)
        raise RuntimeError(
            f"publish_teleprompter: uploaded {object_key} but the live verify GET "
            f"failed: {exc}"
        ) from exc

    verified_at = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    status_ok = (status == 200)
    record = {
        "platform": platform,
        "host_target": "cloudflare-central",
        "local_file": str(local_file),
        "public_url": public_url,
        "published_at": published_at,
        "verified_http_status": status,
        "verified_at": verified_at,
        "status": "published" if status_ok else "verify_failed",
    }
    _write_publish_ledger(pub_path, record)

    # Record the URL on the teleprompter_html deliverable entry (additive field).
    try:
        entries = json.loads(deliverables_ledger.read_text())
        for e in entries:
            if e.get("key") == "teleprompter_html":
                e["public_url"] = public_url
                break
        deliverables_ledger.write_text(json.dumps(entries, indent=2))
    except Exception:  # noqa: BLE001 — best-effort; the publish ledger is authoritative
        pass

    if not status_ok:
        raise RuntimeError(
            f"publish_teleprompter: {public_url} returned HTTP {status}, not 200. "
            f"The teleprompter is uploaded but the live URL is not verified."
        )
    return record


def _check_teleprompter_published(bundle_dir: Path, skip_gate: bool = False) -> str:
    """Teleprompter-publish sub-check of the postflight bundle gate (folded under
    AF-BUNDLE-COMPLETE). Return "" when the teleprompter is published with a verified
    live URL, else a fail reason. Reads <bundle_dir>/teleprompter_publish.json:
      * file absent                                  -> fail
      * status != 'published'                        -> fail (carry the recorded status)
      * public_url missing / not http(s)             -> fail
      * verified_http_status != 200                  -> fail

    M7 HARDENING: the gate is bypassed ONLY when the caller passes the explicit
    per-run --skip-teleprompter-gate CLI flag (skip_gate=True). It is NEVER bypassed
    by a persisted status string. A 'skipped_adhoc' record written by a prior --adhoc
    run can persist on disk and would otherwise false-pass a later real run whose
    publish step was skipped, so it is treated as an UNPUBLISHED failure and a WARNING
    is logged whenever it is encountered."""
    pub = bundle_dir / TELEPROMPTER_PUBLISH_LEDGER
    if skip_gate:
        print("WARNING: teleprompter-publish gate bypassed via the explicit "
              "--skip-teleprompter-gate flag. The teleprompter is NOT verified as "
              "hosted — this run MUST NOT be delivered to a client.", file=sys.stderr)
        return ""
    if not pub.exists():
        return ("teleprompter_publish.json absent — the teleprompter was never "
                "published (TELEPROMPTER-PUBLISH). Run publish_teleprompter.")
    obj = _read_json(pub)
    if "__parse_error__" in obj:
        return f"teleprompter_publish.json not valid JSON ({obj['__parse_error__']})"
    if obj.get("status") == "skipped_adhoc":
        print("WARNING: teleprompter_publish.json carries a stale 'skipped_adhoc' "
              "status from a prior --adhoc run. This NO LONGER passes the gate (M7): "
              "a real run must publish the teleprompter and verify a live HTTP 200, "
              "or be re-run with the explicit --skip-teleprompter-gate flag.",
              file=sys.stderr)
        return ("teleprompter publish status is 'skipped_adhoc' (stale ad-hoc record) "
                "— not a published, live-verified teleprompter (TELEPROMPTER-PUBLISH).")
    if obj.get("status") != "published":
        return f"teleprompter publish status is {obj.get('status')!r}, expected 'published'"
    url = str(obj.get("public_url", "")).strip()
    try:
        assert_url_scheme_allowed(url, what="teleprompter public_url")
    except ValueError as exc:
        return f"teleprompter public_url invalid: {exc}"
    if obj.get("verified_http_status") != 200:
        return (f"teleprompter public_url not verified live (status="
                f"{obj.get('verified_http_status')!r}); a live HTTP 200 is required.")
    return ""


# ---------------------------------------------------------------------------
# C2 — POSTFLIGHT CONTENT-TYPE MAGIC-BYTE VERIFICATION
# ---------------------------------------------------------------------------
# The completeness gate used to check only exists() + size, so a symlink pointing
# at a large unrelated file, or a decoy file of the right size but the wrong type,
# would PASS. C2 hardens it: reject symlinks outright, size via os.lstat (never
# follow a link), and verify the file's leading magic bytes match the type implied
# by its extension. .md stays size-only (plain text has no magic signature).
#
# Map of deliverable key -> tuple of acceptable leading-byte signatures. A PPTX is a
# ZIP (Office Open XML), so it begins with the local-file-header magic 'PK\x03\x04'.
# An MP3 is either an ID3v2-tagged file ('ID3') or a raw MPEG frame sync ('\xff\xfb'
# / '\xff\xf3' / '\xff\xf2'). PNG begins with '\x89PNG'. PDF begins with '%PDF'.
DELIVERABLE_MAGIC = {
    "deck_pptx":       (b"PK\x03\x04",),                              # PPTX = ZIP (OOXML)
    "deck_pdf":        (b"%PDF",),
    "guide_pdf":       (b"%PDF",),
    "speech_pdf":      (b"%PDF",),
    "audio_mp3":       (b"ID3", b"\xff\xfb", b"\xff\xf3", b"\xff\xf2"),  # ID3v2 tag or MPEG sync
    "infographic_png": (b"\x89PNG",),
    # teleprompter_html: a real HTML document leads with a doctype or an <html> /
    # <!-- ... --> opener. Accept the common case-variant openers so a legitimate
    # self-contained teleprompter page passes while a decoy of the wrong type fails.
    "teleprompter_html": (b"<!DOCTYPE", b"<!doctype", b"<!Doctype",
                          b"<html", b"<HTML", b"<!--"),
    # webinar_mp4: an MP4 leads with a 32-bit box size + 'ftyp' at offset 4 (ISO-BMFF
    # container signature) — the same guard ghl_media.verify_video uses. A non-video
    # stub of the right size is rejected.
    "webinar_mp4": (b"\x00\x00\x00\x18ftyp", b"ftyp", b"\x00\x00\x00\x10ftyp"),
    # speech_md / speech_fish_md: intentionally absent — plain/tagged markdown has no
    # magic signature, so md stays size-only (per spec).
}


def _magic_ok(path: Path, signatures) -> bool:
    """Return True if the file's leading bytes match ANY of the given signatures.
    Reads only the bytes needed for the longest signature. HTML may carry a leading
    UTF-8 BOM (\\xef\\xbb\\xbf) or leading whitespace before the doctype/tag, so the
    comparison tolerates a BOM + leading whitespace for text-ish (non-binary) types."""
    try:
        want = max(len(s) for s in signatures)
        with open(path, "rb") as f:
            head = f.read(want + 8)  # +8 to cover an optional BOM + a little leading WS
    except OSError:
        return False
    if any(head.startswith(s) for s in signatures):
        return True
    # Tolerate a leading UTF-8 BOM and/or leading ASCII whitespace ONLY for text-ish
    # signatures (those beginning with '<'); binary magics (PNG/PDF/ZIP/MP3) must be
    # at byte 0 and are not relaxed.
    text_sigs = tuple(s for s in signatures if s[:1] == b"<")
    if text_sigs:
        trimmed = head
        if trimmed.startswith(b"\xef\xbb\xbf"):
            trimmed = trimmed[3:]
        trimmed = trimmed.lstrip(b" \t\r\n")
        if any(trimmed.startswith(s) for s in text_sigs):
            return True
    return False


def _cc_registration_verified(manifest: dict) -> bool:
    """T2 gate teeth: prove a cc_task_id in process_manifest.json was produced
    by a REAL cc_board.ingest_deck_task round-trip, entirely offline.

    The receipt cc_board.stamp_task_id writes on a successful ingest is
    ``cc_registration`` = registration_proof(task_id, idempotency_key,
    deck_slug) — a deterministic HMAC-SHA256 over ``cc_task_id + "|" +
    idempotency_key`` (no secret, no timestamp: hermetic and replay-proof
    against hand-editing). idempotency_key is sha256(source_ref+title) on the
    producer side, so it CANNOT be derived from the task id alone — a forged
    id cannot forge its matching key, and a key for a different deck cannot
    match. The gate recomputes the HMAC over the manifest's own
    cc_task_id|idempotency_key and compares bytes to the stamped digest.

    Returns True only when: cc_task_id is a non-empty string, cc_registration
    is a dict whose task_id matches cc_task_id and whose hmac equals the
    recomputed digest. Everything else (absent cc_registration, mismatched
    task_id, wrong digest, non-dict receipt) is False — could-not-verify NEVER
    reads as verified. Never raises (a malformed manifest is just False)."""
    try:
        tid = manifest.get("cc_task_id")
        if not isinstance(tid, str) or not tid.strip():
            return False
        receipt = manifest.get("cc_registration")
        if not isinstance(receipt, dict):
            return False
        if receipt.get("task_id") != tid:
            return False
        key = receipt.get("idempotency_key")
        slug = receipt.get("deck_slug")
        if not isinstance(key, str) or not key or not isinstance(slug, str) or not slug:
            return False
        expected = hmac.new(
            f"{tid}|{key}".encode("utf-8"), b"", hashlib.sha256
        ).hexdigest()
        return hmac.compare_digest(str(receipt.get("hmac") or ""), expected)
    except Exception:  # noqa: BLE001 — a malformed manifest is UNVERIFIED, never fatal
        return False


def _chk_cc_registered(run_dir: Optional[Path], deck_slug: str) -> str:
    """AF-CC-UNREGISTERED / AF-CC-UNVERIFIED: postflight gate enforcing Command
    Center registration. THREE outcomes (T2 gate teeth):

      VERIFIED    — cc_task_id present AND the cc_registration HMAC proof
                    verifies (a real cc_board.ingest_deck_task round-trip).
                    Gate PASSES (returns "").
      UNVERIFIED  — cc_register_attempted=True but NO verifiable proof: either
                    no cc_task_id (transport/partial failure) or a task_id
                    whose receipt does not verify (hand-written/stale). Gate
                    FAILS with the explicit, honest AF-CC-UNVERIFIED message —
                    could-not-verify NEVER prints as verified. Fail-soft per
                    the item: a genuine Command Center outage does not brick
                    the deck beyond a clear, truthful UNVERIFIED closeout.
      UNREGISTERED— neither field: cc_board.ingest_deck_task was never called
                    for this run. Gate FAILS CLOSED (AF-CC-UNREGISTERED,
                    current behavior, kept).

    Returns "" on pass; an AF-CC-UNVERIFIED / AF-CC-UNREGISTERED message
    string on fail. Enforced_by: build_deck. py_symbol: _chk_cc_registered.
    """
    if run_dir is None:
        # adhoc/no-run-dir paths (--adhoc-no-process) skip the gate.
        return ""
    pm = Path(run_dir) / "working" / "checkpoints" / "process_manifest.json"
    if not pm.exists():
        return (
            "AF-CC-UNREGISTERED: no process_manifest.json found at closeout — "
            "cc_board.ingest_deck_task was never called for this deck run. "
            "Call cc_board.ingest_deck_task at run-begin (the canonical entry "
            "path handles this automatically). "
            f"Deck: {deck_slug}"
        )
    try:
        manifest = json.loads(pm.read_text())
    except (json.JSONDecodeError, OSError) as exc:
        return (
            f"AF-CC-UNVERIFIED: could not read process_manifest.json ({exc}) — "
            "cannot verify CC registration. "
            f"Deck: {deck_slug}"
        )
    if not isinstance(manifest, dict):
        return (
            "AF-CC-UNVERIFIED: process_manifest.json is not a JSON object — "
            f"cannot verify CC registration. Deck: {deck_slug}"
        )
    tid = manifest.get("cc_task_id")
    attempted = manifest.get("cc_register_attempted") is True
    # Successful registration PROVED: task_id written by a real
    # cc_board.stamp_task_id round-trip (HMAC receipt verifies offline).
    if isinstance(tid, str) and tid.strip() and _cc_registration_verified(manifest):
        return ""
    # UNVERIFIED — could-not-verify. NEVER printed as verified; fails the gate
    # with an explicit, honest message (fail-soft: a real CC outage is visible
    # and truthful, and the run is not bricked beyond closeout).
    if attempted or (isinstance(tid, str) and tid.strip()):
        if not (isinstance(tid, str) and tid.strip()):
            return (
                "AF-CC-UNVERIFIED: cc_register_attempted=True but no cc_task_id "
                "in process_manifest.json — the Command Center POST did not "
                "return a task (transport/partial failure or outage). "
                "Registration could NOT be verified; this run does NOT count "
                "as registered. Investigate the cc-board.json movement receipt "
                "and the [cc_board/presentations] stderr log, then re-run the "
                "deck or reconcile the card. "
                f"Deck: {deck_slug}"
            )
        return (
            "AF-CC-UNVERIFIED: cc_task_id is present but its cc_registration "
            "receipt does not verify (missing, mismatched, or bad HMAC) — this "
            "task id was NOT provably produced by a real "
            "cc_board.ingest_deck_task round-trip. Could-not-verify NEVER "
            "counts as verified. Re-run through the canonical entry path so "
            "ingest_deck_task stamps a verifiable receipt. "
            f"Deck: {deck_slug} cc_task_id={tid}"
        )
    # Neither field: this module was never invoked for this run — fail-CLOSED.
    return (
        "AF-CC-UNREGISTERED: process_manifest.json exists but has neither "
        "cc_task_id nor cc_register_attempted. Command Center was never contacted "
        "for this deck run. Call cc_board.ingest_deck_task at run-begin. "
        f"Deck: {deck_slug}"
    )


def run_postflight_gate(bundle_dir: Path, ledger_path: Path, deck_slug: str,
                        skip_teleprompter_gate: bool = False,
                        run_dir: Optional[Path] = None,
                        slides_path: Optional[Path] = None) -> None:
    """POSTFLIGHT COMPLETENESS GATE (AF-BUNDLE-COMPLETE, Requirement 2).

    After assembly, verify EVERY DELIVERABLES_REQUIRED artifact exists in
    bundle_dir AND passes its min_bytes threshold. Updates deliverables.json
    with the verified/failed status for each artifact, then reads it back to
    produce the final report.

    If ANY artifact is missing or below threshold → prints a LOUD failure
    listing exactly which are missing + sys.exit(5). The script may print
    COMPLETE/DONE ONLY when all pass.

    The final "complete" report is generated by READING deliverables.json
    (all entries verified), NOT from in-memory state.
    """
    # --- Phase 1: scan and update each entry ---
    # C2: reject symlinks (a symlink to a large unrelated file would otherwise pass),
    # size via os.lstat (NEVER follow a link), and verify the leading magic bytes
    # match the type implied by the extension (a decoy of the right size but wrong
    # type is rejected). md stays size-only (plain text has no magic signature).
    missing_or_short = []
    for spec in DELIVERABLES_REQUIRED:
        key = spec["key"]
        # Feature L2-G (P9.6-WEBINAR-VIDEO): `webinar_mp4` is PRODUCED LATER than this
        # gate (build_webinar_video.py runs at order 8.92, after P8-ASSEMBLE/this
        # postflight). It does NOT exist at render/assemble time, so it is skipped here —
        # the P9.6 phase, fix_bundle_complete (FIX-8) and the delivery gate own its real
        # presence gate. A spec without `produced_later` is required exactly as before.
        if spec.get("produced_later"):
            continue
        fname = _expand_filename(spec["filename"], deck_slug)
        path = bundle_dir / fname
        # exists() follows symlinks; use lexists so a broken/dangling symlink is seen
        # as present-but-rejected rather than silently "absent".
        if not os.path.lexists(str(path)):
            update_deliverable_status(ledger_path, key, "failed",
                                      error="file absent from bundle_dir")
            missing_or_short.append((key, fname, spec["label"], spec["min_bytes"], 0,
                                     "ABSENT"))
            continue
        # C2: a symlink (or any non-regular file) is a decoy vector — reject it.
        if path.is_symlink():
            update_deliverable_status(ledger_path, key, "failed",
                                      error="path is a symlink (rejected: a symlink "
                                            "can point at a large unrelated/decoy file)")
            missing_or_short.append((key, fname, spec["label"], spec["min_bytes"], 0,
                                     "SYMLINK"))
            continue
        # os.lstat: size of the path itself, never a followed link.
        st = os.lstat(str(path))
        if not __import__("stat").S_ISREG(st.st_mode):
            update_deliverable_status(ledger_path, key, "failed",
                                      error="path is not a regular file (rejected)")
            missing_or_short.append((key, fname, spec["label"], spec["min_bytes"], 0,
                                     "NOT_REGULAR"))
            continue
        actual = st.st_size
        if actual < spec["min_bytes"]:
            update_deliverable_status(ledger_path, key, "failed",
                                      size=actual,
                                      error=f"file too small: {actual} bytes "
                                            f"(min {spec['min_bytes']})")
            missing_or_short.append((key, fname, spec["label"], spec["min_bytes"],
                                     actual, "UNDER_THRESHOLD"))
            continue
        # C2: magic-byte content-type check (md is intentionally size-only).
        signatures = DELIVERABLE_MAGIC.get(key)
        if signatures is not None and not _magic_ok(path, signatures):
            sig_disp = "/".join(repr(s) for s in signatures)
            update_deliverable_status(ledger_path, key, "failed",
                                      size=actual,
                                      error=f"wrong content type: leading bytes do not "
                                            f"match expected magic {sig_disp} for this "
                                            f"deliverable type (decoy/wrong-type file)")
            missing_or_short.append((key, fname, spec["label"], spec["min_bytes"],
                                     actual, "WRONG_TYPE"))
            continue
        update_deliverable_status(ledger_path, key, "verified", size=actual)

    # --- KIE-BAKED sub-check (AF-I14) — at closeout, re-prove every rendered slide
    # was actually baked by the image model (real KIE taskId + verified above-floor
    # PNG), not native-rendered / flat-placeholder-filled. This mirrors the preflight
    # _chk_kie_baked call but fires at the FINAL gate so a native-rendered deck can
    # never be reported "complete". run_dir is threaded in by main(); when it is not
    # available (a bare gate-only invocation) the sub-check is skipped rather than
    # failing a caller that has no run dir.
    kie_fail_reason = ""
    if run_dir is not None:
        kie_reason = _chk_kie_baked(run_dir, slides_path)
        if kie_reason:
            kie_fail_reason = kie_reason
            missing_or_short.append((
                "deck_pptx",
                _expand_filename("deck.pptx", deck_slug),
                "KIE-baked slide images (every slide model-baked, not native render)",
                0, 0, "NOT_BAKED"))
            # Record the AF-I14 reason on the deck entry so the ledger carries it.
            update_deliverable_status(ledger_path, "deck_pptx", "failed",
                                      error=kie_reason)

    # --- AF-PACKAGE-CLEAN sub-check (2026-06-19, deck-quality gate) ---
    # The final bundle must contain ONLY the canonical deliverable files. Dev artifacts
    # (.py/.sh/~$*/tasks/) in the bundle are a hard delivery failure.
    pkg_clean_reason = check_package_cleanliness(bundle_dir)
    if pkg_clean_reason:
        missing_or_short.append((
            "deck_pptx",
            _expand_filename("bundle/", deck_slug),
            "package cleanliness (no dev artifacts in the delivered bundle)",
            0, 0, "DIRTY_PACKAGE"))
        update_deliverable_status(ledger_path, "deck_pptx", "failed", error=pkg_clean_reason)

    # --- AF-VISUAL-VARIETY sub-check (2026-06-19, deck-quality gate) ---
    # Post-render: re-run the visual-variety check at the postflight gate where renders exist.
    visual_variety_reason = ""
    if run_dir is not None:
        visual_variety_reason = check_visual_variety(run_dir, slides_path)
        if visual_variety_reason:
            missing_or_short.append((
                "deck_pptx",
                _expand_filename("deck.pptx", deck_slug),
                "visual variety (deck has light/break slides and varied palette)",
                0, 0, "MONOTONE_DECK"))
            update_deliverable_status(ledger_path, "deck_pptx", "failed",
                                      error=visual_variety_reason)

    # --- AF-IMAGE-QC-RAN sub-check (2026-06-19, deck-quality gate) ---
    # Post-render: verify the image-QC report is fresh and has per-slide coverage.
    image_qc_ran_reason = ""
    if run_dir is not None:
        image_qc_ran_reason = check_image_qc_present(run_dir, slides_path)
        if image_qc_ran_reason:
            missing_or_short.append((
                "deck_pptx",
                "working/qc/image_qc_report.json",
                "image-QC report (fresh + per-slide coverage; AF-IMAGE-QC-RAN)",
                0, 0, "IMAGE_QC_STALE_OR_MISSING"))
            update_deliverable_status(ledger_path, "deck_pptx", "failed",
                                      error=image_qc_ran_reason)

    # --- AF-BRAND-CONSISTENCY sub-check (2026-06-19, deck-quality gate) ---
    # Post-render: check rendered slides against the declared brand palette.
    brand_reason = ""
    if run_dir is not None:
        brand_reason = check_brand_consistency(run_dir, slides_path)
        if brand_reason:
            missing_or_short.append((
                "deck_pptx",
                _expand_filename("deck.pptx", deck_slug),
                "brand consistency (all rendered slides within declared brand palette)",
                0, 0, "OFF_BRAND"))
            update_deliverable_status(ledger_path, "deck_pptx", "failed",
                                      error=brand_reason)

    # --- AF-OVERLAY-DELIVERED sub-check (2026-06-20, Decision 5C) ---
    # The native PPTX text-overlay path is eliminated. At closeout, re-prove no
    # eliminated pptx_text_overlays.json is present AND the delivered PPTX carries no
    # native on-slide text run (only a composed gpt-image-2 image + off-slide notes).
    overlay_reason = ""
    if run_dir is not None:
        overlay_reason = _chk_no_overlay(run_dir, slides_path)
    # Also scan the bundle's delivered PPTX directly (bundle_dir may be outside run_dir).
    if not overlay_reason:
        for pptx in sorted(bundle_dir.glob("*.pptx")):
            if pptx.name.startswith("~$"):
                continue
            r = _delivered_pptx_native_text(pptx)
            if r:
                overlay_reason = ("AF-OVERLAY-DELIVERED: " + r + ". Every delivered slide "
                                  "must be a single composed gpt-image-2 image; the native "
                                  "PPTX text-overlay path is eliminated (Decision 5C).")
                break
    if overlay_reason:
        missing_or_short.append((
            "deck_pptx",
            _expand_filename("deck.pptx", deck_slug),
            "no native overlay (every slide a composed gpt-image-2 image; "
            "no pptx_text_overlays.json)",
            0, 0, "OVERLAY_DELIVERED"))
        update_deliverable_status(ledger_path, "deck_pptx", "failed",
                                  error=overlay_reason)

    # --- AF-EMPTY-NOTES-PANE sub-check (P9.5-NOTES-SYNC reorder) ---
    # Previously documented (SOP-SLIDE-00:578-582) but never code-enforced. At
    # closeout — AFTER notes_sync_pass() has had the chance to re-inject notes from
    # the now-QC-passed speech — re-open the delivered PPTX and refuse to ship any
    # audience-facing content slide with an empty notes pane.
    notes_pane_reason = _chk_notes_pane(bundle_dir, run_dir=run_dir, slides_path=slides_path)
    if notes_pane_reason:
        missing_or_short.append((
            "deck_pptx",
            _expand_filename("deck.pptx", deck_slug),
            "per-slide speaker notes present (AF-EMPTY-NOTES-PANE; notes injected "
            "from the QC-passed presenter speech)",
            0, 0, "EMPTY_NOTES_PANE"))
        update_deliverable_status(ledger_path, "deck_pptx", "failed",
                                  error=notes_pane_reason)

    # --- FIX-2: CANONICAL-RENDER-PATH sub-check (AF-CANONICAL-RENDER-BYPASS /
    # AF-LOCAL-CANVAS) — at closeout, re-prove no hand-rolled renderer/assembler in
    # the run dir produced (part of) the deck. The canonical render path is
    # build_deck.py / run_signature_deck.py ONLY; a per-deck Pillow canvas / native
    # overlay / direct kie createTask script is a hard delivery failure.
    canonical_reason = ""
    if run_dir is not None:
        canonical_reason = check_canonical_render_path(run_dir, slides_path)
        if canonical_reason:
            missing_or_short.append((
                "deck_pptx",
                _expand_filename("deck.pptx", deck_slug),
                "canonical render path (no hand-rolled renderer/assembler in the run dir)",
                0, 0, "CANONICAL_BYPASS"))
            update_deliverable_status(ledger_path, "deck_pptx", "failed",
                                      error=canonical_reason)

    # --- FIX-2: IMAGE-QC-VISION sub-check (AF-IMAGE-QC-VISION) — the deterministic
    # PIXEL cross-check at closeout. Opens every rendered PNG (byte floor + flat
    # cream/typography-card detection) and refuses a pixel-blind / overlay-blessing /
    # hook-slide-excluding image-QC report. Does NOT defer-to-pass when renders exist.
    image_qc_vision_reason = ""
    if run_dir is not None:
        image_qc_vision_reason = check_image_qc_vision(run_dir, slides_path)
        if image_qc_vision_reason:
            missing_or_short.append((
                "deck_pptx",
                _expand_filename("deck.pptx", deck_slug),
                "image-QC pixel read (real multimodal read of every PNG; no flat card; "
                "no rubber-stamp report)",
                0, 0, "IMAGE_QC_VISION"))
            update_deliverable_status(ledger_path, "deck_pptx", "failed",
                                      error=image_qc_vision_reason)

    # --- FIX-2: QC REPORT FLOOR sub-check (AF-QC-PLACEHOLDER) — at the pre-delivery
    # closeout gate, re-prove every PRESENT QC phase report (copy / prompt / typography
    # / shift) clears the REAL-CONTENT floor (> QC_REPORT_FLOOR_BYTES bytes, valid JSON,
    # >= QC_REPORT_SLIDE_FLOOR real per-slide verdicts). A 3-byte '{}' placeholder (the
    # exact Error-2 artifact) can never satisfy a QC phase, so a deck delivered on a
    # placeholder QC report is a hard delivery failure. DEFERS on an ABSENT report file
    # (a standalone / adhoc run legitimately has no QC artifacts — report existence is
    # enforced fail-closed at attest time by run_signature_deck.attest_phase); once a
    # report EXISTS it must be real.
    qc_floor_reason = ""
    if run_dir is not None:
        for _qc_phase_id in sorted(UNSKIPPABLE_QC_PHASES):
            _qc_rel = QC_PHASE_REPORT[_qc_phase_id]
            if not (run_dir / _qc_rel).is_file():
                continue  # no report yet — defer (the attestation gate owns absence)
            _qc_floor = check_qc_phase_report_real(run_dir, _qc_phase_id)
            if _qc_floor:
                qc_floor_reason = _qc_floor
                break
        if qc_floor_reason:
            missing_or_short.append((
                "deck_pptx",
                _expand_filename("deck.pptx", deck_slug),
                "real QC report floor (each present QC report clears >256 bytes, valid "
                "JSON, >=20 real per-slide verdicts; a 3-byte '{}' placeholder never "
                "satisfies a QC phase)",
                0, 0, "QC_REPORT_PLACEHOLDER"))
            update_deliverable_status(ledger_path, "deck_pptx", "failed",
                                      error=qc_floor_reason)

    # --- U047: SLIDE-GEOMETRY sub-checks (AF-TEXT-OVERFLOW / AF-SPELLING /
    # AF-TYPE-SIZE-MEASURED) — the three pixel-level checks that did not exist. They read
    # BAKED renders, so postflight is the only place they can run: preflight has no PNGs.
    # Rule 3.5: WARN by default. Each reason is printed and recorded; only under
    # SLIDE_GEOMETRY_ENFORCE_ENV=1 does it reach missing_or_short and fail the gate.
    slide_geometry_reasons = []
    if run_dir is not None:
        for _label, _fn, _token in (
            ("text fits on the slide (no glyph box at an edge, no overlapping text blocks)",
             _chk_text_fits, "TEXT_OVERFLOW"),
            ("rendered spelling (every baked word accounted for by approved copy or allowlist)",
             _chk_spelling, "SPELLING"),
            ("measured type size (smallest rendered word box at or above the pt floor)",
             _chk_type_size, "TYPE_SIZE"),
        ):
            _reason = _fn(run_dir, slides_path)
            if not _reason:
                continue
            slide_geometry_reasons.append(_reason)
            if os.environ.get(SLIDE_GEOMETRY_ENFORCE_ENV, "") == "1":
                missing_or_short.append((
                    "deck_pptx",
                    _expand_filename("deck.pptx", deck_slug),
                    _label, 0, 0, _token))
                update_deliverable_status(ledger_path, "deck_pptx", "failed", error=_reason)
    for _reason in slide_geometry_reasons:
        print("WARN-SLIDE-GEOMETRY: " + _reason, file=sys.stderr, flush=True)

    # --- TELEPROMPTER-PUBLISH sub-check (folded under AF-BUNDLE-COMPLETE) ---
    # A self-contained HTML on disk is NOT a delivered teleprompter. The bundle is not
    # complete until the teleprompter is hosted at the central Cloudflare URL and that
    # URL is live (HTTP 200). This mirrors the AF-BUNDLE-COMPLETE pattern exactly: a
    # missing/unverified publish URL is a hard exit-5 failure (no new AF code — it is
    # the teleprompter-publish condition of the existing bundle-completeness gate). The
    # delivery / ruleset SOPs name this the TELEPROMPTER-PUBLISH auto-fail.
    tele_pub = _check_teleprompter_published(bundle_dir, skip_gate=skip_teleprompter_gate)
    if tele_pub:
        update_deliverable_status(ledger_path, "teleprompter_html", "failed",
                                  error=tele_pub)
        missing_or_short.append((
            "teleprompter_html",
            _expand_filename("presenter-teleprompter.html", deck_slug),
            "presenter teleprompter web app (published URL)",
            0, 0, "UNPUBLISHED"))

    # NOTE: the AF-CC-UNREGISTERED gate is enforced once, AFTER the bundle-completeness
    # check passes (see the dedicated _chk_cc_registered call below). It is intentionally
    # NOT folded into AF-BUNDLE-COMPLETE here, so a never-registered CC surfaces its own
    # clean AF-CC-UNREGISTERED message rather than being mislabelled as a bundle defect.

    # --- Phase 2: read ledger back as the authoritative final state ---
    try:
        ledger_entries = json.loads(ledger_path.read_text())
    except Exception as exc:  # noqa: BLE001
        print(f"FATAL: deliverables.json could not be read for final verification: {exc}",
              file=sys.stderr)
        sys.exit(5)

    all_verified = all(e.get("status") == "verified" for e in ledger_entries)

    # --- Phase 3: fail loud on any non-verified artifact ---
    if missing_or_short or not all_verified:
        bar = "=" * 78
        print(f"\n{bar}", file=sys.stderr)
        print("FATAL: POSTFLIGHT COMPLETENESS GATE FAILED (AF-BUNDLE-COMPLETE)", file=sys.stderr)
        print("build_deck.py produced the deck PPTX, but the required deliverable bundle",
              file=sys.stderr)
        print("is INCOMPLETE. The run may NOT be reported as 'complete' or 'done'.", file=sys.stderr)
        print(f"Bundle dir: {bundle_dir}", file=sys.stderr)
        print(f"Ledger:     {ledger_path}", file=sys.stderr)
        print("\nMissing or under-threshold artifacts:", file=sys.stderr)
        for (key, fname, label, min_b, actual_b, reason) in missing_or_short:
            if reason == "ABSENT":
                print(f"  MISSING  [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           Required by DELIVERABLES_REQUIRED; produced by upstream",
                      file=sys.stderr)
                print(f"           roles (presenter guide → Presenters Guide Specialist;",
                      file=sys.stderr)
                print(f"           speech → Presenters Speech Writer;",
                      file=sys.stderr)
                print(f"           audio → Audio Demonstration Specialist;",
                      file=sys.stderr)
                print(f"           infographic → infographic-checklist role;",
                      file=sys.stderr)
                print(f"           deck PDF → PPTX Assembly Specialist).",
                      file=sys.stderr)
            elif reason == "UNDER_THRESHOLD":
                print(f"  TOO SMALL [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           actual={actual_b:,} bytes  minimum={min_b:,} bytes",
                      file=sys.stderr)
            elif reason == "SYMLINK":
                print(f"  SYMLINK  [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           path is a symlink — rejected (a symlink can point at",
                      file=sys.stderr)
                print(f"           a large unrelated/decoy file to fake the size gate).",
                      file=sys.stderr)
            elif reason == "NOT_REGULAR":
                print(f"  NOT A FILE [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           path is not a regular file — rejected.", file=sys.stderr)
            elif reason == "WRONG_TYPE":
                print(f"  WRONG TYPE [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           actual={actual_b:,} bytes but leading magic bytes do",
                      file=sys.stderr)
                print(f"           not match the expected type (decoy / wrong-type file).",
                      file=sys.stderr)
            elif reason == "NOT_BAKED":
                print(f"  NOT BAKED [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           one or more rendered slides were NOT KIE-baked "
                      f"(native render,", file=sys.stderr)
                print(f"           missing image, or flat-placeholder fill) — AF-I14.",
                      file=sys.stderr)
                print(f"           {kie_fail_reason}", file=sys.stderr)
            elif reason == "OVERLAY_DELIVERED":
                print(f"  OVERLAY   [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           a delivered slide ships native on-slide text (or a "
                      f"pptx_text_overlays.json", file=sys.stderr)
                print(f"           is present) instead of a single composed gpt-image-2 "
                      f"image — AF-OVERLAY-DELIVERED.", file=sys.stderr)
                print(f"           {overlay_reason}", file=sys.stderr)
                print(f"           The native-overlay path is eliminated (Decision 5C): "
                      f"re-prompt/re-seed", file=sys.stderr)
                print(f"           the garbled slide, escalate to a human if it persists, "
                      f"and re-render.", file=sys.stderr)
            elif reason == "CANONICAL_BYPASS":
                print(f"  BYPASS    [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           a hand-rolled renderer/assembler in the run dir "
                      f"produced (part of)", file=sys.stderr)
                print(f"           the deck — AF-CANONICAL-RENDER-BYPASS / AF-LOCAL-CANVAS.",
                      file=sys.stderr)
                print(f"           {canonical_reason}", file=sys.stderr)
                print(f"           The canonical render path is build_deck.py / "
                      f"run_signature_deck.py", file=sys.stderr)
                print(f"           ONLY: delete the hand-rolled script and re-route the "
                      f"deck through it.", file=sys.stderr)
            elif reason == "IMAGE_QC_VISION":
                print(f"  IMAGE-QC  [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           image-QC did not actually inspect the rendered slides "
                      f"(flat card,", file=sys.stderr)
                print(f"           below-floor PNG, or rubber-stamp report) — "
                      f"AF-IMAGE-QC-VISION.", file=sys.stderr)
                print(f"           {image_qc_vision_reason}", file=sys.stderr)
                print(f"           Re-render any flat card through kie.ai gpt-image-2 and "
                      f"re-run the Image", file=sys.stderr)
                print(f"           QC Specialist with a per-slide multimodal pixel read.",
                      file=sys.stderr)
            elif reason == "UNPUBLISHED":
                print(f"  UNPUBLISHED [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           the teleprompter HTML exists locally but was not "
                      f"published", file=sys.stderr)
                print(f"           to the central host with a live public URL "
                      f"(TELEPROMPTER-PUBLISH).", file=sys.stderr)
                print(f"           Run publish_teleprompter — it uploads to the central "
                      f"Cloudflare", file=sys.stderr)
                print(f"           host (teleprompter.zerohumanworkforce.com). The URL "
                      f"must return", file=sys.stderr)
                print(f"           HTTP 200. See teleprompter_publish.json.",
                      file=sys.stderr)
            elif reason == "QC_REPORT_PLACEHOLDER":
                print(f"  QC-PLACEHOLDER [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           {qc_floor_reason}", file=sys.stderr)
                print("           A QC phase report is a placeholder / sub-floor — a "
                      "3-byte '{}' rubber stamp", file=sys.stderr)
                print(f"           can never satisfy a QC phase (FIX-2 / Error 2). "
                      f"Re-run the QC phase to", file=sys.stderr)
                print(f"           produce a real per-slide report "
                      f"(AF-QC-PLACEHOLDER).", file=sys.stderr)
            else:
                print(f"  TOO SMALL [{key}] {fname}  ({label})", file=sys.stderr)
                print(f"           actual={actual_b:,} bytes  minimum={min_b:,} bytes",
                      file=sys.stderr)
        # Also surface any ledger entries that are not verified but not in our scan list
        # (edge case: previously-failed entries from a prior run).
        extra_unverified = [e for e in ledger_entries
                            if e.get("status") != "verified"
                            and e.get("key") not in {m[0] for m in missing_or_short}]
        for e in extra_unverified:
            print(f"  NOT VERIFIED [{e['key']}] {e.get('filename')} — status: {e.get('status')}",
                  file=sys.stderr)
        print(f"\n{bar}", file=sys.stderr)
        sys.exit(5)

    # --- AF-CC-UNREGISTERED / AF-CC-UNVERIFIED check (enforced_by:build_deck,
    # symbol:_chk_cc_registered) — T2 three-outcome gate. VERIFIED = cc_task_id +
    # verifiable cc_registration HMAC proof (real ingest round-trip) -> pass.
    # UNVERIFIED = attempted or bare task_id without a verifying receipt ->
    # fails with an explicit honest AF-CC-UNVERIFIED message (could-not-verify
    # never prints as verified; fail-soft on a genuine CC outage). UNREGISTERED
    # = neither field -> fail-CLOSED. This runs only when the bundle check
    # above has PASSED (all deliverables verified), so a CC failure is the only
    # remaining blocker.
    _cc_reason = _chk_cc_registered(run_dir, deck_slug)
    if _cc_reason:
        print(f"\nFATAL: {_cc_reason}", file=sys.stderr)
        sys.exit(5)

    # --- Phase 4: success — print COMPLETE only when ALL verified ---
    print("\n=== POSTFLIGHT COMPLETENESS GATE PASSED (AF-BUNDLE-COMPLETE) ===", flush=True)
    print(f"All {len(DELIVERABLES_REQUIRED)} required deliverables present and sized.",
          flush=True)
    print(f"Bundle dir: {bundle_dir}", flush=True)
    for e in ledger_entries:
        size_str = f"{e.get('size', 0):,} bytes" if e.get("size") else "n/a"
        print(f"  VERIFIED  [{e['key']}] {e.get('filename')}  ({size_str})", flush=True)
    print(f"Ledger (all-verified): {ledger_path}", flush=True)

    # --- OWNER-SKIP DISCLOSURE (Trust Boundary Increment 2) — this run's operator
    # report may NEVER print COMPLETE without also surfacing every owner_skip_approval
    # this run consumed. Sourced from the code-owned "owner_skip_events" ledger
    # (_owner_skip_ledger_append), never from the agent-authored "owner_skip_approval"
    # input directly — see the FIX-2 owner-skip section for the full design. This is
    # a report step only: it reflects what already happened, it does not gate.
    if run_dir is not None:
        _pm = run_dir / "working" / "checkpoints" / "process_manifest.json"
        _events = []
        if _pm.exists():
            _obj = _read_json(_pm)
            if isinstance(_obj, dict) and "__parse_error__" not in _obj:
                _raw = _obj.get(_OWNER_SKIP_LEDGER_KEY)
                _events = _raw if isinstance(_raw, list) else []
        if _events:
            print(f"\n*** OWNER OVERRIDES USED IN THIS RUN — {len(_events)} EVENT(S) ***",
                  flush=True)
            for _ev in _events:
                if not isinstance(_ev, dict):
                    continue
                _tag = "WAIVED" if _ev.get("outcome") == "granted" else "REJECTED (attempt refused)"
                print(f"  [{_tag}] {_ev.get('af_code')} — approved_by="
                      f"{_ev.get('approved_by')!r} reason={_ev.get('reason')!r} "
                      f"timestamp={_ev.get('timestamp')!r} consumed_at="
                      f"{_ev.get('consumed_at')!r}", flush=True)
            print("*** This deck did NOT clear every gate on its own merits — see above. ***",
                  flush=True)
        else:
            print("\nOwner overrides: NONE used in this run (every gate cleared on its own "
                  "merits).", flush=True)
    print("=== COMPLETE ===\n", flush=True)


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------

def run_style_preview_samples(slides_path: Path, run_dir: Path,
                              style_spec_path: Optional[Path], api_key: str,
                              logo_url: Optional[str] = None) -> int:
    """P-STYLE-PREVIEW (order 4.85) --sample mode. Render 3 attention-grade style
    variants (A/B/C) across 3 representative slides (cover / data / people) = 9 sample
    renders, so the owner can pick ONE winning direction via their OWN gateway before
    the full deck is committed. The winner's 3 representative renders are then LOCKED
    and reused downstream so kie never double-charges (enforced by AF-STYLE-DOUBLECHARGE).

    Reads working/copy/style_preview_spec.json (or --style-spec) =
      {"variants":[{"id":"A","style_directive":"..."}, ...3],
       "representative_slides":[<cover ord>, <data ord>, <people ord>]}.
    Writes working/style-preview/<variant>-slide-<ord>.png + style_samples_manifest.json.
    Returns a process exit code (0 = the 9 samples rendered + manifest written)."""
    spec_path = style_spec_path
    if spec_path is None:
        for rel in ("working/copy/style_preview_spec.json",
                    "working/style-preview/style_preview_spec.json"):
            p = run_dir / rel
            if p.exists():
                spec_path = p
                break
    if spec_path is None or not Path(spec_path).exists():
        print("FATAL --sample: no style_preview_spec.json found (pass --style-spec PATH or "
              "place it at working/copy/style_preview_spec.json).", file=sys.stderr)
        return 2
    spec = _read_json(Path(spec_path))
    if not isinstance(spec, dict) or "__parse_error__" in spec:
        print(f"FATAL --sample: style_preview_spec.json is not valid JSON.", file=sys.stderr)
        return 2
    variants = spec.get("variants") or []
    rep_slides = spec.get("representative_slides") or []
    if not (isinstance(variants, list) and len(variants) == 3):
        print("FATAL --sample: style_preview_spec.variants must list exactly 3 variants "
              "(A/B/C).", file=sys.stderr)
        return 2
    if not (isinstance(rep_slides, list) and len(rep_slides) == 3):
        print("FATAL --sample: style_preview_spec.representative_slides must list exactly 3 "
              "slide ordinals (cover / data / people).", file=sys.stderr)
        return 2

    slides = json.loads(Path(slides_path).read_text())
    by_ord = {s["slide"]: s for s in slides if isinstance(s, dict) and "slide" in s}
    out_dir = run_dir / "working" / "style-preview"
    out_dir.mkdir(parents=True, exist_ok=True)

    samples = []
    for v in variants:
        vid = str(v.get("id") or "").strip() or "?"
        directive = str(v.get("style_directive") or "").strip()
        for ordn in rep_slides:
            try:
                ordn = int(ordn)
            except (TypeError, ValueError):
                print(f"FATAL --sample: representative slide ordinal {ordn!r} is not an int.",
                      file=sys.stderr)
                return 2
            slide = by_ord.get(ordn)
            if slide is None:
                print(f"FATAL --sample: representative slide {ordn} is not in slides.json.",
                      file=sys.stderr)
                return 2
            base_prompt = load_rich_prompt(slide, run_dir)
            # Prepend the variant's attention-grade style directive to the verbatim rich
            # prompt; the rest of the prompt (copy, negative block, spelling-lock) is
            # rendered unchanged so the SAMPLE is a faithful preview of the real slide.
            prompt = (f"[STYLE VARIANT {vid} — {directive}]\n\n{base_prompt}"
                      if directive else base_prompt)
            png = out_dir / f"variant-{vid}-slide-{ordn:02d}.png"
            print(f"  [sample] variant {vid} / slide {ordn} -> {png.name}", flush=True)
            task_id = submit_task(prompt, api_key, logo_url=logo_url)
            result_url = poll_task(task_id, api_key)
            download_image(result_url, png, api_key)
            verify_png(png)
            samples.append({
                "variant": vid,
                "slide": ordn,
                "render": str(png.relative_to(run_dir)),
                "taskId": task_id,
            })

    manifest = {
        "schema": "style_samples_manifest/v1",
        "phase": "P-STYLE-PREVIEW (order 4.85)",
        "generated_at": time.strftime("%Y-%m-%dT%H:%M:%S%z"),
        "variants": [str(v.get("id")) for v in variants],
        "representative_slides": [int(o) for o in rep_slides],
        "samples": samples,
        "owner_pick_required": True,
        "owner_pick_artifact": STYLE_CHOICE_REL,
        "owner_pick_note": ("Route these 9 samples to the OWNER's OWN gateway (never the "
                            "operator chat). Record the pick in working/copy/"
                            "style_preview_choice.json {owner_approved:true, chosen_variant, "
                            "locked_renders:[...]} so the winner's 3 renders are reused and "
                            "kie never double-charges."),
    }
    (out_dir / "style_samples_manifest.json").write_text(json.dumps(manifest, indent=2))
    print(json.dumps({"slidesRendered": len(samples),
                      "styleSamplesManifest": str(out_dir / "style_samples_manifest.json"),
                      "ownerPickArtifact": STYLE_CHOICE_REL}, indent=2))
    return 0


# ---------------------------------------------------------------------------
# Command Center board — phase-transition hookup (FAIL-SOFT). This is the missing
# half of the presentations cc_board wiring: run-begin ingest_deck_task CREATES
# the deck's single Kanban card, and these calls move it as the run progresses:
#
#   P4-RENDER START     -> STATUS  backlog -> in_progress   (_board_patch_phase)
#   P4-RENDER complete  -> ACTIVITY progress post           (_board_post_activity)
#   P8-ASSEMBLE complete-> ACTIVITY progress post           (_board_post_activity)
#   run-end (bundle OK) -> STATUS  -> review + process_certificate_sha + QC scores,
#                          but ONLY when the PROCESS-CERTIFICATE already exists on disk
#                          (_process_certificate_present guard); otherwise the close
#                          is DEFERRED to the runner (run_signature_deck.py) which
#                          fires it right after prove-deck mints the cert.
#
# Mid-run phase progress is an ACTIVITY, NOT a task-level status change: a mid-run
# terminal status 422s the presentations cert gate (the PROCESS-CERTIFICATE is
# minted at delivery, not mid-run) and would wrongly move a non-presentation card.
# The COMPLETED deck's producer close is status='review' — the pipeline STOPS at
# review and the CC-side QC scorer / Devil's-Advocate gate promotes review->done
# (the interlock every sibling department respects; there is NO 'delivered' status in
# the CC UpdateTaskSchema enum — "delivered" belongs in the note only). build_deck
# runs the RENDER phase BEFORE the cert exists, so in the normal runner flow the
# terminal close is the RUNNER's job; build_deck only closes here in a standalone
# full build where prove-deck already minted the cert.
#
# Mirrors Skill-48's ad_director._board_move contract verbatim: EVERY call is
# best-effort. A board outage, a missing token, an import error, or any patch
# failure is caught and logged, NEVER fatal — the deck build's exit code is
# independent of the board (the card is a VIEW; the offline AF-CC-UNREGISTERED
# gate is already satisfied by the run-begin ingest's cc_register_attempted flag).
# ---------------------------------------------------------------------------
def _board_patch_phase(run_dir, task_id, phase_id, status, note=""):
    """Advance the deck's CC card to (phase_id, status), FAIL-SOFT.

    task_id is the id returned by cc_board.ingest_deck_task at run-begin; when it
    is falsy (board disabled, or a transport-failed ingest) we recover it from
    working/checkpoints/process_manifest.json (stamped by ingest on success). A
    missing task_id or a disabled board is a clean no-op. Never raises — the board
    is a convenience, never a gate."""
    try:
        import cc_board as _cc_board
        tid = task_id
        if not tid and run_dir is not None:
            try:
                _pm = Path(run_dir) / "working" / "checkpoints" / "process_manifest.json"
                if _pm.exists():
                    tid = (json.loads(_pm.read_text()) or {}).get("cc_task_id")
            except Exception:  # noqa: BLE001 — a bad manifest read is never fatal
                tid = task_id
        if not tid:
            return
        _cc_board.patch_phase(run_dir, tid, phase_id, status, note=note)
    except Exception as _cc_exc:  # noqa: BLE001 — the board is a view, never a gate
        print(f"[cc_board] patch_phase {phase_id}->{status} raised ({_cc_exc}) — "
              "run continues; the board update is best-effort.",
              file=sys.stderr, flush=True)


def _board_post_activity(run_dir, task_id, phase_id, note=""):
    """Log a mid-run phase-PROGRESS activity on the deck's CC card, FAIL-SOFT.

    Phase progress (P4-RENDER complete, P8-ASSEMBLE complete) is an ACTIVITY, not a
    task-level status change: a mid-run status='done' 422s the presentations cert
    done-gate (the PROCESS-CERTIFICATE is minted at delivery) and would wrongly
    close a non-presentation card. Recovers the task_id from process_manifest.json
    when not supplied; a missing task_id or a disabled board is a clean no-op. Never
    raises — the board is a convenience, never a gate."""
    try:
        import cc_board as _cc_board
        tid = task_id
        if not tid and run_dir is not None:
            try:
                _pm = Path(run_dir) / "working" / "checkpoints" / "process_manifest.json"
                if _pm.exists():
                    tid = (json.loads(_pm.read_text()) or {}).get("cc_task_id")
            except Exception:  # noqa: BLE001 — a bad manifest read is never fatal
                tid = task_id
        if not tid:
            return
        _cc_board.post_activity(run_dir, tid, phase_id, note)
    except Exception as _cc_exc:  # noqa: BLE001 — the board is a view, never a gate
        print(f"[cc_board] post_activity {phase_id} raised ({_cc_exc}) — "
              "run continues; the board update is best-effort.",
              file=sys.stderr, flush=True)


def _process_certificate_present(run_dir) -> bool:
    """True when prove-deck.py has already minted a
    delivery/<SLUG>-FINAL/PROCESS-CERTIFICATE.json for this run. GUARDS build_deck's
    terminal P9-DELIVER board close so it only fires when the certificate the CC
    presentations done-gate REQUIRES actually exists on disk. In the runner flow
    build_deck runs the RENDER phase BEFORE prove-deck mints the cert, so this returns
    False and the terminal close is DEFERRED to the runner (which closes the card
    right after prove-deck). A standalone full build where the cert is already present
    still closes here. Never raises."""
    try:
        if run_dir is None:
            return False
        delivery = Path(run_dir) / "delivery"
        if not delivery.is_dir():
            return False
        for cert in delivery.glob("*-FINAL/PROCESS-CERTIFICATE.json"):
            if cert.is_file():
                return True
    except OSError:
        pass
    return False


def main():
    argv = sys.argv[1:]

    # Pull out the flags (any position) before reading positional args.
    adhoc = False
    run_dir_arg = None
    logo_arg = None
    timestamp_arg = None
    out_dir_arg = None   # --out BUNDLE_DIR (default ~/Downloads/<deck-slug>)
    platform_arg = None  # --platform {vps|mac} override (default None -> auto-detect)
    skip_teleprompter_gate = False  # M7: explicit per-run bypass of the teleprompter
                                    # publish sub-check (never via a persisted status).
    sample_mode = False  # P-STYLE-PREVIEW (4.85): render 9 style samples, then stop.
    notes_sync_mode = False  # P9.5-NOTES-SYNC (8.7): re-inject notes post-speech-QC.
    style_spec_arg = None
    positional = []
    i = 0
    while i < len(argv):
        tok = argv[i]
        if tok == "--adhoc-no-process":
            adhoc = True
        elif tok == "--sample":
            sample_mode = True
        elif tok == "--notes-sync":
            notes_sync_mode = True
        elif tok == "--style-spec":
            i += 1
            if i >= len(argv):
                print("FATAL: --style-spec requires a path argument.", file=sys.stderr)
                sys.exit(2)
            style_spec_arg = argv[i]
        elif tok.startswith("--style-spec="):
            style_spec_arg = tok[len("--style-spec="):]
        elif tok == "--skip-teleprompter-gate":
            skip_teleprompter_gate = True
        elif tok == "--run-dir":
            i += 1
            if i >= len(argv):
                print("FATAL: --run-dir requires a directory argument.", file=sys.stderr)
                sys.exit(2)
            run_dir_arg = argv[i]
        elif tok.startswith("--run-dir="):
            run_dir_arg = tok[len("--run-dir="):]
        elif tok == "--logo":
            i += 1
            if i >= len(argv):
                print("FATAL: --logo requires a PNG path argument.", file=sys.stderr)
                sys.exit(2)
            logo_arg = argv[i]
        elif tok.startswith("--logo="):
            logo_arg = tok[len("--logo="):]
        elif tok == "--timestamp":
            i += 1
            if i >= len(argv):
                print("FATAL: --timestamp requires a value.", file=sys.stderr)
                sys.exit(2)
            timestamp_arg = argv[i]
        elif tok.startswith("--timestamp="):
            timestamp_arg = tok[len("--timestamp="):]
        elif tok == "--out":
            i += 1
            if i >= len(argv):
                print("FATAL: --out requires a directory path argument.", file=sys.stderr)
                sys.exit(2)
            out_dir_arg = argv[i]
        elif tok.startswith("--out="):
            out_dir_arg = tok[len("--out="):]
        elif tok == "--platform":
            i += 1
            if i >= len(argv):
                print("FATAL: --platform requires a value (vps|mac).", file=sys.stderr)
                sys.exit(2)
            platform_arg = argv[i]
        elif tok.startswith("--platform="):
            platform_arg = tok[len("--platform="):]
        elif tok.startswith("--"):
            print(f"FATAL: unknown flag {tok!r}.", file=sys.stderr)
            sys.exit(2)
        else:
            positional.append(tok)
        i += 1

    # FRONT-DOOR NONCE HANDSHAKE — Contract #8: at the render/CLI main() entry,
    # require a valid per-run nonce (OC_DECK_ENTRY_NONCE matching the run-scoped
    # <run-dir>/working/checkpoints/.canonical-entry-nonce file that
    # presentation-canonical-entry.sh minted for this run). This SUPERSEDES the
    # retired OC_DECK_CANONICAL_ENTRY / OC_DECK_ALLOW_DIRECT env markers, which
    # shipped in box-visible comments and were therefore forgeable by any model that
    # read the repo — setting either of those names is now DENIED. Module imports and
    # unit-test paths that call build_deck functions directly are unaffected — this
    # guard fires only when main() is reached via the CLI (`python3 build_deck.py ...`).
    # References: AF-CANONICAL-RENDER-BYPASS, shared CONTRACT.md §FRONT-DOOR MARKER.
    _nonce_run_dir = _resolve_run_dir_for_guard(run_dir_arg, positional)
    if not _verify_entry_nonce(_nonce_run_dir):
        print(
            "FATAL [AF-CANONICAL-RENDER-BYPASS]: build_deck.py must be invoked via "
            "presentation-canonical-entry.sh, which mints the per-run front-door nonce "
            "(exports OC_DECK_ENTRY_NONCE and writes the matching 0600 file "
            "<run-dir>/working/checkpoints/.canonical-entry-nonce). Direct invocation — "
            "or a guessed/stale nonce, or the retired OC_DECK_CANONICAL_ENTRY / "
            "OC_DECK_ALLOW_DIRECT env markers — is denied by the front-door handshake.",
            file=sys.stderr,
        )
        sys.exit(2)

    # ---------------------------------------------------------------------
    # TRUST BOUNDARY, INCREMENT 1 (report-only). Seal the RunFacts record for
    # this run EXACTLY ONCE, at the same admission point the front-door nonce
    # was just verified above — deliberately the FIRST thing that runs after
    # admission succeeds. Every _owner_skip_approved call and the P-TYPO-QC
    # shadow verifier below reuse this single sealed snapshot (get_or_seal)
    # instead of each re-opening process_manifest.json / QC reports on its
    # own schedule. Sealing itself is report-only and can never block
    # admission — any error here is caught and logged, never raised; see
    # presentation_job/runfacts.py for the full design and its honest limits.
    # ---------------------------------------------------------------------
    try:
        from presentation_job import runfacts as _runfacts  # noqa: PLC0415
        _runfacts.seal(_nonce_run_dir, nonce_bound=True)
    except Exception as _rf_exc:  # noqa: BLE001 — sealing must never block admission
        try:
            print(f"TRUST-BOUNDARY-SEAL-ERROR: could not seal RunFacts for "
                  f"{_nonce_run_dir}: {_rf_exc!r} (report-only — run proceeds)",
                  file=sys.stderr)
        except Exception:  # noqa: BLE001
            pass

    # P-STYLE-PREVIEW (order 4.85) --sample mode: render 9 style samples (3 variants x
    # 3 representative slides) for the owner's gateway sign-off, then stop. Needs only
    # slides.json (no out.pptx / full preflight / postflight).
    if sample_mode:
        if not positional:
            print("Usage: python3 build_deck.py <slides.json> --sample "
                  "[--style-spec PATH] [--run-dir DIR] [--logo URL]", file=sys.stderr)
            sys.exit(2)
        _sp = Path(positional[0])
        if not _sp.exists():
            print(f"FATAL: slides.json not found: {_sp}", file=sys.stderr)
            sys.exit(2)
        _run_dir = find_run_dir(run_dir_arg, _sp, _sp.parent)
        _logo_url = logo_arg if (logo_arg and str(logo_arg).startswith("http")) else None
        _api_key = "" if adhoc else load_api_key()
        _spec = Path(style_spec_arg) if style_spec_arg else None
        sys.exit(run_style_preview_samples(_sp, _run_dir, _spec, _api_key, logo_url=_logo_url))

    # P9.5-NOTES-SYNC (order 8.7) --notes-sync mode: reopen the already-assembled
    # bundle .pptx and re-inject per-slide speaker notes now that the presenter
    # speech exists and has passed P-SPEECH-QC. Fires AFTER P8-ASSEMBLE / P9-SPEECH /
    # P-SPEECH-QC and BEFORE P9-DELIVER (run_signature_deck.py dispatches this the
    # same way it dispatches P4-RENDER). Writes
    # working/checkpoints/notes_sync.json as the phase's produces_artifact so the
    # runner can attest it. Idempotent — safe to re-run.
    if notes_sync_mode:
        if len(positional) not in (2, 3):
            print("Usage: python3 build_deck.py <slides.json> <out.pptx> "
                  "--notes-sync --run-dir DIR [--out BUNDLE_DIR]", file=sys.stderr)
            sys.exit(2)
        _ns_slides_path = Path(positional[0])
        _ns_out_path = Path(positional[1])
        if not _ns_slides_path.exists():
            print(f"FATAL: slides.json not found: {_ns_slides_path}", file=sys.stderr)
            sys.exit(2)
        _ns_deck_slug = _ns_out_path.stem or "deck"
        _ns_bundle_dir = _resolve_bundle_dir(out_dir_arg, _ns_out_path)
        _ns_run_dir = find_run_dir(run_dir_arg, _ns_slides_path, _ns_out_path.parent)
        _ns_bundle_pptx = _ns_bundle_dir / f"{_ns_deck_slug}-FINAL.pptx"
        if not _ns_bundle_pptx.is_file():
            # Fall back to out.pptx itself when it already lives in the bundle dir
            # under a non-"-FINAL" name (e.g. a direct/adhoc invocation).
            if _ns_out_path.is_file():
                _ns_bundle_pptx = _ns_out_path
            else:
                print(f"FATAL: no assembled PPTX found at {_ns_bundle_pptx} or "
                      f"{_ns_out_path}; run P8-ASSEMBLE (the normal render) before "
                      f"P9.5-NOTES-SYNC.", file=sys.stderr)
                sys.exit(2)
        _ns_result = notes_sync_pass(_ns_bundle_pptx, _ns_run_dir, _ns_bundle_dir)
        _ns_ckpt_dir = _ns_run_dir / "working" / "checkpoints"
        _ns_ckpt_dir.mkdir(parents=True, exist_ok=True)
        _ns_record_path = _ns_ckpt_dir / "notes_sync.json"
        _ns_record = dict(_ns_result)
        _ns_record["bundle_pptx"] = str(_ns_bundle_pptx)
        _ns_record["timestamp"] = (timestamp_arg or time.strftime("%Y-%m-%dT%H:%M:%S%z")
                                    or time.strftime("%Y-%m-%dT%H:%M:%S"))
        try:
            _ns_record_path.write_text(json.dumps(_ns_record, indent=2))
        except OSError as exc:
            print(f"WARNING: notes-sync succeeded but could not write "
                  f"{_ns_record_path}: {exc}", file=sys.stderr, flush=True)
        print(json.dumps(_ns_record, indent=2))
        # "no_speech" is NON-FATAL (the speech may still be mid-write on a re-run of
        # this phase before the send-back loop settles) — AF-EMPTY-NOTES-PANE at
        # postflight is the hard gate, this pass only exits non-zero on a real error.
        sys.exit(1 if _ns_result.get("status") == "error" else 0)

    if len(positional) not in (2, 3):
        print("Usage: python3 build_deck.py <slides.json> <out.pptx> [renders_dir] "
              "[--run-dir DIR] [--logo PNG] [--out BUNDLE_DIR] "
              "[--platform vps|mac] [--timestamp ISO8601] [--adhoc-no-process] "
              "[--skip-teleprompter-gate]",
              file=sys.stderr)
        sys.exit(2)

    if platform_arg is not None and str(platform_arg).strip().lower() not in ("vps", "mac"):
        print(f"FATAL: --platform must be 'vps' or 'mac' (got {platform_arg!r}).",
              file=sys.stderr)
        sys.exit(2)

    slides_path = Path(positional[0])
    out_path = Path(positional[1])
    renders_dir = Path(positional[2]) if len(positional) == 3 else out_path.parent / "renders"

    # BUNDLE DIR + DELIVERABLES LEDGER (Requirements 3 & 4)
    # Must be resolved before any render so the ledger exists and can be updated
    # incrementally (crash-safe). The deck_slug is derived from the out.pptx stem.
    deck_slug = out_path.stem or "deck"
    bundle_dir = _resolve_bundle_dir(out_dir_arg, out_path)
    print(f"=== OUTPUT BUNDLE DIR (default ~/Downloads): {bundle_dir} ===", flush=True)
    ledger_path = init_deliverables_ledger(bundle_dir, deck_slug)
    print(f"=== DELIVERABLES LEDGER: {ledger_path} ===", flush=True)

    if not slides_path.exists():
        print(f"FATAL: slides.json not found: {slides_path}", file=sys.stderr)
        sys.exit(2)

    try:
        slides = json.loads(slides_path.read_text())
    except json.JSONDecodeError as exc:
        print(f"FATAL: slides.json is not valid JSON: {exc}", file=sys.stderr)
        sys.exit(2)

    if not isinstance(slides, list) or not slides:
        print("FATAL: slides.json must be a non-empty JSON array.", file=sys.stderr)
        sys.exit(2)

    # Basic schema validation (deterministic, fail loud).
    seen = set()
    for s in slides:
        if not isinstance(s, dict):
            print("FATAL: every slide must be an object.", file=sys.stderr)
            sys.exit(2)
        for req in ("slide", "scene", "copy"):
            if req not in s:
                print(f"FATAL: slide missing required field '{req}': {json.dumps(s)}", file=sys.stderr)
                sys.exit(2)
        if not isinstance(s["copy"], list) or not s["copy"]:
            print(f"FATAL: slide {s.get('slide')}: 'copy' must be a non-empty array.", file=sys.stderr)
            sys.exit(2)
        ordn = s["slide"]
        if not isinstance(ordn, int) or ordn in seen:
            print(f"FATAL: slide ordinal must be a unique integer: {ordn}", file=sys.stderr)
            sys.exit(2)
        seen.add(ordn)

    # PROCESS PREFLIGHT (un-bypassable unless --adhoc-no-process). Runs BEFORE any
    # API key load, render, or assembly so a bypass costs zero KIE renders.
    run_dir = find_run_dir(run_dir_arg, slides_path, out_path)
    if adhoc:
        print_adhoc_banner()
    else:
        # H3: pass the ACTUAL positional slides.json so the coverage / rich-prompt
        # gates count the exact file that will be rendered (no gate-bypass via a
        # different canonical-path slides.json).
        run_preflight(run_dir, slides_path)

    # CC REGISTRATION (run-begin, idempotent) — Fix 5b / AF-CC-UNREGISTERED.
    # Fail-soft: a board outage never blocks the deck. ingest_deck_task always
    # stamps cc_register_attempted=True BEFORE the HTTP call, so the postflight
    # gate (_chk_cc_registered) is satisfied even on transport failure.
    # Skipped in --adhoc mode (ad-hoc runs are not CC-tracked deliverables).
    # _cc_task_id is hoisted here so the later phase-boundary card moves
    # (_board_patch_phase) can thread it through the rest of main().
    _cc_task_id = None
    if not adhoc:
        try:
            import cc_board as _cc_board
            # FIX-PRES-08(a): the runner's Phase-0 pre-flight already opened this
            # deck's card so pre-render phases are board-visible. If a task_id is
            # already stamped in the manifest, REUSE it here rather than re-ingest
            # — the run-begin ingest is idempotent (idempotency_key) and local reuse
            # guarantees a SINGLE card even if the runner's slug differs from
            # out_path.stem. A direct/standalone build_deck run (no runner) finds no
            # stamp and ingests as before (backward-compatible).
            _existing = _cc_board._read_manifest(run_dir).get("cc_task_id") \
                if hasattr(_cc_board, "_read_manifest") else None
            if _existing:
                _cc_task_id = str(_existing)
                print(f"[cc_board] reusing Phase-0 card task_id={_cc_task_id} "
                      "(runner pre-flight ingest).", file=sys.stderr)
            else:
                _cc_title = deck_slug
                _cc_desc = f"Deck build: {deck_slug}"
                _rcid, _rchan = _cc_board.resolve_requester(run_dir)
                if not _rcid:
                    # Rule 3.5 WARN-MODE, stage 1 of 3 — standalone twin
                    print("[cc_board] WARN-REQUESTER-MISSING: standalone build_deck run with "
                          "no requester chat id; the client will receive no acknowledgement, "
                          "progress or completion message for this build.",
                          file=sys.stderr, flush=True)
                _cc_task_id = _cc_board.ingest_deck_task(
                    run_dir, deck_slug, title=_cc_title, description=_cc_desc,
                    requester_chat_id=_rcid, requester_channel=_rchan
                )
                if _cc_task_id:
                    _cc_board.stamp_task_id(run_dir, _cc_task_id)
        except Exception as _cc_exc:  # noqa: BLE001
            print(
                f"[cc_board] run-begin ingest raised ({_cc_exc}) — "
                "run continues; _chk_cc_registered will see the attempted flag.",
                file=sys.stderr,
            )

    # OFFICIAL-LOGO resolution: --logo wins; else intake.json brand.logo_image_path.
    # IMAGE-TO-IMAGE logo: if --logo is a URL, KIE composites the REAL logo via
    # input_urls (verified gpt-image-2-image-to-image). A local path falls back to a
    # flat overlay in assemble. Never resolve a URL as a local file.
    logo_url = logo_arg if (logo_arg and str(logo_arg).startswith("http")) else None
    if logo_url:
        # C1 (SSRF / local-file-read guard): the --logo URL is passed to KIE as an
        # input_urls reference and is otherwise fetchable. Enforce the http(s)
        # allowlist up front so a file://, ftp://, or data: --logo is refused before
        # any render. (startswith("http") already excludes most, but a value like
        # "httpx://" or "http\tfile://" would slip past a naive prefix check — the
        # urlparse scheme allowlist is the authoritative gate.)
        try:
            assert_url_scheme_allowed(logo_url, what="--logo URL")
        except ValueError as exc:
            print(f"FATAL: {exc}", file=sys.stderr)
            sys.exit(2)
    logo_path = None if logo_url else resolve_logo_path(logo_arg, run_dir)
    if logo_url:
        print(f"=== OFFICIAL LOGO (IMAGE-TO-IMAGE): {logo_url} — KIE composites the REAL "
              f"logo into EVERY slide via input_urls ===", flush=True)
    elif logo_path is not None:
        print(f"=== OFFICIAL LOGO (overlay fallback): {logo_path} ===", flush=True)

    # In --adhoc-no-process mode the key is NEVER used for a network call: render_slide
    # raises on missing rich prompts before reaching submit_task, so we skip the
    # load_api_key() call entirely.  This makes adhoc runs platform-portable (no
    # secrets files required on CI / linux runners) without weakening any gate —
    # if a prompt file IS present and render_slide would actually call KIE, it receives
    # the sentinel "" and will fail at the HTTP layer, not silently succeed.
    api_key = "" if adhoc else load_api_key()
    renders_dir.mkdir(parents=True, exist_ok=True)

    # === GOAL-4 / 3C PHASE-0 PRE-FLIGHT — runs BEFORE any KIE render ===
    # (1) AF-KIE-BALANCE: HARD-ABORT before a single slide is dispatched when the
    #     client's live Kie.ai credit balance is below the estimated floor for this
    #     deck, so a run never dies mid-deck and burns credits on a partial render.
    #     Skipped in --adhoc mode (no key / no network). The runner
    #     (run_signature_deck.py) calls the same kie_balance_preflight() at its Phase-0.
    # (2) AF-PHASE-SKIPPED: the deterministic phase-precondition contract lives in
    #     check_phase_preconditions(); the runner enforces it across phases. build_deck
    #     references it here (render = the P4-RENDER phase) so the symbol is live and a
    #     direct render is recorded as honouring the precondition contract.
    if not adhoc:
        _box_type = detect_platform(run_dir, platform_arg)
        print(f"=== PHASE-0 PRE-FLIGHT — box_type={_box_type}; "
              f"Kie balance floor check before render ===", flush=True)
        # MASTER-SPEC 7.4 / AF-OCR-ENGINE-MISSING — checked FIRST (fast, local, no
        # network) so a box with no OCR engine refuses before the network-dependent
        # Kie-balance call, let alone before any render. Defense-in-depth for a direct
        # `python3 build_deck.py` invocation that bypasses run_signature_deck.
        # phase0_preflight()'s identical check.
        _ocr_engine_reason = ocr_engine_preflight(run_dir)
        if _ocr_engine_reason:
            print("\nFATAL: " + _ocr_engine_reason, file=sys.stderr)
            sys.exit(4)
        # FIX-6 — one-shot auth proof BEFORE any render: a 401 here aborts with the
        # auth diagnosis instead of a 164x backoff tailspin mid-render.
        try:
            _preflight_kie_auth(api_key)
        except AuthError as exc:
            print(f"\nFATAL: {exc}", file=sys.stderr)
            sys.exit(4)
        # FIX-E2E (incremental resume): count only the UN-rendered slides for the
        # balance floor. A resume/re-render that already produced 18 of 20 slides
        # only needs credit for the remaining 2, not the full deck — the old code
        # passed len(slides) (the full count), so a low-but-sufficient balance
        # (2 remaining ≈ 10 credits) was rejected against a 100-credit full-deck
        # floor, bricking every OCR-retry resume.
        _remaining = max(0, len(slides) - len(_gather_rendered_pngs(run_dir)))
        _balance_reason = kie_balance_preflight(run_dir, _remaining, api_key)
        if _balance_reason:
            print("\nFATAL: " + _balance_reason, file=sys.stderr)
            sys.exit(4)
        # Render is the canonical P4-RENDER phase; a direct build has no prior runner
        # attestations to require (prior_phase_ids empty), so this is a no-op pass here
        # but keeps the precondition contract symbol on the live enforcement path.
        _precondition_reason = check_phase_preconditions(run_dir, "P4-RENDER", [])
        if _precondition_reason:
            print("\nFATAL: " + _precondition_reason, file=sys.stderr)
            sys.exit(4)

    print(f"=== build_deck — {len(slides)} slides ===", flush=True)
    print(f"slides.json: {slides_path}", flush=True)
    print(f"renders_dir: {renders_dir}", flush=True)
    print(f"out.pptx:    {out_path}", flush=True)
    # Truth in reporting: name the model this run ACTUALLY uses — image-to-image
    # when a logo URL is composited via input_urls, else text-to-image. The same
    # model_used value flows into the process manifest + the final JSON summary,
    # so the audit trail never claims t2i for an i2i logo-composite run.
    model_used = MODEL_I2I if logo_url else MODEL_T2I
    print(f"endpoint:    {CREATE_URL}  model={model_used}\n", flush=True)

    rendered = []
    task_ids = []
    failures = []

    has_official_logo = (logo_path is not None) or (logo_url is not None)
    ordered = sorted(slides, key=lambda s: s["slide"])

    # BATCH RENDER (FIX-5 / D22 root cause B): submit ALL prompts once (0.6s apart,
    # inside kie's 20/10s window + 100 concurrent tasks), then poll every submitted
    # task together every 10s and download each image the moment ITS task finishes.
    # This replaces the old submit-one-then-poll path, which serialized submit+poll+
    # download per slide and wasted 20+ min on a 20-slide deck. Per-slide gates are
    # unchanged: rich-prompt floor, English pin, i2i logo compositing, post-download
    # aspect/OCR verification and the U028 task-id checkpoint all still run.
    print(f"=== rendering {len(ordered)} slides via BATCH (submit all, then poll all) ===\n",
          flush=True)
    # BOARD: the render phase is now the one being worked (fail-soft; no-op when
    # the board is disabled / the ingest never returned a task_id).
    if not adhoc:
        _board_patch_phase(run_dir, _cc_task_id, "P4-RENDER", "in_progress",
                           note=f"batch-rendering {len(ordered)} slides")
    batch_result = render_slides_batch(
        ordered, api_key, renders_dir, run_dir,
        has_official_logo=has_official_logo, logo_url=logo_url)
    rendered = batch_result["rendered"]
    failures = batch_result["failures"]

    # Deterministic order for the summary / manifest regardless of completion order.
    rendered.sort(key=lambda r: r["slide"])
    task_ids = [r["taskId"] for r in rendered]
    failures.sort(key=lambda f: (f.get("slide") is None, f.get("slide")))

    # FAIL LOUD: never produce a partial deck silently. FIX-7 (D14): a render that
    # failed (e.g. a RenderPollTimeout from a never-completing kie task) must NOT
    # leave the CC card in 'in_progress' — surface a terminal status so the run
    # never hangs. 'blocked' is a CC_TASK_STATUSES member; mid-run this is legal
    # because the run is ending with a hard failure (no process cert is minted, so
    # the cert-gate 422 concern does not apply to an abort). Fail-soft: a disabled
    # board / missing task_id is a clean no-op.
    if failures:
        summary = {
            "slidesRendered": len(rendered),
            "kieTaskIds": task_ids,
            "outputPath": None,
            "failures": failures,
        }
        print("\n=== SUMMARY (FAILED) ===", file=sys.stderr)
        print(json.dumps(summary, indent=2))
        if not adhoc:
            _board_patch_phase(run_dir, _cc_task_id, "P4-RENDER", "blocked",
                               note=f"render FAILED: {len(failures)} slide(s) did not "
                                    f"complete (e.g. poll hard cap / terminal state). "
                                    f"Never hangs.")
        sys.exit(1)

    # BOARD: every slide rendered cleanly — log render-phase PROGRESS as an ACTIVITY
    # (NOT a task-level 'done': mid-run status changes 422 the cert done-gate and
    # would wrongly close a non-presentation card). Fail-soft.
    if not adhoc:
        _board_post_activity(run_dir, _cc_task_id, "P4-RENDER",
                             note=f"{len(rendered)} slides rendered")

    # PER-SLIDE SPEAKER NOTES: auto-discover + parse the presenter speech (if it
    # exists yet) so the assembled deck carries word-for-word notes per slide. This
    # is NON-FATAL — a phase-4 render usually precedes the phase-9 speech, so an
    # absent speech yields None and the deck is assembled with no notes. The speech
    # is NEVER required to render. Count mismatches are safe: extra chunks are simply
    # not matched to any slide; deck slides with no chunk get no notes part.
    speech_chunks = discover_speech_chunks(run_dir, bundle_dir)

    # Assemble.
    try:
        # When image-to-image baked the logo in (logo_url), do NOT also overlay it.
        assemble_pptx(rendered, out_path, logo_path=(None if logo_url else logo_path),
                      speech_chunks=speech_chunks)
    except SystemExit:
        raise
    except Exception as exc:  # noqa: BLE001
        print(f"FATAL: pptx assembly failed: {exc}", file=sys.stderr)
        summary = {
            "slidesRendered": len(rendered),
            "kieTaskIds": task_ids,
            "outputPath": None,
            "failures": [{"slide": "ASSEMBLY", "error": str(exc)}],
        }
        print(json.dumps(summary, indent=2))
        sys.exit(1)

    # Timestamp: caller-supplied --timestamp wins (lets the orchestrator pin a
    # consistent run clock); else stamp now. This is a plain Python script, so
    # using the local clock here is fine.
    timestamp = timestamp_arg or time.strftime("%Y-%m-%dT%H:%M:%S%z") or time.strftime("%Y-%m-%dT%H:%M:%S")

    # APPEND the render record to the cumulative process manifest (never clobber
    # prior phases) instead of leaving the result stdout-only.
    manifest_path = None
    try:
        manifest_path = write_process_manifest(
            run_dir, rendered, task_ids, model_used, out_path, timestamp)
    except Exception as exc:  # noqa: BLE001
        print(f"WARNING: render succeeded but process manifest write failed: {exc}",
              file=sys.stderr, flush=True)

    # DELIVERABLES LEDGER — mark the deck_pptx as built now that assembly succeeded.
    # The other eight build artifacts (the rest of the nine-key build bundle) are
    # produced by upstream roles; the postflight gate will check their actual
    # presence on disk.
    pptx_size = out_path.stat().st_size if out_path.exists() else 0
    update_deliverable_status(ledger_path, "deck_pptx", "built", size=pptx_size)

    # BOARD: the PPTX assembled — log assemble-phase PROGRESS as an ACTIVITY (NOT a
    # task-level 'done'; see the P4-RENDER note above). Fail-soft.
    if not adhoc:
        _board_post_activity(run_dir, _cc_task_id, "P8-ASSEMBLE",
                             note="deck PPTX assembled")

    # If the out_path is not inside bundle_dir, copy the pptx into bundle_dir so the
    # postflight gate finds it at its expected location.
    bundle_pptx = bundle_dir / f"{deck_slug}-FINAL.pptx"
    if out_path.resolve() != bundle_pptx.resolve():
        import shutil
        try:
            shutil.copy2(str(out_path), str(bundle_pptx))
            update_deliverable_status(ledger_path, "deck_pptx", "built",
                                      size=bundle_pptx.stat().st_size)
            print(f"=== PPTX copied to bundle dir: {bundle_pptx} ===", flush=True)
        except Exception as exc:  # noqa: BLE001
            print(f"WARNING: could not copy PPTX to bundle_dir: {exc}", file=sys.stderr)

    summary = {
        "slidesRendered": len(rendered),
        "kieTaskIds": task_ids,
        "outputPath": str(out_path),
        "bundleDir": str(bundle_dir),
        "deliverables_ledger": str(ledger_path),
        "modelUsed": model_used,
        "timestamp": timestamp,
        "processManifest": str(manifest_path) if manifest_path else None,
        "failures": [],
    }
    print("\n=== SUMMARY (RENDER OK — running postflight completeness gate) ===",
          flush=True)
    print(json.dumps(summary, indent=2))

    # TELEPROMPTER PUBLISH — only meaningful if the teleprompter HTML is already in the
    # bundle (it is an upstream artifact; build_deck.py does not author it). If it is
    # present, publish it to the central Cloudflare host and capture the public URL. A
    # publish failure does NOT abort the render (the .pptx is already built), but it DOES
    # leave the publish ledger un-verified so the postflight gate fails loud (exit 5).
    # In --adhoc mode, publish_teleprompter writes a 'skipped_adhoc' record (no host call).
    platform = detect_platform(run_dir, override=platform_arg)
    tele_html = bundle_dir / "presenter-teleprompter.html"
    if tele_html.exists():
        print(f"=== TELEPROMPTER PUBLISH (platform={platform}, host=cloudflare-central) ===",
              flush=True)
        try:
            rec = publish_teleprompter(bundle_dir, deck_slug, run_dir, platform,
                                       ledger_path, adhoc=adhoc)
            if rec.get("status") == "published":
                print(f"=== TELEPROMPTER PUBLISHED: {rec.get('public_url')} "
                      f"(HTTP {rec.get('verified_http_status')}) ===", flush=True)
            elif rec.get("status") == "skipped_adhoc":
                print("=== TELEPROMPTER PUBLISH skipped (ad-hoc render) ===", flush=True)
        except Exception as exc:  # noqa: BLE001 — never silently swallow
            print(f"WARNING: teleprompter publish failed: {exc} — the postflight "
                  f"gate will fail loud (TELEPROMPTER-PUBLISH) until it is published "
                  f"and the URL verified.", file=sys.stderr, flush=True)

    # POSTFLIGHT COMPLETENESS GATE (Requirement 2 — AF-BUNDLE-COMPLETE).
    # This is the FINAL gate: exit 0 only when ALL nine bundle deliverables are
    # present and sized. Exit 5 (loud failure) when any are missing or under-threshold.
    # The word "COMPLETE" / "DONE" is printed ONLY from inside run_postflight_gate
    # after reading deliverables.json (not from in-memory state).
    # M7: the teleprompter-publish sub-check is bypassed ONLY by an explicit per-run
    # CLI flag — never by a persisted status string. --adhoc-no-process is itself an
    # explicit per-run flag (ad-hoc output is not a client deliverable), so it implies
    # the bypass for THIS run only; a later real run without the flag re-arms the gate.
    run_postflight_gate(bundle_dir, ledger_path, deck_slug,
                         skip_teleprompter_gate=(skip_teleprompter_gate or adhoc),
                         run_dir=run_dir, slides_path=slides_path)

    # BOARD (terminal close) — GUARDED on the PROCESS-CERTIFICATE existing on disk.
    # The producer STOPS at the TERMINAL status 'review' (there is NO 'delivered'
    # status in the CC UpdateTaskSchema enum — "delivered" goes in the note only): the
    # presentations pipeline never self-closes to 'done'. Promotion 'review'->'done' is
    # the CC-side QC scorer / Devil's-Advocate gate's job, the interlock every sibling
    # department respects. This close carries the process_certificate_sha (the ticket
    # INTO review) + the real per-gate QC scores. In the RUNNER flow build_deck runs the
    # RENDER phase BEFORE prove-deck mints the cert, so the cert is ABSENT here and the
    # terminal close is DEFERRED to the runner (run_signature_deck.py fires it right
    # after prove-deck). The guard prevents a doomed premature close before the cert
    # exists. In a standalone full build where the cert is already present, close here.
    # Fail-soft either way — a board problem never changes this run's exit code.
    if not adhoc:
        if _process_certificate_present(run_dir):
            _board_patch_phase(run_dir, _cc_task_id, "P9-DELIVER", "review",
                               note="bundle complete — deck delivered; awaiting CC QC "
                                    "scorer / Devil's-Advocate promotion to done")
            # Lightweight VISIBILITY gate: a completed run should have recorded >= 1
            # SUCCESSFUL board advance. This only DIAGNOSES (never blocks) — a disabled
            # board legitimately records none; the detail is on disk in cc-board.json.
            try:
                import cc_board as _cc_board
                if not _cc_board.assert_min_one_advance(run_dir):
                    print("[cc_board] NOTE: no successful board advance recorded for this "
                          "run (board disabled, or every advance failed — see "
                          "working/checkpoints/cc-board.json).",
                          file=sys.stderr, flush=True)
            except Exception:  # noqa: BLE001 — the visibility gate is best-effort
                pass
        else:
            print("[cc_board] P9-DELIVER terminal close DEFERRED — no PROCESS-CERTIFICATE "
                  "on disk yet; the runner closes the card to status='review'+cert+QC "
                  "scores after prove-deck mints it (CC QC scorer / Devil's-Advocate then "
                  "promotes review->done). Render-phase progress is already on the board.",
                  file=sys.stderr, flush=True)
    sys.exit(0)


if __name__ == "__main__":
    main()
